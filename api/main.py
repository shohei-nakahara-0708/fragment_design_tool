# app.py
from __future__ import annotations

import base64
import json
import os, subprocess
import re
import signal
import shlex
import uuid
from dataclasses import dataclass, field as dataclass_field
from pathlib import Path
from typing import List, Optional, Dict, Any, Literal, Tuple
from urllib.parse import quote

import httpx
from dotenv import load_dotenv
from fastapi import APIRouter,FastAPI, UploadFile, File, HTTPException, Body,Form, BackgroundTasks,Request,Response
from fastapi.responses import FileResponse, JSONResponse,StreamingResponse,RedirectResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, Field

import time
import sqlite3
import zipfile
import tempfile
from datetime import datetime, timezone, date
import traceback

import gspread
from google.auth.transport.requests import Request as GoogleAuthRequest
from google.oauth2.service_account import Credentials
from collections import Counter

from difflib import SequenceMatcher

import random
from gspread.exceptions import APIError
import psycopg
from psycopg.rows import dict_row
from psycopg import OperationalError

import fitz  # pymupdf

import shutil

import math

from PIL import Image, ImageOps

import mimetypes
import requests

import io
import asyncio
import queue
import threading
from concurrent.futures import ThreadPoolExecutor

import logging

logger = logging.getLogger(__name__)
load_dotenv()

DATABASE_URL = os.getenv("DATABASE_URL", "")
API_BASE_URL = os.getenv("API_BASE_URL", "")

# データベース接続設定
DB_CONNECT_TIMEOUT = int(os.getenv("DB_CONNECT_TIMEOUT", "30"))  # 接続タイムアウト（秒）
DB_QUERY_TIMEOUT = int(os.getenv("DB_QUERY_TIMEOUT", "60"))      # クエリタイムアウト（秒）
DB_RETRY_ATTEMPTS = int(os.getenv("DB_RETRY_ATTEMPTS", "3"))     # リトライ回数
DB_RETRY_DELAY = float(os.getenv("DB_RETRY_DELAY", "2.0"))        # リトライ間隔（秒）

SUPABASE_URL = os.environ["SUPABASE_URL"]
SUPABASE_SERVICE_ROLE_KEY = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
SUPABASE_BUCKET = os.environ.get("SUPABASE_BUCKET", "jobs")

APP_DIR = Path(__file__).resolve().parent

def resolve_data_dir() -> Path:
    v = (os.getenv("DATA_DIR") or "").strip()
    if v:
        p = Path(v)
        try:
            p.mkdir(parents=True, exist_ok=True)
            return p
        except PermissionError:
            pass

    p = APP_DIR / "_data"
    p.mkdir(parents=True, exist_ok=True)
    return p

DATA_DIR = resolve_data_dir()

VM_DIFF_PREVIEW_DIR = DATA_DIR / "vm_diff_previews"
VM_DIFF_PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

TEMPLATE_PATH = APP_DIR / "template.html"
_cached_template: Optional[str] = None

# Postgres移行したなら不要。残すならローカル専用に。
DB_PATH = DATA_DIR / "index.sqlite"

EXPORT_DIR = DATA_DIR / "_exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)

MAX_HEIGHT = 2000
BASE_VIEWPORT = {"width": 600, "height": 800}

# より高性能なモデルを使用して精度向上
AI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o")  # gpt-4o-mini → gpt-4o
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
OPENAI_BASE_URL = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1")
AI_TIMEOUT = 180  # より複雑な処理に対してタイムアウトを延長

# EMU -> pt (pptx uses EMU units for font size)
EMU_PER_PT = 12700

TIME_PAT = re.compile(r"(\d{1,2}:\d{2}\s*[～〜\-ー~]\s*\d{1,2}:\d{2})")

# 医療用語の標準化辞書（精度向上のため）
MEDICAL_ORGANIZATION_NORMALIZATION = {
    "東京大学医学部附属病院": ["東大病院", "東大医学部", "東京大医学部", "東大附属病院"],
    "慶應義塾大学病院": ["慶應病院", "慶大病院", "慶応大学病院"],
    "順天堂大学医学部附属順天堂病院": ["順天堂病院", "順天堂大病院"],
    "国立がん研究センター": ["がん研", "国がん"],  # 「がんセンター」を削除（他の機関との混同を避ける）
    "日本医科大学付属病院": ["日医大病院", "日本医大病院"],
    "大阪国際がんセンター": ["大阪国際がんセンター"],  # 正式名称を追加
}

MEDICAL_TITLE_NORMALIZATION = {
    "教授": ["prof", "Prof", "Professor", "教授"],
    "准教授": ["准教授", "準教授", "Associate Professor"],
    "講師": ["講師", "Lecturer"],
    "助教": ["助教", "Assistant Professor"],
    "部長": ["部長", "Director"],
    "主任": ["主任", "Chief"],
    "センター長": ["センター長", "Center Director"],
}

def normalize_medical_terms(text: str) -> str:
    """医療用語を標準化（日本語テキスト対応の境界検出）"""
    import re
    result = text

    # 日本語対応の境界パターン（\b は日本語文字で機能しないため独自実装）
    # 組織名・役職名の前後が漢字でないことを境界とする（助詞「の」「は」等は境界OK）
    def _jp_boundary_pattern(variant: str) -> str:
        escaped = re.escape(variant)
        # CJK漢字を含むかチェック
        has_cjk = bool(re.search(r'[\u4e00-\u9fff]', variant))
        if has_cjk:
            # 日本語の場合: 前後が漢字でないことを確認（ひらがな・カタカナ・記号は境界扱い）
            return f'(?<![\u4e00-\u9fff]){escaped}(?![\u4e00-\u9fff])'
        else:
            # 英語の場合: 通常の \b を使用
            return r'\b' + escaped + r'\b'

    # 組織名の標準化
    for standard, variants in MEDICAL_ORGANIZATION_NORMALIZATION.items():
        for variant in variants:
            pattern = _jp_boundary_pattern(variant)
            result = re.sub(pattern, standard, result)

    # 役職名の標準化
    for standard, variants in MEDICAL_TITLE_NORMALIZATION.items():
        for variant in variants:
            pattern = _jp_boundary_pattern(variant)
            result = re.sub(pattern, standard, result)

    return result

ORG_CANON = {
    "MSD": "MSD株式会社",
    "MSD KK": "MSD株式会社",
    "MSD K.K.": "MSD株式会社",
    "Merck": "MSD株式会社",
}

SESSION_TIME_RE = re.compile(r"[①②③④⑤⑥⑦⑧⑨⑩]?\s*(\d{1,2}[:：]\d{2}\s*[～〜\-ー~]\s*\d{1,2}[:：]\d{2})")
TYPESET_JS = r"""
({ data }) => {
  try {
    console.log('=== TYPESET_JS EXECUTION START ===');
    console.log('JavaScript is running successfully');
    console.log('Data object exists:', !!data);
    console.log('Talks array exists:', !!data.talks);
    console.log('Talks count:', data.talks?.length || 0);
  } catch (e) {
    console.log('CRITICAL ERROR in initial log:', e);
  }
  
  // =========================================================
  // helpers: normalize
  // =========================================================
  const norm = (s) => String(s ?? "")
    .replace(/\u3000/g, " ")
    .replace(/[ \t\r\f\v]+/g, " ")
    .replace(/\n+/g, "\n")
    .trim();

  const oneLine = (s) => norm(String(s ?? "").replace(/\n/g, " "));

  const unifyTilde = (s) => String(s ?? "")
    .replace(/～/g, "〜")
    .trim();

  const stripEdgePunct = (s) => String(s ?? "")
    .replace(/^[\s,，、:：/／]+/, "")
    .replace(/[\s,，、:：/／]+$/, "")
    .trim();

  const isMostlyAscii = (s) => {
    const t = String(s ?? "");
    const visible = [...t].filter(ch => !/\s/.test(ch));
    if (!visible.length) return false;
    const ascii = visible.filter(ch => ch.charCodeAt(0) < 128).length;
    return ascii / visible.length >= 0.7;
  };

  const wordCount = (s) => oneLine(s).split(/\s+/).filter(Boolean).length;

  const isOnlySymbols = (s) => {
    const t = oneLine(s);
    if (!t) return false;
    return /^[~〜～\-–—―−－・･:：,，/／()（）]+$/.test(t);
  };

  const startsWithWeakSymbol = (s) => {
    const t = oneLine(s);
    return /^[~〜～\-–—―−－:：,，/／]+/.test(t);
  };

  // title_lines を join するときに記号前の余計な空白を作らない
  const isJpChar = (ch) => /[ぁ-んァ-ヶー一-龠々〇\u4E00-\u9FFF]/.test(ch || "");

  const joinTitleLinesSmart = (lines) => {
    const arr = (lines || []).map(x => oneLine(x)).filter(Boolean);
    if (!arr.length) return "";

    let out = arr[0];
    for (let i = 1; i < arr.length; i++) {
      const cur = arr[i];

      // 記号始まりの行は前に空白を入れない
      if (/^[~〜～\-–—―−－・･:：/／]/.test(cur)) {
        out += cur;
      // 日本語同士の結合はスペース不要
      } else if (isJpChar(out[out.length - 1]) && isJpChar(cur[0])) {
        out += cur;
      } else {
        out += " " + cur;
      }
    }
    return oneLine(out);
  };

  // =========================================================
  // measurer element
  // =========================================================
  const getMeasurer = () => {
    let el = document.getElementById("__measurer__");
    if (!el) {
      el = document.createElement("div");
      el.id = "__measurer__";
      el.style.position = "fixed";
      el.style.left = "-10000px";
      el.style.top = "-10000px";
      el.style.whiteSpace = "pre";
      el.style.padding = "0";
      el.style.margin = "0";
      el.style.border = "0";
      document.body.appendChild(el);
    }
    return el;
  };

  const measure = (text, style) => {
    const el = getMeasurer();
    el.style.fontFamily = style.fontFamily;
    el.style.fontWeight = String(style.fontWeight);
    el.style.fontSize = style.fontSize;
    el.style.letterSpacing = style.letterSpacing ?? "normal";
    el.textContent = String(text ?? "");
    return el.scrollWidth;
  };

  // =========================================================
  // break candidates
  // =========================================================
  const hasJapanese = (s) => /[ぁ-んァ-ヶ一-龠々]/.test(String(s ?? ""));

  // 医学・専門用語の境界を考慮した改行位置決定
  const breakPositions = (s) => {
    const out = [];
    const jp = hasJapanese(s);

    // 基本的な区切り文字
    const breakers = new Set(
      jp
        ? [" ", "、", "。", ",", "，", ":", "：", "/", "／"]
        : [" ", "、", "。", ",", "，", ":", "：", "/", "／"]
    );

    // 英語の単語境界を検出する関数
    const isWordBoundary = (str, index) => {
      if (index <= 0 || index >= str.length) return false;
      const before = str[index - 1];
      const after = str[index];
      // アルファベットから非アルファベット、またはその逆の境界
      const beforeIsAlpha = /[a-zA-Z]/.test(before);
      const afterIsAlpha = /[a-zA-Z]/.test(after);
      return (beforeIsAlpha && !afterIsAlpha) || (!beforeIsAlpha && afterIsAlpha);
    };

    // 英語の单語途中かどうかチェック
    const isMidWord = (str, index) => {
      if (index <= 0 || index >= str.length) return false;
      const before = str[index - 1];
      const after = str[index];
      return /[a-zA-Z]/.test(before) && /[a-zA-Z]/.test(after);
    };

    // 医学用語・専門用語のパターン（途中で改行を避けるべき）
    const medicalTerms = [
      /GLP-1受容体作動薬/, /SGLT-2阻害薬/, /DPP-4阻害薬/, 
      /ACE阻害薬/, /ARB/, /β遮断薬/, /Ca拮抗薬/,
      /2型糖尿病/, /1型糖尿病/, /糖尿病性腎症/, /糖尿病性網膜症/,
      /心血管疾患/, /循環器疾患/, /虚血性心疾患/, /心不全/,
      /エビデンス/, /ガイドライン/, /プロトコル/, /アルゴリズム/,
      /バイオマーカー/, /リスクファクター/, /アウトカム/,
      /webセミナー/, /Webセミナー/, /WEBセミナー/,
    ];

    // 意味のある境界で改行候補を追加（優先順位付き）
    const meaningfulBreaks = [
      // 最高優先度：サブタイトル記号（医学文書でよく使用）
      { pattern: /－/g, priority: 1 },
      { pattern: /−/g, priority: 1 },
      { pattern: /–/g, priority: 1 },
      { pattern: /—/g, priority: 1 },
      { pattern: /―/g, priority: 1 },
      // 高優先度：長い文章の自然な区切り
      { pattern: /における/g, priority: 2 },
      { pattern: /について/g, priority: 2 },
      { pattern: /に関する/g, priority: 2 },
      { pattern: /による/g, priority: 3 },
      { pattern: /からの/g, priority: 3 },
      { pattern: /への/g, priority: 3 },
      { pattern: /との/g, priority: 3 },
      { pattern: /での/g, priority: 3 },
      { pattern: /としての/g, priority: 3 },
      { pattern: /という/g, priority: 3 },
      // 中優先度：接続詞的な表現
      { pattern: /〜/g, priority: 4 },
      { pattern: /～/g, priority: 4 },
      { pattern: /から/g, priority: 4 },
      { pattern: /まで/g, priority: 4 },
      { pattern: /より/g, priority: 4 },
      { pattern: /など/g, priority: 4 },
    ];

    // 基本的な区切り文字での改行候補
    for (let i = 0; i < s.length; i++) {
      const ch = s[i];
      if (breakers.has(ch)) out.push(i + 1);
    }

    // 意味のある境界での改行候補（医学用語の途中でなければ）
    for (const breakItem of meaningfulBreaks) {
      const pattern = breakItem.pattern;
      let match;
      while ((match = pattern.exec(s)) !== null) {
        const pos = match.index + match[0].length;
        
        // 医学用語の途中でないかチェック
        let inMedicalTerm = false;
        for (const medPattern of medicalTerms) {
          const medMatch = s.match(medPattern);
          if (medMatch && pos > medMatch.index && pos < medMatch.index + medMatch[0].length) {
            inMedicalTerm = true;
            break;
          }
        }
        
        if (!inMedicalTerm && !isMidWord(s, pos)) {
          out.push({ pos, priority: breakItem.priority });
          // 全角ダッシュはダッシュの前でも改行できるようにする
          if (/[–—―−－]/.test(match[0]) && match.index > 0) {
            out.push({ pos: match.index, priority: breakItem.priority });
          }
        }
      }
    }

    // 医学用語の途中かどうかチェックするヘルパー関数
    const isMedicalTermMiddle = (str, index) => {
      for (const medPattern of medicalTerms) {
        const medMatch = str.match(medPattern);
        if (medMatch && index > medMatch.index && index < medMatch.index + medMatch[0].length) {
          return true;
        }
      }
      return false;
    };

    // 前後に空白がある dash だけ候補にする
    const dashAroundSpaceRe = /\s[-–—―−－]\s/g;
    let m;
    while ((m = dashAroundSpaceRe.exec(s)) !== null) {
      out.push({ pos: m.index + 1, priority: 2 });
      out.push({ pos: m.index + m[0].length - 1, priority: 2 });
    }

    // 助詞の後（ただし医学用語の途中でなければ）
    const particleRe = /(を|の|に|と|へ|や|で|が|は|も)/g;
    while ((m = particleRe.exec(s)) !== null) {
      const pos = m.index + m[0].length;
      
      // 医学用語の途中でないかチェック
      let inMedicalTerm = false;
      for (const medPattern of medicalTerms) {
        const medMatch = s.match(medPattern);
        if (medMatch && pos > medMatch.index && pos < medMatch.index + medMatch[0].length) {
          inMedicalTerm = true;
          break;
        }
      }
      
      if (!inMedicalTerm && !isMidWord(s, pos)) {
        out.push({ pos, priority: 4 }); // 助詞は低優先度
      }
    }

    // 基本的な区切り文字での改行候補を優先度付きで追加
    for (let i = 0; i < s.length; i++) {
      const ch = s[i];
      if (breakers.has(ch)) {
        const pos = i + 1;
        // 英語の単語途中での改行を避ける
        if (!isMidWord(s, pos)) {
          out.push({ pos, priority: ch === '、' || ch === '，' ? 1 : (ch === ' ' ? 2 : 3) });
        }
      }
    }

    // 英語の単語境界を明示的に追加（日本語混在テキスト用）
    for (let i = 1; i < s.length; i++) {
      if (isWordBoundary(s, i) && !isMedicalTermMiddle(s, i)) {
        out.push({ pos: i, priority: 3 });
      }
    }

    // priorityでソート、同じpriorityなら位置順
    out.sort((a, b) => {
      if (typeof a === 'number') a = { pos: a, priority: 5 }; // 既存の数値は低優先度に
      if (typeof b === 'number') b = { pos: b, priority: 5 };
      return a.priority - b.priority || a.pos - b.pos;
    });

    // 位置のみを返す（重複除去）
    const positions = [...new Set(out.map(item => typeof item === 'number' ? item : item.pos))];
    return positions.sort((a, b) => a - b);
  };

  const shouldIgnoreParenSubtitle = (s, idx) => {
    const open = s[idx];
    const close = open === "（" ? "）" : ")";
    const endIdx = s.indexOf(close, idx + 1);
    if (endIdx === -1) return false;

    const inside = s.slice(idx + 1, endIdx).trim();
    const ignoreWords = new Set([
      "仮", "予定", "案", "再", "改", "新",
      "案1", "案2", "案3"
    ]);

    if (ignoreWords.has(inside)) return true;
    if (inside.length <= 2) return true;

    return false;
  };

  // =========================================================
  // subtitle split
  // =========================================================
  const splitBySubtitle = (s) => {
    s = unifyTilde(oneLine(s));
    if (!s) return null;

    // " - " / " – " / " — "
    let m = /\s[-–—―−－]\s/.exec(s);
    if (m && m.index > 0) {
      const p = m.index + 1;
      return [s.slice(0, p).trimEnd(), s.slice(p).trimStart()];
    }

    // 〜 / ~
    let idx = s.indexOf("〜");
    if (idx > 0) {
      return [s.slice(0, idx).trimEnd(), s.slice(idx).trimStart()];
    }

    idx = s.indexOf("~");
    if (idx > 0) {
      return [s.slice(0, idx).trimEnd(), s.slice(idx).trimStart()];
    }

    // full-width dash (スペースなしでもサブタイトル区切りとみなす)
    const fwDashSubRe = /[–—―−－]/g;
    let fwm;
    while ((fwm = fwDashSubRe.exec(s)) !== null) {
      const p = fwm.index;
      if (p <= 1) continue;
      const before = s.slice(0, p).trim();
      const after  = s.slice(p).trim();
      // after が閉じダッシュだけ(eg "―")でなく実質的内容を持つか
      const afterCore = after.replace(/^[–—―−－\s]+/, "").replace(/[–—―−－\s]+$/, "");
      if (before.length > 2 && afterCore.length > 2) {
        return [before, after];
      }
    }

    // dash fallback: 語中ハイフンは避ける
    const dashRe = /[-–—―−－]/g;
    while ((m = dashRe.exec(s)) !== null) {
      const p = m.index;
      const prev = s[p - 1] || "";
      const next = s[p + 1] || "";
      if (/\s/.test(prev) || /\s/.test(next)) {
        return [s.slice(0, p + 1).trimEnd(), s.slice(p + 1).trimStart()];
      }
    }

    // 括弧
    idx = s.indexOf("（");
    if (idx > 0 && !shouldIgnoreParenSubtitle(s, idx)) {
      return [s.slice(0, idx).trimEnd(), s.slice(idx).trimStart()];
    }

    idx = s.indexOf("(");
    if (idx > 0 && !shouldIgnoreParenSubtitle(s, idx)) {
      return [s.slice(0, idx).trimEnd(), s.slice(idx).trimStart()];
    }

    return null;
  };

  // =========================================================
  // line cleanup / penalties
  // =========================================================
  const mergeDanglingDotLines = (lines) => {
    const out = [];
    for (const line of lines) {
      const s = oneLine(line);
      if (!s) continue;

      if (out.length && /[・･]$/.test(out[out.length - 1])) {
        out[out.length - 1] = `${out[out.length - 1]}${s}`;
      } else {
        out.push(s);
      }
    }
    return out;
  };

  const mergeTinyParenTail = (lines) => {
    if (!Array.isArray(lines) || lines.length < 2) return lines || [];
    const out = [...lines];

    for (let i = out.length - 1; i >= 1; i--) {
      const cur = oneLine(out[i]);
      if (/^[（(][^)）]{1,3}[)）]$/.test(cur)) {
        out[i - 1] = `${oneLine(out[i - 1])}${cur}`;
        out.splice(i, 1);
      }
    }
    return out;
  };

  const mergeSymbolOnlyTail = (lines) => {
    if (!Array.isArray(lines) || lines.length < 2) return lines || [];
    const out = [...lines];

    for (let i = out.length - 1; i >= 1; i--) {
      const cur = oneLine(out[i]);
      if (!cur) {
        out.splice(i, 1);
        continue;
      }

      if (isOnlySymbols(cur)) {
        out[i - 1] = `${oneLine(out[i - 1])}${cur}`;
        out.splice(i, 1);
      }
    }
    return out;
  };

  const linePenalty = (line, idx, total, opts = {}) => {
    let p = 0;
    const t = oneLine(line);
    const wc = wordCount(t);

    if (!t) p += 9999;

    // 医学用語・専門用語の途中で改行されている場合は重いペナルティ
    const medicalTerms = [
      'GLP-1', 'SGLT-2', 'DPP-4', 'ACE阻害', 'β遮断', 'Ca拮抗',
      '糖尿病性', '心血管', '循環器', '虚血性', 'バイオマーカー',
      'エビデンス', 'ガイドライン', 'プロトコル', 'アルゴリズム',
      'webセミナー', 'Webセミナー', 'WEBセミナー'
    ];
    
    for (const term of medicalTerms) {
      if (t.includes(term.slice(0, -1)) && !t.includes(term)) {
        p += 5000; // 医学用語の途中改行は避ける
      }
    }

    // 最終行1語だけ（英語の場合）
    if (idx === total - 1 && wc === 1 && isMostlyAscii(t)) p += 1600;

    // 最終行が短すぎる（ただし意味のある単語なら許容）
    if (idx === total - 1 && t.length <= 8 && !medicalTerms.some(term => t.includes(term))) {
      p += 900;
    }

    // 最終行が記号だけはほぼ禁止
    if (idx === total - 1 && isOnlySymbols(t)) p += 100000;

    // 最終行が弱い記号スタートで極端に短いのも避ける
    if (idx === total - 1 && startsWithWeakSymbol(t) && t.length <= 2) p += 50000;

    // 先頭が記号っぽい（ただし医学用語の一部なら軽減）
    if (/^[-–—―−－:：,，/／]+/.test(t)) {
      const hasImportantContent = medicalTerms.some(term => t.includes(term));
      p += hasImportantContent ? 200 : 700;
    }

    // GLP-1 みたいな医学用語で終わる場合は許容
    if (/^[A-Za-z0-9]+-$/.test(t)) {
      const isMedicalAbbrev = /^(GLP|SGLT|DPP|ACE|ARB)-$/.test(t);
      p += isMedicalAbbrev ? 300 : 1800;
    }
    if (/[A-Za-z0-9]-$/.test(t) && !/^(GLP|SGLT|DPP|ACE|ARB)-/.test(t)) {
      p += 1200;
    }

    // 行末が中途半端（助詞で終わる）
    if (/[のにとへやではがも]$/.test(t)) p += 120;

    // カタカナ語の途中分断ペナルティ
    // 行末がカタカナ（ー含む）で終わり、次行頭もカタカナの場合
    if (idx < total - 1 && /[ァ-ヶー]$/.test(t)) p += 3000;

    // ・終わり
    if (/[・･]$/.test(t)) p += 800;

    // 意味のある医学用語を保持している行は優先
    if (medicalTerms.some(term => t.includes(term))) {
      p -= 100; // ボーナス
    }

    if (opts.preferBalancedAscii && isMostlyAscii(t) && wc === 1) {
      p += 250;
    }

    return p;
  };

  const scoreLines = (lines, maxPx, style, opts = {}) => {
    const widths = lines.map(line => measure(line, style));
    const slackScore = widths.reduce((acc, w) => acc + (maxPx - w) ** 2, 0);
    const penalties = lines.reduce((acc, line, i) => {
      return acc + linePenalty(line, i, lines.length, opts);
    }, 0);

    // より均等な行に対するボーナス（医学コンテンツでは読みやすさ重視）
    let balanceBonus = 0;
    if (widths.length >= 2) {
      const avg = widths.reduce((sum, w) => sum + w, 0) / widths.length;
      const variance = widths.reduce((sum, w) => sum + (w - avg) ** 2, 0) / widths.length;
      
      // 分散が小さい（均等）ほどボーナス
      if (variance < 2000) balanceBonus = -200;
      else if (variance < 5000) balanceBonus = -100;
    }

    let raggedPenalty = 0;
    if (widths.length >= 2) {
      const mx = Math.max(...widths);
      const mn = Math.min(...widths);
      raggedPenalty = (mx - mn) * 1.5; // ペナルティを少し軽減（内容重視）
    }

    // 医学用語が適切に配置されている場合のボーナス
    let medicalTermBonus = 0;
    const medicalTerms = ['GLP-1', 'SGLT-2', 'DPP-4', '糖尿病', '循環器', 'エビデンス'];
    for (const line of lines) {
      if (medicalTerms.some(term => line.includes(term))) {
        medicalTermBonus -= 50; // 医学用語を含む行にボーナス
      }
    }

    return slackScore + penalties + raggedPenalty + balanceBonus + medicalTermBonus;
  };

  // =========================================================
  // generic wrap
  // =========================================================
  const wrapPx = (s, maxPx, style, maxLines, opts = {}) => {
    if (String(s ?? "").length > 400) return [oneLine(s)];

    s = unifyTilde(oneLine(s));
    if (!s) return [];

    const {
      forceSubtitle2ndHead = false,
      preferBalancedAscii = false,
      avoidSingleWordLastLine = false,
      enableEarlyBreak = false,
    } = opts;

        // 不自然な語中分割（例: つ|いて）を抑制
        const isAwkwardPhraseSplit = (left, right) => {
            const l = String(left || "");
            const r = String(right || "");
            const badPairs = [
                ["につ", "いて"],
                ["にお", "ける"],
                ["に関", "する"],
                ["によ", "る"],
            ];
            return badPairs.some(([a, b]) => l.endsWith(a) && r.startsWith(b));
        };

        // 意味的に自然な分割位置に小さなボーナス
        const isPreferredPhraseBoundary = (left) => {
            const l = String(left || "");
            return (
                l.endsWith("について") ||
                l.endsWith("における") ||
                l.endsWith("に関する") ||
                l.endsWith("による")
            );
        };

    // subtitle rule
    if (forceSubtitle2ndHead) {
      const sp = splitBySubtitle(s);
      if (sp) {
        const a = stripEdgePunct(sp[0]);
        const b = stripEdgePunct(sp[1]);

        if (a && b) {
          if (!isOnlySymbols(b) && measure(a, style) <= maxPx && measure(b, style) <= maxPx) {
            return mergeTinyParenTail(mergeDanglingDotLines([a, b])).slice(0, maxLines);
          }

          const aLines = wrapPx(a, maxPx, style, maxLines, {
            ...opts,
            forceSubtitle2ndHead: false,
          });
          const remain = Math.max(1, maxLines - aLines.length);
          const bLines = wrapPx(b, maxPx, style, remain, {
            ...opts,
            forceSubtitle2ndHead: false,
          });

          return mergeTinyParenTail(
            mergeDanglingDotLines([...aLines, ...bLines].filter(Boolean))
          ).slice(0, maxLines);
        }
      }
    }

    // one line - 実際にmaxPxに収まるなら1行で返す
    const currentWidth = measure(s, style);
    
    if (currentWidth <= maxPx) return [s];

    const cand = breakPositions(s);
    try {
      console.log(`[WRAP] Break candidates: ${cand.length} positions found`);
    } catch (e) {
      console.log('[WRAP] Candidates debug failed:', e);
    }
    
    let best = null;
    let bestScore = null;

    // try 2 lines
    for (const p of cand) {
      const a = stripEdgePunct(s.slice(0, p).trim());
      const b = stripEdgePunct(s.slice(p).trim());
      if (!a || !b) continue;
      if (/[・･]$/.test(a)) continue;
      if (isOnlySymbols(b)) continue;

      const wa = measure(a, style);
      const wb = measure(b, style);
      
      if (wa > maxPx || wb > maxPx) continue;

      // 早期改行判定：カラム落ちリスクがある場合は優先的に採用
      const isEarlyBreakCandidate = enableEarlyBreak && shouldEarlyBreak(s, p, maxPx, style);

      if (avoidSingleWordLastLine && isMostlyAscii(b) && wordCount(b) === 1) continue;

      const lines = [a, b];
      let score = scoreLines(lines, maxPx, style, { preferBalancedAscii });

            // 語中分割は強くペナルティ
            if (isAwkwardPhraseSplit(a, b)) {
                score += 2200;
            }

            // 句として自然な境界はやや優遇
            if (isPreferredPhraseBoundary(a)) {
                score -= 180;
            }
      
      // カタカナ語の途中分断を回避（フォーラム→フォー|ラム 等）
      if (/[ァ-ヶー]$/.test(a) && /^[ァ-ヶー]/.test(b)) {
        score += 8000;
      }
      
      // 早期改行候補は大幅にスコアを改善（カラム落ち防止優先）
      if (isEarlyBreakCandidate) {
        score -= 1000;
      }
      
      // ダッシュ記号での改行は高スコア（自然なサブタイトル区切り）
      if (/[－−–—―]/.test(s.slice(Math.max(0, p-3), p+1))) {
        score -= 500;
      }

      console.log(`Break at ${p}: "${a}" | "${b}" | Score: ${score} | Widths: ${wa}/${wb}`);

      if (bestScore === null || score < bestScore) {
        bestScore = score;
        best = lines;
      }
    }
    if (best) return mergeTinyParenTail(mergeDanglingDotLines(best)).slice(0, maxLines);

    // try 3 lines
    if (maxLines >= 3) {
      for (const p of cand) {
        const a = stripEdgePunct(s.slice(0, p).trim());
        const rest = stripEdgePunct(s.slice(p).trim());
        if (!a || !rest) continue;
        if (/[・･]$/.test(a)) continue;
        if (measure(a, style) > maxPx) continue;

        const cand2 = breakPositions(rest);
        for (const q of cand2) {
          const b = stripEdgePunct(rest.slice(0, q).trim());
          const c = stripEdgePunct(rest.slice(q).trim());
          if (!b || !c) continue;
          if (/[・･]$/.test(b)) continue;
          if (isOnlySymbols(c)) continue;

          const wb = measure(b, style);
          const wc = measure(c, style);
          if (wb > maxPx || wc > maxPx) continue;

          if (avoidSingleWordLastLine && isMostlyAscii(c) && wordCount(c) === 1) continue;

                    const lines = [a, b, c];
                    let score = scoreLines(lines, maxPx, style, { preferBalancedAscii });

                    if (isAwkwardPhraseSplit(a, b) || isAwkwardPhraseSplit(b, c)) {
                        score += 2200;
                    }
                    if (isPreferredPhraseBoundary(a) || isPreferredPhraseBoundary(b)) {
                        score -= 180;
                    }

          if (bestScore === null || score < bestScore) {
            bestScore = score;
            best = lines;
          }
        }
      }
    }
    if (best) return mergeTinyParenTail(mergeDanglingDotLines(best)).slice(0, maxLines);

    // last resort
    const out = [];
    let cur = "";

    for (let i = 0; i < s.length; i++) {
      const ch = s[i];
      const nxt = cur + ch;

      const prevCh = s[i - 1] || "";
      const nextCh = s[i + 1] || "";

      const isHyphenInsideToken =
        ch === "-" &&
        /[A-Za-z0-9]/.test(prevCh) &&
        /[A-Za-z0-9]/.test(nextCh);

      if (!cur || measure(nxt, style) <= maxPx || isHyphenInsideToken) {
        cur = nxt;
      } else {
        out.push(cur.trim());
        cur = ch;
        if (out.length >= maxLines - 1) break;
      }
    }
    if (cur.trim()) out.push(cur.trim());

    return mergeTinyParenTail(mergeDanglingDotLines(out)).slice(0, maxLines);
  };

  // =========================================================
  // specialized wrappers
  // =========================================================
  // スマートな幅計算：カラム落ちを防ぐ安全マージン付き
  const getSmartMaxWidth = (maxPx, context = 'normal') => {
    // コンテキストに応じた安全マージン（改行時の各行の上限）
    const margins = {
      'hero': 0.96,    // イベントタイトル：4%マージン
            'talk': 0.95,    // 演題タイトル：5%マージン（カラム落ち防止優先）
      'normal': 0.97   // 通常：3%マージン
    };
    
    const margin = margins[context] || margins.normal;
    return Math.floor(maxPx * margin);
  };

  // 早期改行判定：長い行になりそうな場合の予防的改行
  const shouldEarlyBreak = (text, currentPos, maxPx, style) => {
    const remaining = text.slice(currentPos);
    const currentWidth = measure(text.slice(0, currentPos), style);
    
    // 現在の幅が70%を超えて、残りテキストが長い場合
    if (currentWidth > maxPx * 0.7 && remaining.length > 15) {
      return true;
    }
    
    // 残りに医学用語が含まれていて分割困難な場合
    const hasMedicalTerm = /GLP-1|SGLT-2|DPP-4|受容体作動薬|阻害薬|糖尿病性/.test(remaining.slice(0, 20));
    if (currentWidth > maxPx * 0.65 && hasMedicalTerm) {
      return true;
    }
    
    return false;
  };

  // =========================================================
  const wrapHeroTitle = (line, maxPx, style) => {
    console.log('[WRAP_HERO] Called with:', { line, maxPx });
    const s = unifyTilde(oneLine(line));
    if (!s) return [];

    const ascii = isMostlyAscii(s);
    const smartMaxPx = getSmartMaxWidth(maxPx, 'hero');

    // 医学セミナータイトルの特別な処理
    const isMedicalSeminar = /セミナー|webinar|symposium|conference/i.test(s);
    const hasMedicalTerms = /GLP-1|SGLT-2|DPP-4|糖尿病|循環器|腎臓|心血管/.test(s);

    const result = wrapPx(s, smartMaxPx, style, 8, {
      forceSubtitle2ndHead: true,
      preferBalancedAscii: ascii && !hasMedicalTerms, // 医学用語がある場合はバランスより内容重視
      avoidSingleWordLastLine: ascii && !isMedicalSeminar,
      enableEarlyBreak: true, // 早期改行を有効化
    });
    
    console.log('[WRAP_HERO] Result:', result);
    return result;
  };

  const wrapTalkTitle = (line, maxPx, style) => {
    try {
      console.log('wrapTalkTitle called with line:', JSON.stringify(line));
      const s = unifyTilde(oneLine(line));
      if (!s) return [];

      const smartMaxPx = getSmartMaxWidth(maxPx, 'talk');
      const hasMedicalTerms = /GLP-1|SGLT-2|DPP-4|糖尿病|循環器|腎臓|心血管|エビデンス|ガイドライン/.test(s);
      const hasLongMedicalPhrase = /受容体作動薬|阻害薬|合併症|バイオマーカー|プロトコル/.test(s);

      const result = wrapPx(s, smartMaxPx, style, 5, {
        forceSubtitle2ndHead: true,
        preferBalancedAscii: isMostlyAscii(s) && !hasMedicalTerms,
        avoidSingleWordLastLine: isMostlyAscii(s) && !hasLongMedicalPhrase,
        enableEarlyBreak: true, // 早期改行を有効化
      });
      
      console.log('wrapTalkTitle result:', JSON.stringify(result));
      return result;
    } catch (e) {
      console.log('ERROR in wrapTalkTitle:', e);
      return [line]; // フォールバック
    }
  };

  const normalizeAffiliation = (s) => {
    return String(s ?? "")
      .split("\n")
      .map(x => oneLine(x))
      .filter(Boolean)
      .join("\n")
      .replace(/\s*演\s*者\s*$/g, "")  // 末尾の「演 者」「演者」を除去
      .replace(/\s*座\s*長\s*$/g, "")  // 末尾の「座 長」「座長」も除去
      .trim();
  };

  // =========================================================
  // compute widths from DOM
  // =========================================================
  const wrapEl = document.querySelector(".wrap");
  const wrapW = wrapEl ? wrapEl.clientWidth : 600;
  // .sheet の padding を考慮して実コンテンツ幅を算出
  const sheetEl = document.querySelector(".sheet");
  let sheetInnerW = wrapW;
  if (sheetEl) {
    const cs = getComputedStyle(sheetEl);
    sheetInnerW = sheetEl.clientWidth - parseFloat(cs.paddingLeft) - parseFloat(cs.paddingRight);
  }
  // pill(74px) + gap(24px) = 98px → content area
  const contentMax = sheetInnerW - 98;
  const talkMax = contentMax;
  console.log('[TYPESET] wrapW:', wrapW, 'sheetInnerW:', sheetInnerW, 'contentMax:', contentMax);

  // =========================================================
  // styles
  // =========================================================
  const heroStyle = {
    fontFamily: "Invention JP",
    fontWeight: 700,
    fontSize: "30px",
    letterSpacing: "normal"
  };

  const talkStyle = {
    fontFamily: "Invention JP",
    fontWeight: 700,
    fontSize: "25px",
    letterSpacing: "normal"
  };

  const affStyle = {
    fontFamily: "Invention JP",
    fontWeight: 700,
    fontSize: "16px",
    letterSpacing: "normal"
  };

  const chairAffStyle = {
    fontFamily: "Invention JP",
    fontWeight: 700,
    fontSize: "16px",
    letterSpacing: "normal"
  };

  // 所属テキストのスマート改行
  const wrapAffiliation = (text, maxPx, style) => {
    const s = oneLine(text);
    if (!s) return "";
    // 安全マージンを適用（カラム落ち防止）
    const safeMaxPx = getSmartMaxWidth(maxPx, 'normal');
    const w = measure(s, style);
    if (w <= safeMaxPx) return s;

    // 所属テキスト用の改行候補を生成
    const lines = wrapPx(s, safeMaxPx, style, 3, {
      forceSubtitle2ndHead: false,
      preferBalancedAscii: false,
      avoidSingleWordLastLine: false,
      enableEarlyBreak: false,
    });
    return lines.join("\n");
  };

  // =========================================================
  // event title
  // =========================================================
  console.log('Processing event title started');
  
const evBaseLines =
  Array.isArray(data.event_title_lines) && data.event_title_lines.length
    ? data.event_title_lines.map(x => oneLine(x)).filter(Boolean)
    : [oneLine(data.event_title || "")].filter(Boolean);

console.log('evBaseLines:', JSON.stringify(evBaseLines));

// ★ 改行処理を適用してスマートな改行位置を決定
console.log('[TYPESET] Processing event title:', evBaseLines);
let evLines = [];
if (evBaseLines.length > 0) {
  const combinedEventTitle = evBaseLines.join(" ");
  console.log('[TYPESET] Combined event title:', combinedEventTitle);
  // heroタイトル用の最大幅を計算（wrapW - マージン）
  const heroMax = wrapW - 60; // 余裕を持ったマージン
  console.log('[TYPESET] Hero max width:', heroMax);
  // wrapHeroTitleを使って改行位置を最適化
  evLines = wrapHeroTitle(combinedEventTitle, heroMax, heroStyle);
  // 空の場合は元の形式にフォールバック
  if (!evLines.length) {
    console.log('[TYPESET] Event title fallback to original');
    evLines = Array.isArray(data.event_title_lines) && data.event_title_lines.length
      ? data.event_title_lines
      : evBaseLines;
  } else {
    console.log('[TYPESET] New event title lines:', evLines);
  }
} else {
  evLines = Array.isArray(data.event_title_lines) && data.event_title_lines.length
    ? data.event_title_lines
    : evBaseLines;
}

data.event_title_lines = evLines;
data.event_title = evLines.join("\n");

console.log('Event title processing completed:', JSON.stringify(evLines));

  // =========================================================
  // chair
  // =========================================================
  if (data.chair) {
    const rawChairAff = normalizeAffiliation(data.chair.affiliation ?? "");
    // 座長所属のスマート改行（pill 74px + gap 24px = 98px を引いた領域）
    data.chair.affiliation = wrapAffiliation(rawChairAff, contentMax, chairAffStyle);
  }

  // =========================================================
  // talks
  // =========================================================
if (Array.isArray(data.talks)) {
  console.log('Processing talks started');
  console.log('Talks length:', data.talks.length);
  
  // 講演タイトル・所属のピクセルベース改行
  const talkTitleMax = getSmartMaxWidth(talkMax, 'talk');
  const talkAffMax = contentMax; // 所属は pill(74) + gap(24) の右側

  data.talks = data.talks.map((t, index) => {
    if (t?.item_type === "chair") {
      const rawChairAff = normalizeAffiliation(String(t?.affiliation ?? ""));
      return {
        ...t,
        role: t?.role || "座長",
        affiliation: wrapAffiliation(rawChairAff, talkAffMax, affStyle),
      };
    }

    const rawTitleLines =
      Array.isArray(t?.title_lines) && t.title_lines.length
        ? t.title_lines.map(x => oneLine(x)).filter(Boolean)
        : [oneLine(t?.title ?? "")].filter(Boolean);

    let debug_info = {};

    // ★ ピクセルベースのスマート改行
    let title_lines = [];
    if (rawTitleLines.length > 0) {
      // タイトル行を結合（記号始まりの行は空白なしで結合）
      const combinedTitle = joinTitleLinesSmart(rawTitleLines);
      debug_info.original_title = combinedTitle;
      debug_info.title_length = combinedTitle.length;

      // wrapTalkTitle でピクセル実測ベースの改行
      title_lines = wrapTalkTitle(combinedTitle, talkMax, talkStyle);
      debug_info.split_method = "pixel_wrap";

      if (!title_lines.length) {
        title_lines = rawTitleLines;
        debug_info.split_method = "fallback";
      }
    } else {
      title_lines = t.title_lines?.length ? t.title_lines : rawTitleLines;
    }

    // ★ 所属テキストのピクセルベース改行
    const rawAff = normalizeAffiliation(String(t?.affiliation ?? ""));
    const affiliation = wrapAffiliation(rawAff, talkAffMax, affStyle);

    return {
      ...t,
      title_lines,
      title: title_lines.join("\n"),
      affiliation,
      _debug: debug_info
    };
  });
  
  console.log('Talks processing completed successfully');
}

  return data;
}
"""

# ---------------- Models ----------------

class DatetimeParts(BaseModel):
    year: str = ""
    month: str = ""
    day: str = ""
    dow: str = ""      # "月" "火" ...
    time: str = ""     # "19:00~20:20"

class TextOverride(BaseModel):
    # どれで対象を特定するか（どれか1つ使えればOK）
    target: Optional[str] = ""    # 対象テキスト（完全一致/部分一致に使う）
    index: Optional[int] = None   # title_lines の行番号

    # 変更したい見た目
    font_size: Optional[int] = None
    font_weight: Optional[int] = None
    color: Optional[str] = None


    
class Talk(BaseModel):
    item_type: str = "talk"
    program_index: int = 0
    role: str = "演者"
    name_display: str = ""
    time: str = ""
    title: str = "" 
    title_lines: List[str] = Field(default_factory=list)  # 改行保持 + ~...~ は別行
    speaker: str = ""
    speaker_display: str = ""
    affiliation: str = ""
    title_overrides: List[TextOverride] = Field(default_factory=list)
    honorific_title: str = "先生"


class Chair(BaseModel):
    role: str = "座長"
    name: str = ""
    name_display: str = ""
    affiliation: str = ""
    honorific_title: str = "先生"


class DesignJSON(BaseModel):
    event_title_lines: List[str] = Field(default_factory=list)  # 改行保持 + ~...~ は別行
    event_title: str = ""  # 互換/テンプレ移行用（event_title_linesを \n で結合）
    title_overrides: List[TextOverride] = Field(default_factory=list)
    datetime: str = ""
    datetime_parts: Optional[DatetimeParts] = None
    datetime_time_newline: bool = False  # datetime_parts.time を改行するか
    datetime_note: str = ""
    datetime_note_font_size: int = 14
    datetime_note_left: int = 5  # datetime_note の左位置
    organizer: str = ""
    chair: Chair = Chair()
    talks: List[Talk] = Field(default_factory=list)
    warnings: List[str] = Field(default_factory=list)
    confidence: float = 0.0
    manual_override: bool = False
    note: str = ""
    locked: bool = False
    title_font_size: int = 30
    region: str = ""  # 追加: 地域
    unit: str = ""    # 追加: 取得単位
    event_id: str = "" # 追加: イベントID


class RenderReq(BaseModel):
    jobId: str
    design: DesignJSON

class VmDiffByEventIdRequest(BaseModel):
    event_id: str

@dataclass
class TimeCand:
    text: str
    top: int
    left: int


# ---------------- Utilities ----------------
BREAK_CHARS = ["／", "/", "・", " ", "　", "～", "~", "-", "－", "—", "–", "（", "(", "）", ")", "、", "。", ":", "："]
DROP_AT_BREAK = set(["－", "—", "–"])

EMU_PER_PT = 12700

TIME_RE = re.compile(r"(\d{1,2}\s*[:\：]\s*\d{2})\s*[~〜～\-–—−－]\s*(\d{1,2}\s*[:\：]\s*\d{2})")

def _merge_blocks_to_rows(blocks: list, top_tolerance: int = 200000):
    """同一行（近い top）のブロックを left 順で結合し (top, left, merged_text) を返す。
    時間ブロックが個別に分割されている場合（"19:00" "～" "19:30"）に結合する。"""
    if not blocks:
        return []
    sorted_blocks = sorted(blocks, key=lambda b: (b.top, b.left))
    rows = []
    for b in sorted_blocks:
        if rows and abs(b.top - rows[-1][0].top) < top_tolerance:
            rows[-1].append(b)
        else:
            rows.append([b])
    out = []
    _hhmm_re = re.compile(r"\d{1,2}[:：]\d{2}")
    for row in rows:
        row.sort(key=lambda b: b.left)
        merged = " ".join((b.text or "").replace("\n", " ") for b in row)
        merged = normalize_space(merged)
        if merged:
            out.append((row[0].top, row[0].left, merged))
        # フォールバック: 通常結合で TIME_RE マッチしない場合、
        # 各ブロックの先頭行から HH:MM を抽出して時間値順で結合
        if merged and not TIME_RE.search(_norm_time(merged)):
            times = []
            for b in row:
                first_line = (b.text or "").split("\n")[0].strip()
                first_line = first_line.translate(str.maketrans("０１２３４５６７８９：", "0123456789:"))
                m = _hhmm_re.search(first_line)
                if m:
                    times.append(m.group().replace("：", ":"))
            if len(times) >= 2:
                times.sort(key=lambda t: int(t.replace(":", "")))
                time_merged = f"{times[0]}～{times[1]}"
                out.append((row[0].top, row[0].left, time_merged))
    return out

def _norm_time(s: str) -> str:
    """時間表記の正規化（全角・半角・記号統一）"""
    s = str(s or "").replace("\u3000", " ")
    s = re.sub(r"\s+", " ", s).strip()
    
    # 全角数字・コロンを半角に変換
    s = s.translate(str.maketrans("０１２３４５６７８９：", "0123456789:"))
    
    # ダッシュ類の統一
    s = re.sub(r"[–—−－\-]", "~", s)
    s = s.replace("～", "~").replace("〜", "~")
    
    return s

def extract_time_cands_with_pos(blocks: list[TextBlock]) -> list[TimeCand]:
    """ブロックから時間範囲を抽出（位置情報付き）"""
    _event_dt_re = re.compile(r"20\d{2}\s*年|\d{1,2}\s*月\s*\d{1,2}\s*日|日時")
    out = []
    seen_tops = set()
    for b in blocks:
        # 日時ラベル・年月日を含むブロックはイベント日時なので除外
        if _event_dt_re.search(b.text or ""):
            continue
        t = _norm_time(b.text)
        
        # 時間範囲のパターンを探す
        m = TIME_RE.search(t)
        if not m:
            continue
            
        start_time = m.group(1).replace(" ", "")
        end_time = m.group(2).replace(" ", "")
        
        # HH:MM形式に正規化
        start_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", start_time)
        end_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", end_time)
        
        time_text = f"{start_norm}~{end_norm}"
        out.append(TimeCand(text=time_text, top=b.top, left=b.left))
        seen_tops.add(b.top)

    # 個別ブロックで取れない場合、結合行からも抽出（"19:00" "～" "19:30" バラバラ対応）
    for top, left, merged in _merge_blocks_to_rows(blocks):
        if top in seen_tops:
            continue
        t = _norm_time(merged)
        m = TIME_RE.search(t)
        if not m:
            continue
        start_time = m.group(1).replace(" ", "")
        end_time = m.group(2).replace(" ", "")
        start_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", start_time)
        end_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", end_time)
        time_text = f"{start_norm}~{end_norm}"
        out.append(TimeCand(text=time_text, top=top, left=left))
    
    out.sort(key=lambda x: (x.top, x.left))
    return out

def parse_event_time_range(dt_text: str) -> tuple[str, str] | None:
    # "2026年4月18日（土）14:00～16:00" みたいなのから取る
    t = _norm_time(dt_text)
    m = TIME_RE.search(t.replace("-", "~"))
    if not m:
        return None
    return (m.group(1), m.group(2))

def time_to_minutes(hhmm: str) -> int | None:
    m = re.match(r"^(\d{1,2}):(\d{2})$", hhmm)
    if not m:
        return None
    h = int(m.group(1)); mi = int(m.group(2))
    return h * 60 + mi

def is_within_event(tc: TimeCand, ev: tuple[str,str] | None) -> bool:
    if not ev:
        return True
    s, e = ev
    s0 = time_to_minutes(s); e0 = time_to_minutes(e)
    m = TIME_RE.search(tc.text)
    if not m or s0 is None or e0 is None:
        return True
    a = time_to_minutes(m.group(1)); b = time_to_minutes(m.group(2))
    if a is None or b is None:
        return True
    # 多少の誤差許容
    return (s0 - 10) <= a and b <= (e0 + 10)

def z2h_digits(s: str) -> str:
    return str(s).translate(str.maketrans({
        "０": "0", "１": "1", "２": "2", "３": "3", "４": "4",
        "５": "5", "６": "6", "７": "7", "８": "8", "９": "9",
    }))

CIRCLED_NUM_MAP = {
    "①": 1, "②": 2, "③": 3, "④": 4, "⑤": 5,
    "⑥": 6, "⑦": 7, "⑧": 8, "⑨": 9, "⑩": 10,
}

ROMAN_JP_MAP = {
    "Ⅰ": 1, "II": 2, "Ⅱ": 2, "III": 3, "Ⅲ": 3,
    "IV": 4, "Ⅳ": 4, "V": 5, "Ⅴ": 5,
}
def normalize_time_text(s: str) -> str:
    s = str(s or "")
    s = s.replace("：", ":")
    s = re.sub(r"[～〜\-ー]", "~", s)
    s = re.sub(r"\s*~\s*", "~", s)
    return s.strip()

def extract_talk_number_and_time_from_text(text: str) -> tuple[int | None, str]:
    raw = normalize_space(text or "").replace("\n", " ")
    s = normalize_key(raw)  # ← 空白を潰した比較用

    talk_no = None

    # 1) 講演1 / 講演１ / 講演① / 講 演 1 / 演題1 / 演題①
    #    rawで先にマッチ（スペースが自然な境界になり時間の数字を食わない）
    m = re.search(r"(?:講演|演題)\s*([0-9０-９]+|[①②③④⑤⑥⑦⑧⑨⑩])", raw)
    if not m:
        # フォールバック: normalize_key版（"講 演 1" のようにスペース入りラベル対応）
        m = re.search(r"(?:講演|演題)([0-9０-９]+|[①②③④⑤⑥⑦⑧⑨⑩])", s)
    if m:
        raw_no = m.group(1)
        if raw_no in CIRCLED_NUM_MAP:
            talk_no = CIRCLED_NUM_MAP[raw_no]
        else:
            talk_no = int(z2h_digits(raw_no))
        # 妥当性チェック: 演題番号が大きすぎる場合は時間数字を食った可能性
        if talk_no is not None and talk_no > 20:
            talk_no = None

    # 2) 一般講演Ⅰ / 一般 講 演 Ⅱ / III
    if talk_no is None:
        m = re.search(r"一般講演([ⅠⅡⅢⅣⅤ]+|I{1,3}|IV|V)", s)
        if m:
            raw_no = m.group(1)
            talk_no = ROMAN_JP_MAP.get(raw_no)

    # 3) 特別講演
    if talk_no is None and "特別講演" in s:
        talk_no = 3

    # time
    t0 = normalize_time_text(raw)
    m = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－]\s*(\d{1,2}:\d{2})", t0)
    time_text = f"{m.group(1)}~{m.group(2)}" if m else ""

    return talk_no, time_text

def extract_talk_times_in_order(blocks: list[TextBlock]) -> list[str]:
    ordered = sorted(blocks, key=lambda b: (b.top, b.left))
    out: list[str] = []
    seen = set()

    for b in ordered:
        txt = normalize_space(getattr(b, "text", "") or "")
        one = normalize_time_text(txt.replace("\n", " "))

        m = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－]\s*(\d{1,2}:\d{2})", one)
        if not m:
            continue

        tm = f"{m.group(1)}~{m.group(2)}"

        if tm not in seen:
            out.append(tm)
            seen.add(tm)

    # 個別ブロックで取れない場合、結合行からも抽出
    if not out:
        for top, left, merged in _merge_blocks_to_rows(blocks):
            one = normalize_time_text(merged)
            m = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－]\s*(\d{1,2}:\d{2})", one)
            if not m:
                continue
            tm = f"{m.group(1)}~{m.group(2)}"
            if tm not in seen:
                out.append(tm)
                seen.add(tm)

    return out

def assign_talk_times_by_order_fallback(
    payload: DesignJSON,
    blocks: list[TextBlock],
) -> DesignJSON:
    if not getattr(payload, "talks", None):
        return payload

    time_list = extract_talk_times_in_order(blocks)
    if not time_list:
        return payload

    full_dt = normalize_time_text(normalize_space(getattr(payload, "datetime", "") or ""))
    m_full = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－]\s*(\d{1,2}:\d{2})", full_dt)
    full_time = f"{m_full.group(1)}~{m_full.group(2)}" if m_full else ""

    # 全体時間を除外
    time_list = [tm for tm in time_list if tm != full_time]

    if not time_list:
        return payload

    targets = [t for t in payload.talks if _is_program_talk_item(t) and not getattr(t, "time", "")]
    if not targets:
        return payload

    for t, tm in zip(targets, time_list):
        t.time = tm

    return payload

def extract_talk_time_map_by_anchor(blocks: list[TextBlock]) -> dict[int, str]:
    ordered = sorted(blocks, key=lambda b: (b.top, b.left))
    # 結合行も含めた時間検索用リスト
    merged_rows = _merge_blocks_to_rows(blocks)
    out: dict[int, str] = {}

    for i, b in enumerate(ordered):
        talk_no, time_text = extract_talk_number_and_time_from_text(b.text)
        if talk_no is not None and time_text:
            out[talk_no] = time_text
            continue

        talk_no, _ = extract_talk_number_and_time_from_text(b.text)
        if talk_no is None:
            continue

        best_time = ""
        best_dist = None
        # 個別ブロックから探す
        for j in range(max(0, i - 6), min(len(ordered), i + 7)):
            if j == i:
                continue

            cand = normalize_space(ordered[j].text or "")
            cand = normalize_time_text(cand)

            tm = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－]\s*(\d{1,2}:\d{2})", cand)
            if not tm:
                continue

            dist = abs(ordered[j].top - b.top) + abs(ordered[j].left - b.left) * 0.15
            if best_dist is None or dist < best_dist:
                best_dist = dist
                best_time = f"{tm.group(1)}~{tm.group(2)}"

        # 個別ブロックで見つからない場合、結合行から探す
        if not best_time:
            for top, left, merged in merged_rows:
                cand = normalize_time_text(merged)
                tm = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－]\s*(\d{1,2}:\d{2})", cand)
                if not tm:
                    continue
                dist = abs(top - b.top) + abs(left - b.left) * 0.15
                if dist > 1500000:
                    continue
                if best_dist is None or dist < best_dist:
                    best_dist = dist
                    best_time = f"{tm.group(1)}~{tm.group(2)}"

        if best_time:
            out[talk_no] = best_time

    return out
    

def assign_talk_times_by_anchor(blocks: list[TextBlock], payload: DesignJSON) -> DesignJSON:
    talks = list(payload.talks or [])
    if not talks:
        return payload

    talk_time_map = extract_talk_time_map_by_anchor(blocks)
    if not talk_time_map:
        return payload

    if any(_is_program_chair_item(t) for t in talks):
        talk_only = [t for t in talks if _is_program_talk_item(t)]
        for idx, t in enumerate(talk_only, start=1):
            tm = talk_time_map.get(idx)
            if tm:
                t.time = tm
                setattr(t, "_talk_index", idx)
        payload.talks = talks
        return payload

    # --- talks を blocks 内のスピーカー位置に基づいて正しい順序に並び替え ---
    # 各 talk のスピーカーが blocks 内のどの位置にいるかを特定
    if len(talks) >= 2 and len(talk_time_map) >= 2:
        def _speaker_top_in_blocks(t) -> int:
            """talk のスピーカー名が blocks 内で最初に出現する top 位置"""
            sp = (getattr(t, "speaker", "") or "").replace(" ", "").replace("\u3000", "")
            if not sp:
                return 10**9
            for b in sorted(blocks, key=lambda b: (b.top, b.left)):
                bt = (b.text or "").replace(" ", "").replace("\u3000", "").replace("\n", "")
                if sp in bt:
                    return b.top
            return 10**9

        speaker_tops = [(i, _speaker_top_in_blocks(t)) for i, t in enumerate(talks)]
        # 全員の位置が取れた場合のみ並び替え
        if all(top < 10**9 for _, top in speaker_tops):
            sorted_indices = [i for i, _ in sorted(speaker_tops, key=lambda x: x[1])]
            # 現在の順番と異なる場合のみ並び替え
            if sorted_indices != list(range(len(talks))):
                talks = [talks[i] for i in sorted_indices]
                print(f"[assign-talk-times] reordered talks by blocks position: {sorted_indices}")

    # アンカーで特定できた講演のみ時間を設定（他はリセットしない）
    for idx, t in enumerate(talks, start=1):
        tm = talk_time_map.get(idx)
        if tm:
            t.time = tm
            setattr(t, "_talk_index", idx)

    payload.talks = talks
    return payload


def assign_talk_times_by_proximity(blocks: list[TextBlock], payload: DesignJSON) -> DesignJSON:
    """講演と時間の近接性による割り当て（既存の時間は保持）"""
    # 1) イベント全体の時間枠
    ev = parse_event_time_range(getattr(payload, "datetime", "") or "")

    # 2) time候補を抽出＆イベント範囲でフィルタ
    cands = [c for c in extract_time_cands_with_pos(blocks) if is_within_event(c, ev)]
    if not cands or not getattr(payload, "talks", None):
        return payload

    # 3) 講演のアンカー位置を精密に特定
    def find_precise_anchor_top(talk) -> int:
        search_terms = []
        
        # 講演者名（複数パターン）
        if getattr(talk, "speaker", ""):
            speaker = talk.speaker.strip()
            search_terms.append(speaker)
            search_terms.append(speaker.replace(" ", ""))  # スペース無し
            if " " in speaker:
                search_terms.append(" ".join(speaker.split()))  # 正規化
        
        if getattr(talk, "speaker_display", ""):
            search_terms.append(talk.speaker_display.strip())
            search_terms.append(talk.speaker_display.replace(" ", ""))
        
        # タイトル（部分マッチも考慮）
        if getattr(talk, "title_lines", []):
            for line in talk.title_lines:
                if line.strip():
                    search_terms.append(line.strip())
        elif getattr(talk, "title", ""):
            title_lines = talk.title.split('\n')
            for line in title_lines:
                if line.strip():
                    search_terms.append(line.strip())
        
        best_top = 10**18
        best_score = 0
        
        for b in blocks:
            bt = normalize_space(b.text or "")
            if not bt:
                continue
                
            score = 0
            # より精密なマッチング
            for term in search_terms:
                term_norm = normalize_space(term).replace(" ", "")
                bt_norm = bt.replace(" ", "").replace("\n", "")
                
                if term_norm and term_norm in bt_norm:
                    # 完全一致は高得点
                    if term_norm == bt_norm:
                        score += 10
                    # 部分一致
                    elif len(term_norm) >= 3:
                        score += 5
                    else:
                        score += 2
            
            if score > best_score:
                best_score = score
                best_top = b.top
        
        return best_top

    talks = list(payload.talks or [])
    talk_infos = []
    for idx, t in enumerate(talks):
        if _is_program_chair_item(t):
            continue
        talk_infos.append((idx, find_precise_anchor_top(t)))
    
    # 上→下に並べる
    talk_infos.sort(key=lambda x: x[1])

    used = set()

    for idx, anchor_top in talk_infos:
        t = talks[idx]
        # 既に時間が設定されている場合はスキップ
        if _norm_time(getattr(t, "time", "")):
            continue

        # アンカーに最も近い未使用の時間候補を探す
        best = None
        best_dist = None

        for ci, c in enumerate(cands):
            if ci in used:
                continue
                
            # 距離計算（縦方向優先、横方向も少し考慮）
            vert_dist = abs(anchor_top - c.top)
            horiz_dist = abs(c.left - 500000)  # 左側を少し優遇
            
            # 時間は講演の少し上にあることが多い
            if c.top <= anchor_top:
                vert_dist *= 0.7  # 上側を優遇
            
            total_dist = vert_dist + horiz_dist * 0.3
            
            if best_dist is None or total_dist < best_dist:
                best_dist = total_dist
                best = (ci, c)

        if best and best_dist is not None and best_dist < 2000000:  # 距離制限
            ci, c = best
            used.add(ci)
            t.time = c.text

    payload.talks = talks
    return payload


def extract_blocks_from_pdf(pdf_path: Path, first_page_only: bool = True) -> List[TextBlock]:
    doc = fitz.open(str(pdf_path))
    blocks: List[TextBlock] = []

    pages = [doc[0]] if (first_page_only and doc.page_count > 0) else [doc[i] for i in range(doc.page_count)]

    for page in pages:
        d = page.get_text("dict")

        for b in d.get("blocks", []):
            if b.get("type") != 0:  # 0=text
                continue

            x0, y0, x1, y1 = b.get("bbox", (0, 0, 0, 0))

            # 段落/行を組み立て（PDFのline/spansを尊重）
            lines = []
            max_font_pt = 0.0

            for line in b.get("lines", []):
                spans = line.get("spans", [])
                # spanをそのまま連結（余計な空白は後でnormalize）
                t = "".join(s.get("text", "") for s in spans)
                if t and t.strip():
                    lines.append(t.strip())

                for s in spans:
                    try:
                        max_font_pt = max(max_font_pt, float(s.get("size") or 0.0))
                    except Exception:
                        pass

            text = normalize_keep_newlines("\n".join(lines))
            if not text:
                continue

            # PDF(pt) -> EMU に変換
            left_emu   = int(round(x0 * EMU_PER_PT))
            top_emu    = int(round(y0 * EMU_PER_PT))
            width_emu  = int(round((x1 - x0) * EMU_PER_PT))
            height_emu = int(round((y1 - y0) * EMU_PER_PT))

            blocks.append(
                TextBlock(
                    text=text,
                    left=left_emu,
                    top=top_emu,
                    width=width_emu,
                    height=height_emu,
                    max_font_pt=float(max_font_pt or 0.0),
                )
            )

    doc.close()
    blocks.sort(key=lambda b: (b.top, b.left))
    return blocks


def extract_blocks_from_pdf2(file_path: str) -> list[dict[str, Any]]:
    """
    PDFから line 単位でテキストブロックを抽出する。
    さらに各 line の spans も返す。
    """
    doc = fitz.open(file_path)
    out: list[dict[str, Any]] = []

    try:
        for page_index, page in enumerate(doc):
            page_no = page_index + 1
            page_width = float(page.rect.width)
            page_height = float(page.rect.height)

            text_dict = page.get_text("dict")

            for block in text_dict.get("blocks", []):
                if block.get("type") != 0:
                    continue

                for line in block.get("lines", []):
                    raw_spans = line.get("spans", [])
                    if not raw_spans:
                        continue

                    line_spans: list[dict[str, Any]] = []
                    parts: list[str] = []
                    max_font_pt = 0.0

                    x0s = []
                    y0s = []
                    x1s = []
                    y1s = []

                    for span in raw_spans:
                        text = str(span.get("text", "") or "")
                        if not text:
                            continue

                        bbox = span.get("bbox")
                        if not bbox or len(bbox) != 4:
                            continue

                        x0, y0, x1, y1 = map(float, bbox)
                        size = float(span.get("size", 0) or 0)

                        parts.append(text)
                        max_font_pt = max(max_font_pt, size)

                        x0s.append(x0)
                        y0s.append(y0)
                        x1s.append(x1)
                        y1s.append(y1)

                        line_spans.append({
                            "text": text,
                            "left": x0,
                            "top": y0,
                            "width": max(0.0, x1 - x0),
                            "height": max(0.0, y1 - y0),
                            "font_size": size,
                        })

                    line_text = "".join(parts).strip()
                    if not line_text or not line_spans:
                        continue

                    left = min(x0s)
                    top = min(y0s)
                    right = max(x1s)
                    bottom = max(y1s)

                    width = max(0.0, right - left)
                    height = max(0.0, bottom - top)

                    if width <= 0 or height <= 0:
                        continue

                    out.append({
                        "text": line_text,
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "max_font_pt": max_font_pt,
                        "page": page_no,
                        "_page_width": page_width,
                        "_page_height": page_height,
                        "_coord_unit": "pdf_page",
                        "spans": line_spans,
                    })

        return out

    finally:
        doc.close()


def blocks_to_dicts(blocks: list[Any]) -> list[dict]:
    out: list[dict] = []

    for b in blocks or []:
        if isinstance(b, dict):
            out.append({
                "text": b.get("text", "") or "",
                "left": b.get("left", 0) or 0,
                "top": b.get("top", 0) or 0,
                "width": b.get("width", 0) or 0,
                "height": b.get("height", 0) or 0,
                "max_font_pt": b.get("max_font_pt", 0) or 0,
                "page": b.get("page", 1) or 1,
                "page_width": b.get("_page_width", 0) or 0,
                "page_height": b.get("_page_height", 0) or 0,
                "coord_unit": b.get("_coord_unit", "") or "",
                "spans": b.get("spans", []) or [],
            })
        else:
            out.append({
                "text": getattr(b, "text", "") or "",
                "left": getattr(b, "left", 0) or 0,
                "top": getattr(b, "top", 0) or 0,
                "width": getattr(b, "width", 0) or 0,
                "height": getattr(b, "height", 0) or 0,
                "max_font_pt": getattr(b, "max_font_pt", 0) or 0,
                "page": getattr(b, "page", 1) or 1,
                "page_width": getattr(b, "_page_width", 0) or 0,
                "page_height": getattr(b, "_page_height", 0) or 0,
                "coord_unit": getattr(b, "_coord_unit", "") or "",
                "spans": getattr(b, "spans", []) or [],
            })

    return out

def merge_event_title_blocks_strict(blocks: list[TextBlock]) -> list[TextBlock]:
    # 上部の大フォントだけ抽出（学習済みmedianをベースに閾値を決定）
    _lp = _get_layout_pattern_cache()
    _etfp = _lp.get("event_title_font_pt", {})
    if _etfp.get("count", 0) >= 3:
        # イベントタイトルとスピーカーの中間を閾値に（ノイズ除去）
        _sp_median = _lp.get("speaker_font_pt", {}).get("median", 0)
        _et_median = _etfp.get("median", 0)
        if _et_median > 0 and _sp_median > 0:
            _threshold = (_et_median + _sp_median) / 2.0
        elif _et_median > 0:
            _threshold = _et_median * 0.75  # イベントタイトルmedianの75%
        else:
            _threshold = 22.0
        _threshold = max(16.0, min(_threshold, 32.0))  # 16〜32ptにクランプ
    else:
        _threshold = 22.0
    candidates = [b for b in blocks if b.max_font_pt >= _threshold]
    if not candidates:
        return blocks

    candidates.sort(key=lambda b: b.top)

    # 上から連続しているものだけ取る
    merged_group = [candidates[0]]

    for b in candidates[1:]:
        # 前のブロックと縦距離が近ければ同グループ
        if abs(b.top - merged_group[-1].top) < 500000:
            merged_group.append(b)
        else:
            break  # 離れたら終了

    merged_text = "\n".join(b.text for b in merged_group)

    merged_block = TextBlock(
        text=merged_text,
        left=min(b.left for b in merged_group),
        top=min(b.top for b in merged_group),
        width=max(b.width for b in merged_group),
        height=sum(b.height for b in merged_group),
        max_font_pt=max(b.max_font_pt for b in merged_group),
    )

    # 元ブロックから除去
    new_blocks = [b for b in blocks if b not in merged_group]
    new_blocks.append(merged_block)
    new_blocks.sort(key=lambda b: (b.top, b.left))

    return new_blocks

def merge_pdfish_blocks(blocks):
    blocks = sorted(blocks, key=lambda b: (b["top"], b["left"]))
    merged = []

    for b in blocks:
        if not merged:
            merged.append(b)
            continue

        prev = merged[-1]

        # 縦距離が近くて左が近いなら同じ行扱い
        if abs(b["top"] - prev["top"]) < 6 and abs(b["left"] - prev["left"]) < 10:
            prev["text"] += b["text"]
            prev["width"] = max(prev["width"], b["width"])
        else:
            merged.append(b)

    return merged

def extract_blocks_any(path: Path, first_only: bool = True) -> List[TextBlock]:
    suf = path.suffix.lower()
    if suf == ".pdf":
        blocks = extract_blocks_from_pdf(path, first_page_only=first_only)
        # PDFは分断が多いので前処理（後で強化できる）
        # blocks = merge_pdfish_blocks(blocks)
        return blocks
    # .ppt は python-pptx では基本ダメなのでここでは弾くか、事前変換前提
    return extract_blocks_from_pptx(path, first_slide_only=first_only)

def pget(obj, key, default=None):
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)

def pset(obj, key, value):
    if isinstance(obj, dict):
        obj[key] = value
        return
    setattr(obj, key, value)

def extract_session_times(s: str) -> list[str]:
    s0 = normalize_datetime_text(s)  # ここで全角コロン等も正規化される :contentReference[oaicite:8]{index=8}
    out = []
    for m in SESSION_TIME_RE.finditer(s0):
        t = normalize_datetime_text(m.group(1)).replace(" ", "")
        if t not in out:
            out.append(t)
    return out

def split_tilde_subtitle(s: str) -> list[str]:
    """
    末尾の ～xxx～ / ~xxx~ を別行に分離（あれば）
    """
    s = normalize_space(s)
    if not s:
        return []
    m = re.search(r"(.*?)(\s*[~～].+[~～]\s*)$", s)
    if m:
        a = normalize_space(m.group(1))
        b = normalize_space(m.group(2))
        return [x for x in [a, b] if x]
    return [s]

def wrap_by_chars(s: str, max_len: int, *, back: int = 8) -> list[str]:
    """
    文字数近似で自然改行（直近の区切り文字を優先）
    """
    s = normalize_space(s)
    if not s:
        return []
    if len(s) <= max_len:
        return [s]

    out = []
    rest = s
    while rest and len(rest) > max_len:
        cut = -1
        start = max(0, max_len - back)
        for i in range(min(max_len, len(rest) - 1), start - 1, -1):
            if rest[i] in BREAK_CHARS:
                cut = i + 1
                break
        if cut == -1:
            cut = max_len

        head = rest[:cut].strip()
        tail = rest[cut:].strip()

        # ハイフンなど “落としたい区切り” を行末/行頭に残さない
        if head and head[-1] in DROP_AT_BREAK:
            head = head[:-1].rstrip()
        if tail and tail[0] in DROP_AT_BREAK:
            tail = tail[1:].lstrip()

        if head:
            out.append(head)
        rest = tail

    if rest:
        out.append(rest)

    return [x for x in out if x]

def join_short_suffix(lines: list[str]) -> list[str]:
    """
    “るために” みたいな短い尻尾行ができたら、前行に戻して繋げる（軽い補正）
    """
    if not lines:
        return lines
    out = [lines[0]]
    for l in lines[1:]:
        if len(l) <= 4 and out:
            out[-1] = (out[-1] + l).strip()
        else:
            out.append(l)
    return out


MEASURE_JS = """
({ text, font }) => {
  const c = document.createElement('canvas');
  const ctx = c.getContext('2d');
  ctx.font = font;
  return ctx.measureText(text).width;
}
"""

async def measure_px(page, text: str, font_css: str) -> float:
    return float(await page.evaluate(MEASURE_JS, {"text": text, "font": font_css}))

def split_tilde_head_2nd(s: str):
    s = (s or "").replace("～", "〜").strip()
    i = s.find("〜")
    if i <= 0:
        return None
    a = s[:i].rstrip()
    b = s[i:].lstrip()  # 2行目は必ず「〜」から
    if not a or not b:
        return None
    return a, b

BREAK_CHARS_PX = [" ", "　", "、", "。", "・", "／", "/", ":", "：", "）", ")", "】", "]"]

def candidate_breaks(s: str) -> list[int]:
    pos = []
    for i, ch in enumerate(s):
        if ch in BREAK_CHARS_PX:
            pos.append(i + 1)  # その直後で折る
    return sorted(set(pos))

async def wrap_px(page, text: str, max_px: int, font_css: str, max_lines: int = 3, force_tilde: bool = False) -> list[str]:
    s = (text or "").replace("\n", " ").strip()
    if not s:
        return []

    # 〜強制（talk用）
    if force_tilde:
        sp = split_tilde_head_2nd(s)
        if sp:
            a, b = sp
            if await measure_px(page, a, font_css) <= max_px and await measure_px(page, b, font_css) <= max_px:
                return [a, b]

    # 1行で入るなら1行
    if await measure_px(page, s, font_css) <= max_px:
        return [s]

    # 2行以上：候補位置で分割して、収まりつつ “余白が少ない” ものを選ぶ
    breaks = candidate_breaks(s)
    best = None
    best_score = None

    # まずは2行を狙う（ダメなら後段で3行）
    for p in breaks:
        a = s[:p].strip()
        b = s[p:].strip()
        if not a or not b:
            continue
        wa = await measure_px(page, a, font_css)
        wb = await measure_px(page, b, font_css)
        if wa <= max_px and wb <= max_px:
            score = (max_px - wa) ** 2 + (max_px - wb) ** 2
            if best_score is None or score < best_score:
                best_score = score
                best = [a, b]
    if best:
        return best

    # 3行まで許す：2回折る（粗いが強い）
    if max_lines >= 3:
        for p in breaks:
            a = s[:p].strip()
            rest = s[p:].strip()
            if not a or not rest:
                continue
            if await measure_px(page, a, font_css) > max_px:
                continue
            # rest を2行にする
            breaks2 = candidate_breaks(rest)
            for q in breaks2:
                b = rest[:q].strip()
                c = rest[q:].strip()
                if not b or not c:
                    continue
                wb = await measure_px(page, b, font_css)
                wc = await measure_px(page, c, font_css)
                if wb <= max_px and wc <= max_px:
                    return [a, b, c]

    # 最後の保険：強制分割（絶対はみ出さない）
    out = []
    cur = ""
    for ch in s:
        nxt = cur + ch
        if not cur or await measure_px(page, nxt, font_css) <= max_px:
            cur = nxt
        else:
            out.append(cur.strip())
            cur = ch
            if len(out) >= max_lines - 1:
                break
    if cur.strip():
        out.append(cur.strip())
    return out[:max_lines]

async def apply_precise_typeset_initial(payload: DesignJSON, page=None) -> DesignJSON:
    # ---- 初期値のみ：編集・ロック・上書き指定がある場合は何もしない ----
    if getattr(payload, "manual_override", False):
        return payload
    if getattr(payload, "locked", False):
        return payload
    if (getattr(payload, "title_overrides", None) or []):
        return payload
    for t in (getattr(payload, "talks", None) or []):
        ov = (t.get("title_overrides") if isinstance(t, dict) else getattr(t, "title_overrides", None)) or []
        if ov:
            return payload

    # ---- payload -> dict ----
    data_json = payload.model_dump_json() if hasattr(payload, "model_dump_json") else payload.json(ensure_ascii=False)
    data_obj = json.loads(data_json)

    async def _run(pg):
        global _cached_template
        if _cached_template is None:
            _cached_template = TEMPLATE_PATH.read_text(encoding="utf-8")

        await pg.set_content(_cached_template, wait_until="domcontentloaded")
        await pg.evaluate("() => document.fonts && document.fonts.ready")

        # TYPESET_JS: data_obj を px実測で event_title_lines/title_lines/affiliation に整形して返す
        print("🔧 [PYTHON] About to execute TYPESET_JS...")
        print("🔧 [PYTHON] data_obj keys:", list(data_obj.keys()) if isinstance(data_obj, dict) else "Not a dict")
        try:
            result = await pg.evaluate(TYPESET_JS, {"data": data_obj})
            print("🔧 [PYTHON] TYPESET_JS execution completed successfully")  
            print("🔧 [PYTHON] Result type:", type(result))
            
            # talks の改行結果を確認
            if isinstance(result, dict) and 'talks' in result:
                print("🔧 [PYTHON] Talks processing results:")
                for i, talk in enumerate(result['talks']):
                    if 'title_lines' in talk:
                        print(f"  Talk {i}: {talk['title_lines']}")
                        # デバッグ情報も表示
                        if '_debug' in talk:
                            debug = talk['_debug']
                            print(f"    Debug: title_length={debug.get('title_length')}, "
                                  f"has_dash={debug.get('has_dash')}, "
                                  f"split_method={debug.get('split_method')}")
                            if 'colon_split_result' in debug:
                                print(f"    Colon split: {debug['colon_split_result']}")
                            if 'dash_split_result' in debug:
                                print(f"    Dash split: {debug['dash_split_result']}")
                            if 'second_line_length' in debug:
                                print(f"    Second line length: {debug['second_line_length']}")
            
            return result
        except Exception as e:
            print("🔧 [PYTHON] ERROR: TYPESET_JS execution failed:", str(e))
            print("🔧 [PYTHON] Error type:", type(e))
            raise

    # ---- page が無ければグローバル _browser から一時ページを作成 ----
    if page is None:
        global _browser
        if _browser is not None:
            # グローバルブラウザから一時ページを作成（2重ブラウザ起動によるOOMを回避）
            ctx = await _browser.new_context(viewport=BASE_VIEWPORT)
            pg = await ctx.new_page()
            try:
                new_obj = await _run(pg)
            finally:
                await ctx.close()
        else:
            # フォールバック: グローバルブラウザが未初期化の場合のみ新規起動
            async with async_playwright() as p:
                browser = await p.chromium.launch()
                pg = await browser.new_page(viewport=BASE_VIEWPORT)
                try:
                    new_obj = await _run(pg)
                finally:
                    await browser.close()
    else:
        new_obj = await _run(page)

    # ---- dict -> payload に戻す ----
    if hasattr(payload.__class__, "model_validate"):
        payload = payload.__class__.model_validate(new_obj)
    else:
        payload = payload.__class__.parse_obj(new_obj)

    # ---- template の参照差でズレないように同期（重要）----
    if getattr(payload, "event_title_lines", None):
        payload.event_title = "\n".join(payload.event_title_lines)

    for t in (getattr(payload, "talks", None) or []):
        if getattr(t, "title_lines", None):
            t.title_lines = fix_title_lines_jp(t.title_lines)
            t.title = "\n".join(t.title_lines)

    # payload.typeset_done = True
    return payload



def format_title_initial(
    raw: str,
    *,
    max_len: int,
    max_lines: int = 3,
    force_tilde_second_line: bool = False,
) -> list[str]:
    """
    talk.title / event_title の初期整形
    """
    raw = normalize_space(raw)
    if not raw:
        return []

    # 表記ゆれ統一
    raw = raw.replace("～", "〜")

    # ★talk用：最初の「〜」を必ず2行目先頭へ
    # 例) 「…治療〜GLP-1…〜」→
    #   1行目「…治療」
    #   2行目「〜GLP-1…〜」
    if force_tilde_second_line and "〜" in raw:
        idx = raw.find("〜")
        if idx > 0:
            first = raw[:idx].rstrip()
            second = raw[idx:].lstrip()  # 2行目は必ず「〜」から
            lines = []
            lines.extend(wrap_by_chars(first, max_len=max_len))
            lines.extend(wrap_by_chars(second, max_len=max_len))
            lines = [x for x in lines if x]
            lines = join_short_suffix(lines)
            if max_lines and len(lines) > max_lines:
                lines = lines[: max_lines - 1] + [" ".join(lines[max_lines - 1 :]).strip()]
            return lines

    # 既存ロジック：末尾の ~xxx~ / ～xxx～ は別行に分離
    lines: list[str] = []
    for part in split_tilde_subtitle(raw):
        lines.extend(wrap_by_chars(part, max_len=max_len))

    lines = [x for x in lines if x]
    lines = join_short_suffix(lines)

    if max_lines and len(lines) > max_lines:
        lines = lines[: max_lines - 1] + [" ".join(lines[max_lines - 1 :]).strip()]
    return lines





def format_affiliation_initial(raw: str, *, max_len: int, max_lines: int = 2) -> str:
    """
    affiliation 初期整形（最大2行）
    """
    raw = normalize_space(raw)
    if not raw:
        return ""
    lines = wrap_by_chars(raw, max_len=max_len)
    return "\n".join(lines[:max_lines]).strip()

def post_format_design_initial(payload):
    # 学習済み1行文字数を取得（learning-based max_len）
    _llen = _get_title_line_len_cache()
    _talk_max_len = int(max(14, min(_llen.get("talk_title_p90") or 18, 26)))
    _event_max_len = int(max(16, min(_llen.get("event_title_p90") or 22, 30)))

    # event title
    if getattr(payload, "event_title_lines", None):
        payload.event_title_lines = [
            normalize_space(x) for x in payload.event_title_lines if normalize_space(x)
        ]
    else:
        base = normalize_space(getattr(payload, "event_title", "") or "")
        if base:
            if len(base) <= 24:
                payload.event_title_lines = [base]
            else:
                payload.event_title_lines = format_title_initial(
                    base,
                    max_len=_event_max_len,
                    max_lines=20,
                    force_tilde_second_line=False,
                )

    payload.event_title = "\n".join(payload.event_title_lines).strip()

    # talks
    for t in getattr(payload, "talks", []) or []:
        if getattr(t, "title_lines", None):
            t.title_lines = [normalize_space(x) for x in t.title_lines if normalize_space(x)]
        else:
            raw_title = normalize_space(getattr(t, "title", "") or "")
            if raw_title:
                t.title_lines = format_title_initial(
                    raw_title,
                    max_len=_talk_max_len,
                    max_lines=20,
                    force_tilde_second_line=True,
                )

        if getattr(t, "title_lines", None):
            t.title = "\n".join(t.title_lines)

        # affiliation だけ軽く整形
        t.affiliation = format_affiliation_initial(t.affiliation or "", max_len=28, max_lines=20)

    return payload

def _get_field(payload, key: str, default=None):
    # dict
    if isinstance(payload, dict):
        return payload.get(key, default)
    # pydantic model
    return getattr(payload, key, default)

def _set_field(payload, key: str, value):
    if isinstance(payload, dict):
        payload[key] = value
    else:
        setattr(payload, key, value)

TIME_RANGE_IN_TEXT_RE = re.compile(
    r"(\d{1,2}[:：]\d{2}\s*[～〜\-ー−－—–~]\s*\d{1,2}[:：]\d{2})"
)

SESSION_LABEL_RE = re.compile(
    r"[①②③④⑤⑥⑦⑧⑨⑩]"
    r"|(?:[0-9０-９]+|[一二三四五六七八九十]+)\s*回目"
    r"|第\s*(?:[0-9０-９]+|[一二三四五六七八九十]+)\s*回"
)

def has_session_label(s: str) -> bool:
    return bool(SESSION_LABEL_RE.search(str(s or "")))

def normalize_time_range(s: str) -> str:
    s = str(s or "")
    s = s.replace("：", ":")

    # 変な連続区切りを 1 個の ~ に寄せる
    s = re.sub(r"\s*[-－ー‐-‒–—―〜～~]+\s*", "~", s)

    # 19:30-~20:15 みたいな中途半端な並びも潰す
    s = re.sub(r"~+", "~", s)

    return s.strip()

def extract_session_times_from_blocks(blocks) -> list[str]:
    if not blocks:
        return []

    ordered = sorted(
        blocks,
        key=lambda b: (
            b.get("top", 0) if isinstance(b, dict) else getattr(b, "top", 0),
            b.get("left", 0) if isinstance(b, dict) else getattr(b, "left", 0),
        )
    )

    out = []

    for i, b in enumerate(ordered):
        t = (b.get("text") if isinstance(b, dict) else getattr(b, "text", "")) or ""

        # 1) 同一block内にセッションラベル + time
        if has_session_label(t):
            for m in TIME_RANGE_IN_TEXT_RE.finditer(t):
                tt = normalize_time_range(m.group(1))
                if tt and tt not in out:
                    out.append(tt)
            continue

        # 2) セッションラベルが前後blockにあるケース
        near_texts = [t]
        if i > 0:
            prev_t = (ordered[i-1].get("text") if isinstance(ordered[i-1], dict) else getattr(ordered[i-1], "text", "")) or ""
            near_texts.append(prev_t)
        if i + 1 < len(ordered):
            next_t = (ordered[i+1].get("text") if isinstance(ordered[i+1], dict) else getattr(ordered[i+1], "text", "")) or ""
            near_texts.append(next_t)

        if not any(has_session_label(x) for x in near_texts):
            continue

        for m in TIME_RANGE_IN_TEXT_RE.finditer(t):
            tt = normalize_time_range(m.group(1))
            if tt and tt not in out:
                out.append(tt)

    return out



def extract_session_times_from_datetime(dt: str) -> list[str]:
    dt = str(dt or "")
    if not has_session_label(dt):
        return []

    out = []
    for m in TIME_RANGE_IN_TEXT_RE.finditer(dt):
        t = normalize_time_range(m.group(1))
        if t and t not in out:
            out.append(t)
    return out

def should_hide_talk_times(payload, blocks=None) -> bool:
    dt = str(_get_field(payload, "datetime", "") or "")
    parts = _get_field(payload, "datetime_parts", None)

    time_str = ""
    if parts:
        time_str = parts.get("time", "") if isinstance(parts, dict) else getattr(parts, "time", "")

    # セッション形式なら非表示
    if has_session_label(dt):
        return True

    if has_session_label(time_str):
        return True

    # イベント時間が blocks に1箇所しかない → 日時ブロック由来のみ → talk には不要
    if time_str and blocks:
        _time_ns = normalize_time_range(time_str).replace(":", "").replace("~", "")
        if _time_ns:
            _count = sum(
                1 for b in blocks
                if _time_ns in (b.text or "").replace(":", "").replace("~", "").replace("～", "").replace("〜", "")
            )
            if _count <= 1:
                # ただし、イベント全体時間とは異なる個別演題時間がブロックに存在する場合は非表示にしない
                full_time_norm = normalize_time_range(time_str)
                all_times = extract_talk_times_in_order(blocks)
                individual_times = [t for t in all_times if normalize_time_range(t) != full_time_norm]
                if individual_times:
                    return False
                return True

    return False

def clear_talk_times(payload):
    talks = _get_field(payload, "talks", []) or []
    for t in talks:
        if isinstance(t, dict):
            t["time"] = ""
        else:
            setattr(t, "time", "")
    _set_field(payload, "talks", talks)
    return payload

def _ensure_datetime_parts(parts):
    if parts is None:
        return DatetimeParts(year="", month="", day="", dow="", time="")
    if isinstance(parts, DatetimeParts):
        return parts
    if isinstance(parts, dict):
        return DatetimeParts(**parts)
    return DatetimeParts(
        year=getattr(parts, "year", "") or "",
        month=getattr(parts, "month", "") or "",
        day=getattr(parts, "day", "") or "",
        dow=getattr(parts, "dow", "") or "",
        time=getattr(parts, "time", "") or "",
    )


def extract_dow_from_blocks(blocks) -> str:
    ordered = sorted(blocks, key=lambda b: (b.top, b.left))
    # 1) 同一ブロック内に「年月日(曜)」があるケース
    for b in ordered:
        s = normalize_datetime_text(b.text or "")
        m = re.search(r"\d{4}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日\s*[（(]\s*([月火水木金土日])\s*[）)]", s)
        if m:
            return m.group(1)

    # 2) 曜日が独立ブロック "(金)" or "(水) 17:00~" or 単独 "水" のケース: 日付ブロック近傍を探す
    date_block = None
    for b in ordered:
        s = normalize_datetime_text(b.text or "")
        if re.search(r"\d{4}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日", s):
            date_block = b
            break
    if date_block:
        for b in ordered:
            s = normalize_datetime_text(b.text or "").strip()
            # "(水)" or "(水) 17:00~" パターン
            m = re.match(r"[（(]\s*([月火水木金土日])\s*[）)]", s)
            if not m:
                # 単独 "水" のような括弧なしブロック
                m = re.fullmatch(r"([月火水木金土日])", s)
            if m:
                # 日付ブロックとの距離が近い（縦方向±500000 emu）
                if abs(b.top - date_block.top) <= 500000:
                    return m.group(1)

    return ""
  

def fill_datetime_parts(payload, blocks=None):
    def pget(obj, key, default=None):
        if isinstance(obj, dict):
            return obj.get(key, default)
        return getattr(obj, key, default)

    def pset(obj, key, value):
        if isinstance(obj, dict):
            obj[key] = value
        else:
            setattr(obj, key, value)

    dt = str(pget(payload, "datetime", "") or "")

    # 全体時間
    full_dt = normalize_datetime_text(dt)
    # マルチセッション（N回目）の場合は全体時間を取らない
    if has_session_label(full_dt):
        full_time = ""
    else:
        m_full = TIME_RANGE_RE.search(full_dt)
        full_time = normalize_time_range(m_full.group(0)) if m_full else ""

    session_times = extract_session_times_from_blocks(blocks)
    if not session_times:
        session_times = extract_session_times_from_datetime(dt)

    # 全体時間は除外
    if full_time:
        session_times = [t for t in session_times if normalize_time_range(t) != full_time]

    # 重複除去
    uniq = []
    seen = set()
    for t in session_times:
        nt = normalize_time_range(t)
        if nt and nt not in seen:
            uniq.append(nt)
            seen.add(nt)
    session_times = uniq

    m = re.search(
        r"(?P<y>\d{4})\s*年\s*(?P<mo>\d{1,2})\s*月\s*(?P<d>\d{1,2})\s*日"
        r"(?:\s*[（(]\s*(?P<dow>[^）)\s]+)\s*[）)])?",
        normalize_datetime_text(dt)
    )

    parts = _ensure_datetime_parts(pget(payload, "datetime_parts", None))

    if m:
        y, mo, d, dow = (m.group("y") or "", m.group("mo") or "", m.group("d") or "", m.group("dow") or "")
        parts.year, parts.month, parts.day = y, mo, d
        parts.dow = dow or extract_dow_from_blocks(blocks)
    else:
        dow = extract_dow_from_blocks(blocks)
        if dow:
            parts.dow = dow

    if session_times:
        # イベント全体の時間（full_time）が優先
        # 個別セッション時間は talks の time フィールドで管理
        if full_time:
            time_joined = full_time
            newline = False
        elif len(session_times) == 1:
            time_joined = session_times[0]
            newline = False
        else:
            # 全体時間が不明で複数セッションがある場合のみ「〇回目」表記
            time_joined = ", ".join([f"{i+1}回目{t}" for i, t in enumerate(session_times)])
            newline = True
    else:
        time_joined = full_time or ""
        newline = False

    # ★ 正解DBで True 比率が70%超なら初期値を True に（案内状の多数派フォーマット優先）
    if not newline:
        _dnc = _get_datetime_newline_cache()
        if _dnc.get("true_ratio", 0) >= 0.70:
            newline = True

    parts.time = time_joined

    pset(payload, "datetime_parts", parts)
    pset(payload, "datetime_time_newline", newline)

    # 不要な talk time を消す（セッション複数回 or 時間blocks1個くだり=日時ブロック由来）
    if should_hide_talk_times(payload, blocks):
        clear_talk_times(payload)

    return payload

def _norm1(s: str) -> str:
    return normalize_space(s).replace("\n", " ").strip()

_TIME_RE = re.compile(
    r"""^\s*
    (?:\d{1,2}\s*[:：]\s*\d{2})      # 19:05
    \s*(?:[〜～~\-－–—]\s*)          # ～ / 〜 / ~ / - 系
    (?:\d{1,2}\s*[:：]\s*\d{2})      # 19:35
    \s*$""",
    re.VERBOSE
)

def looks_like_affil_line(s: str) -> bool:
    s = _norm1(s)
    if not s:
        return False

    # 正解DBに登録済みの所属・施設なら即 True
    _s_ns = re.sub(r'[\s\u3000]+', '', s)
    if _s_ns and len(_s_ns) >= 3 and _s_ns in _get_facility_name_dict_cache():
        return True

    # ラベル/案内っぽいのは除外
    if any(k in s for k in ["演題", "演者", "座長", "日時", "会場", "共催", "主催", "提供", "企画", "運営", "詳細は"]):
        return False
    
    if _TIME_RE.match(s):
        return False

    # 敬称入りは名前行の可能性
    if "先生" in s:
        return False

    # ★演題っぽい記号がある行は基本除外（所属に入ることは稀）
    #   - ダッシュ/波線/カギカッコ があると演題率が高い
    if any(p in s for p in ["「", "」", "～", "—", "－", "–"]):
        return False
    # ハイフンは施設名にも混ざり得るので「両側に空白がない長文」の時だけ除外
    if "-" in s and len(s) >= 18:
        return False

    # ★所属として強い接頭辞（法人格など）
    entity_prefix = ("医療法人", "一般財団法人", "一般社団法人", "公益財団法人", "公益社団法人",
                     "独立行政法人", "国立病院機構", "学校法人")
    if s.startswith(entity_prefix):
        return True

    # ★施設キーワード（これが無いなら所属扱いしない）
    facility_kw = ["病院", "クリニック", "医院", "診療所", "大学", "機構", "センター", "医師会", "総合病院"]
    has_facility = any(k in s for k in facility_kw)

    # 役職・部署っぽい語（施設名と一緒に出ると所属確度UP）
    role_kw = ["内科", "外科", "科", "部", "室", "教授", "准教授", "講師", "主任", "部長", "院長", "理事長"]
    has_role = any(k in s for k in role_kw)

    # ★施設名があるなら所属。施設名がなく役職だけはNG。
    if has_facility:
        return True

    # 施設名が無い場合でも「県立中央」など施設っぽい固有パターンを少し救う（任意）
    if ("県立" in s or "市立" in s) and has_role:
        return True

    return False



def wrap_vm_rows_for_rank(vm_rows: list[dict]) -> list[dict]:
    out = []
    for r in (vm_rows or []):
        # すでに {"data": ...} 形式ならそのまま
        if isinstance(r, dict) and "data" in r and isinstance(r["data"], dict):
            out.append(r)
            continue

        # rowdict -> {"data": rowdict} に包む
        if isinstance(r, dict):
            out.append({"data": r})
            continue

    return out

def _key(v: Any) -> str:
    return str(v or "").strip()

def _col_to_a1(col_idx_1based: int) -> str:
    s = ""
    n = col_idx_1based
    while n:
        n, r = divmod(n - 1, 26)
        s = chr(65 + r) + s
    return s

def _find_col(headers: List[str], name: str) -> int:
    for i, h in enumerate(headers):
        if h == name:
            return i + 1
    raise KeyError(f"column not found: {name}")

def _rowdict(headers: List[str], row: List[Any], *, sheet: str, rownum: int) -> Dict[str, Any]:
    if len(row) < len(headers):
        row = list(row) + [""] * (len(headers) - len(row))
    row = row[:len(headers)]
    d = {headers[i]: row[i] for i in range(len(headers))}
    d["_sheet"] = sheet
    d["_row"] = rownum
    return d

def _extract_1col(vr: dict) -> List[Any]:
    """
    values_batch_get の valueRanges から「1列データ」を取り出す（空行/欠損に強い）
    返り値は [v0, v1, ...]
    """
    vals = vr.get("values") or []
    if not vals:
        return []

    # ケース1: [["a","b","c"]] (1行に横並びで入ってくる)
    if len(vals) == 1 and isinstance(vals[0], list):
        return vals[0]

    # ケース2: [["a"], ["b"], [], ["c"]] (縦で返る＋空行が混ざる)
    if isinstance(vals[0], list):
        out: List[Any] = []
        for r in vals:
            if not r:          # [] をスキップ（または "" を入れるのでもOK）
                out.append("") # ← 行番号ズレ防止のため空を入れるのがおすすめ
                continue
            out.append(r[0])
        return out

    # ケース3: まれに ["a","b"] のような形（そのまま）
    return vals

def _pad_list(xs: List[Any], need_len: int) -> List[Any]:
    if len(xs) < need_len:
        xs = list(xs) + [""] * (need_len - len(xs))
    return xs

def pick_last_rownum(rownums: List[int]) -> int:
    return max(rownums) if rownums else 0

def _build_id_index_from_column(values: List[Any], *, start_row: int) -> Dict[str, List[int]]:
    idx: Dict[str, List[int]] = {}
    for offset, v in enumerate(values):
        k = str(v or "").strip()
        if not k:
            continue
        rownum = start_row + offset
        idx.setdefault(k, []).append(rownum)
    return idx


def _retry_gspread(fn, *, tries=6, base=0.5, jitter=0.2):
    last = None
    for i in range(tries):
        try:
            return fn()
        except APIError as e:
            last = e
            # INTERNAL(500) とかのときだけリトライ
            msg = str(getattr(e, "response", "")) + " " + str(e)
            if "INTERNAL" not in msg and "'code': 500" not in msg:
                raise
            sleep = base * (2 ** i) + random.uniform(0, jitter)
            time.sleep(sleep)
    raise last

def batch_fetch_system_and_vm_rows(
    workbook,
    *,
    ws_map,  # 
    event_ids: List[str],
    presence_sheets: List[str],
    presence_header_row: int,
    presence_id_col: str,
    vm_sheet: str,
    vm_header_row: int,
    vm_id_col_candidates: List[str],
    col_end: str = "N",
    chunk_size: int = 200,
) -> Tuple[Dict[str, List[Dict]], Dict[str, List[Dict]], str]:

    # ---- headers: 1回ずつ ----
    presence_headers_by_sheet: Dict[str, List[str]] = {}
    presence_id_col_letter_by_sheet: Dict[str, str] = {}

    presence_col_end_by_sheet: Dict[str, str] = {}

    for s in presence_sheets:
        ws = ws_map[s]
        headers = make_unique(ws.row_values(presence_header_row))
        presence_headers_by_sheet[s] = headers

        id_col_idx = _find_col(headers, presence_id_col)
        presence_id_col_letter_by_sheet[s] = _col_to_a1(id_col_idx)

        # ★行取得で「後ろ列が空」にならないよう、ヘッダの最終列を col_end にする
        presence_col_end_by_sheet[s] = _col_to_a1(len(headers))

    ws_vm = ws_map[vm_sheet]
    vm_headers = make_unique(ws_vm.row_values(vm_header_row))

    vm_id_col_used = ""
    for c in vm_id_col_candidates:
        if c in vm_headers:
            vm_id_col_used = c
            break
    if not vm_id_col_used:
        raise KeyError(f"VM id column not found. candidates={vm_id_col_candidates}")

    vm_id_col_letter = _col_to_a1(_find_col(vm_headers, vm_id_col_used))

    # ---- batchGet: ID列だけ（まとめて）----
    ranges = []
    # presence
    for s in presence_sheets:
        col_letter = presence_id_col_letter_by_sheet[s]
        start = presence_header_row + 1
        # 「列の下全部」(末尾空はAPIが返さないことがある)
        ranges.append(f"{s}!{col_letter}{start}:{col_letter}")
    # vm
    start_vm = vm_header_row + 1
    ranges.append(f"{vm_sheet}!{vm_id_col_letter}{start_vm}:{vm_id_col_letter}")

    print(f"[INFO] batch fetching ID columns: {ranges}")
    resp = _retry_gspread(lambda: workbook.values_batch_get(ranges))

    # ---- index化（event_id -> rownums）----
    presence_index_by_sheet: Dict[str, Dict[str, List[int]]] = {}
    for i, s in enumerate(presence_sheets):
        ws = ws_map[s]
        start = presence_header_row + 1
        col_values = _extract_1col(resp["valueRanges"][i])

        # ★row_countで膨らませない。返ってきた範囲内だけで十分
        # ただし「途中に空行がある」ケースは _extract_1col が "" を入れてくれるのでズレにくい
        presence_index_by_sheet[s] = _build_id_index_from_column(col_values, start_row=start)

    # VM index
    start = vm_header_row + 1
    end_row = ws_vm.row_count

    vm_col_values = _extract_1col(resp["valueRanges"][-1])
    need_len = end_row - start + 1
    vm_col_values = _pad_list(vm_col_values, need_len)

    vm_index = _build_id_index_from_column(vm_col_values, start_row=start)

    uniq_event_ids = [str(e or "").strip() for e in event_ids if str(e or "").strip()]
    presence_rows_by_event: Dict[str, List[Dict]] = {eid: [] for eid in uniq_event_ids}
    vm_rows_by_event: Dict[str, List[Dict]] = {eid: [] for eid in uniq_event_ids}

    # ---- 必要行rangeを組み立てて batchGet（行だけ）----
    presence_row_ranges = []
    presence_row_meta = []  # (eid, sheet, rownum)
    vm_row_ranges = []
    vm_row_meta = []        # (eid, rownum)

    for eid in uniq_event_ids:
        for s in presence_sheets:
            rownums = presence_index_by_sheet.get(s, {}).get(eid, [])
            last_row = pick_last_rownum(rownums)
            if last_row:
                end_letter = presence_col_end_by_sheet[s]  # ★ここ
                presence_row_ranges.append(f"{s}!A{last_row}:{end_letter}{last_row}")
                presence_row_meta.append((eid, s, last_row))
    
    for eid in uniq_event_ids:
        # ... presence ...
        for rownum in vm_index.get(eid, []):
            vm_row_ranges.append(f"{vm_sheet}!A{rownum}:{col_end}{rownum}")
            vm_row_meta.append((eid, rownum))

    def _chunks(xs, n):
        for i in range(0, len(xs), n):
            yield xs[i:i+n]

    # presence 行取得
    for rchunk, mchunk in zip(_chunks(presence_row_ranges, chunk_size), _chunks(presence_row_meta, chunk_size)):
        if not rchunk:
            continue
        rresp = _retry_gspread(lambda: workbook.values_batch_get(rchunk))
        for vr, meta in zip(rresp["valueRanges"], mchunk):
            rowvals = vr.get("values") or [[]]
            row = rowvals[0] if rowvals else []
            eid, sheet, rownum = meta
            headers = presence_headers_by_sheet[sheet]
            presence_rows_by_event[eid].append(_rowdict(headers, row, sheet=sheet, rownum=rownum))

    # VM 行取得
    for rchunk, mchunk in zip(_chunks(vm_row_ranges, chunk_size), _chunks(vm_row_meta, chunk_size)):
        if not rchunk:
            continue
        rresp = _retry_gspread(lambda: workbook.values_batch_get(rchunk))
        for vr, meta in zip(rresp["valueRanges"], mchunk):
            rowvals = vr.get("values") or [[]]
            row = rowvals[0] if rowvals else []
            eid, rownum = meta
            vm_rows_by_event[eid].append(_rowdict(vm_headers, row, sheet=vm_sheet, rownum=rownum))
    
    for eid in uniq_event_ids:
        presence_rows = presence_rows_by_event.get(eid, [])
        vm_rows = vm_rows_by_event.get(eid, [])

        # Presenceからsheet一覧を取得
        sheets = {
            normalize_space(r.get("_sheet", ""))
            for r in presence_rows
            if isinstance(r, dict) and r.get("_sheet")
        }

        # vm_rowsに付与
        for r in vm_rows:
            if isinstance(r, dict):
                r["_presence_sheets"] = list(sheets)

    return presence_rows_by_event, vm_rows_by_event, vm_id_col_used

def make_unique(headers):
    counter = Counter()
    out = []
    for h in headers:
        key = (h or "").strip()
        if key == "":
            key = "EMPTY"
        counter[key] += 1
        out.append(key if counter[key] == 1 else f"{key}_{counter[key]}_{counter[key]}")
    return out

def pick_last_presence_hit(hits: list[dict]) -> dict | None:
    if not hits:
        return None
    return max(hits, key=lambda h: int(h.get("row") or 0))

def normalize_day(v) -> str:
    """
    'YYYY-MM-DD' に正規化
    """
    if v is None:
        return ""

    if isinstance(v, datetime):
        return v.date().isoformat()
    if isinstance(v, date):
        return v.isoformat()

    s = str(v).strip()
    if not s:
        return ""

    m = re.search(r"(\d{4})[\/\-.](\d{1,2})[\/\-.](\d{1,2})", s)
    if not m:
        return ""

    try:
        return date(int(m[1]), int(m[2]), int(m[3])).isoformat()
    except Exception:
        return ""
    
def pick_best_presence_row(rows: list[dict]) -> dict | None:
    if not rows: return None
    def filled_count(d): return sum(1 for k,v in d.items() if not k.startswith("_") and str(v).strip())
    return max(rows, key=filled_count)

def fetch_row_dict(workbook, sheet_name: str, row_num: int, *, header_row: int = 2, col_end: str = "Z") -> dict:
    """
    指定シートの指定行(row_num)を、ヘッダ(header_row)の列名でdict化して返す。
    """
    ws = workbook.worksheet(sheet_name)

    # ヘッダ
    headers = make_unique(ws.row_values(header_row))
    # 1行分（A{row}:Z{row}）
    row_values = ws.get(f"A{row_num}:{col_end}{row_num}")
    row = (row_values[0] if row_values else [])

    # 列数合わせ
    if len(row) < len(headers):
        row = row + [""] * (len(headers) - len(row))
    elif len(row) > len(headers):
        row = row[:len(headers)]

    data = {headers[i]: (row[i] if i < len(row) else "") for i in range(len(headers))}

    # 便利情報として行番号も入れる（不要なら消してOK）
    data["_sheet"] = sheet_name
    data["_row"] = row_num
    return data

def build_system_id_to_rows(workbook, sheet_name: str, header_row: int = 2, id_col_name: str = "講演会ID"):
    ws = workbook.worksheet(sheet_name)

    headers = make_unique(ws.row_values(header_row))
    if id_col_name not in headers:
        raise RuntimeError(f"'{id_col_name}' not found in '{sheet_name}' header_row={header_row}")

    id_col = headers.index(id_col_name) + 1
    id_values = ws.col_values(id_col)  # 列だけ（軽め）

    index: dict[str, list[int]] = {}
    for row_num, v in enumerate(id_values, start=1):
        if row_num <= header_row:
            continue
        key = str(v).strip()
        if not key:
            continue
        index.setdefault(key, []).append(row_num)  # 上から順に溜まる

    return index

def fetch_rows_for_system_id_fast(
    workbook,
    sheet_name: str,
    index: dict[str, list[int]],
    system_id: str,
    header_row: int = 2,
    col_end: str = "Z",
):
    ws = workbook.worksheet(sheet_name)
    headers = make_unique(ws.row_values(header_row))

    rows = index.get(str(system_id).strip(), [])
    if not rows:
        return []

    # 該当範囲を一括取得（例: A10:Z80）
    min_row, max_row = min(rows), max(rows)
    values = ws.get(f"A{min_row}:{col_end}{max_row}")  # まとめて取る

    out = []
    for row_num in rows:
        rel = row_num - min_row  # values内のindex
        row_values = values[rel] if 0 <= rel < len(values) else []

        if len(row_values) < len(headers):
            row_values += [""] * (len(headers) - len(row_values))

        out.append({"sheet": sheet_name, "row": row_num, "data": dict(zip(headers, row_values))})

    return out





def _norm(s: str) -> str:
    return " ".join((s or "").replace("　"," ").split()).strip().lower()

def sim(a: str, b: str) -> float:
    a = _norm(a); b = _norm(b)
    if not a or not b: return 0.0
    return SequenceMatcher(None, a, b).ratio()

def rank_vm_candidates(pptx_title: str, pptx_speaker: str, vm_rows: list[dict], k: int = 5):
    scored = []
    for r in vm_rows:
        d = r["data"]
        s_title = d.get("演題","")
        s_name  = d.get("案内状掲載 医師名","")
        # タイトル重視＋名前少し
        sc = 0.75*sim(pptx_title, s_title) + 0.25*sim(pptx_speaker, s_name)
        scored.append((sc, r))
    scored.sort(key=lambda x: x[0], reverse=True)
    return scored[:k]


def apply_vm_correction_no_ai(payload, vm_rows: list[dict], *, hi=0.90, gap=0.06, k=5):
    """
    AIなし：スコアが高い時だけ talk.title を補正。
    演者はシート抜けもあるので基本 keep（触らない）。
    """
    if not vm_rows or not getattr(payload, "talks", None):
        return payload

    warnings = getattr(payload, "warnings", None) or []
    payload.warnings = warnings

    for talk in payload.talks:
        pptx_title = (getattr(talk, "title", "") or "").strip()
        pptx_speaker = (getattr(talk, "speaker", "") or "").strip()

        if not pptx_title and getattr(talk, "title_lines", None):
            pptx_title = "\n".join(talk.title_lines).strip()

        top = rank_vm_candidates(pptx_title, pptx_speaker, vm_rows, k=k)
        if not top:
            continue

        best_score = top[0][0]
        second_score = top[1][0] if len(top) > 1 else 0.0

        if best_score >= hi and (best_score - second_score) >= gap:
            chosen = top[0][1]["data"]
            # 演題のみ補正（必要ならここで項目追加）
            if chosen.get("演題"):
                talk.title = chosen["演題"]
        else:
            payload.warnings.append("vm_match_not_confident")

    return payload

def get_gsa_credentials(scopes):
    gsa_json = (os.getenv("GSA_JSON") or "").strip()
    if gsa_json:
        return Credentials.from_service_account_info(
            json.loads(gsa_json),
            scopes=scopes,
        )

    # ローカル用 fallback
    sa_path = (os.getenv("GOOGLE_SA_PATH") or str(APP_DIR / "_master" / "client_secret.json")).strip()
    if Path(sa_path).exists():
        return Credentials.from_service_account_file(
            sa_path,
            scopes=scopes,
        )

    raise RuntimeError("GSA_JSON (or GOOGLE_SA_PATH) is not set")


def get_gsa_client_email() -> str:
    gsa_json = (os.getenv("GSA_JSON") or "").strip()
    if gsa_json:
        return json.loads(gsa_json).get("client_email", "")

    sa_path = (os.getenv("GOOGLE_SA_PATH") or str(APP_DIR / "_master" / "client_secret.json")).strip()
    if Path(sa_path).exists():
        return json.loads(Path(sa_path).read_text(encoding="utf-8")).get("client_email", "")

    return ""



def dump_json(obj) -> str:
    if hasattr(obj, "model_dump_json"):
        return obj.model_dump_json(indent=2)
    return obj.json(ensure_ascii=False, indent=2)


def new_session_id() -> str:
    return uuid.uuid4().hex


def normalize_space(s: str) -> str:
    s = (s or "").replace("\u3000", " ")
    s = re.sub(r"\s+", " ", s).strip()
    return s


def normalize_keep_newlines(s: str) -> str:
    """
    改行は保持したまま、各行の空白だけ正規化する
    """
    s = (s or "").replace("\u3000", " ")
    lines = []
    for line in s.splitlines():
        line = re.sub(r"[ \t]+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def normalize_key(s: str) -> str:
    s = (s or "").replace("\u3000", " ")
    s = re.sub(r"\s+", "", s)
    return s


ROLE_PREFIXES = [
    # 司会・進行系
    "総合司会", "座長", "司会", "進行",
    # 学術・医療肩書き
    "教授", "准教授", "準教授", "客員教授", "名誉教授", "特任教授", "臨床教授",
    "講師", "助教", "特任講師", "臨床講師", "助手",
    # 医療機関肩書き
    "院長", "副院長", "名誉院長", "理事長", "副理事長",
    "部長", "副部長", "医長", "副医長", "課長", "副課長", "主任", "副主任",
    "センター長", "科長", "室長", "技師長", "看護師長",
    "診療部長", "診療科長", "診療科部長",
    # 学会・組織肩書き
    "会長", "副会長", "理事", "副理事", "監事", "評議員", "幹事", "委員長", "副委員長",
    # 医師資格・専門医
    "医師", "歯科医師", "薬剤師", "看護師", "専門医", "認定医", "指導医",
    # その他
    "博士", "Dr.", "Professor", "Prof.", "講 師", "教 授", "部 長"
]


def strip_role_prefix(s: str) -> str:
    """名前から肩書きを除去（前置・後置両方対応）"""
    x = normalize_space(s or "").replace("\n", "")
    
    # 前置肩書きの除去（繰り返し処理）
    changed = True
    while changed and x:
        changed = False
        for p in ROLE_PREFIXES:
            if x.startswith(p):
                x = x[len(p):].strip()
                changed = True
    
    # 後置肩書きの除去
    changed = True
    while changed and x:
        changed = False
        for p in ROLE_PREFIXES:
            if x.endswith(p):
                x = x[:-len(p)].strip()
                changed = True
    
    return x


def detect_chair_role(text: str) -> str:
    s = normalize_space(text or "").replace("\n", "")
    if not s:
        return ""
    for p in ["総合司会", "座長"]:
        if p in s:
            return p
    return ""


def normalize_person_name(s: str) -> str:
    x = strip_role_prefix(s)
    x = normalize_space(x).replace("先生", "").strip()
    x = x.replace("　", " ")
    x = x.replace(" ", "")
    return x


def normalize_affiliation(s: str) -> str:
    """所属から肩書きを除去し、適切に正規化"""
    x = normalize_space(s or "").replace("\n", " ")
    
    # 先頭の人名（先生付き）を除去 - より包括的な日本語文字パターン
    original_x = x
    
    # パターン1: より包括的な日本語文字セット（西根広樹等に対応）
    # \u4e00-\u9faf: CJK統合漢字
    # \u3041-\u3096: ひらがな  
    # \u30a1-\u30fa: カタカナ
    # \u3005: 々
    # \u3006: 〆
    # \u3024: 〤  
    x = re.sub(r'^[\u4e00-\u9faf\u3041-\u3096\u30a1-\u30fa\u3005\u3006\u3024]{2,8}\s*先生\s*[（(]?\s*', '', x)
    
    # 除去後に空または短すぎる場合は、スペース込みでより長い名前も除去
    if not x.strip() or len(x.strip()) < 3:
        x = original_x
        # スペース込みのより長い名前パターンも除去（西根　広樹等）
        x = re.sub(r'^[\u4e00-\u9faf\u3041-\u3096\u30a1-\u30fa\u3005\u3006\u3024\s]{2,16}先生\s*[（(]?\s*', '', x)
    
    # まだ短すぎる場合は、括弧内の機関情報だけを抽出
    if not x.strip() or len(x.strip()) < 3:
        x = original_x
        # 括弧内の機関情報だけを抽出
        bracket_match = re.search(r'[（(]([^)）]+)[）)]', x)
        if bracket_match:
            x = bracket_match.group(1)
        else:
            # 括弧がない場合は「先生」以降の部分を取得
            match = re.search(r'先生\s*(.+)', original_x)
            if match:
                x = match.group(1)
    
    # 先頭が ( で始まる場合、( を除去  
    x = re.sub(r'^\s*[（(]\s*', '', x)
    # 末尾が ) で終わる場合、) を除去
    x = re.sub(r'\s*[）)]\s*$', '', x)
    
    # 先頭の肩書き除去
    changed = True
    while changed and x:
        changed = False
        for p in ROLE_PREFIXES:
            if x.startswith(p):
                x = x[len(p):].strip()
                changed = True
                break
    
    # 「講 師」のようなスペース入り肩書きも除去
    x = re.sub(r'^(講\s*師|教\s*授|部\s*長|院\s*長|課\s*長)\s*', '', x)
    
    # 結果が短すぎる場合は元の文字列（ただし改行は除去）
    if len(x.strip()) < 3:
        return normalize_space(original_x)
    
    return x.strip()


def normalize_person_display(s: str) -> str:
    x = strip_role_prefix(s)
    x = normalize_space(x).replace("先生", "").replace("\n", " ").strip()
    return build_speaker_display(x)


def norm_name(s: str) -> str:
    return normalize_person_name(s)
def split_tilde_subtitle_lines(line: str) -> List[str]:
    """
    1行内の ~...~ / ～...～ を「別行扱い」にする
    """
    s = normalize_space(line)
    if not s:
        return []

    if re.fullmatch(r"[~～].+[~～]", s):
        return [s]

    m = re.search(r"([~～].+[~～])", s)
    if not m:
        return [s]

    before = normalize_space(s[: m.start()])
    sub = normalize_space(m.group(1))
    after = normalize_space(s[m.end() :])

    out: List[str] = []
    if before:
        out.append(before)
    if sub:
        out.append(sub)
    if after:
        out.append(after)
    return out

def clean_event_title_line(s: str) -> str:
    s = normalize_space(s or "")

    # 「〇〇のご案内」→後ろだけ残す
    m = re.match(r"^(.+?)のご案内\s*(.*)$", s)
    if m:
        after = m.group(2).strip()
        if after:
            return after
        return ""

    # 「開催のご案内」系
    m = re.match(r"^(.+?)開催のご案内\s*(.*)$", s)
    if m:
        after = m.group(2).strip()
        if after:
            return after
        return ""

    return s

def normalize_lines_keep_order(lines: List[str]) -> List[str]:
    out: List[str] = []
    seen = set()
    for l in lines:
        for x in split_tilde_subtitle_lines(l):
            x = normalize_space(x)
            if not x:
                continue
            if x not in seen:
                out.append(x)
                seen.add(x)
    cleaned = [clean_event_title_line(l) for l in out]
    cleaned = [l for l in cleaned if l]  # 空削除

    if cleaned:
        out = cleaned
    return out


def job_paths(job_id: str):
    d = DATA_DIR / job_id
    d.mkdir(parents=True, exist_ok=True)
    return {
        "dir": d,
        "input": d / "input.bin",
        "pptx": d / "input.pptx",
        "json": d / "latest.json",
        "jpg": d / "preview.jpg",
        "debug_html": d / "debug.html",
        "debug_blocks": d / "blocks.json",
    }


def fix_warnings(payload: DesignJSON) -> None:
    w = set(payload.warnings or [])
    if payload.organizer:
        w.discard("missing_organizer")

    # ★ほぼ埋まっていて confidence が高いなら ai_refined は外す（運用用）
    core_ok = bool(payload.event_title) and bool(payload.datetime) and bool(payload.organizer)
    if core_ok and (payload.confidence or 0) >= 0.98:
        w.discard("ai_refined")

    payload.warnings = sorted(w)


ORG_LABEL_PAT = re.compile(r"^(主催|共催|提供|企画|運営)\s*[:：]\s*(.+)$")

ORG_LABEL_PAT = re.compile(r"^(主催|共催|提供|企画|運営)\s*[:：]\s*(.+)$")
ORG_BRACKET_PAT = re.compile(r"^[【\[]\s*(主催|共催|提供|企画|運営)\s*[】\]]\s*(.+)$")

def _organizer_seps_to_space(s: str) -> str:
    # ／,、・ などは全部スペースに寄せる
    s = s.replace("／", " ").replace("/", " ")
    s = s.replace(",", " ").replace("，", " ")
    s = s.replace("、", " ").replace("・", " ")
    s = normalize_space(s)
    return s

def normalize_organizer(org: str) -> str:
    s = normalize_space(org).replace("（", "(").replace("）", ")")

    m = ORG_LABEL_PAT.match(s)
    if not m:
        m = ORG_BRACKET_PAT.match(s)

    if m:
        label = m.group(1)
        body = m.group(2)
        body = ORG_CANON.get(body, body)
        body = _organizer_seps_to_space(body)
        return f"{label}: {body}"   # ★半角コロン + 半角スペース

    # ラベル無し
    s2 = ORG_CANON.get(s, s)
    return _organizer_seps_to_space(s2)

KANJI_NAME_PAT = re.compile(r"^[\u4E00-\u9FFF]{2,6}$")

COMMON_LASTNAMES = {
    # 上位頻出
    "佐藤","鈴木","高橋","田中","伊藤","渡辺","山本","中村","小林","加藤",
    "吉田","山田","佐々木","山口","松本","井上","木村","林","清水","山崎",
    "森","阿部","池田","橋本","山下","石川","中島","前田","藤田","後藤",
    "小川","岡田","長谷川","村上","近藤","石井","斉藤","坂本","遠藤","青木",
    "藤井","西村","福田","太田","三浦","岡本","松田","中川","中野","原田",
    "小野","田村","竹内","金子","和田","中山","石田","上田","森田","原",
    "酒井","工藤","横山","宮崎","宮本","内田","高木","安藤","谷口","大野",
    "今井","丸山","高田","藤本","武田","村田","上野","杉山","増田","小島",
    "大塚","平野","菅原","久保","松井","千葉","岩崎","桜井","木下","野口",
    "松尾","野村","菊地","佐野","杉本","新井","浜田","市川","古川","小松",
    "高野","水野","吉川","島田","小山","大西","西田","西川","土屋","飯田",
    "渡部","川口","関","川村","永井","齋藤","本田","佐久間","松岡","山中",
    "川上","北村","西山","五十嵐","福島","安田","平田","中田","川崎","飯塚",
    "荒木","河野","田口","星野","岡崎","荒井","大久保","浅野","野田","松下",
    "小池","山内","中西","篠原","須藤","広瀬","吉岡","長田","本間","川島",
    "藤原","熊谷","片山","小沢","成田","宮田","大橋","石原","岡","富田",
    "大島","大谷","西岡","児玉","馬場","矢野","田辺","秋山","松浦","堀",
    "大川","宮下","吉村","岩田","奥田","松原","栗原","大石","中井","尾崎",
    "横田","岡村","三宅","松村","岩本","菊池","早川","吉野","中谷","片岡",
    "内藤","中尾","奥村","松永","望月","岩下","福井","村井","大森","片桐",
    "石橋","黒田","堀内","大竹","大場","高山","宮内","西本","矢島","川田",
    "松崎","徳永","川辺","平山","大沢","吉沢","横井","奥野","柳沢","大村",
    "宮原","三好","大島","藤川","北川","川端","本多","福本","石塚","古田",
    "長尾","永田","江口","杉浦","高井","大山","神田","森本","土井","水谷",
    "小倉","柴田","山岸","川合","三輪","西尾","谷","村松","高岡","白石",
    "大槻","小泉","坂井","岸本","松山","安部","宮川","岩井","金田","藤岡",
    "大崎","岡野","杉田","島崎","浜口","村山","黒川","中沢","江藤","武藤",
    "上原","津田","大内","森山","菅野","高見","柴山","坂田","矢口","川本",
    "坂上","石黒","高松","石野","黒木","大原","宮崎","木原","宮沢","島村",
    "松谷","平井","今村","吉本","石垣","川原","小関","宮沢","西谷","杉原",

    # 医療系でよく出る拡張（重要）
    "石和田","田邉","渡邊","齊藤","齋藤","髙橋","髙田","髙木","髙野",
    "長谷部","長谷川","佐々木","小笠原","宇都宮","上野山","久保田",
    "川井田","川井","川瀬","川原田","西條","西條","西脇","西尾",
    "森下","森川","森岡","森口","森山","森元",
    "林田","林原","林本","林崎",
    "石川","石原","石橋","石山","石丸",
    "藤田","藤原","藤村","藤野","藤沢","藤岡","藤本",
    "高橋","高木","高田","高野","高山","高岡","高井","米田","樋口"
}

def split_name_by_dictionary(name: str) -> str:
    core = norm_name(name)
    if not core:
        return ""

    # 最大4文字まででマッチ（長い姓優先）
    for i in range(min(4, len(core)), 0, -1):
        last = core[:i]
        if last in COMMON_LASTNAMES:
            if i < len(core):
                return f"{last} {core[i:]}"
            return core

    return ""

ONE_CHAR_LASTNAMES = {
    "森", "林", "原", "堀", "関", "郭", "秦", "東", "西", "南", "北",
    "辻", "堤", "岸", "今", "岡", "萩", "星", "楊", "呉", "文", "李"
}

COMMON_GIVEN_3_SUFFIXES = {
    "一郎", "二郎", "三郎", "四郎", "五郎",
    "太郎", "次郎",
    "子", "美", "香", "菜", "奈", "乃", "江", "恵"
}

def _looks_like_one_char_lastname_case(core: str) -> bool:
    """
    4文字名を 1+3 にしてよさそうかを雑に判定
    例:
      森啓一郎 -> True
      西田育功 -> False
    """
    if len(core) != 4:
        return False

    if core[0] not in ONE_CHAR_LASTNAMES:
        return False

    given = core[1:]  # 3文字

    # 典型: 啓一郎 / 健太郎 / 恒一郎 みたいに末尾2文字がよくある名の終わり
    if given[1:] in COMMON_GIVEN_3_SUFFIXES:
        return True

    # 1文字目が「啓/健/裕/智/和/直/信/良/洋/雅/孝/祐/雄/達/亮」あたりで
    # かつ末尾が 郎 のようなケースを救う
    if given.endswith("郎"):
        return True

    return False


def add_space_to_jp_name(name: str) -> str:
    raw = str(name or "")
    s = raw.replace("\u3000", " ").strip()
    if not s:
        return ""

    # すでに空白ありなら尊重
    parts = [p for p in re.split(r"\s+", s) if p]
    if len(parts) >= 2:
        clean_parts = [re.sub(r"^(座長|演者|司会|講師|:|：)\s*", "", p) for p in parts]
        clean_parts = [re.sub(r"(先生)\s*$", "", p) for p in clean_parts]
        clean_parts = [p for p in clean_parts if p]

        if len(clean_parts) >= 2 and len(clean_parts[0]) >= 1 and len(clean_parts[1]) >= 1:
            return f"{clean_parts[0]} {clean_parts[1]}"

    core = re.sub(r"\s+", "", s)
    core = re.sub(r"^(座長|演者|司会|講師|:|：)", "", core)
    core = re.sub(r"(先生)$", "", core)

    if not all(("\u3040" <= ch <= "\u30ff") or ("\u4e00" <= ch <= "\u9fff") or (ch == "々") for ch in core):
        return core

    n = len(core)
    if n <= 2:
        return core

    if n == 3:
        return core[:2] + " " + core[2:]

    if n == 4:
        # 1文字姓救済は条件付き
        if _looks_like_one_char_lastname_case(core):
            return core[:1] + " " + core[1:]
        return core[:2] + " " + core[2:]

    if n == 5:
        return core[:2] + " " + core[2:]

    if n == 6:
        return core[:3] + " " + core[3:]

    return core

def build_speaker_display(name: str) -> str:
    core = norm_name(name)
    if not core:
        return ""

    # ① 正解DBから学習した分割位置（最優先）
    cache = _get_speaker_display_cache()
    if core in cache:
        cached_val = cache[core]
        # キャッシュ値の文字（スペース除去）がnameと一致する場合のみ適用
        # 異なる場合は誤ったDB登録の可能性があるためスキップ
        if cached_val.replace(" ", "").replace("\u3000", "") == core:
            return cached_val

    # ② 辞書
    v = split_name_by_dictionary(core)
    if v:
        return v

    # ③ fallback（既存ロジック）
    return add_space_to_jp_name(core)

def extend_lastname_dict_from_vm(vm_rows):
    for row in vm_rows:
        d = row["data"] if isinstance(row, dict) and "data" in row else row
        name = norm_name(d.get("案内状掲載 医師名", ""))
        if len(name) >= 2:
            COMMON_LASTNAMES.add(name[:2])
            COMMON_LASTNAMES.add(name[:3])

def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()

def safe_title_for_list(payload: DesignJSON) -> str:
    # 一覧用（短く）
    if payload.event_title_lines:
        return payload.event_title_lines[0]
    return (payload.event_title or "").splitlines()[0] if payload.event_title else ""


def join_lines(lines: list[str]) -> str:
    lines = [str(l).rstrip() for l in (lines or [])]
    lines = [l for l in lines if l != ""]
    return "\n".join(lines)

def pull_affil_out_of_title_lines(talk) -> None:
    if getattr(talk, "affiliation", ""):
        return

    lines = list(getattr(talk, "title_lines", []) or [])
    if not lines:
        return

    affils = [ln for ln in lines if looks_like_affil_line(ln)]
    if not affils:
        return

    new_lines = [ln for ln in lines if ln not in affils]

    # まず affiliation は確定させる（全部所属でもOK）
    talk.affiliation = " / ".join(_norm1(a) for a in affils)

    # ★全部所属だった場合：title_lines は空にするが、title は保持する（または空なら所属を残す）
    if not new_lines:
        talk.title_lines = []
        if hasattr(talk, "title"):
            # title が空なら affiliation を見える場所に残す（運用しやすい）
            if not (talk.title or "").strip():
                talk.title = join_lines(lines)  # 元の表示を維持
        return

    # 通常ケース
    talk.title_lines = new_lines
    if hasattr(talk, "title"):
        talk.title = join_lines(new_lines)


def normalize_for_render(payload: DesignJSON) -> DesignJSON:
    # event title 合成
    if hasattr(payload, "event_title_lines") and payload.event_title_lines:
        payload.event_title = join_lines(payload.event_title_lines)

    # ★chair ラベルの掃除（role は別フィールドへ）
    if getattr(payload, "chair", None):
        if not getattr(payload.chair, "role", ""):
            payload.chair.role = detect_chair_role((getattr(payload.chair, "name", "") or "") + " " + (getattr(payload.chair, "name_display", "") or ""))
        if getattr(payload.chair, "name", ""):
            payload.chair.name = normalize_person_name(payload.chair.name)
        if getattr(payload.chair, "name_display", ""):
            payload.chair.name_display = normalize_person_display(payload.chair.name_display)
        elif getattr(payload.chair, "name", ""):
            payload.chair.name_display = build_speaker_display(payload.chair.name)

    # talk title 合成（titleフィールドが存在する時だけ代入）
    for t in (payload.talks or []):
        # ★ここで「所属が title_lines に混ざったやつ」を剥がす
        pull_affil_out_of_title_lines(t)

        if hasattr(t, "title_lines") and t.title_lines:
            if hasattr(t, "title"):
                t.title = join_lines(t.title_lines)

    return payload


DT_RE = re.compile(
    r"(20\d{2}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日.*?\d{1,2}\s*[:：]\s*\d{2}\s*[～〜\-ー~]\s*\d{1,2}\s*[:：]\s*\d{2})"
)

DATE_ONLY_RE = re.compile(r"(20\d{2}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日)")

TIME_RANGE_RE = re.compile(
    r"(\d{1,2}:\d{2})\s*[-－ー‐-‒–—―〜～~]+\s*(\d{1,2}:\d{2})"
)

def normalize_time_colon(s: str) -> str:
    # "19：00" -> "19:00"
    return (s or "").replace("：", ":")

_ZEN2HAN = str.maketrans("０１２３４５６７８９：", "0123456789:")

def normalize_datetime_text(s: str) -> str:
    s = (s or "").translate(_ZEN2HAN)
    s = normalize_space(s)
    s = s.replace("：", ":")
    # 年月日まわりの空白を削除
    s = re.sub(r"\s*年\s*", "年", s)
    s = re.sub(r"\s*月\s*", "月", s)
    s = re.sub(r"\s*日\s*", "日", s)
    # 月日のゼロ埋め除去: 04月→4月, 02日→2日
    s = re.sub(r"(\d{4}年)0(\d月)", r"\1\2", s)
    s = re.sub(r"(\d{1,2}月)0(\d日)", r"\1\2", s)
    # コロン前後
    s = re.sub(r"\s*:\s*", ":", s)
    # 20 :00 → 20:00
    s = re.sub(r"(\d)\s+(\d{2})", r"\1\2", s)
    return s.strip()

def looks_like_datetime_text(s: str) -> bool:
    s0 = normalize_datetime_text(s)
    if DT_RE.search(s0):
        return True
    # 「年/月/日 + ～」があれば日時濃厚
    if ("年" in s0 and "月" in s0 and "日" in s0 and ("～" in s0 or "〜" in s0 or "-" in s0)):
        return True
    return False

def looks_like_label(s: str) -> bool:
    k = normalize_key(s)
    # 「日 時」みたいな分割にも強い
    return any(x in k for x in ["日時", "日時", "日", "時", "座長", "演者", "主催", "共催", "提供", "企画", "運営", "会場", "形式", "登録", "視聴"])

def looks_like_talk_anchor(s: str) -> bool:
    k = normalize_key(s)
    # 講演１ / 講演1 / 演題1 等
    return bool(re.search(r"(講演|演題)([0-9]|[１-９])", k))



# ---------------- Jobs list / filter ----------------

def parse_warnings(warnings_json: str) -> List[str]:
    try:
        return json.loads(warnings_json or "[]")
    except Exception:
        return []

# def row_to_job_item(r: sqlite3.Row) -> Dict[str, Any]:
#     job_id = r["job_id"]
#     session_id = r["session_id"]
#     event_id = r["event_id"]
#     return {
#         "jobId": job_id,
#         "filename": r["filename"],
#         "session_id": session_id, 
#         "event_id": event_id, 
#         "status": r["status"],
#         "createdAt": r["created_at"],
#         "updatedAt": r["updated_at"],
#         "title": r["title"] or "",
#         "organizer": r["organizer"] or "",
#         "datetime": r["datetime"] or "",
#         "confidence": float(r["confidence"] or 0.0),
#         "warnings": parse_warnings(r["warnings_json"]),
#         "manualOverride": bool(r["manual_override"]),
#         "note": r["note"] or "",
#         "locked": bool(r["locked"]),
#         "errorMessage": r["error_message"],
#         "previewUrl": f"/preview/{job_id}.jpg",
#         "downloadUrl": f"/download/{job_id}.jpg",
#     }

def row_to_job_item(r) -> Dict[str, Any]:
    job_id = r["job_id"]
    return {
        "jobId": job_id,
        "filename": r.get("filename") or "",
        "session_id": r.get("session_id") or "",
        "event_id": r.get("event_id") or "",
        "status": r["status"],
        "createdAt": r["created_at"].isoformat(),
        "updatedAt": r["updated_at"].isoformat(),
        "title": r.get("title") or "",
        "organizer": r.get("organizer") or "",
        "datetime": r.get("datetime") or "",
        "confidence": float(r.get("confidence") or 0.0),
        "warnings": r.get("warnings_json") or [],
        "manualOverride": bool(r.get("manual_override")),
        "note": r.get("note") or "",
        "locked": bool(r.get("locked")),
        "errorMessage": r.get("error_message"),
        "previewUrl": f"/preview/{job_id}.jpg",
        "downloadUrl": f"/download/{job_id}.jpg",
    }




# ---------------- Database Connection with Retry ----------------

def db_connect():
    """
    PostgreSQL接続をタイムアウト設定とリトライ機構付きで行う
    """
    if not DATABASE_URL:
        raise ValueError("DATABASE_URL is not set")
        
    for attempt in range(DB_RETRY_ATTEMPTS):
        try:
            logger.info(f"Database connection attempt {attempt + 1}/{DB_RETRY_ATTEMPTS}")
            
            # 接続パラメータにタイムアウト設定を追加
            conn = psycopg.connect(
                DATABASE_URL, 
                row_factory=dict_row,
                connect_timeout=DB_CONNECT_TIMEOUT,
                options=f"-c statement_timeout={DB_QUERY_TIMEOUT}s"
            )
            
            # 接続テスト
            with conn.cursor() as cur:
                cur.execute("SELECT 1")
                cur.fetchone()
            
            logger.info("Database connection established successfully")
            return conn
            
        except OperationalError as e:
            error_msg = str(e)
            logger.warning(f"Database connection attempt {attempt + 1} failed: {error_msg}")
            
            if attempt < DB_RETRY_ATTEMPTS - 1:
                # 最後の試行でなければ待機
                delay = DB_RETRY_DELAY * (2 ** attempt)  # 指数バックオフ
                logger.info(f"Retrying in {delay} seconds...")
                time.sleep(delay)
            else:
                # 最後の試行でも失敗した場合
                logger.error(f"Database connection failed after {DB_RETRY_ATTEMPTS} attempts")
                if "timed out" in error_msg.lower():
                    raise OperationalError(
                        f"Database connection timeout after {DB_RETRY_ATTEMPTS} attempts. "
                        f"Please check network connectivity and database server status. "
                        f"Original error: {error_msg}"
                    )
                else:
                    raise
                    
        except Exception as e:
            logger.error(f"Unexpected database connection error: {e}")
            if attempt < DB_RETRY_ATTEMPTS - 1:
                delay = DB_RETRY_DELAY * (2 ** attempt)
                time.sleep(delay)
            else:
                raise

def safe_db_operation(operation_func, *args, **kwargs):
    """
    データベース操作を安全に実行するヘルパー関数
    """
    max_retries = 2
    for attempt in range(max_retries):
        try:
            return operation_func(*args, **kwargs)
        except OperationalError as e:
            if "connection" in str(e).lower() and attempt < max_retries - 1:
                logger.warning(f"Database operation failed, retrying: {e}")
                time.sleep(1)
                continue
            raise
        except Exception:
            raise

class SafeDBConnection:
    """
    安全なデータベース接続のコンテキストマネージャー
    """
    def __init__(self, max_retries=2):
        self.max_retries = max_retries
        self.connection = None
        
    def __enter__(self):
        def _get_connection():
            self.connection = db_connect()
            return self.connection
            
        return safe_db_operation(_get_connection)
        
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.connection:
            try:
                if exc_type is None:
                    self.connection.commit()
                else:
                    self.connection.rollback()
            except Exception as e:
                logger.error(f"Error during connection cleanup: {e}")
            finally:
                try:
                    self.connection.close()
                except Exception as e:
                    logger.error(f"Error closing connection: {e}")

def safe_db_connect():
    """
    安全なデータベース接続のコンテキストマネージャーを返す
    """
    return SafeDBConnection()


def init_db():
    """データベース初期化（リトライ機構付き）"""
    def _init_operation():
        con = db_connect()
        try:
            logger.info("Initializing database tables...")
            con.execute("""
            CREATE TABLE IF NOT EXISTS jobs (
                job_id TEXT PRIMARY KEY,
                filename TEXT,
                session_id TEXT,
                event_id TEXT NOT NULL,     
                status TEXT NOT NULL,                 -- ok / error
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,

                title TEXT,
                organizer TEXT,
                datetime TEXT,
                
                region TEXT,
                unit TEXT,

                confidence REAL,
                warnings_json TEXT,

                manual_override INTEGER NOT NULL DEFAULT 0,
                note TEXT NOT NULL DEFAULT '',
                locked INTEGER NOT NULL DEFAULT 0,

                error_message TEXT
            );
            """)
            con.execute("CREATE INDEX IF NOT EXISTS idx_jobs_updated_at ON jobs(updated_at);")
            con.execute("CREATE INDEX IF NOT EXISTS idx_jobs_status ON jobs(status);")
            con.execute("CREATE INDEX IF NOT EXISTS idx_jobs_created_at ON jobs(created_at DESC);")
            con.execute("CREATE INDEX IF NOT EXISTS idx_jobs_event_id ON jobs(event_id);")

            con.execute("""
            CREATE TABLE IF NOT EXISTS correct_answers (
                id SERIAL PRIMARY KEY,
                job_id TEXT UNIQUE NOT NULL,
                event_title TEXT NOT NULL DEFAULT '',
                blocks_text TEXT NOT NULL DEFAULT '',
                keywords TEXT[] NOT NULL DEFAULT '{}',
                correct_json JSONB NOT NULL DEFAULT '{}',
                embedding DOUBLE PRECISION[],
                created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
            );
            """)
            con.execute("CREATE INDEX IF NOT EXISTS idx_ca_job_id ON correct_answers(job_id);")
            # embedding カラムが無ければ追加（既存テーブル互換）
            # CREATE TABLE に embedding を含めているため、新規作成時は不要。
            # 既存テーブルのみ対象とし、軽量に存在チェックしてから ALTER する。
            row = con.execute("""
                SELECT 1 FROM information_schema.columns
                WHERE table_name = 'correct_answers' AND column_name = 'embedding'
            """).fetchone()
            if not row:
                con.execute("ALTER TABLE correct_answers ADD COLUMN embedding DOUBLE PRECISION[];")
            # layout_hints カラムが無ければ追加（レイアウトパターン学習用）
            row = con.execute("""
                SELECT 1 FROM information_schema.columns
                WHERE table_name = 'correct_answers' AND column_name = 'layout_hints'
            """).fetchone()
            if not row:
                con.execute("ALTER TABLE correct_answers ADD COLUMN layout_hints JSONB DEFAULT '{}';")
            con.commit()
            logger.info("Database tables initialized successfully")
        finally:
            con.close()

    # 安全な実行でリトライ機構を利用
    safe_db_operation(_init_operation)

# def upsert_job_ok(
#     job_id: str,
#     filename: str,
#     payload: DesignJSON,
#     session_id: Optional[str] = None,
#     event_id: Optional[str] = None,
# ):
#     con = db_connect()
#     try:
#         created_at = now_iso()
#         updated_at = created_at

#         warnings_json = json.dumps(payload.warnings or [], ensure_ascii=False)

#         # 既存があれば created_at 維持 + session_id/event_id を未指定なら引き継ぐ
#         row = con.execute(
#             "SELECT created_at, session_id, event_id FROM jobs WHERE job_id=?",
#             (job_id,),
#         ).fetchone()
#         if row:
#             created_at = row["created_at"]
#             updated_at = now_iso()
#             if session_id is None:
#                 session_id = row["session_id"]
#             if event_id is None:
#                 event_id = row["event_id"]

#         # 新規insertで None が残ると DB制約で死ぬなら空文字に落とす（安全策）
#         if session_id is None:
#             session_id = ""
#         if event_id is None:
#             event_id = ""

#         con.execute(
#             """
#             INSERT INTO jobs (
#                 job_id, filename, session_id, event_id, status, created_at, updated_at,
#                 title, organizer, datetime,
#                 region, unit,
#                 confidence, warnings_json,
#                 manual_override, note, locked,
#                 error_message
#             ) VALUES (?, ?, ?, ?, 'ok', ?, ?,
#                       ?, ?, ?,
#                       ?, ?,
#                       ?, ?,
#                       ?, ?, ?,
#                       NULL)
#             ON CONFLICT(job_id) DO UPDATE SET
#                 filename=excluded.filename,
#                 session_id=excluded.session_id,
#                 event_id=excluded.event_id,
#                 status='ok',
#                 updated_at=excluded.updated_at,
#                 title=excluded.title,
#                 organizer=excluded.organizer,
#                 datetime=excluded.datetime,
#                 region=excluded.region,
#                 unit=excluded.unit,
#                 confidence=excluded.confidence,
#                 warnings_json=excluded.warnings_json,
#                 manual_override=excluded.manual_override,
#                 note=excluded.note,
#                 locked=excluded.locked,
#                 error_message=NULL
#             """,
#             (
#                 job_id,
#                 filename,
#                 session_id,
#                 event_id,  # ★ここはそのまま正しい並び
#                 created_at,
#                 updated_at,
#                 safe_title_for_list(payload),
#                 payload.organizer,
#                 payload.datetime,
#                 payload.region,
#                 payload.unit,
#                 float(payload.confidence or 0.0),
#                 warnings_json,
#                 1 if getattr(payload, "manual_override", False) else 0,
#                 getattr(payload, "note", "") or "",
#                 1 if getattr(payload, "locked", False) else 0,
#             ),
#         )
#         con.commit()
#     finally:
#         con.close()

def now_ts():
    return datetime.now(timezone.utc)

def upsert_job_ok(job_id: str, filename: str, payload, session_id: str = "", event_id: str = ""):
    created_at = now_ts()
    updated_at = created_at
    warnings = payload.warnings or []

    # jsonb array として入れる（["missing_organizer", ...]）
    warnings_jsonb = json.dumps(warnings, ensure_ascii=False)

    with db_connect() as con:
        row = con.execute(
            "SELECT created_at, session_id, event_id FROM jobs WHERE job_id=%s",
            (job_id,),
        ).fetchone()

        if row:
            created_at = row["created_at"]
            updated_at = now_ts()
            if not session_id:
                session_id = row.get("session_id") or ""
            if not event_id:
                event_id = row.get("event_id") or ""

        con.execute(
            """
            INSERT INTO jobs (
              job_id, filename, session_id, event_id, status, created_at, updated_at,
              title, organizer, datetime, region, unit, confidence, warnings_json,
              manual_override, note, locked, error_message
            )
            VALUES (%s,%s,%s,%s,'ok',%s,%s,
                    %s,%s,%s,%s,%s,%s,%s::jsonb,
                    %s,%s,%s,NULL)
            ON CONFLICT (job_id) DO UPDATE SET
              filename=excluded.filename,
              session_id=excluded.session_id,
              event_id=excluded.event_id,
              status='ok',
              updated_at=excluded.updated_at,
              title=excluded.title,
              organizer=excluded.organizer,
              datetime=excluded.datetime,
              region=excluded.region,
              unit=excluded.unit,
              confidence=excluded.confidence,
              warnings_json=excluded.warnings_json,
              manual_override=excluded.manual_override,
              note=excluded.note,
              locked=excluded.locked,
              error_message=NULL
            """,
            (
                job_id, filename, session_id or "", event_id or "",
                created_at, updated_at,
                safe_title_for_list(payload),
                payload.organizer or "",
                payload.datetime or "",
                payload.region or "",
                payload.unit or "",
                float(payload.confidence or 0.0),
                warnings_jsonb,
                bool(getattr(payload, "manual_override", False)),
                getattr(payload, "note", "") or "",
                bool(getattr(payload, "locked", False)),
            ),
        )
        con.commit()  # ★必須


def upsert_job_error(job_id: str, filename: str, error_message: str, event_id: str = ""):
    created_at = now_ts()
    updated_at = created_at

    with db_connect() as con:
        row = con.execute(
            "SELECT created_at, event_id FROM jobs WHERE job_id=%s",
            (job_id,),
        ).fetchone()

        if row:
            created_at = row["created_at"]
            updated_at = now_ts()
            if not event_id:
                event_id = (row.get("event_id") or "")

        con.execute(
            """
            INSERT INTO jobs (
              job_id, filename, session_id, event_id, status, created_at, updated_at,
              title, organizer, datetime, region, unit, confidence, warnings_json,
              manual_override, note, locked, error_message
            )
            VALUES (%s,%s,%s,%s,'error',%s,%s,
                    %s,%s,%s,%s,%s,%s,%s::jsonb,
                    %s,%s,%s,%s)
            ON CONFLICT (job_id) DO UPDATE SET
              filename=excluded.filename,
              status='error',
              event_id=excluded.event_id,
              updated_at=excluded.updated_at,
              error_message=excluded.error_message
            """,
            (
                job_id,
                filename,
                "",                 # session_id は error の時は空でOK（必要なら保持）
                event_id or "",
                created_at,
                updated_at,
                "", "", "",         # title/organizer/datetime
                "", "",             # region/unit
                0.0,
                "[]",               # warnings_json
                False,
                "",
                False,
                error_message or "",
            ),
        )
        con.commit()

# ---------------- PPTX (Blocks) ----------------
@dataclass
class TextBlock:
    text: str
    left: int
    top: int
    width: int
    height: int
    max_font_pt: float


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        # GROUP = 6
        if getattr(sh, "shape_type", None) == 6:
            for sub in iter_shapes(sh.shapes):
                yield sub




def extract_blocks_from_pptx(pptx_path: Path, first_slide_only: bool = True) -> List[TextBlock]:
    prs = Presentation(str(pptx_path))
    blocks: List[TextBlock] = []

    slides = [prs.slides[0]] if (first_slide_only and len(prs.slides) > 0) else prs.slides
    for slide in slides:
        for sh in iter_shapes(slide.shapes):
            if not getattr(sh, "has_text_frame", False):
                continue
            tf = sh.text_frame
            if not tf:
                continue

            paras = []
            max_font = 0.0
            for p in tf.paragraphs:
                t = (p.text or "").strip()
                if t:
                    paras.append(t)
                for run in p.runs:
                    if run.font and run.font.size:
                        max_font = max(max_font, run.font.size / EMU_PER_PT)

            # ★改行保持する
            text = normalize_keep_newlines("\n".join(paras))
            if not text:
                continue

            blocks.append(
                TextBlock(
                    text=text,
                    left=int(sh.left),
                    top=int(sh.top),
                    width=int(sh.width),
                    height=int(sh.height),
                    max_font_pt=float(max_font or 0.0),
                )
            )

    blocks.sort(key=lambda b: (b.top, b.left))
    return blocks


def blocks_to_lines(blocks: List[TextBlock]) -> List[str]:
    # 改行は潰して良い用途（datetime/organizer検出など）向け
    out: List[str] = []
    seen = set()
    for b in blocks:
        s = normalize_space(b.text.replace("\n", " "))
        if not s:
            continue
        if s not in seen:
            out.append(s)
            seen.add(s)
    return out


def in_region(b: TextBlock, x0: float, y0: float, x1: float, y1: float) -> bool:
    cx = b.left + b.width / 2.0
    cy = b.top + b.height / 2.0
    return (x0 <= cx <= x1) and (y0 <= cy <= y1)

def looks_like_body_text_for_title(s: str) -> bool:
    s2 = normalize_space(s)
    if not s2:
        return False
    # long polite-body sentences
    if len(s2) >= 30:
        if any(k in s2 for k in ["謹啓", "謹白", "時下", "平素", "ご高配", "ご多用", "厚く御礼", "お慶び"]):
            return True
    # explicitly exclude these keywords even if short
    if any(k in s2 for k in ["謹啓", "謹白"]):
        return True
    
    body_kw = [
        "本会は", "事前参加登録", "参加をご希望", "担当者へご連絡",
        "医療従事者", "医療系資格", "学生", "受付", "医療事務",
        "ご参加はご遠慮", "お願い申し上げます", "ご了承ください",
        "ご視聴には", "事前参加予約", "芳名録", "個人情報",
    ]
    if any(k in s2 for k in body_kw):
        return True
    
    return False


# def looks_like_body_text_for_title(s: str) -> bool:
#     s2 = normalize_space(s)
#     if not s2:
#         return False
        

#     # 既存: 挨拶文
#     if any(k in s2 for k in ["謹啓", "謹白", "時下", "平素", "ご高配", "厚く御礼"]):
#         return True

#     # ★追加: 案内・注意文（今回の混入パターン）
#     body_kw = [
#         "本会は", "事前参加登録", "参加をご希望", "担当者へご連絡",
#         "医療従事者", "医療系資格", "学生", "受付", "医療事務",
#         "ご参加はご遠慮", "お願い申し上げます", "ご了承ください",
#         "ご視聴には", "事前参加予約", "芳名録", "個人情報",
#     ]
#     if any(k in s2 for k in body_kw):
#         return True

#     # 長文は本文率高い（ただしタイトル長めもあるので閾値は控えめに）
#     if len(s2) >= 40:
#         return True

#     return False

def looks_like_format_value(s: str) -> bool:
    s2 = normalize_space(s)
    return (
        "Live配信" in s2
        or "Web（" in s2
        or s2.endswith("による開催")
        or s2.endswith("による配信")
        or s2.startswith("Web")
    )




def extract_event_title_lines_from_blocks(blocks: List[TextBlock]) -> List[str]:
    if not blocks:
        return []

    tops = [b.top for b in blocks]
    min_top, max_top = min(tops), max(tops)
    mid_top = min_top + (max_top - min_top) * 0.35  # 上1/3

    def looks_like_date_title_block(s: str) -> bool:
        s0 = normalize_space(s)
        s1 = re.sub(r"\s+", "", s0)

        # 典型: 2026年 5月19日(火), 2026年\n5月19日(火)
        if re.search(r"\d{4}年\d{1,2}月\d{1,2}日", s1):
            return True

        # 年と月日が改行で分かれているケース
        if re.search(r"\d{4}年", s1) and re.search(r"\d{1,2}月\d{1,2}日", s1):
            return True

        # 曜日つき
        if re.search(r"\d{1,2}月\d{1,2}日[（(][月火水木金土日][)）]", s1):
            return True

        # 時間だけ大きく出るケースも一応除外
        if re.fullmatch(r"\d{1,2}:\d{2}[~〜～\-－]\d{1,2}:\d{2}", s1):
            return True

        return False

    def kw_norm(s: str) -> str:
        return re.sub(r"\s+", "", s)

    def is_title_excluded(s: str) -> bool:
        k = kw_norm(s)

        # 単独ラベルだけ除外
        if k in {"主催", "共催", "座長", "演者", "演題", "会場", "形式", "日時", "PROGRAM"}:
            return True

        # これは強めに除外してよい
        if "共催" in k or "主催" in k:
            return True

        # 「座長 池上達義先生」「演者 二宮貴一朗先生」等の人名付き役職ブロックを除外
        if "先生" in k and any(lbl in k for lbl in ["座長", "演者", "総合司会", "司会"]):
            return True

        # 「添田周先生福島県立医科大学産婦人科学講座教授」等の人名+先生+所属ブロックを除外
        if "先生" in k and any(kw in k for kw in ["大学", "病院", "センター", "クリニック", "教授", "部長", "院長", "医長", "講師", "准教授"]):
            return True

        return False

    def contains_japanese(s: str) -> bool:
        return bool(re.search(r"[ぁ-んァ-ン一-龯]", s))

    def english_ratio(s: str) -> float:
        if not s:
            return 0.0
        latin = sum(1 for ch in s if ("A" <= ch <= "Z") or ("a" <= ch <= "z"))
        return latin / max(1, len(s))

    def trim_mixed_english_title(s: str) -> str:
        """
        例: "Seminar さてこの度…" → "Seminar" に切る
        """
        # 最初の日本語文字が出た位置でカット
        m = re.search(r"[ぁ-んァ-ン一-龯]", s)
        if m:
            s = s[: m.start()].rstrip()
        return s.strip()
    
    def title_head_score(b: TextBlock) -> tuple:
        s = normalize_space(b.text)
        dt_penalty = 1 if looks_like_date_title_block(s) else 0
        # フォントサイズは1pt単位で丸める（36.0 vs 36.03 の誤差吸収）
        rounded_pt = round(b.max_font_pt or 0)
        # penalty が小さい方が優先、次にフォント（大きい方優先）、次に上側
        return (dt_penalty, -rounded_pt, b.top)

    def clean_title_text(s: str) -> str:
        s = normalize_space(s)

        # 改行区切りで末尾ラベルを落とす
        parts = [x.strip() for x in str(s).split("\n") if x.strip()]
        while parts and parts[-1] in {"日時", "会場", "座長", "演者", "演題", "開催形式"}:
            parts.pop()

        s = " ".join(parts)

        # 念のため末尾にぶら下がった単独ラベルも落とす
        s = re.sub(r"\s*(日時|会場|座長|演者|演題|開催形式)\s*$", "", s)

        return s.strip()
    

    cand = []
    for b in blocks:
        s = normalize_space(b.text)
        if not s:
            continue
        if looks_like_datetime_text(s) or looks_like_date_title_block(s) or looks_like_format_value(s) or looks_like_body_text_for_title(s):
            continue
        if is_title_excluded(s):
            continue
        if b.top <= mid_top and (b.max_font_pt or 0) >= 18:
            cand.append(b)

    if not cand:
        # 上1/3に候補なし → 全体から同じフィルタ（top制限なし）で最大フォントを選ぶ
        cand2 = []
        for b in blocks:
            s = normalize_space(b.text)
            if not s:
                continue
            if looks_like_datetime_text(s) or looks_like_date_title_block(s) or looks_like_format_value(s) or looks_like_body_text_for_title(s):
                continue
            if is_title_excluded(s):
                continue
            cand2.append(b)
        if not cand2:
            return []

        # 最大フォントを特定し、同等フォントの全ブロックを top でグループ化
        # （タイトルが画像化されて単語ごとに分割されているケースに対応）
        max_font = max(b.max_font_pt or 0 for b in cand2)
        large_blocks = [b for b in cand2 if (b.max_font_pt or 0) >= max_font - 2]
        large_blocks.sort(key=lambda b: b.top)

        # top が近いブロックを同一行にグループ化
        rows: list = []
        for b in large_blocks:
            placed = False
            for row in rows:
                rep = row[0]
                tol = max(rep.height or 400000, b.height or 400000) * 0.5
                if abs(b.top - rep.top) <= tol:
                    row.append(b)
                    placed = True
                    break
            if not placed:
                rows.append([b])

        if not rows:
            return []

        # 最上段の行を1行目、直下の行を2行目として返す
        rows.sort(key=lambda r: min(x.top for x in r))
        out_lines = []
        for row in rows[:2]:
            row.sort(key=lambda b: b.left or 0)
            merged = "".join(normalize_space(b.text) for b in row)
            if merged:
                out_lines.append(merged)
        return out_lines if out_lines else []

    # head（タイトル本体）
    # head = sorted(cand, key=lambda b: ((b.max_font_pt or 0), -b.top), reverse=True)[0]
    head = sorted(cand, key=title_head_score)[0]
    head_text = normalize_space(head.text)

    # head が英語っぽいなら “混ざり” を切る
    head_is_english = english_ratio(head_text) >= 0.45 and not contains_japanese(head_text)
    if head_is_english:
        head_text = trim_mixed_english_title(head_text)

    # ✅ near 範囲（タイトルは2行分程度）
    # head の height を基準にして、同等サイズのブロックが2行目にある場合を拾う
    title_line_height = max(head.height, 300000)  # 最低30万EMU(約0.5行分)
    x0 = head.left - 400000
    x1 = head.left + head.width + 400000
    y0 = head.top - title_line_height
    y1 = head.top + title_line_height * 4  # 約4行分の余裕

    near = [b for b in blocks if in_region(b, x0, y0, x1, y1)]
    near = sorted(near, key=lambda b: (b.top, b.left))

    lines: List[str] = []
    for b in near:
        s = normalize_space(b.text)
        s = clean_title_text(b.text)
        if not s:
            continue
        if looks_like_datetime_text(s) or looks_like_format_value(s) or looks_like_body_text_for_title(s):
            continue
        if is_title_excluded(s):
            continue

        # ✅ 本文（小さいフォント）を落とす：head との差でフィルタ
        if (b.max_font_pt or 0) < max(14, (head.max_font_pt or 0) - 4):
            continue

        # ✅ 英語タイトルモードなら、日本語混在をカットして本文側は捨てる
        if head_is_english:
            s2 = trim_mixed_english_title(s)
            if not s2:
                continue
            s = s2

        lines.append(s)
        # ✅ タイトルは最大2行まで（Web Seminar系の暴走防止）
        if len(lines) >= 2:
            break

    # head が拾えてない時の保険
    if not lines:
        lines = [head_text] if head_text else [normalize_space(head.text)]

    # 重複除去（順序維持）
    out, seen = [], set()
    for s in lines:
        if s not in seen:
            out.append(s)
            seen.add(s)

    return out





DATE_RE = re.compile(r"(20\d{2})\s*年\s*(\d{1,2})\s*月\s*(\d{1,2})\s*日(?:\s*[（(]?\s*([月火水木金土日])\s*[）)]?)?")
# 月・日の数字がテンプレートとは別テキストに分かれているケース:
# 「３ 26 日時 2026年 月 日（木）12:30~13:10」→ month=3, day=26, year=2026
SPLIT_DATE_RE = re.compile(r"(\d{1,2})\s+(\d{1,2})\s+[^\d]+?(20\d{2})\s*年\s*月\s*日(?:\s*[（(]\s*([月火水木金土日])\s*[）)])?")  
TIME_RE2 = re.compile(r"(\d{1,2}[:：]\d{2})\s*[～〜\-ー~]\s*(\d{1,2}[:：]\d{2})")

def _norm_time2(s: str) -> str:
    s = (s or "").strip()
    s = s.replace("：", ":")
    s = s.replace("〜", "～").replace("~", "～").replace("-", "～").replace("ー", "～")
    return s

def _normalize_cjk_compat(s: str) -> str:
    """CJK互換部首文字を通常の漢字に変換: ⽇→日, ⽉→月, ⽕→火, ⽔→水, ⽊→木, ⾦→金, ⼟→土, ⽇→日"""
    table = str.maketrans({
        "\u2F47": "\u65E5",  # ⽇ → 日
        "\u2F49": "\u6708",  # ⽉ → 月
        "\u2F51": "\u65E5",  # ⽇ → 日 (duplicate radical)
        "\u2F6B": "\u706B",  # ⽕ → 火
        "\u2F54": "\u6C34",  # ⽔ → 水
        "\u2F32": "\u571F",  # ⼟ → 土
        "\u2F46": "\u65E5",  # ⽆ variant
        "\u2F4E": "\u6728",  # ⽊ → 木
        "\u2F91": "\u91D1",  # ⾦ → 金
        "\u2F44": "\u5E74",  # ⽄ → 年 (if needed)
    })
    return s.translate(table)


def _merge_scattered_blocks_to_lines(blocks: List[TextBlock], top_tolerance: int = 200000) -> List[str]:
    """同一行（近い top）のブロックを left 順で結合して行リストを返す。
    1文字ずつバラバラに配置された日時等を復元するため。"""
    if not blocks:
        return []
    sorted_blocks = sorted(blocks, key=lambda b: (b.top, b.left))
    rows: List[List[TextBlock]] = []
    for b in sorted_blocks:
        if rows and abs(b.top - rows[-1][0].top) < top_tolerance:
            rows[-1].append(b)
        else:
            rows.append([b])
    out = []
    for row in rows:
        row.sort(key=lambda b: b.left)
        merged = "".join((b.text or "").replace("\n", " ") for b in row)
        merged = _normalize_cjk_compat(merged)
        merged = normalize_space(merged)
        if merged:
            out.append(merged)
    return out


def extract_datetime_from_blocks(blocks: List[TextBlock]) -> str:
    # 1) まず全テキスト（行）を作る
    lines = [normalize_space(x) for x in blocks_to_lines(blocks)]
    lines = [x for x in lines if x]

    # 1b) CJK互換部首文字を正規化
    lines = [_normalize_cjk_compat(l) for l in lines]

    # 1c) バラバラブロックを結合した行も追加（1文字ずつ分離されたケース対応）
    merged_lines = _merge_scattered_blocks_to_lines(blocks)
    for ml in merged_lines:
        if ml not in lines:
            lines.append(ml)

    # 2) 日付を探す（どこか1行にあることが多い）
    y = m = d = None
    dow = ""
    date_line = ""
    for l in lines:
        mm = DATE_RE.search(l)
        if mm:
            y, m, d = mm.group(1), mm.group(2), mm.group(3)
            dow = (mm.group(4) or "").strip()
            date_line = l
            break

    # 2b) split-date fallback: 数字が年月日のテンプレートと別テキストに分離しているケース
    if not (y and m and d):
        for l in lines:
            mm = SPLIT_DATE_RE.search(l)
            if mm:
                m, d, y = mm.group(1), mm.group(2), mm.group(3)
                dow = (mm.group(4) or "").strip()
                date_line = l
                break

    # 3) 時間を探す（別行のことが多いので全行から探す）
    # 「1回目：12:30～13:00  2回目：13:10～13:40」のような複数回パターンを優先
    multi_session_re = re.compile(r"(\d+回目\s*[：:]\s*\d{1,2}[:：]\d{2}\s*[～〜~\-ー]\s*\d{1,2}[:：]\d{2})")
    t0 = t1 = ""
    multi_time = ""
    for l in lines:
        sessions = multi_session_re.findall(_norm_time2(l))
        if len(sessions) >= 2:
            # 複数回パターン検出 → そのまま結合して返す
            multi_time = " ".join(normalize_space(s) for s in sessions)
            break
    if not multi_time:
        for l in lines:
            if "日時" in l and TIME_RE2.search(_norm_time2(l)) is None:
                # 「日時：2026年5月...」みたいに日付専用行もあるのでスルー
                pass
            mm = TIME_RE2.search(_norm_time2(l))
            if mm:
                t0, t1 = _norm_time2(mm.group(1)), _norm_time2(mm.group(2))
                break

    # 4) 「日時:」行のフォールバック（文字列をそのまま返す用途）
    if not (y and m and d) and any("日時" in l for l in lines):
        for l in lines:
            if "日時" in l:
                s = normalize_space(l)
                s = re.sub(r"^.*日時\s*[:：]?\s*", "", s).strip()
                return s

    if not (y and m and d):
        return ""

    # 5) 表示文字列を組み立て（時間が無いなら日付だけ）
    date_str = f"{int(y)}年{int(m)}月{int(d)}日"
    if dow:
        date_str += f"（{dow}）"
    if multi_time:
        return f"{date_str} {multi_time}"
    if t0 and t1:
        return f"{date_str} {t0}～{t1}"
    return date_str


def _normalize_organizer(s: str) -> str:
    """organizer テキストを正規化: ラベル内スペース除去 + コロン前後整形 + MSD追加"""
    if not s:
        return s
    # "共 催" → "共催" 等: ラベル内のスペースを除去
    s = re.sub(r"^(主)\s*(催)", r"\1\2", s)
    s = re.sub(r"^(共)\s*(催)", r"\1\2", s)
    s = re.sub(r"^(提)\s*(供)", r"\1\2", s)
    s = re.sub(r"^(企)\s*(画)", r"\1\2", s)
    s = re.sub(r"^(運)\s*(営)", r"\1\2", s)
    # "：" 前後の無駄スペース除去 → "共催：XXX"
    s = re.sub(r"\s*[:：]\s*", "：", s)
    # MSD株式会社 が含まれていなければ末尾に追加
    s_ns = re.sub(r'[\s\u3000]+', '', s)
    if "MSD" not in s_ns and "MSD" not in s:
        s = s.rstrip() + " MSD株式会社"
    return s


def extract_organizer_from_blocks(blocks: List[TextBlock]) -> str:
    lines = blocks_to_lines(blocks)

    # ラベル行をそのまま返す（主催/共催/提供/企画/運営）
    # "共 催" のようにラベル内にスペースが入るケースに対応
    pat = re.compile(r"^(主\s*催|共\s*催|提\s*供|企\s*画|運\s*営)\s*[:：]\s*(.+)$")
    for l in lines:
        s = normalize_space(l)
        m = pat.match(s)
        if m:
            return _normalize_organizer(s)

    # ★ 正解DBの主催者辞書と照合（ラベルなし行でも検出）
    org_dict = _get_organizer_dict_cache()
    if org_dict:
        for l in lines:
            s = normalize_space(l)
            s_ns = re.sub(r'[\s\u3000]+', '', s)
            if len(s_ns) < 4:
                continue
            if s_ns in org_dict:
                return _normalize_organizer(s)
            for known in org_dict:
                if len(known) >= 6 and known in s_ns:
                    return _normalize_organizer(s)

    # fallback（会社名っぽい行）
    corp_pat = re.compile(r"(株式会社|有限会社|合同会社|Inc\.|LLC|Ltd\.|Co\.,?\s*Ltd\.|GmbH)")
    for l in reversed(lines):
        s = normalize_space(l)
        if corp_pat.search(s):
            return _normalize_organizer(s)
    return ""





def extract_time_candidates_from_blocks(blocks: List[TextBlock]) -> List[str]:
    out: List[str] = []
    seen = set()
    for b in blocks:
        m = TIME_PAT.search(b.text.replace("\n", " "))
        if m:
            t = normalize_space(m.group(1))
            if t not in seen:
                out.append(t)
                seen.add(t)
    return out

def split_name_affil_inline(text: str) -> tuple[str, str]:
    """
    例:
    '永井　英明　先生　（ Web 講演） 独立行政法人...感染症センター長'
    '國島広之先生 （聖マリアンナ医科大学 感染症学講座 主任教授）'
    -> ('永井英明', '独立行政法人...感染症センター長')
    -> ('國島広之', '聖マリアンナ医科大学 感染症学講座 主任教授')
    """
    s = normalize_space(text)
    if "先生" not in s:
        return "", ""
    
    # パターン1: 括弧内に所属がある場合「名前先生（所属）」
    m1 = re.search(r"(.+?)\s*先生\s*（([^）]+)）", s)
    if m1:
        raw_name = normalize_space(m1.group(1))
        raw_aff = normalize_space(m1.group(2))
        key = norm_name(raw_name)
        clean_aff = normalize_affiliation(raw_aff)

        return key, clean_aff
    
    # パターン2: 先生より前を名前、後ろを所属（括弧の補足は捨てる）
    m2 = re.search(r"(.+?)\s*先生\s*(?:（.*?）\s*)?(.*)$", s)
    if m2:
        raw_name = normalize_space(m2.group(1))
        raw_aff = normalize_space(m2.group(2))
        key = norm_name(raw_name)
        clean_aff = normalize_affiliation(raw_aff)

        return key, clean_aff
    
    return "", ""

def is_overall_datetime(tm: str, overall: str) -> bool:
    if not tm or not overall:
        return False
    return normalize_space(tm) in normalize_space(overall)


def _ns(s: str) -> str:
    return re.sub(r"\s+", " ", str(s or "").replace("\u3000", " ")).strip()

def _is_greeting(s: str) -> bool:
    s = _ns(s)
    return any(k in s for k in ["謹啓","謹白","平素は","厚く御礼","さてこの度","ご清祥","お慶び","ご多用"])

def extract_chair_from_blocks(blocks, speaker_map):
    ordered = sorted(blocks, key=lambda b: (b.top, b.left))
    filtered = [b for b in ordered if not _is_greeting(b.text)]

    def norm_key_for_map(name_disp: str) -> str:
        return norm_name(_ns(name_disp))

    def looks_like_affil_line(s: str) -> bool:
        s = _ns(s).replace("\n", " ")
        if not s:
            return False
        if "先生" in s:
            return False
        # 時間パターンを除外（HH:MM～HH:MM形式）
        time_pattern = r'\d{1,2}[:：]\d{2}\s*[～〜~\-－]\s*\d{1,2}[:：]\d{2}'
        if re.search(time_pattern, s):
            return False
        if any(w in s for w in ["座長", "演者", "講演", "日時", "会場", "主催", "共催", "提供", "視聴", "登録", "お願い", "ご注意"]):
            return False
        kw = ["大学", "病院", "クリニック", "センター", "科", "部", "教授", "講師", "医師", "部長", "院長", "医療"]
        return any(w in s for w in kw) or len(s) >= 10

    def find_affil_right_same_row(target_block) -> str:
        # ★同じ高さ帯の右側ブロックを拾う（今回の blocks にドンピシャ）
        cand = []
        for b in filtered:
            if b is target_block:
                continue
            if b.left <= target_block.left:
                continue
            # 高さ帯が近い
            if abs(b.top - target_block.top) > 450000:
                continue
            s = _ns(b.text)
            if not looks_like_affil_line(s):
                continue
            dx = b.left - target_block.left
            dy = abs(b.top - target_block.top)
            score = dx + dy * 0.3
            cand.append((score, s))
        cand.sort(key=lambda x: x[0])
        return cand[0][1] if cand else ""

    def pick_affil_near_lines(lines, i, fallback_key):
        for j in range(i + 1, min(i + 5, len(lines))):
            if "【講演" in lines[j] or "講演" == lines[j].replace(" ", ""):
                break
            if looks_like_affil_line(lines[j]):
                return _ns(lines[j])
        return (speaker_map.get(fallback_key) or "").strip()

    # (1) 最優先：座長：◯◯先生 が同一ブロックにある
    for b in filtered:
        t = _ns(b.text)
        if "座長" in t and "先生" in t:
            m = re.search(r"座長[：:\s]*([^\n]+?)\s*先生", t)
            if m:
                name_disp = _ns(m.group(1))
                key = norm_key_for_map(name_disp)
                # ★まず横（右側）→ ダメなら speaker_map
                aff = find_affil_right_same_row(b) or (speaker_map.get(key) or "").strip()
                return {"name": key, "name_display": name_disp, "affiliation": aff}

    # (2) 次：座長 ラベル単独
    chair_labels = []
    for b in filtered:
        t = _ns(b.text).replace("：", "").replace(":", "").replace(" ", "")
        if t == "座長":
            chair_labels.append(b)

    if chair_labels:
        lbl = sorted(chair_labels, key=lambda b: (b.top, b.left))[0]
        x0 = lbl.left - 200000
        x1 = lbl.left + 6500000
        y0 = lbl.top - 200000
        y1 = lbl.top + 1600000

        cands = []
        for b in filtered:
            if b is lbl:
                continue
            if not in_region(b, x0, y0, x1, y1):
                continue
            if "先生" not in (b.text or ""):
                continue
            if "【講演" in (b.text or ""):
                continue
            cands.append(b)

        cands.sort(key=lambda b: (abs(b.top - lbl.top) + abs(b.left - lbl.left), b.top, b.left))

        for b in cands[:5]:
            lines = [_ns(x) for x in str(b.text).split("\n") if _ns(x)]
            for i, line in enumerate(lines):
                if "先生" in line:
                    name_disp = _ns(line.replace("先生", ""))
                    key = norm_key_for_map(name_disp)
                    if not key:
                        continue
                    aff = pick_affil_near_lines(lines, i, key)
                    return {"name": key, "name_display": name_disp, "affiliation": aff}

    # (3) fallback：巨大ブロックの「講演より前の最初の先生」
    has_chair_label = any(_ns(b.text).replace("：", "").replace(":", "").replace(" ", "") == "座長" for b in filtered)
    if has_chair_label:
        bigs = [b for b in filtered if "先生" in (b.text or "")]
        bigs.sort(key=lambda b: (-(b.width * b.height), b.top, b.left))
        for b in bigs[:3]:
            lines = [_ns(x) for x in str(b.text).split("\n") if _ns(x)]
            for i, line in enumerate(lines):
                if "【講演" in line:
                    break
                if "先生" in line:
                    name_disp = _ns(line.replace("先生", ""))
                    key = norm_key_for_map(name_disp)
                    if not key:
                        continue
                    aff = pick_affil_near_lines(lines, i, key)
                    return {"name": key, "name_display": name_disp, "affiliation": aff}

    # (4) 最後：上から最初の先生
    for b in filtered:
        if "先生" not in (b.text or ""):
            continue
        lines = [_ns(x) for x in str(b.text).split("\n") if _ns(x)]
        for i, line in enumerate(lines):
            if "先生" in line and "【講演" not in line:
                name_disp = _ns(line.replace("先生", ""))
                key = norm_key_for_map(name_disp)
                if key:
                    aff = pick_affil_near_lines(lines, i, key)
                    return {"name": key, "name_display": name_disp, "affiliation": aff}

    return None

def ensure_display_fields(payload: DesignJSON) -> DesignJSON:
    def _compact_person(value: str) -> str:
        return norm_name(value or "").replace("先生", "")

    # chair
    if getattr(payload, "chair", None):
        c = payload.chair
        name = (getattr(c, "name", "") or "").strip()
        display = (getattr(c, "name_display", "") or "").strip()
        if name and (not display or _compact_person(display) != _compact_person(name)):
            c.name_display = build_speaker_display(c.name) or c.name

    # talks
    for t in (payload.talks or []):
        if _is_program_chair_item(t):
            name_display = (getattr(t, "name_display", "") or "").strip()
            if name_display:
                t.name_display = build_speaker_display(name_display) or name_display
            continue

        # speaker_display を必ず作る（speaker優先）
        sp = (getattr(t, "speaker", "") or "").strip()
        disp = (getattr(t, "speaker_display", "") or "").strip()
        if sp and (not disp or _compact_person(disp) != _compact_person(sp)):
            t.speaker_display = build_speaker_display(sp) or sp

        # speaker が空で display だけある場合は speaker を作る（逆補完）
        disp = (getattr(t, "speaker_display", "") or "").strip()
        if (not sp) and disp:
            t.speaker = norm_name(disp) or disp.replace(" ", "").replace("\u3000", "")

    return payload

def ensure_display_fields_in_dict(data: dict) -> dict:
    try:
        payload = DesignJSON(**(data or {}))
        payload = ensure_display_fields(payload)
        return (
            payload.model_dump(exclude_none=True)
            if hasattr(payload, "model_dump")
            else json.loads(payload.json(ensure_ascii=False))
        )
    except Exception:
        return data

def is_valid_person_name(name: str) -> bool:
    """有効な人名かどうかをチェック（トップレベル関数）"""
    if not name or len(name) < 2:
        return False

    # 正解DBに登録済みの人名なら即 True（旧字体・珍しい姓も対応）
    _name_ns = name.replace(" ", "").replace("\u3000", "")
    if _name_ns and _name_ns in _get_person_name_dict_cache():
        return True

    # 明らかに人名ではない単語を除外
    invalid_words = {
        "男子", "女子", "学生", "医師", "看護", "患者", "症例", 
        "治療", "診断", "手術", "検査", "薬剤", "病院", "クリニック",
        "大学", "学会", "講演", "演題", "座長", "司会", "質問", "回答",
        "時間", "場所", "会場", "参加", "登録", "視聴", "配信",
        "キーワード", "ポイント", "重要", "注意", "対象", "方法", "結果", "考察",
        "14歳", "15歳", "16歳", "17歳", "18歳", "19歳", "20歳", "30歳", "40歳",
        "年齢", "代表", "担当", "責任", "監修", "編集", "翻訳", "協力", "支援"
    }
    
    # 正規化したname（すべて小文字）で比較
    normalized_name = normalize_key(name).lower()
    if normalized_name in invalid_words:
        return False
        
    # 数字のみの文字列は人名ではない
    if re.match(r'^[\d\s]+$', name):
        return False
        
    # 長すぎる名前は所属情報の可能性
    if len(name) > 15:
        return False
        
    # アルファベットのみの短い文字列は除外
    if re.match(r'^[A-Za-z\s]{1,3}$', name):
        return False
        
    return True


def enrich_speaker_map_with_vm(
    speaker_map: Dict[str, str],
    blocks: List[TextBlock],
    vm_rows: List[dict],
) -> Dict[str, str]:
    """
    VMの医師名をblocksから直接検索し、speaker_mapに追加・補強する。
    - VMに載っている演者名がblocks内に存在すれば speaker_map に追加
    - blocks内の位置情報から所属(affiliation)も上下近傍ブロックを探索して取得
    - 既存の speaker_map エントリが所属空の場合も補完する
    """
    if not vm_rows or not blocks:
        return speaker_map

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    # VM から演者名リストを抽出（役職=="演者" のみ）
    vm_speakers: List[dict] = []
    for r in vm_rows:
        name_raw = (r.get("案内状掲載 医師名") or "").strip()
        if not name_raw:
            continue
        name_norm = _norm_person_name(name_raw)
        if not name_norm or len(name_norm) < 2:
            continue
        vm_speakers.append({
            "name_norm": name_norm,
            "name_raw": name_raw,
            "facility": (r.get("案内状掲載 施設名") or "").strip(),
            "dept": (r.get("案内状掲載 所属科") or "").strip(),
            "role_title": (r.get("案内状掲載 役職") or "").strip(),
            "is_speaker": (r.get("役職") or "") == "演者",
        })

    if not vm_speakers:
        return speaker_map

    def _find_name_in_blocks(name_norm: str) -> Optional[TextBlock]:
        """blocks内にVM名がテキストとして存在するか検索"""
        for b in ordered:
            txt_ns = b.text.replace(" ", "").replace("\u3000", "").replace("先生", "")
            if name_norm in txt_ns:
                return b
        return None

    def _extract_affil_near_block(anchor: TextBlock, facility_hint: str = "") -> str:
        """anchor(名前ブロック)の近傍から所属情報を抽出"""
        # 上方向を優先的に探索（医療セミナーでは name の上に affiliation が多い）
        cands_above: List[tuple] = []
        cands_below: List[tuple] = []

        for b in ordered:
            if b is anchor:
                continue
            txt = normalize_space(b.text.replace("\n", " "))
            if not txt or len(txt) < 3:
                continue

            # 人名単体は除外
            if txt.endswith("先生") and len(txt) <= 16:
                continue
            # ラベル行は除外
            txt_key = normalize_key(txt)
            if txt_key in {"演者", "座長", "講演1", "講演2", "講演3", "講演4"}:
                continue
            if any(kw in txt_key for kw in ["座長", "演者", "主催", "共催", "日時", "会場", "形式", "視聴", "登録"]):
                continue

            # 所属っぽいか確認
            is_aff = any(k in txt for k in [
                "大学", "病院", "クリニック", "センター", "内科", "外科",
                "部", "科", "教授", "准教授", "講師", "助教", "医長",
                "部長", "院長", "研究", "機構", "医院", "診療",
            ])
            if not is_aff:
                continue

            # 横方向の距離制限
            dx = abs((b.left + b.width / 2) - (anchor.left + anchor.width / 2))
            if dx > 5000000:
                continue

            dy = anchor.top - b.top  # 正 = bがanchorの上
            dist_y = abs(dy)
            if dist_y > 2000000:
                continue

            # facility_hint があれば優先度を上げる
            boost = 0
            if facility_hint:
                fac_ns = facility_hint.replace(" ", "").replace("\u3000", "")
                txt_ns = txt.replace(" ", "").replace("\u3000", "")
                if fac_ns in txt_ns:
                    boost = -10000000  # 最優先

            if dy > 0:
                # bはanchorの上
                cands_above.append((dist_y + boost, txt))
            else:
                # bはanchorの下
                cands_below.append((dist_y + boost, txt))

        # 上方向を優先（ただし近いものが一番）
        cands_above.sort(key=lambda x: x[0])
        cands_below.sort(key=lambda x: x[0])

        # 上にあるaffが1200000以内なら優先
        if cands_above and cands_above[0][0] < 1200000:
            return cands_above[0][1]
        # 下にあるaff
        if cands_below and cands_below[0][0] < 1200000:
            return cands_below[0][1]
        # どちらかあれば返す
        if cands_above:
            return cands_above[0][1]
        if cands_below:
            return cands_below[0][1]
        return ""

    def _clean_affil_text(s: str) -> str:
        """所属テキストから演題名（「...」）部分を除去"""
        s = re.sub(r'\s*[「「][^」」]*[」」]?\s*$', '', s).strip()
        return normalize_space(s) if s else ""

    added = 0
    enriched = 0
    for vs in vm_speakers:
        name_norm = vs["name_norm"]
        name_block = _find_name_in_blocks(name_norm)
        if not name_block:
            continue

        existing_aff = speaker_map.get(name_norm, "")

        if name_norm not in speaker_map:
            # speaker_map に無い → 新規追加
            aff = _extract_affil_near_block(name_block, vs["facility"])
            speaker_map[name_norm] = _clean_affil_text(aff) if aff else ""
            added += 1
            print(f"[VM→speaker_map] 追加: {name_norm} aff='{speaker_map[name_norm][:40]}'")
        elif not existing_aff:
            # speaker_map にあるがaff空 → 補完
            aff = _extract_affil_near_block(name_block, vs["facility"])
            if aff:
                speaker_map[name_norm] = _clean_affil_text(aff)
                enriched += 1
                print(f"[VM→speaker_map] 補完: {name_norm} aff='{speaker_map[name_norm][:40]}'")

    if added or enriched:
        print(f"[VM→speaker_map] 合計: 追加={added}, 補完={enriched}")

    return speaker_map


def extract_speaker_affil_map_by_blocks(blocks: List[TextBlock]) -> Dict[str, str]:
    mp: Dict[str, str] = {}

    # ① 同一ブロック内（先生 + 所属）優先（既存のまま）
    for b in blocks:
        key, aff = split_name_affil_inline(b.text)
        if key:
            mp.setdefault(key, "")
        if key and aff:
            mp[key] = aff

    def looks_like_affil(s: str) -> bool:
        s = normalize_space(s)
        if not s:
            return False
        # 正解DBから学習した役割を優先
        learned = lookup_text_role(s)
        if learned == "affiliation":
            return True
        if learned in ("person_name", "event_title", "talk_title"):
            return False
        if s.startswith("※"):
            return False
        if "先生" in s:
            return False
        k = normalize_key(s)
        if any(x in k for x in ["座長","演者","講演","日時","会場","主催","共催","提供","企画","運営","登録","視聴"]):
            return False
        kw = ["大学", "病院", "クリニック", "センター", "内科", "外科", "部", "科",
              "講師", "教授", "医師", "部長", "院長", "研究", "機構"]
        return any(w in s for w in kw) or len(s) >= 10

    # 座長エリアにあるブロックを除外する関数
    def is_block_in_chair_area(block: TextBlock) -> bool:
        """ブロックが座長エリア（座長ラベルより上）にあるかどうかをチェック"""
        for b in blocks:
            if "座長" in normalize_key(b.text):
                # 座長ラベル近傍だが、座長ラベルより上のエリアのみ
                chair_x0 = b.left - 200000
                chair_x1 = b.left + 5000000
                chair_y0 = b.top - 200000
                chair_y1 = b.top - 50000  # 座長ラベルより上のみ（修正）
                
                if in_region(block, chair_x0, chair_y0, chair_x1, chair_y1):
                    return True
        return False

    # 座長エリア以外の「先生」ブロックのみを対象とする
    name_blocks = [
        b for b in blocks 
        if ("先生" in b.text and "座長" not in normalize_key(b.text) and "演者" not in normalize_key(b.text))
        and not is_block_in_chair_area(b)
    ]

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    for nb in name_blocks:
        raw = normalize_space(nb.text.replace("先生", ""))
        key = norm_name(raw)
        
        # 人名として有効かチェック
        if not key or not is_valid_person_name(key):
            continue
            
        mp.setdefault(key, "")  # ★所属が見つからなくてもキーだけ作る（後段の照合が安定）

        # ★まず「直下」を最優先（このテンプレで一番多い）
        below = []
        x0 = nb.left - 400000
        x1 = nb.left + nb.width + 400000
        y0 = nb.top + nb.height - 100000
        y1 = nb.top + nb.height + 1200000
        for b in ordered:
            if b is nb:
                continue
            if not in_region(b, x0, y0, x1, y1):
                continue
            s = normalize_space(b.text.replace("\n", " "))
            if looks_like_affil(s):
                # 所属から肩書き除去を適用
                clean_s = normalize_affiliation(s)
                # 所属情報から人名を除去（座長名の重複を防ぐ）
                clean_s = _remove_person_names_from_affiliation(clean_s, key)
                below.append((abs(b.top - nb.top), clean_s))
        below.sort(key=lambda x: x[0])
        if below:
            mp[key] = below[0][1]
            continue

        # ★ レイアウトパターンキャッシュから「上方向」を検索（学習済みパターンが上を示す場合）
        lp_cache = _get_layout_pattern_cache()
        lp_affil = lp_cache.get("affil_rel_to_speaker_y", {})
        if lp_affil.get("above_ratio", 0) > 0.3 and lp_affil.get("count", 0) >= 3:
            above = []
            x0a = nb.left - 400000
            x1a = nb.left + nb.width + 400000
            y0a = nb.top - 1200000
            y1a = nb.top + 100000
            for b in ordered:
                if b is nb:
                    continue
                if not in_region(b, x0a, y0a, x1a, y1a):
                    continue
                s = normalize_space(b.text.replace("\n", " "))
                if looks_like_affil(s):
                    clean_s = normalize_affiliation(s)
                    clean_s = _remove_person_names_from_affiliation(clean_s, key)
                    dy = abs(b.top - nb.top)
                    # 学習済みmedianに近いほどスコア優遇
                    median_y = lp_affil.get("median", -300000)
                    rel_y = b.top - nb.top
                    dist_from_median = abs(rel_y - median_y)
                    above.append((dist_from_median, clean_s))
            above.sort(key=lambda x: x[0])
            if above:
                mp[key] = above[0][1]
                continue

        # ★次に「右+下」広め（2カラム/右寄せ対策）— レイアウトパターンで上方向も含める
        cand = []
        x0 = nb.left - 200000
        x1 = nb.left + 6500000
        lp_above_ratio = lp_affil.get("above_ratio", 0) if lp_affil else 0
        y0 = nb.top - (1200000 if lp_above_ratio > 0.3 else 200000)
        y1 = nb.top + 1800000

        for b in ordered:
            if b is nb:
                continue
            if not in_region(b, x0, y0, x1, y1):
                continue
            s = normalize_space(b.text.replace("\n", " "))
            if not looks_like_affil(s):
                continue

            # 所属から肩書き除去を適用
            clean_s = normalize_affiliation(s)
            # 所属情報から人名を除去（座長名の重複を防ぐ）
            clean_s = _remove_person_names_from_affiliation(clean_s, key)

            cx = b.left + b.width / 2.0
            cy = b.top + b.height / 2.0
            nx = nb.left + nb.width / 2.0
            ny = nb.top + nb.height / 2.0

            # ★「下方向」を強く優遇（所属は下に来ることが多い）
            # レイアウトパターンがある場合はmedianに近いブロックを優遇
            dy = max(0, cy - ny)
            dx = abs(cx - nx)
            dist = dx + dy * 0.6  # 下を優遇
            if lp_affil.get("count", 0) >= 3:
                rel_y = b.top - nb.top
                median_y = lp_affil.get("median", 0)
                # medianに近いほどボーナス（最大50%減）
                max_range = abs(lp_affil.get("q75", 0) - lp_affil.get("q25", 0)) or 1000000
                closeness = 1.0 - min(abs(rel_y - median_y) / max_range, 1.0)
                dist *= (1.0 - closeness * 0.5)

            cand.append((dist, clean_s))

        cand.sort(key=lambda x: x[0])
        if cand:
            mp[key] = cand[0][1]

    # ★ 正解DBの人名辞書で「先生なし」ブロックも候補に追加
    person_dict = _get_person_name_dict_cache()
    if person_dict:
        for b in ordered:
            if "先生" in b.text:
                continue
            bt_ns = normalize_space(b.text).replace(" ", "").replace("\u3000", "")
            if not bt_ns or bt_ns not in person_dict:
                continue
            # DB登録済み人名と一致 → name_blocksと同様に所属を探索
            key = bt_ns
            if key in mp:
                continue
            if not is_valid_person_name(key):
                continue
            mp[key] = ""
            # 直下優先探索
            x0 = b.left - 400000
            x1 = b.left + b.width + 400000
            y0 = b.top + b.height - 100000
            y1 = b.top + b.height + 1200000
            for ab in ordered:
                if ab is b:
                    continue
                if not in_region(ab, x0, y0, x1, y1):
                    continue
                s = normalize_space(ab.text.replace("\n", " "))
                if looks_like_affil(s):
                    clean_s = normalize_affiliation(s)
                    clean_s = _remove_person_names_from_affiliation(clean_s, key)
                    mp[key] = clean_s
                    break

    # ★ 学習済み所属フォーマットを適用（スペース位置の正規化）
    aff_cache = _get_affiliation_format_cache()
    if aff_cache:
        for name_key, aff_val in mp.items():
            if not aff_val:
                continue
            ck = _aff_cache_key(aff_val)
            if ck and ck in aff_cache:
                learned = aff_cache[ck]
                if learned != aff_val:
                    mp[name_key] = learned

    return mp


def _remove_person_names_from_affiliation(affiliation: str, person_name: str) -> str:
    """所属情報から人名を除去する"""
    if not affiliation or not person_name:
        return affiliation
    
    # 人名を正規化（スペースや全角スペース除去）
    normalized_person_name = person_name.replace(" ", "").replace("\u3000", "")
    
    # 対象人名のバリエーションを生成（空白の有無を考慮）
    name_variants = []
    
    # 基本の名前パターン
    base_patterns = [
        person_name,
        normalized_person_name,
        person_name.replace(" ", "\u3000"),  # 全角スペース版
    ]
    
    # 各パターンに「先生」を追加したバリエーション
    for pattern in base_patterns:
        if pattern:
            name_variants.extend([
                pattern,
                f"{pattern} 先生",
                f"{pattern}\u3000先生",
                f"{pattern}先生"
            ])
    
    # 名前を構成する各文字コンポーネントも追加
    if len(normalized_person_name) >= 4:  # 姓名が2文字ずつの場合など
        # 可能な姓名分割パターンを試す
        for i in range(2, len(normalized_person_name)-1):
            surname = normalized_person_name[:i]
            given_name = normalized_person_name[i:]
            if len(surname) >= 2 and len(given_name) >= 2:
                # 姓 名の形で含まれる場合
                name_variants.extend([
                    f"{surname} {given_name}",
                    f"{surname}\u3000{given_name}",
                    f"{surname} {given_name} 先生",
                    f"{surname}\u3000{given_name} 先生",
                ])
    
    cleaned = affiliation
    
    # 各バリエーションで除去処理
    for variant in name_variants:
        if not variant:
            continue
        
        # より厳密なパターンマッチング
        patterns = [
            # 前後に区切り文字がある場合（最も安全）
            f"\\s+{re.escape(variant)}\\s+",
            f"^{re.escape(variant)}\\s+",
            f"\\s+{re.escape(variant)}$",
            f"^{re.escape(variant)}$",
            # 先生付きパターン
            f"\\s*{re.escape(variant)}\\s*先生\\s*",
            # より広いパターン（他の文字に続く場合）
            f"([^一-龠々ぁ-ゟァ-ヶ]){re.escape(variant)}([^一-龠々ぁ-ゟァ-ヶ])",
        ]
        
        for pattern in patterns:
            if "([^" in pattern:  # グループパターンの場合
                cleaned = re.sub(pattern, r"\1\2", cleaned)
            else:
                cleaned = re.sub(pattern, " ", cleaned)
    
    # 職位＋人名パターンも除去（「会長 田中太郎」など）
    position_name_patterns = [
        f"(院長|部長|科長|センター長|主任|会長|副会長|理事長|教授|准教授|講師|助教)\\s*{re.escape(normalized_person_name)}",
        f"(院長|部長|科長|センター長|主任|会長|副会長|理事長|教授|准教授|講師|助教)\\s*{re.escape(person_name)}",
    ]
    
    for pattern in position_name_patterns:
        cleaned = re.sub(pattern, r"\1", cleaned)
    
    # 連続する空白・記号を整理
    cleaned = re.sub(r'[\s\u3000]+', ' ', cleaned)  # 全角・半角スペース統一
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()  # 連続空白除去
    
    # 末尾に残った「先生」のみの場合も除去
    cleaned = re.sub(r'\s*先生\s*$', '', cleaned).strip()
    
    # 意味のない短すぎる結果は空文字に
    if len(cleaned.replace(' ', '')) < 3 or cleaned in ['先生', '院長', '部長', '会長', '理事長']:
        return ""
    
    return cleaned

def extract_chair_by_blocks(blocks: List[TextBlock], speaker_map: Dict[str, str], heading_words=None, debug=False) -> Chair:
    """
    「座長」ラベル近傍から座長情報を抽出する
    - 座長ラベルより上のブロックを優先検索（レイアウト修正）
    - 「先生」付きテキストから名前部分を抽出
    - speaker_mapから対応する所属情報を取得
    """
    if heading_words is None:
        heading_words = {"PROGRAM", "P R O G R A M", "AGENDA", "SCHEDULE", "TIME TABLE", "タイムテーブル", "プログラム"}
    if speaker_map is None:
        speaker_map = {}

    chair_anchor = None
    _chair_label_words = _get_chair_label_words()
    for b in blocks:
        b_key = normalize_key(b.text)
        if any(lbl in b_key for lbl in _chair_label_words):
            chair_anchor = b
            break
    if not chair_anchor:
        return Chair()

    # 座長ラベル近傍（下方向を優先）
    x0 = chair_anchor.left - 200000
    x1 = chair_anchor.left + 5000000
    y0 = chair_anchor.top - 200000
    y1 = chair_anchor.top + 1200000

    near = [b for b in blocks if in_region(b, x0, y0, x1, y1)]
    # 座長ラベルより上にあるブロックを優先（レイアウト修正）
    above_chair = [b for b in near if b.top < chair_anchor.top]
    below_chair = [b for b in near if b.top > chair_anchor.top]

    # ① まず座長ラベル「より上」のブロックから人名らしい行を探す
    for b in sorted(above_chair, key=lambda x: -x.top):  # 上から順
        text = normalize_space(b.text)
        print(f"[CHAIR DEBUG] 座長候補（上）: {b.text.strip()} (top: {b.top})")
        print(f"[CHAIR DEBUG] 処理テキスト: '{text}'")
        for line in text.split("\n"):
            line = line.strip()
            # 「先生」付きなら前方の部分を候補に
            if "先生" in line:
                name_part = line.split("先生")[0].strip()
                # スペース・記号で分割
                words = re.split(r'[\s　,、，・/／（）\(\)\[\]【】]+', name_part)
                words = [w for w in words if w]
                affiliation_words = []
                # 所属ワードリスト
                aff_words = {"大学", "病院", "センター", "外科", "内科", "教授", "医長", "診療科", "部", "科", "クリニック", "Department", "Hospital", "Center", "Clinic", "Professor"}
                # 役職ワードリスト（名前ではなく肩書き）
                role_words = {"院長", "副院長", "会長", "副会長", "部長", "副部長", "課長", "理事長", "所長", "室長", "准教授", "講師", "助教", "助手", "名誉院長"}
                # 多語（>3語）のとき: 末尾2語を名前候補として優先（"...所属... 姓 名 先生" パターン）
                if len(words) > 3:
                    last_two = " ".join(words[-2:])
                    last_two_has_aff = any(w in last_two for w in aff_words) or any(w in last_two for w in role_words)
                    if not last_two_has_aff and is_valid_person_name(last_two):
                        name_candidate = last_two
                        affiliation_words = words[:-2]
                    else:
                        # 末尾2語が名前でなければ先頭2語で試行
                        if len(words) >= 2 and (words[1] in aff_words or words[1] in role_words or any(w in words[1] for w in aff_words)):
                            name_candidate = words[0]
                            affiliation_words = words[1:]
                        else:
                            name_candidate = " ".join(words[:2])
                            affiliation_words = words[2:]
                # 2語目が所属/役職ワードなら1語目だけを人名候補に
                elif len(words) >= 2 and (words[1] in aff_words or words[1] in role_words or any(w in words[1] for w in aff_words)):
                    name_candidate = words[0]
                    affiliation_words = words[1:]
                else:
                    name_candidate = " ".join(words[:2]) if len(words) >= 2 else (words[0] if words else name_part)
                    affiliation_words = words[2:] if len(words) > 2 else []
                affiliation_candidate = " ".join(affiliation_words)
                # 見出しワードやアルファベットは除外
                if name_candidate.upper() in heading_words or re.fullmatch(r'[A-Za-z\s\.\-]+', name_candidate):
                    continue
                print(f"[CHAIR DEBUG] 先生付き: '{name_candidate}' → is_valid_person_name={is_valid_person_name(name_candidate)}")
                if is_valid_person_name(name_candidate):
                    key = norm_name(name_candidate)
                    aff = speaker_map.get(key, affiliation_candidate) if speaker_map else affiliation_candidate
                    print(f"[CHAIR DEBUG] 座長選択（上・先生付き）: {key} / aff='{aff}'")
                    return Chair(
                        name=key,
                        name_display=build_speaker_display(key),
                        affiliation=normalize_space(aff),
                    )
            # 「先生」なしでも人名らしい行を候補に
            elif is_valid_person_name(line):
                key = norm_name(line)
                aff = speaker_map.get(key, "") if speaker_map else ""
                print(f"[CHAIR DEBUG] 座長選択（上・人名判定）: {key}")
                return Chair(
                    name=key,
                    name_display=build_speaker_display(key),
                    affiliation=normalize_space(aff),
                )

    # ② 次に座長ラベル「より下」から探す
    for b in sorted(below_chair, key=lambda x: x.top):
        if "先生" in b.text:
            # テキストから名前部分だけを抽出（所属は除去）
            text = normalize_space(b.text)
            # 「先生」以降のテキストから所属候補を抽出
            after_sensei = text.split("先生", 1)[1].strip() if "先生" in text else ""
            after_sensei = re.sub(r'^[（(]\s*', '', after_sensei)
            after_sensei = re.sub(r'\s*[）)]$', '', after_sensei)
            aff_from_text = after_sensei.strip()

            # 「先生」前のテキストからも aff_words で名前と所属を分離
            name_part_raw = text.split("先生")[0].strip()
            words_below = re.split(r'[\s　,、，・/／（）\(\)\[\]【】]+', name_part_raw)
            aff_words_below = {"大学", "病院", "センター", "外科", "内科", "教授", "医長", "診療科", "部", "科", "クリニック", "Department", "Hospital", "Center", "Clinic", "Professor"}
            aff_candidate_below = ""
            name_candidate_below = name_part_raw
            if len(words_below) >= 2 and (words_below[1] in aff_words_below or any(w in words_below[1] for w in aff_words_below)):
                name_candidate_below = words_below[0]
                aff_candidate_below = " ".join(words_below[1:])

            # 所属候補: 先生前から分離した所属 > 先生後テキスト > 空
            fallback_aff = aff_candidate_below or aff_from_text

            # パターン1: "名前先生\n（所属）" または "名前先生 （所属）" 形式から名前だけ抽出
            name_match = re.match(r'^([^\n（(]+)先生[\n\s]*[（(]?', text)
            if name_match:
                name_only = name_match.group(1).strip()
                # aff_words分離が行われた場合はそちらの名前を使う
                if aff_candidate_below:
                    name_only = name_candidate_below
                key = norm_name(name_only)
                if key:
                    aff = speaker_map.get(key, fallback_aff) if speaker_map else fallback_aff
                    return Chair(
                        name=key,
                        name_display=build_speaker_display(key),
                        affiliation=normalize_space(aff),
                    )

            # パターン2: 単純に先生前の部分だけを抽出
            if "先生" in text:
                if aff_candidate_below:
                    key = norm_name(name_candidate_below)
                else:
                    name_part = text.split("先生")[0].strip()
                    key = norm_name(name_part)
                if key and len(key) >= 2:  # 名前として妥当な長さ
                    aff = speaker_map.get(key, fallback_aff) if speaker_map else fallback_aff
                    return Chair(
                        name=key,
                        name_display=build_speaker_display(key),
                        affiliation=normalize_space(aff),
                    )

    # ③ フォールバック：speaker_map のキーが含まれるか
    if speaker_map:
        joined = normalize_key("\n".join(b.text for b in near))
        for key, aff in speaker_map.items():
            if key and key in joined:
                return Chair(
                    name=key,
                    name_display=build_speaker_display(key),
                    affiliation=normalize_space(aff),
                )

    return Chair()


def pick_time(texts: List[str], time_candidates: List[str]) -> str:
    for t in texts:
        m = TIME_PAT.search((t or "").replace("\n", " "))
        if m:
            tm = normalize_space(m.group(1))
            if tm in time_candidates or not time_candidates:
                return tm
    return ""

def pick_time_from_near_texts(texts: List[str]) -> Optional[str]:
    """
    講演近傍に明示された時間のみ拾う。
    全体時間は絶対に拾わない。
    """
    for t in texts:
        m = TIME_PAT.search(t)
        if m:
            return normalize_space(m.group(1))
    return ""


def pick_speaker(texts: List[str], speaker_map: Dict[str, str]) -> tuple[str, str]:
    """
    return (speaker_key, speaker_display)
    """
    for t in texts:
        if "先生" in t:
            raw = normalize_space(t.replace("先生", ""))
            key = norm_name(raw)
            if speaker_map and key in speaker_map:
                return key, raw  # ★ raw を保持
    return "", ""


def pick_title_lines(texts: List[str]) -> List[str]:
    """
    近傍テキスト群から「演題行」を改行保持で抽出し、
    - 講演1 / 講 演１ などは完全除外
    - ~...~ は必ず別行扱い
    """
    skip_keywords = [
        "講演", "演者", "座長", "日時", "会場", "開催",
        "主催", "共催", "提供", "企画", "運営",
        "登録", "視聴", "Web", "Live"
    ]

    lines: List[str] = []

    for t in texts:
        for raw in (t or "").split("\n"):
            s = normalize_space(raw)
            if not s:
                continue

            sk = normalize_key(s)

            # ★★★ 講演1 / 講 演１ / 講演２… を完全除外 ★★★
            if re.fullmatch(r"講演[0-9１-９]+", sk):
                continue

            # ラベル系除外（正規化キーで判定）
            if any(k in sk for k in skip_keywords):
                continue

            # 時間除外
            if TIME_PAT.search(s):
                continue

            # 人名除外
            if s.endswith("先生") and len(s) <= 14:
                continue

            # 短すぎるノイズ除外
            if len(s) < 4:
                continue

            # ~...~ は別行扱い
            lines.extend(split_tilde_subtitle_lines(s))

    # 重複排除（順序維持）
    uniq, seen = [], set()
    for l in lines:
        if l not in seen:
            uniq.append(l)
            seen.add(l)

    return uniq

KANJI = r"[一-龠々]"

TITLE_WORDS = [
    "主任教授","教授","准教授","講師","助教",
    "部長","医長","院長","センター長","科長","室長",
]

ORG_WORDS = ["大学","病院","センター","クリニック","医院","診療所","機構","学部","講座","科","部","外科","内科"]

def _extract_name_anywhere(s: str) -> str:
    s = normalize_space(s or "").replace("先生", "").strip()
    if not s:
        return ""
    # 文字列中の「漢字2〜8連続」を全部拾って、組織語を含む候補を落とす
    cands = re.findall(rf"{KANJI}{{2,8}}", s)
    if not cands:
        return ""
    filtered = []
    for t in cands:
        if any(w in t for w in ORG_WORDS):  # "滋賀医科大学" みたいなのを避ける
            continue
        filtered.append(t)
    return (filtered[-1] if filtered else cands[-1]).strip()
    

def split_speaker_affiliation_fuzzy(s: str) -> tuple[str, str]:
    s = normalize_space(s or "")
    if not s:
        return "", ""

    # 末尾に「姓 名」or「姓名」が付いてるパターンが強い
    parts = s.replace("\u3000", " ").split()
    if len(parts) >= 2:
        a, b = parts[-2], parts[-1]
        if re.fullmatch(rf"{KANJI}{{1,4}}", a) and re.fullmatch(rf"{KANJI}{{1,4}}", b):
            name = (a + b)
            aff = " ".join(parts[:-2]).strip()
            return name, aff

    # 役職語で分割して右側から人名を拾う
    for w in TITLE_WORDS:
        if w in s:
            left, right = s.split(w, 1)
            left = normalize_space(left + " " + w)
            name = _extract_name_anywhere(right)
            if name:
                return name, left

    # 最後の保険：全体から名前を拾って、残りを所属扱い
    name = _extract_name_anywhere(s)
    if name:
        aff = normalize_space(s.replace(name, " ", 1))
        return name, aff

    return "", s

SPEAKER_ROLE_WORDS = [
    "教授", "特任教授", "准教授", "特任准教授",
    "講師", "特任講師", "助教", "特任助教",
    "医長", "部長", "院長", "副院長", "センター長", "主任"
]

def extract_role_name_block(text: str) -> tuple[str, str]:
    lines = [normalize_space(x) for x in str(text or "").split("\n") if normalize_space(x)]
    if not lines:
        return "", ""

    # 2行: 講師 / 森啓一郎先生
    if len(lines) >= 2 and lines[0] in SPEAKER_ROLE_WORDS and "先生" in lines[1]:
        return lines[0], normalize_space(lines[1]).replace("先生", "").strip()

    # 1行: 講師 森啓一郎先生
    one = normalize_space(" ".join(lines))
    for role in SPEAKER_ROLE_WORDS:
        m = re.match(rf"^{re.escape(role)}\s*(.+?)先生$", one)
        if m:
            return role, normalize_space(m.group(1))

    return "", ""


def append_role_to_affiliation(affiliation: str, role: str) -> str:
    aff = normalize_space(affiliation or "")
    role = normalize_space(role or "")
    if not role:
        return aff
    if not aff:
        return role
    if role in aff:
        return aff
    return f"{aff} {role}".strip()


def extract_talks_by_blocks(blocks: List[TextBlock], speaker_map: Dict[str, str], chair: "Chair | None" = None, heading_words=None, debug=False) -> List[Talk]:
    if heading_words is None:
        heading_words = {"PROGRAM", "P R O G R A M", "AGENDA", "SCHEDULE", "TIME TABLE", "タイムテーブル", "プログラム"}
    # debugは引数のデフォルト値False
    """
    まず「講演1/演題1」等のアンカーを優先。
    ただし、時間帯（HH:MM～HH:MM）が複数あるテンプレでは
    “時間行を起点” にセグメント分割して talk を構築する（EM2512対策）。
    """

    
    talks: List[Talk] = []
    if not blocks:
        return talks
    
    # speaker_mapがNoneの場合は空辞書をセット
    if speaker_map is None:
        speaker_map = {}

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))
    lines = [normalize_space(b.text) for b in ordered if normalize_space(b.text)]

    

    def is_time_line(s: str) -> str:
        s2 = normalize_time_colon(normalize_space(s))
        m = TIME_RANGE_RE.search(s2)
        return normalize_space(m.group(1)) if m else ""

    # def is_aff_line(s: str) -> bool:
    #     if not s:
    #         return False
    #     if is_time_line(s):
    #         return False
    #     return any(k in s for k in s for s in [])  # dummy to keep mypy calm (ignored)

    # ↑ 上のダミーは不要なら削除してOK。ここから本物:

    def is_label_only(s: str) -> bool:
        k = normalize_key(s or "")
        return k in {"演者", "座長", "演題", "演題演者"}
    def is_aff_line(s: str) -> bool:
        if not s:
            return False
        if is_time_line(s):
            return False

        # 役職 + 氏名ブロックは affiliation 扱いしない
        role2, name2 = extract_role_name_block(s)
        if role2 and name2:
            return False

        return any(k in s for k in [
            "病院", "クリニック", "医院", "診療所", "大学", "センター", "機構", "総合病院",
            "内科", "外科", "部", "科",
            "教授", "准教授", "講師", "医長", "部長", "院長", "主任", "理事長"
        ])

    def strip_label(prefixes, s: str) -> str:
        s2 = normalize_space(s or "")
        s2_key = normalize_key(s2)  # スペースなどを潰した比較用

        for p in prefixes:
            p_key = normalize_key(p)
            if s2_key.startswith(p_key):
                # 先頭の「演\s*者」みたいな形も含めて消す
                # p が "演者" なら ^演\s*者\s*[:：]?\s* を消す
                chars = list(p_key)  # "演者" -> ["演","者"]
                pat = r"^" + r"\s*".join(map(re.escape, chars)) + r"\s*[:：]?\s*"
                s2 = re.sub(pat, "", s2).strip()
                return s2
        return s2
    
    def _key_variants(name: str) -> List[str]:
        n = normalize_space(name or "")
        n2 = n.replace("\u3000", " ").replace(" ", "")
        return [n, n2, normalize_key(n), normalize_key(n2)]

    speaker_map_norm: Dict[str, str] = {}
    for k, v in (speaker_map or {}).items():
        for kk in _key_variants(k):
            if kk and v and kk not in speaker_map_norm:
                speaker_map_norm[kk] = v

    def aff_from_speaker_map(name: str) -> str:
        for kk in _key_variants(name):
            if kk in speaker_map_norm:
                return speaker_map_norm[kk]
        return ""

    # ------------------------------------------------------------------
    # ★ 先に time_idxs を計算（anchors を使うか判定するため）
    #    日時行（2026/03/06 ... 19:00～20:20）を time として数えない
    # ------------------------------------------------------------------
    time_idxs: List[int] = []
    for i, s in enumerate(lines):
        if looks_like_datetime_text(s):
            continue
        if is_time_line(s):
            time_idxs.append(i)

    # ------------------------------------------------------------------
    # 1) 講演アンカー（ただし time が複数あるテンプレでは使わない）
    # ------------------------------------------------------------------
    anchors = [b for b in ordered if looks_like_talk_anchor(b.text)]
    anchors = sorted(anchors, key=lambda b: (b.top, b.left))

    def near_texts(a: TextBlock):
        x0 = a.left - 200000
        x1 = a.left + 6500000
        y0 = a.top - 200000
        y1 = a.top + 2000000
        near = [b for b in ordered if in_region(b, x0, y0, x1, y1)]
        near = sorted(near, key=lambda b: (b.top, b.left))
        return near, [normalize_space(b.text) for b in near if normalize_space(b.text)]

    def pick_local_time(near_texts_list: List[str]) -> str:
        for t in near_texts_list:
            if looks_like_datetime_text(t):
                continue
            tm = is_time_line(t)
            if tm:
                return tm
        return ""

    def pick_speaker_from_texts(texts: List[str]) -> str:
        """
        テキスト群から演者名を抽出する
        - 「演者：」ラベル付きテキストを優先
        - 座長エリアの人名は演者から除外（レイアウトを考慮）
        - speaker_mapのキーから該当する演者を選択
        """
        # “演者：” 優先
        for t in texts:
            if "演者" in normalize_key(t):
                cand = strip_label(["演者", "演者:", "演者："], t)
                cand = norm_name(cand)
                if cand and is_valid_person_name(cand):
                    return cand
        # 座長エリアの人を特定して除外
        chair_area_names = set()
        
        # 座長ラベルがあるかチェック（スペース入り「座 長」にも対応）
        has_chair_label = any("座長" in normalize_key(t) for t in texts)
        if has_chair_label:
            # 座長ラベルより上にいる人の名前を特定（正しいレイアウト対応）
            chair_texts = []
            seat_label_found = False
            
            for t in texts:
                if "座長" in normalize_key(t):
                    seat_label_found = True
                    continue
                    
                # 座長ラベルより前に出現する「先生」付きテキストを座長候補とする
                if not seat_label_found and "先生" in t:
                    chair_texts.append(t)
            
            # 座長候補から名前を抽出
            for t in chair_texts:
                text_normalized = normalize_space(t)
                if "（" in text_normalized or "(" in text_normalized:
                    name_part = re.split(r'[（(]', text_normalized)[0]
                    name_part = name_part.replace("先生", "").strip()
                    if len(name_part) >= 2:
                        chair_area_names.add(norm_name(name_part))
                elif "先生" in text_normalized:
                    # 括弧なしフォーマット（「名前 先生 所属」）にも対応
                    name_part = text_normalized.split("先生")[0].strip()
                    if len(name_part) >= 2:
                        chair_area_names.add(norm_name(name_part))
        
        if chair_area_names:
            pass  # デバッグログ削除
        
        # 明示的に演者名を探す（「先生」付きテキストから）
        explicit_speakers = []
        for t in texts:
            if "先生" in t and "座長" not in normalize_key(t):
                # 先生から人名を抽出
                name_match = re.search(r'([^\n（(\s]{2,8})\s*先生', t)
                if name_match:
                    potential_name = name_match.group(1).strip()
                    # より厳密な検証を追加
                    if (is_valid_person_name(potential_name) and 
                        potential_name not in chair_area_names and
                        speaker_map and potential_name in speaker_map):
                        explicit_speakers.append(norm_name(potential_name))
        
        # 明示的な演者名が見つかった場合は優先
        for speaker in explicit_speakers:
            if speaker_map and speaker in speaker_map:
                return speaker
        
        if not speaker_map:
            return ""
        
        def is_speaker_in_title_context(speaker_name: str, texts: List[str]) -> bool:
            """講演者名がタイトル文脈内にあるかチェック"""
            title_context_indicators = [
                'キーワード', 'ここからのキーワード', '～ここからのキーワード',
                'ポイント', '対象', '重要', '注意', '方法', '結果', '考察',
                '男子', '女子', '14歳', '15歳', '16歳', '17歳', '18歳', '症例', '治療'
            ]
            
            for text in texts:
                text_normalized = normalize_space(text)
                # スピーカー名がタイトル文脈指示詞と同じテキスト内にある場合
                if speaker_name in text_normalized:
                    for indicator in title_context_indicators:
                        if indicator in text_normalized:
                            # 「～ここからのキーワードは「14歳」と「男子」」のような文脈
                            if ('キーワード' in text_normalized and 
                                ('「' in text_normalized or '『' in text_normalized)):
                                return True
                            # その他の文脈指示詞
                            elif indicator in text_normalized:
                                return True
            return False
        
        joined = normalize_key("\n".join(texts))
        for k in speaker_map.keys():
            if (k and k in joined and k not in chair_area_names and 
                is_valid_person_name(k)):  # 追加検証
                # 強化されたタイトル文脈チェック
                if not is_speaker_in_title_context(k, texts):
                    return k
        
        # フォールバック：座長エリア外の最初の人（ただし講演タイトル文脈を除外）
        if not speaker_map:
            return ""
        
        for k in speaker_map.keys():
            if (k and k not in chair_area_names and 
                is_valid_person_name(k)):  # 追加検証
                # 同様に強化されたタイトル文脈チェック
                if not is_speaker_in_title_context(k, texts):
                    return k
        return ""

    def clean_title_from_time_and_labels(text: str) -> str:
        """演題タイトルから時間情報と講演種別ラベルを除去"""
        s = normalize_space(text)
        if not s:
            return s
            
        # 時間パターンを除去 (XX:XX～XX:XX形式)
        time_pattern = r'\d{1,2}[:：]\d{2}\s*[～〜~\-－]\s*\d{1,2}[:：]\d{2}'
        s = re.sub(time_pattern, '', s)
        
        # 講演種別ラベルを除去
        talk_type_patterns = [
            r'^基調講演\s*',
            r'^特別講演\s*',
            r'^一般演題\s*',
            r'^一般講演\s*',
            r'^教育講演\s*',
            r'^招待講演\s*',
            r'^講演[０-９\d]+\s*',
            r'^演題[０-９\d]*\s*[:：]?\s*'
        ]
        
        for pattern in talk_type_patterns:
            s = re.sub(pattern, '', s)
        
        # 連続する空白を整理
        s = re.sub(r'\s+', ' ', s).strip()
        
        return s
    
    def pick_title_lines(texts: List[str]) -> List[str]:
        out: List[str] = []
        for t in texts:
            s = normalize_space(t)
            if not s:
                continue
            if looks_like_talk_anchor(s):
                continue
            if looks_like_datetime_text(s):
                continue
            if any(x in s for x in ["座長", "演者", "主催", "共催", "会場", "形式", "登録", "視聴"]):
                continue
            if is_time_line(s):
                continue
            # 「演題：」のときはラベル除去
            if "演題" in normalize_key(s):
                s = strip_label(["演題", "演題:", "演題："], s)
            
            # 時間情報と講演種別ラベルを除去
            s = clean_title_from_time_and_labels(s)
            
            # 人名単体は除外
            if s.endswith("先生") and len(s) <= 16:
                continue
            if len(s) >= 6:
                out.append(s)

        res: List[str] = []
        seen = set()
        for s in out:
            if s not in seen:
                res.append(s)
                seen.add(s)
        return res[:4]

    # ★ time が複数あるなら anchors を使わず 2) に任せる（EM2512対策）
    if anchors and len(time_idxs) < 2:
        for a in anchors[:4]:
            _, texts = near_texts(a)
            time = pick_local_time(texts)
            speaker = pick_speaker_from_texts(texts)
            aff = speaker_map.get(speaker, "") if (speaker_map and speaker) else ""
            title_lines = pick_title_lines(texts)
            if title_lines or speaker or time:
                talks.append(Talk(time=time, title_lines=title_lines, speaker=speaker, affiliation=aff))
        return talks[:4]

    # ------------------------------------------------------------------
    # 2) 時間行が複数ある場合：時間起点で分割
    # ------------------------------------------------------------------
    if len(time_idxs) >= 2:
        # 時間ブロックを収集（日時行は除外）
        time_blocks: List[tuple[TextBlock, str]] = []
        for b in ordered:
            if looks_like_datetime_text(b.text):
                continue
            tm = is_time_line(b.text)
            if tm:
                time_blocks.append((b, tm))

        time_blocks.sort(key=lambda x: (x[0].left, x[0].top))

        def _content_left_for_time(tb: TextBlock) -> int:
            """時間ブロック(tb)に紐づく“本文側”の left を推定する。"""
            y0 = tb.top - 250000
            y1 = tb.top + 900000

            cand: List[TextBlock] = []
            # まず「演題」ラベルを探す
            for b in ordered:
                if b.top < y0 or b.top > y1:
                    continue
                s = normalize_space(b.text)
                if not s:
                    continue
                if looks_like_datetime_text(s):
                    continue
                if is_time_line(s):
                    continue
                if "演題" in normalize_key(s):
                    cand.append(b)

            # 無ければ“それっぽい本文”を探す（長めで、所属/ラベル/名前ではない）
            if not cand:
                for b in ordered:
                    if b.top < y0 or b.top > y1:
                        continue
                    s = normalize_space(b.text)
                    if not s:
                        continue
                    if looks_like_datetime_text(s):
                        continue
                    if is_time_line(s) or is_aff_line(s):
                        continue
                    if any(x in normalize_key(s) for x in ["演者", "座長", "主催", "共催", "会場", "形式", "登録", "視聴"]):
                        continue
                    if len(s) >= 10:
                        cand.append(b)

            if not cand:
                return tb.left

            cand.sort(key=lambda b: (abs(b.left - tb.left), b.left))
            return cand[0].left

        # time_blocks それぞれに対して“本文側left”を計算
        time_blocks2: List[tuple[TextBlock, str, int]] = []
        for tb, tm in time_blocks:
            time_blocks2.append((tb, tm, _content_left_for_time(tb)))

        # カラム境界推定
        col_lefts: List[int] = []
        for _, _, left in time_blocks2:
            if not col_lefts:
                col_lefts.append(left)
                continue
            if min(abs(left - x) for x in col_lefts) > 900000:
                col_lefts.append(left)
        col_lefts = sorted(col_lefts)

        def col_right_bound(left: int) -> int:
            for x in col_lefts:
                if x > left + 900000:
                    return x - 300000
            return left + 6500000

        def looks_like_name_line(s: str) -> bool:
            s2 = normalize_space(s).replace("先生", "").strip()
            if TIME_RANGE_RE.search(normalize_time_colon(s2)):
                return False
            # 名前候補が取れればOK（所属が混ざってても良い）
            return bool(_extract_name_anywhere(s2))

        def is_chair_area_name(s: str) -> bool:
            """座長エリアの人名かどうかをチェック"""            
            lines = [normalize_space(x) for x in str(s or "").split("\n") if normalize_space(x)]
            # 同じブロック内に「座長」が含まれているかチェック
            for line in lines:
                if "座長" in normalize_key(line):
                    return True
            
            # 座長ラベル近傍ブロックの人名もチェック
            for b in ordered:
                if "座長" in normalize_key(b.text or ""):
                    # 座長ラベルブロック近傍かチェック
                    chair_x0 = b.left - 200000  
                    chair_x1 = b.left + 5000000
                    chair_y0 = b.top - 200000
                    chair_y1 = b.top + 1200000
                    
                    # 現在のテキストを含むブロックを探す
                    for check_b in ordered:
                        if s.strip() in (check_b.text or ""):
                            if in_region(check_b, chair_x0, chair_y0, chair_x1, chair_y1):
                                return True
            
            return False

        # 本文側leftで安定ソート
        time_blocks2.sort(key=lambda x: (x[2], x[0].top))

        used = set()

        def looks_like_affiliation(s: str) -> bool:
            s = normalize_space(s or "")
            if not s:
                return False
            # 正解DBから学習した役割を優先
            learned = lookup_text_role(s)
            if learned == "affiliation":
                return True
            if learned in ("talk_title", "event_title"):
                return False
            # 施設・所属・役職っぽい語が入ってたら「タイトル継続」ではない
            keywords = [
                "大学", "病院", "センター", "研究科", "学部", "診療科", "内科", "外科",
                "教授", "准教授", "講師", "部長", "科長", "主任", "医長",
                "先生", "MD", "PhD"
            ]
            return any(k in s for k in keywords)

        for idx_tb, (tb, tm, base_left) in enumerate(time_blocks2):
            if id(tb) in used:
                continue

            # 次の同カラム時間を探して下限にする
            next_top = None
            for j in range(idx_tb + 1, len(time_blocks2)):
                tb2, _, base_left2 = time_blocks2[j]
                if tb2.top <= tb.top:
                    continue
                if abs(base_left - base_left2) <= 900000:
                    next_top = tb2.top
                    break

            x0 = base_left - 300000
            x1 = col_right_bound(base_left)
            y0 = tb.top - 200000
            y1 = (next_top - 200000) if next_top is not None else (tb.top + 3500000)

            near = [b for b in ordered if in_region(b, x0, y0, x1, y1)]
            # 座長エリアのブロックを除外
            near = [b for b in near if not is_chair_area_name(b.text)]
            near = sorted(near, key=lambda b: (b.top, b.left))
            seg_lines = [normalize_space(b.text) for b in near if normalize_space(b.text)]

            # ★重要：同一セグメントに「次の時間行」が混ざったらそこで打ち切る
            seg2: List[str] = []
            started = False
            for s in seg_lines:
                if looks_like_datetime_text(s):
                    continue
                tm2 = is_time_line(s)
                if tm2:
                    if not started:
                        started = True
                        seg2.append(s)
                        continue
                    if tm2 != tm:
                        break
                if started:
                    seg2.append(s)
            if seg2:
                seg_lines = seg2

            title_lines: List[str] = []
            speaker = ""
            affiliation = ""
            aff_candidates: List[str] = []
            pending_role = ""

            for i, block in enumerate(blocks):
                text = block.text.strip()
                if not text:
                    continue
                # 先生が含まれる場合、その前を人名＋所属とみなす
                if '先生' in text:
                    pre, *_ = text.split('先生', 1)
                    pre = pre.strip()
                    if not pre:
                        continue
                    # 日本語の場合、空白で分割されないことが多いので、正規表現で分割
                    pre_words = re.findall(r'\S+', pre)
                    # 最初の1～2語を人名候補とする
                    if len(pre_words) >= 2:
                        # 1語目と2語目を連結（間に空白を入れる: ex. "高橋將人 北海道大学病院乳腺外科" → "高橋將人"）
                        name_candidate = pre_words[0] + pre_words[1] if len(pre_words[0]) == 1 else pre_words[0] + ' ' + pre_words[1]
                    else:
                        name_candidate = pre_words[0]
                    # 見出しワードやアルファベットのみは除外
                    if heading_words and name_candidate.upper() in heading_words:
                        if debug:
                            print(f"[DEBUG] Chair name '{name_candidate}' is heading word, skip.")
                        continue
                    if re.fullmatch(r'[A-Z ]+', name_candidate):
                        if debug:
                            print(f"[DEBUG] Chair name '{name_candidate}' is all alpha, skip.")
                        continue
                    # 人名判定（1語目＋2語目の連結 or 1語目のみ）
                    name_for_check = name_candidate.replace(' ', '')
                    if not is_valid_person_name(name_for_check):
                        if debug:
                            print(f"[DEBUG] Chair name '{name_candidate}' is not valid person name, skip.")
                        continue
                    # 所属候補
                    affiliation_candidate = ''.join(pre_words[2:]) if len(pre_words) > 2 else ''
                    continue

                if not speaker and "演者" in normalize_key(s):
                    sp = strip_label(["演者", "演者:", "演者："], s)
                    sp = normalize_space(sp)

                    # 「演者」の直後が演題なら speaker ではなく title とみなす
                    if looks_like_title_text(sp):
                        if sp and not title_lines:
                            title_lines.append(sp)
                        continue

                    # 末尾人名分離を先に試す
                    aff_tail, name_tail = split_affiliation_and_name_tail(sp)
                    if name_tail:
                        speaker = name_tail
                        if aff_tail and not affiliation:
                            affiliation = aff_tail
                        continue

                    # 既存ロジック fallback
                    sp2, aff2 = split_speaker_affiliation_fuzzy(sp)
                    if sp2 and not speaker:
                        speaker = sp2
                    if aff2 and not affiliation:
                        affiliation = aff2
                    continue

                # 「演者」ラベルが無いテンプレ用：名前っぽい行
                if not speaker and looks_like_name_line(s):
                    # 座長エリアの人名は演者から除外
                    if is_chair_area_name(s):
                        continue
                    role2, name2 = extract_role_name_block(s)
                    if role2 and name2:
                        speaker = norm_name(name2)
                        pending_role = role2
                    else:
                        sp2, aff2 = split_speaker_affiliation_fuzzy(s)
                        if sp2:
                            speaker = norm_name(sp2)
                            if not affiliation and aff2 and is_aff_line(aff2):
                                affiliation = aff2

                if is_aff_line(s):
                    # affiliation候補として積む前に、末尾人名が付いてたら分離
                    sp2, aff2 = split_speaker_affiliation_fuzzy(s)
                    if sp2 and not speaker:
                        # 座長の名前は演者として使用しない
                        sp2_key = normalize_key(sp2)
                        chair_name_key = normalize_key(chair.name if chair else "")
                        if chair_name_key and sp2_key == chair_name_key:
                            # 座長の所属情報も演者所属には使わない（混入防止）
                            continue
                        else:
                            speaker = norm_name(sp2)
                    if aff2 and sp2:
                        clean_aff = normalize_affiliation(aff2)
                        if clean_aff:
                            aff_candidates.append(clean_aff)
                    else:
                        clean_s = normalize_affiliation(s)
                        if clean_s:
                            aff_candidates.append(clean_s)

            # affiliation確定
            if not affiliation and aff_candidates:
                def aff_score(a: str) -> int:
                    score = 0
                    if any(k in a for k in ["病院", "クリニック", "大学", "センター", "総合病院", "機構"]):
                        score += 2
                    if any(k in a for k in ["内科", "外科", "科", "部"]):
                        score += 2
                    if any(k in a for k in ["教授", "准教授", "講師", "部長", "院長", "理事長", "主任", "医長"]):
                        score += 2
                    return score
                aff_candidates.sort(key=lambda a: (-aff_score(a), len(a)))
                affiliation = aff_candidates[0]

            # speaker_map から補完（欠損のみ）
            if speaker and not affiliation:
                affiliation = aff_from_speaker_map(speaker)

            if pending_role:
                affiliation = append_role_to_affiliation(affiliation, pending_role)

            print(affiliation)

            # speaker確定（座長との重複チェックは後の処理で実行）
            # この段階では座長情報がないため、単純に演者として設定
          
            # タイトル fallback: 「演題」ラベルが無いテンプレ用
            if not title_lines:
                for s in seg_lines:
                    if looks_like_datetime_text(s):
                        continue
                    if is_time_line(s):
                        continue
                    if is_aff_line(s):
                        continue
                    if "演者" in normalize_key(s) or "座長" in normalize_key(s):
                        continue
                    if looks_like_name_line(s):
                        continue
                    
                    # 時間情報と講演種別ラベルを除去してからチェック
                    cleaned_s = clean_title_from_time_and_labels(s)
                    if len(cleaned_s) >= 10:
                        title_lines.append(cleaned_s)
                        break

            talks.append(Talk(time=tm, title_lines=title_lines[:4], speaker=speaker, affiliation=affiliation))

            used.add(id(tb))
            if len(talks) >= 4:
                break

        talks = [t for t in talks if t.time or t.title_lines or t.speaker or t.affiliation]

        def _is_notice_lines(tl: List[str]) -> bool:
            if not tl:
                return False
            tl2 = [normalize_space(x) for x in tl if normalize_space(x)]
            if not tl2:
                return False
            joined = normalize_key("\n".join(tl2))
            # 箇条書き + 注意語
            if all(x.startswith("・") for x in tl2[:min(3, len(tl2))]):
                if any(k in joined for k in ["事前", "参加", "登録", "ご遠慮", "医療従事者", "資格", "担当者へご連絡"]):
                    return True
            return False

        def _is_empty_chair_only(t: Talk) -> bool:
            # タイトルなし/演者なし で affiliation だけ（=座長/施設だけ拾ったゴミ）を落とす
            if (t.speaker or "").strip():
                return False
            if (t.title_lines or []) and any(normalize_space(x) for x in t.title_lines):
                return False
            # affiliation だけがあるケース
            if (t.affiliation or "").strip():
                return True
            return False

        def _should_drop(t: Talk) -> bool:
            # 注意書きtalk
            if _is_notice_lines(t.title_lines or []):
                return True
            # speaker/title無しの “座長/施設だけ”
            if _is_empty_chair_only(t):
                return True
            return False

        talks = [t for t in talks if (t.time or t.title_lines or t.speaker or t.affiliation)]
        talks = [t for t in talks if not _should_drop(t)]
  

        return talks[:4]

    # ------------------------------------------------------------------
    # 3) 最後のfallback：演題ラベル周辺を複数件拾う（2) に入らない時の保険）
    # ------------------------------------------------------------------
    label_idxs: List[int] = []
    for i, s in enumerate(lines):
        if "演題" in normalize_key(s):
            label_idxs.append(i)

    seen = set()

    for label_idx in label_idxs[:6]:
        seg = lines[max(0, label_idx - 3): min(len(lines), label_idx + 25)]

        tm = ""
        for s in seg:
            if looks_like_datetime_text(s):
                continue
            tm2 = is_time_line(s)
            if tm2:
                tm = tm2
                break

        title_lines: List[str] = []
        speaker = ""
        affiliation = ""

        for j, s in enumerate(seg):
            k = normalize_key(s)

            if not title_lines and "演題" in k:
                t = strip_label(["演題", "演題:", "演題："], s).strip()
                if t:
                    for ln in t.split("\n"):
                        ln = normalize_space(ln)
                        if not ln or is_aff_line(ln):
                            continue
                        title_lines.append(ln)

                if j + 1 < len(seg):
                    nxt = normalize_space(seg[j + 1])
                    if nxt and (nxt.startswith("～") or nxt.startswith("~")):
                        title_lines.append(nxt)

            if not speaker and "演者" in k:
                speaker = norm_name(strip_label(["演者", "演者:", "演者："], s))

        # affiliation は time 行の次の所属っぽい行を優先
        for j, s in enumerate(seg):
            if looks_like_datetime_text(s):
                continue
            if is_time_line(s):
                for kk in range(j + 1, min(len(seg), j + 6)):
                    ss = normalize_space(seg[kk])
                    if not ss:
                        continue
                    if "演題" in normalize_key(ss) or "演者" in normalize_key(ss) or "座長" in normalize_key(ss):
                        continue
                    if is_aff_line(ss):
                        affiliation = ss
                        break
                break

        if speaker and not affiliation:
            affiliation = speaker_map.get(speaker, "") or ""

        if title_lines or speaker or tm or affiliation:
            key = (normalize_space(tm), normalize_space(speaker), join_lines(title_lines))
            if key not in seen:
                talks.append(Talk(time=tm, title_lines=title_lines[:4], speaker=speaker, affiliation=affiliation))
                seen.add(key)

        if len(talks) >= 4:
            break

    

        

    return talks[:4]


def _json_item_type(item: Any) -> str:
    if isinstance(item, dict):
        v = normalize_space(str(item.get("item_type") or "talk")).lower()
        return "chair" if v == "chair" else "talk"
    v = normalize_space(str(getattr(item, "item_type", "talk") or "talk")).lower()
    return "chair" if v == "chair" else "talk"


def _json_is_chair_item(item: Any) -> bool:
    return _json_item_type(item) == "chair"


def _json_is_talk_item(item: Any) -> bool:
    return not _json_is_chair_item(item)


def _json_get(item: Any, field: str, default: str = "") -> Any:
    if isinstance(item, dict):
        return item.get(field, default)
    return getattr(item, field, default)


def _json_person_name_value(item: Any) -> str:
    if _json_is_chair_item(item):
        fields = ("name", "name_display", "speaker_display", "speaker")
    else:
        fields = ("speaker", "speaker_display", "name", "name_display")
    for field in fields:
        val = normalize_space(str(_json_get(item, field, "") or ""))
        val = re.sub(r"\s*先生\s*$", "", val).strip()
        if val:
            return val
    return ""


def _json_person_display_value(item: Any) -> str:
    if _json_is_chair_item(item):
        fields = ("name_display", "speaker_display", "name", "speaker")
    else:
        fields = ("speaker_display", "name_display", "speaker", "name")
    for field in fields:
        val = normalize_space(str(_json_get(item, field, "") or ""))
        if val:
            return val
    return ""


def _json_person_name_key(item: Any) -> str:
    return re.sub(r"[\s\u3000]+", "", _json_person_name_value(item) or "")


def _is_program_chair_item(t: Any) -> bool:
    return _json_is_chair_item(t)


def _is_program_talk_item(t: Any) -> bool:
    return _json_is_talk_item(t)


def _looks_like_program_affiliation(s: str) -> bool:
    s = normalize_space(s or "").replace("\n", " ")
    if not s:
        return False
    k = normalize_key(s)
    if TIME_RANGE_RE.search(normalize_time_colon(s)):
        return False
    if looks_like_datetime_text(s):
        return False
    if looks_like_talk_anchor(s):
        return False
    if any(x in k for x in ["日時", "開催形式", "形式", "演者", "座長", "講演", "演題", "主催", "共催", "提供", "登録", "視聴"]):
        return False
    return any(x in s for x in [
        "大学", "病院", "センター", "クリニック", "医院", "診療所", "医療センター",
        "内科", "外科", "科", "部", "講師", "教授", "准教授", "部長", "医長", "院長",
    ])


def _strip_program_role_prefix(s: str) -> str:
    s = normalize_space(s or "").replace("\n", " ")
    return re.sub(r"^(座\s*長|演\s*者|総\s*合\s*司\s*会|司\s*会)\s*[:：]?\s*", "", s).strip()


def _split_program_person_text(text: str) -> tuple[str, str]:
    one = normalize_space(str(text or "").replace("\n", " "))
    one = _strip_program_role_prefix(one)
    if not one:
        return "", ""

    name_src = one
    aff = ""
    if "先生" in one:
        name_src, aff = one.split("先生", 1)
        aff = normalize_affiliation(aff)
        if not _looks_like_program_affiliation(aff):
            aff = ""

    name = _extract_name_anywhere(name_src) or norm_name(name_src)
    name = norm_name(name)
    return name, aff


def _affiliation_after_block(blocks: list[TextBlock], start: TextBlock, *, y_limit: int = 850000, x_limit: int = 2800000) -> str:
    candidates: list[tuple[int, str]] = []
    for b in blocks:
        if b is start:
            continue
        if b.top < start.top:
            continue
        dy = b.top - start.top
        if dy > y_limit:
            continue
        if abs(b.left - start.left) > x_limit and b.left < start.left:
            continue
        s = normalize_space((b.text or "").replace("\n", " "))
        if not _looks_like_program_affiliation(s):
            continue
        candidates.append((dy + abs(b.left - start.left) // 4, normalize_affiliation(s)))
    if not candidates:
        return ""
    candidates.sort(key=lambda x: x[0])
    return candidates[0][1]


def _parse_program_person_block(
    block: TextBlock,
    *,
    role: str,
    ordered: list[TextBlock],
    segment_bottom: int | None = None,
) -> tuple[str, str, str]:
    text = normalize_space((block.text or "").replace("\n", " "))
    role_label = role
    if role == "chair":
        role_label = "総合司会" if "総合司会" in normalize_key(text) else "司会" if normalize_key(text).startswith("司会") else "座長"
    elif role == "speaker":
        role_label = "演者"

    name, aff = _split_program_person_text(text)
    if name and not aff:
        y_limit = 700000
        if segment_bottom is not None:
            y_limit = max(250000, min(900000, segment_bottom - block.top))
        aff = _affiliation_after_block(ordered, block, y_limit=y_limit)

    return role_label, name, aff


def _extract_inline_program_from_blocks(
    blocks: list[TextBlock],
    *,
    chair: Chair | None = None,
) -> tuple[Chair | None, list[Talk]]:
    """講演番号/時間/座長/演者ラベルの縦順が明確な案内状から program items を作る。

    `talks` に item_type="chair" を混ぜるための補助抽出。通常の抽出が苦手な
    PDF（座長が講演の途中に入るレイアウト）だけで採用する。
    """
    ordered = sorted(blocks or [], key=lambda b: (b.top, b.left))
    if not ordered:
        return None, []

    anchors = [b for b in ordered if looks_like_talk_anchor(b.text or "")]
    anchors = sorted(anchors, key=lambda b: (b.top, b.left))
    if not anchors:
        return None, []

    def is_time_line_local(s: str) -> str:
        s2 = normalize_time_colon(normalize_space(s or ""))
        m = TIME_RANGE_RE.search(s2)
        if not m:
            return ""
        return f"{m.group(1)}~{m.group(2)}"

    def is_speaker_label(b: TextBlock) -> bool:
        return "演者" in normalize_key(b.text or "")

    def is_chair_label(b: TextBlock) -> bool:
        return "座長" in normalize_key(b.text or "") or "総合司会" in normalize_key(b.text or "")

    def is_non_lecture_chair_context(b: TextBlock) -> bool:
        """講演ではない短い進行項目に付く座長は、途中座長として追加しない。"""
        non_lecture_words = [
            "会のご挨拶", "ご挨拶", "挨拶", "開会", "閉会", "開会の辞", "閉会の辞",
            "休憩", "質疑応答", "Q&A", "総合討論", "総合討議", "ディスカッション",
            "事務連絡", "諸連絡", "注意事項", "ご案内",
        ]
        window_top = b.top - 900000
        window_bottom = b.top + 250000
        nearby = [
            x for x in ordered
            if x is not b and window_top <= x.top <= window_bottom
        ]
        nearby.sort(key=lambda x: (abs(x.top - b.top), abs(x.left - b.left)))
        for x in nearby[:8]:
            raw = normalize_space(x.text or "")
            if not raw or "先生" in raw:
                continue
            key = normalize_key(raw)
            if any(normalize_key(word) in key for word in non_lecture_words):
                return True
        return False

    def person_blocks_in_segment(y0: int, y1: int, *, speaker_only: bool = False) -> list[TextBlock]:
        out = []
        for b in ordered:
            if b.top < y0 or b.top >= y1:
                continue
            s = normalize_space(b.text or "")
            if "先生" not in s:
                continue
            if speaker_only and is_chair_label(b):
                continue
            out.append(b)
        return sorted(out, key=lambda b: (b.top, b.left))

    def speaker_block_for_segment(seg: list[TextBlock], y0: int, y1: int) -> TextBlock | None:
        labels = [b for b in seg if is_speaker_label(b)]
        for label in labels:
            if "先生" in (label.text or ""):
                return label
            people = person_blocks_in_segment(max(y0, label.top - 550000), min(y1, label.top + 650000), speaker_only=True)
            if people:
                people.sort(key=lambda b: (abs(b.top - label.top), abs(b.left - label.left)))
                return people[0]

        people = person_blocks_in_segment(y0, y1, speaker_only=True)
        if people:
            return people[-1]
        return None

    def title_lines_for_segment(seg: list[TextBlock], *, after_top: int, before_top: int | None) -> list[str]:
        lines: list[str] = []
        seen: set[str] = set()
        for b in seg:
            if b.top <= after_top:
                continue
            if before_top is not None and b.top >= before_top:
                continue
            raw = normalize_space(b.text or "")
            if not raw:
                continue
            if is_time_line_local(raw) or looks_like_talk_anchor(raw) or looks_like_datetime_text(raw):
                continue
            if is_chair_label(b) or is_speaker_label(b):
                continue
            if "先生" in raw or _looks_like_program_affiliation(raw):
                continue
            if (b.max_font_pt or 0) < 18 and not looks_like_title_text(raw):
                continue
            for line in str(b.text or "").split("\n"):
                s = normalize_space(line)
                if not s or s in seen:
                    continue
                lines.append(s)
                seen.add(s)
        return lines[:4]

    talk_items: list[tuple[int, Talk]] = []
    first_title_top: int | None = None

    for i, anchor in enumerate(anchors[:6]):
        next_anchor = anchors[i + 1] if i + 1 < len(anchors) else None
        seg_bottom = next_anchor.top if next_anchor else anchor.top + 3600000
        seg = [b for b in ordered if anchor.top - 250000 <= b.top < seg_bottom]

        time_text = ""
        time_top = anchor.top
        for b in seg:
            if looks_like_datetime_text(b.text or ""):
                continue
            tm = is_time_line_local(b.text or "")
            if tm:
                time_text = tm
                time_top = b.top
                break

        sp_block = speaker_block_for_segment(seg, anchor.top, seg_bottom)
        speaker_top = sp_block.top if sp_block else None
        title_lines = title_lines_for_segment(seg, after_top=time_top, before_top=speaker_top)
        if title_lines and first_title_top is None:
            first_title_top = min(b.top for b in seg if any(normalize_space(x) in normalize_space(b.text or "") for x in title_lines))

        speaker = ""
        speaker_aff = ""
        if sp_block:
            _, speaker, speaker_aff = _parse_program_person_block(
                sp_block,
                role="speaker",
                ordered=ordered,
                segment_bottom=seg_bottom,
            )

        if not (time_text or title_lines or speaker or speaker_aff):
            continue

        talk = Talk(
            item_type="talk",
            role="演者",
            program_index=len(talk_items),
            time=time_text,
            title="\n".join(title_lines),
            title_lines=title_lines,
            speaker=speaker,
            speaker_display=build_speaker_display(speaker) if speaker else "",
            affiliation=speaker_aff,
            honorific_title="先生",
        )
        setattr(talk, "_talk_index", i + 1)
        talk_items.append((anchor.top, talk))

    if not talk_items:
        return None, []

    first_anchor_top = anchors[0].top
    top_chair: Chair | None = None
    inline_chairs: list[tuple[int, Talk]] = []
    seen_chairs: set[tuple[str, int]] = set()

    for b in ordered:
        if not is_chair_label(b):
            continue
        role, name, aff = _parse_program_person_block(
            b,
            role="chair",
            ordered=ordered,
            segment_bottom=first_title_top if b.top < first_anchor_top and first_title_top else None,
        )
        if not name:
            continue
        key = (normalize_key(name), b.top)
        if key in seen_chairs:
            continue
        seen_chairs.add(key)

        if b.top < first_anchor_top:
            top_chair = Chair(
                role=role or "座長",
                name=name,
                name_display=build_speaker_display(name),
                affiliation=aff,
                honorific_title="先生",
            )
            continue

        if is_non_lecture_chair_context(b):
            continue

        item = Talk(
            item_type="chair",
            role=role or "座長",
            program_index=0,
            name_display=build_speaker_display(name),
            speaker="",
            speaker_display="",
            affiliation=aff,
            honorific_title="先生",
        )
        inline_chairs.append((b.top, item))

    combined = sorted([*talk_items, *inline_chairs], key=lambda x: x[0])
    out: list[Talk] = []
    for idx, (_, item) in enumerate(combined):
        item.program_index = idx
        setattr(item, "_talk_index", idx)
        out.append(item)

    # 途中座長が無く、既存抽出のほうが十分なら採用しないため、呼び出し側で判定する。
    return top_chair, out[:6]


def apply_inline_program_extraction(payload: DesignJSON, blocks: list[TextBlock]) -> DesignJSON:
    top_chair, items = _extract_inline_program_from_blocks(blocks, chair=getattr(payload, "chair", None))
    if not items:
        return payload

    extracted_talk_count = sum(1 for t in items if _is_program_talk_item(t))
    extracted_chair_count = sum(1 for t in items if _is_program_chair_item(t))
    current_talk_count = sum(1 for t in (payload.talks or []) if _is_program_talk_item(t))

    strong_program = extracted_talk_count >= 2 and (
        extracted_chair_count > 0
        or current_talk_count < extracted_talk_count
        or any(not _talk_title_text(t) or not getattr(t, "speaker", "") for t in (payload.talks or []) if _is_program_talk_item(t))
    )
    if not strong_program:
        return payload

    if top_chair:
        if not getattr(payload, "chair", None):
            payload.chair = top_chair
        else:
            current_name = normalize_key(getattr(payload.chair, "name", "") or "")
            top_name = normalize_key(top_chair.name or "")
            if (not current_name) or current_name == top_name:
                payload.chair.role = top_chair.role or payload.chair.role or "座長"
                payload.chair.name = top_chair.name or payload.chair.name
                payload.chair.name_display = top_chair.name_display or payload.chair.name_display
                if top_chair.affiliation:
                    payload.chair.affiliation = top_chair.affiliation
                payload.chair.honorific_title = top_chair.honorific_title or payload.chair.honorific_title or "先生"

    payload.talks = items
    warnings = set(payload.warnings or [])
    warnings.discard("talks_pruned_by_vm_hint")
    warnings.discard("talks_pruned_heuristic_only")
    if extracted_chair_count > 0:
        warnings.add("inline_chair_extracted")
    payload.warnings = sorted(warnings)
    return payload



def find_sponsor_logo_blobs(pptx_path: Path) -> list[bytes]:
    prs = Presentation(str(pptx_path))
    if len(prs.slides) == 0:
        return []

    slide = prs.slides[0]

    sponsor_text_shapes = []
    for sh in iter_shapes(slide.shapes):
        if not getattr(sh, "has_text_frame", False):
            continue
        txt = normalize_space(getattr(sh.text_frame, "text", "") or "")
        if not txt:
            continue
        if "主催" in normalize_key(txt):
            sponsor_text_shapes.append(sh)

    if not sponsor_text_shapes:
        return []

    anchor = sorted(sponsor_text_shapes, key=lambda s: int(getattr(s, "top", 0)), reverse=True)[0]

    x0 = int(anchor.left + anchor.width) - 200000
    x1 = int(anchor.left + anchor.width) + 5000000
    y0 = int(anchor.top) - 400000
    y1 = int(anchor.top + anchor.height) + 900000

    blobs: list[bytes] = []
    for sh in iter_shapes(slide.shapes):
        if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            cx = int(sh.left + sh.width / 2)
            cy = int(sh.top + sh.height / 2)
            if x0 <= cx <= x1 and y0 <= cy <= y1:
                try:
                    blobs.append(sh.image.blob)
                except Exception:
                    pass

    return blobs

async def post_with_retry(client, url, *, headers, json_body, retries=3):
    last_exc = None
    for attempt in range(retries):
        try:
            return await client.post(url, headers=headers, json=json_body)
        except httpx.ReadTimeout as e:
            last_exc = e
            if attempt == retries - 1:
                raise
            await asyncio.sleep(1.5 * (attempt + 1))
    raise last_exc

async def ocr_company_name_with_openai(image_bytes: bytes) -> str:
    if not OPENAI_API_KEY:
        return ""

    b64 = base64.b64encode(image_bytes).decode("utf-8")
    data_url = f"data:image/jpg;base64,{b64}"

    prompt = """この画像はセミナー案内の「主催」ロゴです。
ロゴから読み取れる会社名/団体名のみを日本語で1行で返してください。
不明なら空文字を返してください。余計な説明は禁止。"""

    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json",
    }

    body = {
        "model": AI_MODEL,
        "messages": [
            {"role": "system", "content": "Return only plain text. No extra words."},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            },
        ],
        "temperature": 0.0,
    }

    async with httpx.AsyncClient(timeout=AI_TIMEOUT) as client:
        r = await post_with_retry(client, f"{OPENAI_BASE_URL}/chat/completions", headers=headers, json_body=body, retries=3)
        r.raise_for_status()
        data = r.json()

    text = (data["choices"][0]["message"]["content"] or "").strip()
    text = text.splitlines()[0].strip()
    return text


async def try_fill_organizer_from_logo(pptx_path: Path) -> str:
    blobs = find_sponsor_logo_blobs(pptx_path)
    if not blobs:
        return ""

    for blob in blobs[:2]:
        name = await ocr_company_name_with_openai(blob)
        if name:
            return name
    return ""


# ---------------- 正解DB (Correct Answer Store) ----------------

def _compute_embedding(text: str) -> list[float] | None:
    """OpenAI text-embedding-3-small でベクトル化（1536次元）"""
    if not OPENAI_API_KEY or not text.strip():
        return None
    try:
        truncated = text[:8000]  # トークン制限安全マージン
        resp = requests.post(
            f"{OPENAI_BASE_URL}/embeddings",
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json",
            },
            json={"model": "text-embedding-3-small", "input": truncated},
            timeout=30,
        )
        if resp.status_code == 200:
            return resp.json()["data"][0]["embedding"]
        print(f"[embedding] API error {resp.status_code}: {resp.text[:200]}")
    except Exception as e:
        print(f"[embedding] error: {e}")
    return None


def _cosine_similarity(a: list[float], b: list[float]) -> float:
    """コサイン類似度"""
    if not a or not b or len(a) != len(b):
        return 0.0
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = sum(x * x for x in a) ** 0.5
    norm_b = sum(x * x for x in b) ** 0.5
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


# 正解DBから学習した speaker_display キャッシュ: {"name_key": "display"}
# 例: {"秋田谷一輝": "秋田谷 一輝", "下川原裕人": "下川原 裕人"}
_speaker_display_cache: dict[str, str] = {}
_speaker_display_cache_loaded: bool = False

def _build_speaker_display_cache() -> dict[str, str]:
    """正解DBの全レコードから speaker/chair の name→display マッピングを構築。
    新しいレコード(日付降順)を優先。"""
    cache: dict[str, str] = {}
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            # talks
            for t in (cj.get("talks") or []):
                sp = _json_person_name_key(t)
                disp = _json_person_display_value(t).strip()
                if sp and disp and " " in disp and sp not in cache:
                    cache[sp] = disp
            # chair
            ch = cj.get("chair") or {}
            cn = _json_person_name_key(ch)
            cd = _json_person_display_value(ch).strip()
            if cn and cd and " " in cd and cn not in cache:
                cache[cn] = cd
    except Exception as e:
        print(f"[speaker-display-cache] build error: {e}")
    if cache:
        print(f"[speaker-display-cache] loaded {len(cache)} entries")
    return cache

def _get_speaker_display_cache() -> dict[str, str]:
    global _speaker_display_cache, _speaker_display_cache_loaded
    if not _speaker_display_cache_loaded:
        _speaker_display_cache = _build_speaker_display_cache()
        _speaker_display_cache_loaded = True
    return _speaker_display_cache

def invalidate_speaker_display_cache():
    """save_correct_answer 後に呼ぶことで次回再構築される"""
    global _speaker_display_cache_loaded
    _speaker_display_cache_loaded = False


# 正解DBから学習した所属フォーマットキャッシュ: {"nospacekey": "formatted"}
# 例: {"東京大学大学院消化器内科学教授": "東京大学大学院 消化器内科学 教授"}
_affiliation_format_cache: dict[str, str] = {}
_affiliation_format_cache_loaded: bool = False

def _aff_cache_key(s: str) -> str:
    """所属テキストからスペース・改行を全除去してキーにする"""
    return re.sub(r'[\s\u3000]+', '', s or '')

def _build_affiliation_format_cache() -> dict[str, str]:
    """正解DBの全レコードから所属の「正しいスペース・改行位置」を学習。
    新しいレコード(日付降順)を優先。"""
    cache: dict[str, str] = {}
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            # talks
            for t in (cj.get("talks") or []):
                raw = (t.get("affiliation") or "").strip()
                sp = _json_person_name_value(t).strip()
                if not raw:
                    continue
                # 通常キー
                key = _aff_cache_key(raw)
                if key and len(key) >= 4 and key not in cache:
                    cache[key] = raw
                # 人名除去バージョンも登録
                if sp:
                    aff_no_name = _remove_person_names_from_affiliation(raw, sp)
                    key2 = _aff_cache_key(aff_no_name)
                    if key2 and len(key2) >= 4 and key2 not in cache:
                        cache[key2] = aff_no_name
            # chair
            ch = cj.get("chair") or {}
            raw = (ch.get("affiliation") or "").strip()
            name = _json_person_name_value(ch).strip()
            if raw:
                key = _aff_cache_key(raw)
                if key and len(key) >= 4 and key not in cache:
                    cache[key] = raw
                # 人名除去バージョンも登録
                if name:
                    aff_no_name = _remove_person_names_from_affiliation(raw, name)
                    key2 = _aff_cache_key(aff_no_name)
                    if key2 and len(key2) >= 4 and key2 not in cache:
                        cache[key2] = aff_no_name
    except Exception as e:
        print(f"[affiliation-format-cache] build error: {e}")
    if cache:
        print(f"[affiliation-format-cache] loaded {len(cache)} entries")
    return cache

def _get_affiliation_format_cache() -> dict[str, str]:
    global _affiliation_format_cache, _affiliation_format_cache_loaded
    if not _affiliation_format_cache_loaded:
        _affiliation_format_cache = _build_affiliation_format_cache()
        _affiliation_format_cache_loaded = True
    return _affiliation_format_cache

def invalidate_affiliation_format_cache():
    """save_correct_answer 後に呼ぶことで次回再構築される"""
    global _affiliation_format_cache_loaded
    _affiliation_format_cache_loaded = False


# ── テキスト→役割（フィールド）キャッシュ ──
# 正解DBから「このテキストは所属」「このテキストは演者名」等を学習
# {"normalizedtext": {"role": "affiliation", "count": 5, "formatted": "東京大学 医学部"}}
_text_role_cache: dict[str, dict] = {}
_text_role_cache_loaded: bool = False

def _text_role_key(s: str) -> str:
    """テキストからスペース・改行・括弧等を除去して正規化キーにする"""
    s = re.sub(r'[\s\u3000\n\r]+', '', s or '')
    # 括弧類も除去して一致率UP
    for ch in '（）()【】[]「」『』':
        s = s.replace(ch, '')
    return s

def _build_text_role_cache() -> dict[str, dict]:
    """正解DBの全レコードからテキスト→フィールド役割マッピングを構築。
    各テキスト断片が「所属」「演者名」「演題」「講演会名」のどれに使われたかを学習。
    出現回数が多い役割を採用。"""
    from collections import Counter
    role_counts: dict[str, Counter] = {}  # key → Counter({"affiliation": 3, ...})
    formatted: dict[str, str] = {}  # key → 最新のフォーマット済みテキスト

    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}

            # event_title
            et = (cj.get("event_title") or "").strip()
            if et:
                key = _text_role_key(et)
                if key and len(key) >= 4:
                    role_counts.setdefault(key, Counter())["event_title"] += 1
                    if key not in formatted:
                        formatted[key] = re.sub(r'[\n\r]+', ' ', et).strip()

            # chair
            ch = cj.get("chair") or {}
            for field, role in [("name", "person_name"), ("affiliation", "affiliation")]:
                val = (ch.get(field) or "").strip()
                if not val:
                    continue
                key = _text_role_key(val)
                min_len = 2 if role == "person_name" else 4
                if key and len(key) >= min_len:
                    role_counts.setdefault(key, Counter())[role] += 1
                    if key not in formatted:
                        formatted[key] = re.sub(r'[\n\r]+', ' ', val).strip()

            # talks
            for t in (cj.get("talks") or []):
                person = _json_person_name_value(t).strip()
                for val, role in [
                    (person, "person_name"),
                    ((t.get("affiliation") or "").strip(), "affiliation"),
                ]:
                    if not val:
                        continue
                    key = _text_role_key(val)
                    min_len = 2 if role == "person_name" else 4
                    if key and len(key) >= min_len:
                        role_counts.setdefault(key, Counter())[role] += 1
                        if key not in formatted:
                            formatted[key] = re.sub(r'[\n\r]+', ' ', val).strip()

                if _json_is_chair_item(t):
                    continue

                # talk title
                title_lines = t.get("title_lines") or []
                title = " ".join(title_lines).strip() if title_lines else (t.get("title") or "").strip()
                if title:
                    key = _text_role_key(title)
                    if key and len(key) >= 6:
                        role_counts.setdefault(key, Counter())["talk_title"] += 1
                        if key not in formatted:
                            formatted[key] = re.sub(r'[\n\r]+', ' ', title).strip()
    except Exception as e:
        print(f"[text-role-cache] build error: {e}")

    # Counter → most_common role
    cache: dict[str, dict] = {}
    for key, counter in role_counts.items():
        role, count = counter.most_common(1)[0]
        if count >= 1:
            cache[key] = {"role": role, "count": count, "formatted": formatted.get(key, "")}
    if cache:
        print(f"[text-role-cache] loaded {len(cache)} entries")
    return cache

def _get_text_role_cache() -> dict[str, dict]:
    global _text_role_cache, _text_role_cache_loaded
    if not _text_role_cache_loaded:
        _text_role_cache = _build_text_role_cache()
        _text_role_cache_loaded = True
    return _text_role_cache

def invalidate_text_role_cache():
    global _text_role_cache_loaded
    _text_role_cache_loaded = False

def lookup_text_role(text: str) -> str | None:
    """テキストの学習済み役割を返す。未学習ならNone。
    役割: "affiliation", "person_name", "event_title", "talk_title" """
    cache = _get_text_role_cache()
    if not cache:
        return None
    key = _text_role_key(text)
    if not key:
        return None
    entry = cache.get(key)
    if entry:
        return entry["role"]
    return None

def lookup_text_formatted(text: str) -> str | None:
    """テキストの学習済みフォーマット版を返す。未学習ならNone。"""
    cache = _get_text_role_cache()
    if not cache:
        return None
    key = _text_role_key(text)
    if not key:
        return None
    entry = cache.get(key)
    if entry and entry.get("formatted"):
        return entry["formatted"]
    return None


def apply_learned_text_roles(payload) -> object:
    """正解DBから学習したテキスト→役割の知識を使って、フィールド割り振りを検証・修正する。
    例: title_lines に所属テキストが混入していたら affiliation へ移動する等。"""
    cache = _get_text_role_cache()
    if not cache:
        return payload

    # ─────────── talks ───────────
    for t in getattr(payload, "talks", []) or []:
        if _is_program_chair_item(t):
            aff = getattr(t, "affiliation", "") or ""
            if aff:
                fmt = lookup_text_formatted(aff)
                if fmt and fmt != aff.replace("\n", " ").strip():
                    t.affiliation = fmt
            continue

        title_lines = list(getattr(t, "title_lines", []) or [])
        current_aff = normalize_space(getattr(t, "affiliation", "") or "")

        # (A) title_lines に所属テキストが混入 → affiliation に移動
        new_title_lines = []
        moved_aff_parts = []
        for line in title_lines:
            key = _text_role_key(line)
            entry = cache.get(key) if key else None
            if entry and entry["role"] == "affiliation" and entry["count"] >= 1:
                moved_aff_parts.append(line)
                print(f"[text-role-fix] moved from title to affiliation: '{line[:40]}'")
            else:
                new_title_lines.append(line)

        if moved_aff_parts:
            t.title_lines = new_title_lines
            t.title = "\n".join(new_title_lines)
            if not current_aff:
                t.affiliation = " ".join(moved_aff_parts)

        # (B) affiliation に演題テキストが混入 → title_lines に移動
        aff = normalize_space(getattr(t, "affiliation", "") or "")
        if aff:
            key = _text_role_key(aff)
            entry = cache.get(key) if key else None
            if entry and entry["role"] == "talk_title" and entry["count"] >= 1:
                current_title = getattr(t, "title_lines", []) or []
                if not current_title or not "".join(current_title).strip():
                    print(f"[text-role-fix] moved from affiliation to title: '{aff[:40]}'")
                    t.title_lines = [aff]
                    t.title = aff
                    t.affiliation = ""

        # (C) 全フィールドのテキストに学習済みフォーマットを適用
        aff = getattr(t, "affiliation", "") or ""
        if aff:
            fmt = lookup_text_formatted(aff)
            if fmt and fmt != aff.replace("\n", " ").strip():
                t.affiliation = fmt

    # ─────────── chair ───────────
    if getattr(payload, "chair", None):
        aff = getattr(payload.chair, "affiliation", "") or ""
        if aff:
            fmt = lookup_text_formatted(aff)
            if fmt and fmt != aff.replace("\n", " ").strip():
                payload.chair.affiliation = fmt

    return payload


def apply_learned_affiliation_format(payload) -> object:
    """正解DBから学習した所属のスペース位置を適用する。
    blocks由来のスペースなしテキストを、ユーザーが過去に確定した
    スペース入りテキストに置換する。"""
    cache = _get_affiliation_format_cache()
    if not cache:
        return payload

    def _lookup(aff: str) -> str | None:
        """完全一致 → 前方一致の順で学習済みフォーマットを検索"""
        key = _aff_cache_key(aff)
        if not key:
            return None
        # 完全一致
        if key in cache:
            return cache[key]
        # 前方一致（施設名は同じだが役職等が追加されたケース）
        for ck, cv in cache.items():
            if ck and key.startswith(ck) and len(ck) >= 6:
                suffix = key[len(ck):]
                return cv + " " + suffix
        return None

    # chair
    if getattr(payload, "chair", None):
        aff = getattr(payload.chair, "affiliation", "") or ""
        learned = _lookup(aff)
        if learned and learned != aff:
            print(f"[affiliation-format-cache] chair: '{aff[:40]}' -> '{learned[:40]}'")
            payload.chair.affiliation = learned

    # talks
    for t in getattr(payload, "talks", []) or []:
        aff = getattr(t, "affiliation", "") or ""
        learned = _lookup(aff)
        if learned and learned != aff:
            print(f"[affiliation-format-cache] talk: '{aff[:40]}' -> '{learned[:40]}'")
            t.affiliation = learned

    return payload


def _load_correct_answers() -> list[dict]:
    """正解DBを読み込む（Postgres）"""
    try:
        with db_connect() as con:
            rows = con.execute(
                "SELECT job_id, event_title, blocks_text, keywords, correct_json, embedding, created_at "
                "FROM correct_answers ORDER BY created_at DESC LIMIT 500"
            ).fetchall()
        return [
            {
                "job_id": r["job_id"],
                "event_title": r["event_title"],
                "blocks_text": r["blocks_text"],
                "keywords": r["keywords"] or [],
                "correct_json": r["correct_json"] or {},
                "embedding": r["embedding"],
                "created_at": str(r["created_at"]),
            }
            for r in rows
        ]
    except Exception as e:
        print(f"[correct_answers] load error: {e}")
        return []


def _extract_keywords(text: str) -> set[str]:
    """テキストからキーワードを抽出（類似度計算用）"""
    text = normalize_space(text or "")
    # 日本語 + 英数字の連続をトークンとして抽出
    tokens = set(re.findall(r'[\u3040-\u9FFF\uF900-\uFAFF]{2,}|[a-zA-Z0-9]{2,}', text))
    return tokens


# ---------- レイアウトパターン学習 ----------

def _compute_layout_hints(blocks_json: list, correct_json: dict) -> dict:
    """正解JSONの各フィールドがblocks内のどの位置にあるかを逆引きし、
    アンカーからの相対位置を記録する。
    blocks_json: [{"text":..., "top":..., "left":..., ...}, ...]
    """
    if not blocks_json or not correct_json:
        return {}

    def _find_blocks(text_query: str) -> list[dict]:
        """テキストを含むブロック候補を検索（スペース除去で比較）"""
        if not text_query:
            return []
        q = text_query.replace(" ", "").replace("\u3000", "").replace("先生", "")
        if len(q) < 2:
            return []
        out = []
        for b in blocks_json:
            bt = (b.get("text", "") if isinstance(b, dict) else "").replace(" ", "").replace("\u3000", "").replace("先生", "")
            if q in bt:
                out.append(b)
        return out

    def _find_block(text_query: str) -> dict | None:
        """テキストを含む最初のブロックを検索（スペース除去で比較）"""
        blocks = _find_blocks(text_query)
        return blocks[0] if blocks else None

    def _find_anchor(label: str) -> dict | None:
        """座長/演者/講演/PROGRAM等のラベルブロックを検索"""
        lbl = label.replace(" ", "")
        for b in blocks_json:
            bt = (b.get("text", "") if isinstance(b, dict) else "").replace(" ", "").replace("\u3000", "")
            if lbl in bt:
                return b
        return None

    def _find_near_anchor(labels: list[str], near_top: int | None = None) -> dict | None:
        candidates = []
        for b in blocks_json:
            if not isinstance(b, dict):
                continue
            bt = (b.get("text", "") or "").replace(" ", "").replace("\u3000", "")
            if not any(label.replace(" ", "") in bt for label in labels):
                continue
            top = b.get("top", 0)
            if near_top is not None and abs(top - near_top) > 1500000:
                continue
            score = abs(top - near_top) if near_top is not None else top
            candidates.append((score, b))
        if not candidates:
            return None
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]

    all_items = correct_json.get("talks") or []
    regular_talks = [t for t in all_items if _json_is_talk_item(t)]
    inline_chair_items = [t for t in all_items if _json_is_chair_item(t)]

    hints: dict = {
        "talks": [],
        "talk_count": len(regular_talks),
        "program_item_count": len(all_items),
        "inline_chair_count": len(inline_chair_items),
    }

    # PROGRAM アンカー
    prog_block = _find_anchor("PROGRAM") or _find_anchor("プログラム")
    if prog_block:
        hints["program_top"] = prog_block.get("top", 0)

    # イベントタイトルのフォントサイズを記録
    _et_lines = correct_json.get("event_title_lines") or []
    _et_text = _et_lines[0] if _et_lines else (correct_json.get("event_title") or "")
    if _et_text:
        _et_block = _find_block(_et_text[:12])  # 先頭12文字で検索
        if _et_block:
            _pt = float(_et_block.get("max_font_pt") or 0)
            if _pt > 0:
                hints["event_title_font_pt"] = _pt

    # 座長
    chair = correct_json.get("chair") or {}
    chair_name = (chair.get("name") or "").replace(" ", "").replace("\u3000", "")
    chair_anchor = _find_anchor("座長")
    if chair_name and chair_anchor:
        name_block = _find_block(chair_name)
        affil_block = _find_block(chair.get("affiliation") or "")
        anchor_top = chair_anchor.get("top", 0)
        ch_hint: dict = {"anchor_top": anchor_top}
        if name_block:
            ch_hint["name_top"] = name_block.get("top", 0)
            ch_hint["name_rel_y"] = name_block.get("top", 0) - anchor_top
        if affil_block:
            ch_hint["affil_top"] = affil_block.get("top", 0)
            ch_hint["affil_rel_y"] = affil_block.get("top", 0) - anchor_top
        if name_block and affil_block:
            ch_hint["affil_rel_to_name_y"] = affil_block.get("top", 0) - name_block.get("top", 0)
        hints["chair"] = ch_hint

    # 講演（talks）
    talk_idx = 0
    for idx, t in enumerate(all_items):
        if _json_is_chair_item(t):
            continue
        speaker = (t.get("speaker") or "").replace(" ", "").replace("\u3000", "")
        affiliation = t.get("affiliation") or ""
        if not speaker:
            continue

        speaker_block = _find_block(speaker)
        if not speaker_block:
            continue

        sp_top = speaker_block.get("top", 0)
        t_hint: dict = {
            "speaker_top": sp_top,
            "talk_index": talk_idx,
            "program_index": t.get("program_index", idx),
        }
        talk_idx += 1
        # 演者ブロックのフォントサイズを記録
        _sp_pt = float(speaker_block.get("max_font_pt") or 0)
        if _sp_pt > 0:
            t_hint["speaker_font_pt"] = _sp_pt

        # 演者ラベルを付近で検索
        enja_anchor = None
        for b in blocks_json:
            bt = (b.get("text", "") if isinstance(b, dict) else "").replace(" ", "").replace("\u3000", "")
            b_top = b.get("top", 0) if isinstance(b, dict) else 0
            if "演者" in bt and abs(b_top - sp_top) < 1500000:
                enja_anchor = b
                break

        if enja_anchor:
            t_hint["enja_anchor_top"] = enja_anchor.get("top", 0)
            t_hint["speaker_rel_to_enja_y"] = sp_top - enja_anchor.get("top", 0)

        # 所属ブロック
        if affiliation:
            affil_block = _find_block(affiliation)
            if affil_block:
                aff_top = affil_block.get("top", 0)
                t_hint["affil_top"] = aff_top
                t_hint["affil_rel_to_speaker_y"] = aff_top - sp_top

        # タイトルブロック
        title_lines = t.get("title_lines") or []
        title_text = title_lines[0] if title_lines else (t.get("title") or "")
        if title_text:
            title_block = _find_block(title_text)
            if title_block:
                t_hint["title_top"] = title_block.get("top", 0)
                t_hint["title_rel_to_speaker_y"] = title_block.get("top", 0) - sp_top
                _tt_pt = float(title_block.get("max_font_pt") or 0)
                if _tt_pt > 0:
                    t_hint["talk_title_font_pt"] = _tt_pt

        hints["talks"].append(t_hint)

    # talks 内に混ぜる途中座長
    chair_like_anchors = []
    for b in blocks_json:
        if not isinstance(b, dict):
            continue
        bt = (b.get("text", "") or "").replace(" ", "").replace("\u3000", "")
        if any(label in bt for label in ["座長", "総合司会", "司会"]):
            chair_like_anchors.append(b)
    chair_like_anchors.sort(key=lambda b: (b.get("top", 0), b.get("left", 0)))
    top_chair_offset = 1 if (correct_json.get("chair") or {}).get("name") else 0

    inline_chair_hints = []
    inline_idx = 0
    for idx, t in enumerate(all_items):
        if not _json_is_chair_item(t):
            continue
        name = _json_person_name_value(t)
        affiliation = t.get("affiliation") or ""
        role = normalize_space(t.get("role") or "座長")
        anchor = None
        anchor_pos = inline_idx + top_chair_offset
        if chair_like_anchors and anchor_pos < len(chair_like_anchors):
            anchor = chair_like_anchors[anchor_pos]
        inline_idx += 1

        name_candidates = _find_blocks(name)
        affil_candidates = _find_blocks(affiliation)
        anchor_top_guess = anchor.get("top", 0) if anchor else None
        if not anchor:
            near_top_guess = name_candidates[0].get("top", 0) if name_candidates else (affil_candidates[0].get("top", 0) if affil_candidates else None)
            anchor = _find_near_anchor([role, "座長", "総合司会", "司会"], near_top_guess)
            anchor_top_guess = anchor.get("top", 0) if anchor else near_top_guess

        def _nearest(candidates: list[dict], near_top: int | None) -> dict | None:
            if not candidates:
                return None
            if near_top is None:
                return candidates[0]
            return sorted(candidates, key=lambda b: abs((b.get("top", 0) or 0) - near_top))[0]

        name_block = _nearest(name_candidates, anchor_top_guess)
        affil_block = _nearest(affil_candidates, name_block.get("top", 0) if name_block else anchor_top_guess)
        near_top = name_block.get("top", 0) if name_block else (affil_block.get("top", 0) if affil_block else None)
        anchor_top = anchor.get("top", 0) if anchor else (near_top or 0)
        ch_hint: dict = {
            "role": role,
            "program_index": t.get("program_index", idx),
            "anchor_top": anchor_top,
        }
        if name_block:
            ch_hint["name_top"] = name_block.get("top", 0)
            ch_hint["name_rel_y"] = name_block.get("top", 0) - anchor_top
        if affil_block:
            ch_hint["affil_top"] = affil_block.get("top", 0)
            ch_hint["affil_rel_y"] = affil_block.get("top", 0) - anchor_top
        if name_block and affil_block:
            ch_hint["affil_rel_to_name_y"] = affil_block.get("top", 0) - name_block.get("top", 0)
        if name_block or affil_block:
            inline_chair_hints.append(ch_hint)

    if inline_chair_hints:
        hints["inline_chairs"] = inline_chair_hints

    if not hints.get("talks") and not hints.get("chair") and not hints.get("inline_chairs"):
        return {}

    return hints


# ---------- レイアウトパターンキャッシュ ----------

_layout_pattern_cache: dict | None = None
_layout_pattern_cache_loaded = False

def _build_layout_pattern_cache() -> dict:
    """正解DBに蓄積されたlayout_hintsを集計し、位置パターンの統計を返す。
    Returns:
        {
            "affil_rel_to_speaker_y": {"median": -300000, "q25": -500000, "q75": -100000, "count": 50},
            "affil_rel_to_name_y_chair": {"median": -200000, ...},
            "speaker_rel_to_enja_y": {"median": 400000, ...},
            "title_rel_to_speaker_y": {"median": -800000, ...},
        }
    """
    affil_rels: list[int] = []
    chair_affil_rels: list[int] = []
    speaker_enja_rels: list[int] = []
    title_rels: list[int] = []
    event_title_font_pts: list[float] = []
    speaker_font_pts: list[float] = []
    talk_title_font_pts: list[float] = []

    try:
        with db_connect() as con:
            rows = con.execute(
                "SELECT layout_hints FROM correct_answers WHERE layout_hints IS NOT NULL AND layout_hints != '{}'"
            ).fetchall()
        for r in rows:
            hints = r["layout_hints"] if isinstance(r["layout_hints"], dict) else {}
            if not hints:
                continue
            # chair
            ch = hints.get("chair") or {}
            if "affil_rel_to_name_y" in ch:
                chair_affil_rels.append(ch["affil_rel_to_name_y"])
            for ich in (hints.get("inline_chairs") or []):
                if "affil_rel_to_name_y" in ich:
                    chair_affil_rels.append(ich["affil_rel_to_name_y"])
            # event title font
            if "event_title_font_pt" in hints:
                event_title_font_pts.append(float(hints["event_title_font_pt"]))
            # talks
            for th in (hints.get("talks") or []):
                if "affil_rel_to_speaker_y" in th:
                    affil_rels.append(th["affil_rel_to_speaker_y"])
                if "speaker_rel_to_enja_y" in th:
                    speaker_enja_rels.append(th["speaker_rel_to_enja_y"])
                if "title_rel_to_speaker_y" in th:
                    title_rels.append(th["title_rel_to_speaker_y"])
                if "speaker_font_pt" in th:
                    speaker_font_pts.append(float(th["speaker_font_pt"]))
                if "talk_title_font_pt" in th:
                    talk_title_font_pts.append(float(th["talk_title_font_pt"]))
    except Exception as e:
        print(f"[layout-pattern-cache] error: {e}")

    def _stats(vals: list[int]) -> dict:
        if len(vals) < 3:
            return {}
        vals_sorted = sorted(vals)
        n = len(vals_sorted)
        return {
            "median": vals_sorted[n // 2],
            "q25": vals_sorted[n // 4],
            "q75": vals_sorted[3 * n // 4],
            "count": n,
            "above_ratio": sum(1 for v in vals if v < 0) / n,
        }

    def _fstats(vals: list[float]) -> dict:
        """float リスト版の統計（フォントサイズ用）"""
        if len(vals) < 3:
            return {}
        vals_sorted = sorted(vals)
        n = len(vals_sorted)
        return {
            "median": vals_sorted[n // 2],
            "q25": vals_sorted[n // 4],
            "q75": vals_sorted[3 * n // 4],
            "count": n,
        }

    cache = {}
    s = _stats(affil_rels)
    if s:
        cache["affil_rel_to_speaker_y"] = s
    s = _stats(chair_affil_rels)
    if s:
        cache["affil_rel_to_name_y_chair"] = s
    s = _stats(speaker_enja_rels)
    if s:
        cache["speaker_rel_to_enja_y"] = s
    s = _stats(title_rels)
    if s:
        cache["title_rel_to_speaker_y"] = s
    # フォントサイズ統計
    s = _fstats(event_title_font_pts)
    if s:
        cache["event_title_font_pt"] = s
    s = _fstats(speaker_font_pts)
    if s:
        cache["speaker_font_pt"] = s
    s = _fstats(talk_title_font_pts)
    if s:
        cache["talk_title_font_pt"] = s

    if cache:
        print(f"[layout-pattern-cache] loaded {sum(v.get('count',0) for v in cache.values())} samples across {len(cache)} patterns")
        for k, v in cache.items():
            print(f"  {k}: median={v.get('median',0)}, above_ratio={v.get('above_ratio',0):.0%}, n={v.get('count',0)}")
    return cache


def _get_layout_pattern_cache() -> dict:
    global _layout_pattern_cache, _layout_pattern_cache_loaded
    if not _layout_pattern_cache_loaded:
        _layout_pattern_cache = _build_layout_pattern_cache()
        _layout_pattern_cache_loaded = True
    return _layout_pattern_cache or {}


def invalidate_layout_pattern_cache():
    global _layout_pattern_cache_loaded
    _layout_pattern_cache_loaded = False


# ---------- 人名辞書（正解DBから学習）----------
# 正解 DB に登録された speaker / chair name（スペース除去キー）の集合。
# 「先生」なしでも人名として確定できる。
_person_name_dict_cache: set[str] | None = None
_person_name_dict_cache_loaded: bool = False

def _build_person_name_dict_cache() -> set[str]:
    """正解DB の全レコードから speaker / chair 人名を収集してセットを返す。"""
    names: set[str] = set()
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            # talks
            for t in (cj.get("talks") or []):
                sp = _json_person_name_key(t)
                if sp and len(sp) >= 2:
                    names.add(sp)
            # chair
            ch = cj.get("chair") or {}
            cn = _json_person_name_key(ch)
            if cn and len(cn) >= 2:
                names.add(cn)
    except Exception as e:
        print(f"[person-name-dict] build error: {e}")
    if names:
        print(f"[person-name-dict] loaded {len(names)} known names")
    return names

def _get_person_name_dict_cache() -> set[str]:
    global _person_name_dict_cache, _person_name_dict_cache_loaded
    if not _person_name_dict_cache_loaded:
        _person_name_dict_cache = _build_person_name_dict_cache()
        _person_name_dict_cache_loaded = True
    return _person_name_dict_cache or set()

def invalidate_person_name_dict_cache():
    global _person_name_dict_cache_loaded
    _person_name_dict_cache_loaded = False


# ---------- 施設名辞書（正解DBから学習）----------
# 正解 DB の affiliation から施設名・所属文字列（スペース除去キー）を収集。
# 一般キーワード（病院/大学等）を含まない固有施設名も検出できる。
_facility_name_dict_cache: set[str] | None = None
_facility_name_dict_cache_loaded: bool = False

# 所属文字列末尾に付く役職語（除去してコア施設名を得る）
_ROLE_SUFFIX_WORDS = [
    "主任教授", "教授", "准教授", "講師", "助教",
    "病院長", "院長", "副院長", "部長", "副部長", "医長",
    "センター長", "科長", "室長", "所長", "理事長",
]

def _strip_role_suffix(s: str) -> str:
    """所属テキスト末尾の役職語を除去して施設・部署名を返す。"""
    s = s.strip()
    for role in _ROLE_SUFFIX_WORDS:
        if s.endswith(role):
            s = s[: -len(role)].strip()
            break
    return s

def _build_facility_name_dict_cache() -> set[str]:
    """正解DB の全 affiliation から施設名キーを収集してセットを返す。
    フルの所属文字列（スペース除去）と、役職語を除いたコア部分の両方を登録。
    3文字以上のものだけ追加してノイズを防ぐ。
    """
    keys: set[str] = set()
    _bad = {"PROGRAM", "AGENDA", "SCHEDULE", ""}

    def _add(raw: str) -> None:
        raw = raw.strip()
        if not raw or raw.upper() in _bad:
            return
        # フル文字列（スペース除去）
        full_key = re.sub(r'[\s\u3000]+', '', raw)
        if len(full_key) >= 3:
            keys.add(full_key)
        # 役職語を除いたコア部分
        core = _strip_role_suffix(raw)
        core_key = re.sub(r'[\s\u3000]+', '', core)
        if len(core_key) >= 3 and core_key != full_key:
            keys.add(core_key)

    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            for t in (cj.get("talks") or []):
                _add(t.get("affiliation") or "")
            ch = cj.get("chair") or {}
            _add(ch.get("affiliation") or "")
    except Exception as e:
        print(f"[facility-name-dict] build error: {e}")
    if keys:
        print(f"[facility-name-dict] loaded {len(keys)} known facility keys")
    return keys

def _get_facility_name_dict_cache() -> set[str]:
    global _facility_name_dict_cache, _facility_name_dict_cache_loaded
    if not _facility_name_dict_cache_loaded:
        _facility_name_dict_cache = _build_facility_name_dict_cache()
        _facility_name_dict_cache_loaded = True
    return _facility_name_dict_cache or set()

def invalidate_facility_name_dict_cache():
    global _facility_name_dict_cache_loaded
    _facility_name_dict_cache_loaded = False


# ---------- タイトル改行文字数キャッシュ（正解DBから学習）----------
# 正解DBの title_lines / event_title_lines の各行文字数を集計し、
# format_title_initial の max_len と fix_title_lines_jp の MAX_MERGE_LEN に利用する。

_title_line_len_cache: dict | None = None
_title_line_len_cache_loaded: bool = False

def _build_title_line_len_cache() -> dict:
    """正解DB の title_lines / event_title_lines から1行あたりの文字数分布を学習する。
    Returns:
        {
            "talk_title_p90": int,      # トーク演題の90パーセンタイル文字数
            "talk_title_median": int,   # 中央値
            "event_title_p90": int,     # イベントタイトルの90パーセンタイル
            "event_title_median": int,
        }
    """
    talk_lens: list[int] = []
    event_lens: list[int] = []
    result: dict = {}
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            # イベントタイトル各行
            for line in (cj.get("event_title_lines") or []):
                ln = len(re.sub(r'[\s\u3000]+', '', normalize_space(line)))
                if ln >= 3:
                    event_lens.append(ln)
            # トーク演題各行
            for t in (cj.get("talks") or []):
                if _json_is_chair_item(t):
                    continue
                for line in (t.get("title_lines") or []):
                    ln = len(re.sub(r'[\s\u3000]+', '', normalize_space(line)))
                    if ln >= 3:
                        talk_lens.append(ln)
    except Exception as e:
        print(f"[title-line-len-cache] build error: {e}")
        return result

    if len(talk_lens) >= 5:
        st = sorted(talk_lens)
        n = len(st)
        result["talk_title_p90"] = st[min(int(n * 0.9), n - 1)]
        result["talk_title_median"] = st[n // 2]
        result["talk_count"] = n
        print(f"[title-line-len] talk: median={result['talk_title_median']}, p90={result['talk_title_p90']}, n={n}")

    if len(event_lens) >= 3:
        se = sorted(event_lens)
        n = len(se)
        result["event_title_p90"] = se[min(int(n * 0.9), n - 1)]
        result["event_title_median"] = se[n // 2]
        result["event_count"] = n
        print(f"[title-line-len] event: median={result['event_title_median']}, p90={result['event_title_p90']}, n={n}")

    return result

def _get_title_line_len_cache() -> dict:
    global _title_line_len_cache, _title_line_len_cache_loaded
    if not _title_line_len_cache_loaded:
        _title_line_len_cache = _build_title_line_len_cache()
        _title_line_len_cache_loaded = True
    return _title_line_len_cache or {}

def invalidate_title_line_len_cache():
    global _title_line_len_cache_loaded
    _title_line_len_cache_loaded = False


# ---------- 主催者名辞書（正解DBから学習）----------
_organizer_dict_cache: set[str] | None = None
_organizer_dict_cache_loaded: bool = False

def _build_organizer_dict_cache() -> set[str]:
    """正解DB の organizer から既知の主催者文字列（スペース除去キー）を収集。"""
    keys: set[str] = set()
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            org = (cj.get("organizer") or "").strip()
            if not org:
                continue
            # フル文字列キー
            full_key = re.sub(r'[\s\u3000]+', '', org)
            if len(full_key) >= 4:
                keys.add(full_key)
            # 「主催：」「共催：」等のラベルを除いたコア部分も追加
            core = re.sub(r'^(主催|共催|提供|企画|運営)\s*[:：]\s*', '', org).strip()
            core_key = re.sub(r'[\s\u3000]+', '', core)
            if len(core_key) >= 4 and core_key != full_key:
                keys.add(core_key)
    except Exception as e:
        print(f"[organizer-dict] build error: {e}")
    if keys:
        print(f"[organizer-dict] loaded {len(keys)} known organizer keys")
    return keys

def _get_organizer_dict_cache() -> set[str]:
    global _organizer_dict_cache, _organizer_dict_cache_loaded
    if not _organizer_dict_cache_loaded:
        _organizer_dict_cache = _build_organizer_dict_cache()
        _organizer_dict_cache_loaded = True
    return _organizer_dict_cache or set()

def invalidate_organizer_dict_cache():
    global _organizer_dict_cache_loaded
    _organizer_dict_cache_loaded = False


# ---------- 座長ラベル語辞書（正解DBから学習）----------
_chair_label_cache: dict | None = None   # {label_word: count}
_chair_label_cache_loaded: bool = False

def _build_chair_label_cache() -> dict:
    """correct_json.chair.role の出現語を集計し {word: count} を返す。"""
    counts: dict[str, int] = {}
    _default = {"座長", "総合司会", "司会"}
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            role = (cj.get("chair") or {}).get("role") or ""
            role = role.strip()
            if role:
                counts[role] = counts.get(role, 0) + 1
            for t in (cj.get("talks") or []):
                if not _json_is_chair_item(t):
                    continue
                role = (t.get("role") or "").strip()
                if role:
                    counts[role] = counts.get(role, 0) + 1
    except Exception as e:
        print(f"[chair-label-cache] build error: {e}")
    # デフォルト語を必ず含める
    for w in _default:
        counts.setdefault(w, 0)
    if counts:
        print(f"[chair-label-cache] labels: {counts}")
    return counts

def _get_chair_label_cache() -> dict:
    global _chair_label_cache, _chair_label_cache_loaded
    if not _chair_label_cache_loaded:
        _chair_label_cache = _build_chair_label_cache()
        _chair_label_cache_loaded = True
    return _chair_label_cache or {}

def _get_chair_label_words() -> list[str]:
    """出現頻度の高い順に座長ラベル語リストを返す（最低でも「座長」を含む）。"""
    c = _get_chair_label_cache()
    if not c:
        return ["座長", "総合司会", "司会"]
    return [w for w, _ in sorted(c.items(), key=lambda x: -x[1]) if w]

def invalidate_chair_label_cache():
    global _chair_label_cache_loaded
    _chair_label_cache_loaded = False


# ---------- 講演数分布（正解DBから学習）----------
_talk_count_cache: dict | None = None
_talk_count_cache_loaded: bool = False

def _build_talk_count_cache() -> dict:
    """correct_json.talks の件数を集計し統計を返す。"""
    counts: list[int] = []
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            n = sum(1 for t in (cj.get("talks") or []) if _json_is_talk_item(t))
            if n > 0:
                counts.append(n)
    except Exception as e:
        print(f"[talk-count-cache] build error: {e}")
    if len(counts) < 3:
        return {}
    counts_sorted = sorted(counts)
    n = len(counts_sorted)
    result = {
        "median": counts_sorted[n // 2],
        "q25": counts_sorted[n // 4],
        "q75": counts_sorted[3 * n // 4],
        "min": counts_sorted[0],
        "max": counts_sorted[-1],
        "count": n,
    }
    print(f"[talk-count-cache] median={result['median']}, q25={result['q25']}, q75={result['q75']}, n={n}")
    return result

def _get_talk_count_cache() -> dict:
    global _talk_count_cache, _talk_count_cache_loaded
    if not _talk_count_cache_loaded:
        _talk_count_cache = _build_talk_count_cache()
        _talk_count_cache_loaded = True
    return _talk_count_cache or {}

def invalidate_talk_count_cache():
    global _talk_count_cache_loaded
    _talk_count_cache_loaded = False


# ---------- datetime_time_newline 分布（正解DBから学習）----------
_datetime_newline_cache: dict | None = None
_datetime_newline_cache_loaded: bool = False

def _build_datetime_newline_cache() -> dict:
    """correct_json.datetime_time_newline の True/False 出現比率を集計。"""
    true_count = 0
    false_count = 0
    try:
        answers = _load_correct_answers()
        for ans in answers:
            cj = ans.get("correct_json") or {}
            val = cj.get("datetime_time_newline")
            if val is True:
                true_count += 1
            elif val is False:
                false_count += 1
    except Exception as e:
        print(f"[datetime-newline-cache] build error: {e}")
    total = true_count + false_count
    if total < 3:
        return {}
    result = {
        "true_count": true_count,
        "false_count": false_count,
        "true_ratio": true_count / total,
    }
    print(f"[datetime-newline-cache] true={true_count}, false={false_count}, true_ratio={result['true_ratio']:.0%}")
    return result

def _get_datetime_newline_cache() -> dict:
    global _datetime_newline_cache, _datetime_newline_cache_loaded
    if not _datetime_newline_cache_loaded:
        _datetime_newline_cache = _build_datetime_newline_cache()
        _datetime_newline_cache_loaded = True
    return _datetime_newline_cache or {}

def invalidate_datetime_newline_cache():
    global _datetime_newline_cache_loaded
    _datetime_newline_cache_loaded = False


def _compute_similarity(kw_a: set[str], kw_b: set[str]) -> float:
    """Jaccard係数で類似度を計算"""
    if not kw_a or not kw_b:
        return 0.0
    intersection = kw_a & kw_b
    union = kw_a | kw_b
    return len(intersection) / len(union) if union else 0.0


def save_correct_answer(
    blocks_text: str,
    correct_json: dict,
    event_title: str = "",
    job_id: str = "",
    blocks_json: list | None = None,
) -> None:
    """確定済みの正解データをPostgresに保存する（embedding付き）"""
    _INVALID_NAMES = {"PROGRAM", "P R O G R A M", "AGENDA", "SCHEDULE", "TIME TABLE", "タイムテーブル", "プログラム"}

    # 保存前バリデーション
    # (A) chair の不正データ除去
    ch = correct_json.get("chair") or {}
    if ch:
        ch_name = (ch.get("name") or "").strip()
        if ch_name.upper() in _INVALID_NAMES:
            print(f"[correct-answer] WARNING: clearing invalid chair name before save: '{ch_name}'")
            ch["name"] = ""
            ch["name_display"] = ""
        ch_aff = (ch.get("affiliation") or "").strip()
        if ch_aff.upper() in _INVALID_NAMES:
            print(f"[correct-answer] WARNING: clearing invalid chair affiliation before save: '{ch_aff}'")
            ch["affiliation"] = ""

    # (B) talks の affiliation に座長情報が混入していたら除去
    if "talks" in correct_json:
        for t in correct_json.get("talks", []):
            aff = t.get("affiliation", "") or ""
            if _json_is_chair_item(t):
                if aff.upper() in _INVALID_NAMES:
                    print(f"[correct-answer] WARNING: clearing invalid inline chair affiliation before save: '{aff}'")
                    t["affiliation"] = ""
                name_key = _json_person_name_key(t)
                if name_key.upper() in _INVALID_NAMES:
                    print(f"[correct-answer] WARNING: clearing invalid inline chair name before save: '{name_key}'")
                    t["name_display"] = ""
                    t["speaker"] = ""
                    t["speaker_display"] = ""
                t["role"] = normalize_space(t.get("role") or "座長")
                continue
            if "座長" in aff:
                print(f"[correct-answer] WARNING: clearing chair-contaminated affiliation in talk before save: '{aff[:60]}'")
                t["affiliation"] = ""
            # speaker の不正値チェック
            sp = (t.get("speaker") or "").strip()
            if sp.upper() in _INVALID_NAMES:
                print(f"[correct-answer] WARNING: clearing invalid speaker before save: '{sp}'")
                t["speaker"] = ""
                t["speaker_display"] = ""

    keywords = list(_extract_keywords(blocks_text + " " + event_title))
    truncated_text = blocks_text[:2000]
    embedding = _compute_embedding(blocks_text + " " + event_title)
    layout_hints = _compute_layout_hints(blocks_json or [], correct_json) if blocks_json else {}
    print(f"[correct-answer] saving job_id={job_id} event_title='{event_title[:60]}' blocks_text={len(blocks_text)} chars, embedding={'yes' if embedding else 'no'}, layout_hints={len(layout_hints)} keys")

    try:
        with db_connect() as con:
            con.execute(
                """
                INSERT INTO correct_answers (job_id, event_title, blocks_text, keywords, correct_json, embedding, layout_hints)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (job_id) DO UPDATE SET
                    event_title = EXCLUDED.event_title,
                    blocks_text = EXCLUDED.blocks_text,
                    keywords = EXCLUDED.keywords,
                    correct_json = EXCLUDED.correct_json,
                    embedding = EXCLUDED.embedding,
                    layout_hints = EXCLUDED.layout_hints,
                    created_at = NOW()
                """,
                (job_id, event_title, truncated_text, keywords,
                 json.dumps(correct_json, ensure_ascii=False),
                 embedding,
                 json.dumps(layout_hints, ensure_ascii=False)),
            )
            # 500件を超えたら古い順に削除
            con.execute(
                """
                DELETE FROM correct_answers WHERE id IN (
                    SELECT id FROM correct_answers
                    ORDER BY created_at DESC
                    OFFSET 500
                )
                """
            )
            con.commit()
    except Exception as e:
        print(f"[correct_answers] save error: {e}")

    # 保存後にキャッシュを無効化（次回再構築）
    invalidate_speaker_display_cache()
    invalidate_affiliation_format_cache()
    invalidate_text_role_cache()
    invalidate_layout_pattern_cache()
    invalidate_person_name_dict_cache()
    invalidate_facility_name_dict_cache()
    invalidate_title_line_len_cache()
    invalidate_organizer_dict_cache()
    invalidate_chair_label_cache()
    invalidate_talk_count_cache()
    invalidate_datetime_newline_cache()


def find_similar_correct_answers(
    blocks_text: str,
    event_title: str = "",
    top_k: int = 2,
    min_similarity: float = 0.15,
) -> list[dict]:
    """類似する正解データを検索して返す（embedding優先、fallback: Jaccard）"""
    answers = _load_correct_answers()
    print(f"[correct-answer-search] loaded {len(answers)} answers, query event_title='{event_title[:60]}', blocks_text={len(blocks_text)} chars")
    if not answers:
        return []

    query_text = blocks_text + " " + event_title
    query_emb = _compute_embedding(query_text)
    query_title_ns = (event_title or "").replace(" ", "").replace("\u3000", "").replace("\n", "").lower()

    scored = []
    for ans in answers:
        # embedding があればコサイン類似度を使う
        ans_emb = ans.get("embedding")
        if query_emb and ans_emb:
            sim = _cosine_similarity(query_emb, ans_emb)
        else:
            # fallback: Jaccard
            query_kw = _extract_keywords(query_text)
            ans_kw = set(ans.get("keywords", []))
            sim = _compute_similarity(query_kw, ans_kw)

        # event_title 完全一致なら最優先（同じイベントの正解DB）
        ans_title = (ans.get("event_title") or "")
        ans_title_ns = ans_title.replace(" ", "").replace("\u3000", "").replace("\n", "").lower()
        if query_title_ns and ans_title_ns and query_title_ns == ans_title_ns:
            sim = max(sim, 1.0)
            print(f"[correct-answer-search] exact title match boost: '{ans_title[:50]}' sim→1.0")

        if sim >= min_similarity:
            ans["_similarity"] = sim
            scored.append((sim, ans))

    # 同一類似度では新しいレコードを優先（created_at DESCでロード済みなので index が小さいほうが新しい）
    for idx, (sim, ans) in enumerate(scored):
        ans["_rank_index"] = idx
    scored.sort(key=lambda x: (-x[0], x[1].get("_rank_index", 0)))

    results = [item[1] for item in scored[:top_k]]
    for r in results:
        print(f"[correct-answer-search] match: sim={r.get('_similarity', 0):.3f} job_id={r.get('job_id', '')} title='{r.get('event_title', '')[:50]}'")
    if not results:
        # スコア上位を表示してデバッグ支援
        all_scored = sorted([(s, a) for s, a in scored] if scored else [], key=lambda x: -x[0])
        if not all_scored and answers:
            print(f"[correct-answer-search] no matches above min_similarity={min_similarity}")
        for s, a in all_scored[:3]:
            print(f"[correct-answer-search] (below threshold) sim={s:.3f} title='{a.get('event_title', '')[:50]}'")

    return results


def _build_dynamic_few_shot(similar_answers: list[dict]) -> list[dict]:
    """正解DBの類似結果をfew-shot messagesに変換（入力コンテキスト付き）"""
    messages = []
    for ans in similar_answers:
        cj = ans.get("correct_json", {})
        # 重要フィールドだけ抜粋（トークン節約）
        summary = {
            "event_title": cj.get("event_title", ""),
            "talks": [],
        }
        for t in (cj.get("talks") or [])[:6]:
            if _json_is_chair_item(t):
                summary["talks"].append({
                    "item_type": "chair",
                    "role": t.get("role", "座長"),
                    "name_display": t.get("name_display", ""),
                    "affiliation": t.get("affiliation", ""),
                    "program_index": t.get("program_index", len(summary["talks"])),
                })
            else:
                summary["talks"].append({
                    "item_type": "talk",
                    "time": t.get("time", ""),
                    "title_lines": t.get("title_lines", []),
                    "speaker": t.get("speaker", ""),
                    "affiliation": t.get("affiliation", ""),
                    "program_index": t.get("program_index", len(summary["talks"])),
                })
        if cj.get("chair"):
            chair = cj["chair"]
            if chair.get("name"):
                summary["chair"] = {
                    "name": chair.get("name", ""),
                    "affiliation": chair.get("affiliation", ""),
                }
        if cj.get("organizer"):
            summary["organizer"] = cj["organizer"]
        if cj.get("datetime"):
            summary["datetime"] = cj["datetime"]

        # 入力コンテキストを含めることでAIが「何に対してこの正解か」理解できるようにする
        blocks_excerpt = (ans.get("blocks_text", "") or "")[:500]
        user_msg = (
            f"過去の確定済み類似スライド（参考）:\n"
            f"イベント: {ans.get('event_title', '(不明)')}\n"
            f"ブロックテキスト抜粋: {blocks_excerpt}"
        )
        assistant_msg = json.dumps(summary, ensure_ascii=False)
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": assistant_msg})

    return messages


@dataclass
class CorrectAnswerHints:
    """正解DBから抽出した構造的ヒント（テキスト内容ではなく構造情報）"""
    expected_talk_count: int = 0
    talk_speaker_hints: list = dataclass_field(default_factory=list)  # normalized speaker names
    chair_name_hint: str = ""       # normalized chair name
    organizer_hint: str = ""        # normalized organizer
    similarity: float = 0.0
    _correct_json: dict = dataclass_field(default_factory=dict)
    _job_id: str = ""


def compute_correct_answer_hints(blocks: list, event_title: str) -> CorrectAnswerHints:
    """正解DBから構造的ヒントを抽出する（テキスト内容は含めない、構造情報のみ）。
    パイプライン序盤で呼び出して抽出精度を向上させる。"""
    try:
        all_blocks_text = " ".join(
            (b.get("text", "") if isinstance(b, dict) else getattr(b, "text", ""))
            for b in blocks
        )
        similar = find_similar_correct_answers(all_blocks_text, event_title, top_k=1)
        if not similar:
            return CorrectAnswerHints()

        best = similar[0]
        sim = best.get("_similarity", 0.0)
        if sim < 0.80:
            return CorrectAnswerHints()

        correct = best.get("correct_json") or {}
        def _ns(s): return (s or "").replace(" ", "").replace("\u3000", "")

        talks = [t for t in (correct.get("talks") or []) if _json_is_talk_item(t)]
        chair = correct.get("chair") or {}

        hints = CorrectAnswerHints(
            expected_talk_count=len(talks),
            talk_speaker_hints=[_ns(t.get("speaker", "")) for t in talks if t.get("speaker")],
            chair_name_hint=_ns(chair.get("name", "")),
            organizer_hint=_ns(correct.get("organizer", "")),
            similarity=sim,
            _correct_json=correct,
            _job_id=best.get("job_id", ""),
        )
        print(f"[correct-answer-hints] sim={sim:.2f} job={hints._job_id} "
              f"talks={hints.expected_talk_count} chair='{hints.chair_name_hint}' "
              f"speakers={hints.talk_speaker_hints}")
        return hints
    except Exception as e:
        print(f"[correct-answer-hints] error: {e}")
        return CorrectAnswerHints()


def fill_empty_fields_from_blocks_with_hints(
    payload: DesignJSON, blocks: list, hints: CorrectAnswerHints
) -> DesignJSON:
    """正解DBヒントを使って、空フィールドを blocks 内のテキストで補完する。
    DBのテキスト自体は使わない（blocks に実在するテキストのみ採用）。
    パイプライン中盤（repair_talks_from_blocks 後）で呼び出す。"""
    if hints.similarity < 0.80:
        return payload

    def _bt(b):
        return b.get("text", "") if isinstance(b, dict) else getattr(b, "text", "")

    all_text_ns = "".join(_bt(b) for b in blocks).replace(" ", "").replace("\u3000", "")

    # ---- organizer: 空なら blocks から探す ----
    if not (payload.organizer or "").strip() and hints.organizer_hint:
        if hints.organizer_hint in all_text_ns:
            for b in blocks:
                bt = _bt(b)
                bt_ns = bt.replace(" ", "").replace("\u3000", "")
                if hints.organizer_hint in bt_ns:
                    # blocksのテキストから organizer 行を抽出
                    for line in bt.split("\n"):
                        line_ns = line.replace(" ", "").replace("\u3000", "")
                        if hints.organizer_hint in line_ns:
                            payload.organizer = normalize_space(line)
                            print(f"[hints-fill] organizer from blocks: '{payload.organizer}'")
                            break
                    if payload.organizer:
                        break

    # ---- chair.name: 空なら blocks から探す ----
    if payload.chair and not (payload.chair.name or "").strip() and hints.chair_name_hint:
        if hints.chair_name_hint in all_text_ns:
            for b in blocks:
                bt = _bt(b)
                bt_joined = normalize_space(bt.replace("\n", " "))
                m = re.search(
                    r"([一-龥々ぁ-んァ-ヶ]{1,5})\s*([一-龥々ぁ-んァ-ヶ]{1,5})\s*先生",
                    bt_joined,
                )
                if m:
                    cand = norm_name(m.group(1) + m.group(2))
                    cand_ns = cand.replace(" ", "").replace("\u3000", "")
                    if cand_ns == hints.chair_name_hint and is_valid_person_name(cand):
                        payload.chair.name = cand
                        payload.chair.name_display = cand
                        print(f"[hints-fill] chair.name from blocks: '{cand}'")
                        break

    # ---- talks speaker: 空なら blocks から探す ----
    if payload.talks and hints.talk_speaker_hints:
        # 既にマッチ済みの speaker を集計
        matched_hints = set()
        for t in payload.talks:
            if _is_program_chair_item(t):
                continue
            sp_ns = (getattr(t, "speaker", "") or "").replace(" ", "").replace("\u3000", "")
            if sp_ns:
                matched_hints.add(sp_ns)

        for t in payload.talks:
            if _is_program_chair_item(t):
                continue
            sp_ns = (getattr(t, "speaker", "") or "").replace(" ", "").replace("\u3000", "")
            if sp_ns:
                continue  # 既に speaker がある

            # 未マッチのヒントから blocks 内に存在するものを探す
            for hint_sp in hints.talk_speaker_hints:
                if hint_sp in matched_hints:
                    continue
                if hint_sp not in all_text_ns:
                    continue
                # blocks から該当名前を抽出
                found = False
                for b in blocks:
                    bt = _bt(b)
                    bt_joined = normalize_space(bt.replace("\n", " "))
                    m = re.search(
                        r"([一-龥々ぁ-んァ-ヶ]{1,5})\s*([一-龥々ぁ-んァ-ヶ]{1,5})\s*先生",
                        bt_joined,
                    )
                    if m:
                        cand = norm_name(m.group(1) + m.group(2))
                        cand_ns = cand.replace(" ", "").replace("\u3000", "")
                        if cand_ns == hint_sp and is_valid_person_name(cand):
                            t.speaker = cand
                            t.speaker_display = build_speaker_display(cand) or cand
                            matched_hints.add(hint_sp)
                            print(f"[hints-fill] talk speaker from blocks: '{cand}'")
                            found = True
                            break
                if found:
                    break

    # ---- talk affiliation: speaker は埋まっているが affiliation が空 ----
    if payload.talks and hints._correct_json:
        def _ns(s): return (s or "").replace(" ", "").replace("\u3000", "")
        ct_list = [ct for ct in (hints._correct_json.get("talks") or []) if _json_is_talk_item(ct)]
        ct_by_sp = {_ns(ct.get("speaker", "")): ct for ct in ct_list if ct.get("speaker")}
        for t in payload.talks:
            if _is_program_chair_item(t):
                continue
            if getattr(t, "affiliation", ""):
                continue
            sp_ns = _ns(getattr(t, "speaker", ""))
            ct = ct_by_sp.get(sp_ns)
            if not ct or not ct.get("affiliation"):
                continue
            ct_aff_ns = _ns(ct["affiliation"])
            if ct_aff_ns and ct_aff_ns in all_text_ns:
                # blocks に実在 → blocks のテキストから抽出
                for b in blocks:
                    bt = _bt(b)
                    for line in bt.split("\n"):
                        line_ns = line.replace(" ", "").replace("\u3000", "")
                        if ct_aff_ns in line_ns:
                            t.affiliation = normalize_space(line)
                            print(f"[hints-fill] talk affiliation from blocks: '{t.affiliation}'")
                            break
                    if t.affiliation:
                        break

    # ---- expected_talk_count vs actual: 警告 (高類似度かつ差が大きい時のみ) ----
    if hints.expected_talk_count > 0 and hints.similarity >= 0.95:
        actual = sum(1 for t in (payload.talks or []) if _is_program_talk_item(t))
        diff = abs(actual - hints.expected_talk_count)
        if diff >= 2:
            print(f"[hints-fill] talk count mismatch: expected={hints.expected_talk_count} actual={actual} diff={diff} (sim={hints.similarity:.2f})")
            _w = list(payload.warnings or [])
            if "talk_count_mismatch" not in _w:
                _w.append("talk_count_mismatch")
            payload.warnings = _w

    return payload


def apply_correct_answer_overlay(payload: DesignJSON, blocks: list, vm_rows: list | None = None) -> DesignJSON:
    """
    正解DBを参照して「フォーマットの精度向上」と「信頼度スコアリング」を行う。
    テキスト内容の上書きは一切行わない（blocks に存在しないテキストを導入しない）。
    空フィールドで DB 値が blocks 内に存在する場合のみ補完する。

    適用対象:
    - 名前の正規化（スペース位置の復元: 髙田慶応 → 髙田 慶応）
    - name_display の復元（同一人物の場合のみ）
    - 改行位置の復元（event_title_lines, title_lines）
    - 所属のフォーマット復元（内容同一の場合のみスペース位置を復元）
    - 空フィールドの補完（DB値が blocks 内に実在する場合のみ）
    - 信頼度スコアリング（DB との一致度に基づく confidence 値の算出）
    """
    try:
        all_blocks_text = " ".join(
            (b.get("text", "") if isinstance(b, dict) else getattr(b, "text", ""))
            for b in blocks
        )
        all_blocks_ns = all_blocks_text.replace(" ", "").replace("\u3000", "")
        event_title = payload.event_title or ""

        similar = find_similar_correct_answers(all_blocks_text, event_title, top_k=1)
        if not similar:
            return payload

        best = similar[0]
        sim = best.get("_similarity", 0.0)

        if sim < 0.80:
            return payload

        correct = best.get("correct_json") or {}
        print(f"[correct-answer-overlay] sim={sim:.2f} job_id={best.get('job_id','')}")

        def _ns(s: str) -> str:
            """正規化: スペース全除去"""
            return (s or "").replace(" ", "").replace("\u3000", "")

        def _format_if_same(cur: str, db: str) -> str:
            """内容が同一ならDBのフォーマット（スペース位置）を採用、異なるなら現在値を維持"""
            if not cur or not db:
                return cur
            if _ns(cur) == _ns(db):
                return db
            return cur

        def _find_text_in_blocks(target_ns: str) -> str | None:
            """blocks 内に target_ns (normalized) が存在すれば元テキスト行を返す"""
            if not target_ns or target_ns not in all_blocks_ns:
                return None
            for b in blocks:
                bt = b.get("text", "") if isinstance(b, dict) else getattr(b, "text", "")
                for line in bt.split("\n"):
                    if target_ns in line.replace(" ", "").replace("\u3000", ""):
                        return normalize_space(line)
            return None

        # ---- chair: 同一人物ならスペース位置 / name_display を復元 ----
        cc = correct.get("chair") or {}
        if cc and payload.chair:
            cc_name = cc.get("name", "")
            cur_name = getattr(payload.chair, "name", "") or ""
            if cc_name and cur_name and _ns(cur_name) == _ns(cc_name):
                # 同一人物 → スペース位置を復元
                payload.chair.name = cc_name
                cc_disp = cc.get("name_display") or ""
                if cc_disp and _ns(cc_disp) == _ns(cc_name):
                    payload.chair.name_display = cc_disp
                # 所属: 内容同一ならフォーマット復元
                if cc.get("affiliation"):
                    payload.chair.affiliation = _format_if_same(
                        getattr(payload.chair, "affiliation", "") or "",
                        cc["affiliation"],
                    )
            elif cc_name and not cur_name:
                # 座長名が空 → DB値が blocks 内に存在すれば補完
                found = _find_text_in_blocks(_ns(cc_name))
                if found:
                    # blocks のテキストから人名を抽出
                    m = re.search(r"([一-龥々ぁ-んァ-ヶ]{1,5})\s*([一-龥々ぁ-んァ-ヶ]{1,5})\s*先生?", found)
                    if m:
                        cand = norm_name(m.group(1) + m.group(2))
                        if _ns(cand) == _ns(cc_name):
                            payload.chair.name = cand
                            payload.chair.name_display = cand
                            print(f"[correct-answer-overlay] chair.name filled from blocks: '{cand}'")

            # chair affiliation: 空 + DB値がblocksにある場合のみ補完
            if cc.get("affiliation") and not (getattr(payload.chair, "affiliation", "") or "").strip():
                found_aff = _find_text_in_blocks(_ns(cc["affiliation"]))
                if found_aff:
                    payload.chair.affiliation = found_aff
                    print(f"[correct-answer-overlay] chair.affiliation filled from blocks: '{found_aff}'")

        # ---- event_title: 内容同一なら改行位置のみ復元 ----
        _ct_etl = correct.get("event_title_lines") or ([correct["event_title"]] if correct.get("event_title") else [])
        if _ct_etl:
            _ct_et_ns = _ns("".join(_ct_etl))
            _cur_et_ns = _ns((payload.event_title or "").replace("\n", ""))
            if _ct_et_ns and _cur_et_ns and _ct_et_ns == _cur_et_ns:
                payload.event_title_lines = _ct_etl
                payload.event_title = "\n".join(_ct_etl)
                print(f"[correct-answer-overlay] event_title_lines linebreak restored: {_ct_etl}")

        # ---- organizer: 内容同一ならフォーマット復元、空なら blocks から補完 ----
        if correct.get("organizer"):
            cur_org = getattr(payload, "organizer", "") or ""
            if cur_org:
                payload.organizer = _format_if_same(cur_org, correct["organizer"])
            else:
                # 空 → DB値が blocks にあれば補完
                found_org = _find_text_in_blocks(_ns(correct["organizer"]))
                if found_org:
                    payload.organizer = found_org
                    print(f"[correct-answer-overlay] organizer filled from blocks: '{found_org}'")

        # ---- talks: 同一演者の場合のみフォーマット復元 + 空フィールド補完 ----
        correct_program_items = correct.get("talks") or []
        ct_list = [ct for ct in correct_program_items if _json_is_talk_item(ct)]
        if ct_list and payload.talks:
            # speaker名ベースでマッチング
            ct_by_speaker = {}
            for ct in ct_list:
                sp = _ns(ct.get("speaker", ""))
                if sp:
                    ct_by_speaker[sp] = ct

            for t in payload.talks:
                if _is_program_chair_item(t):
                    continue
                sp = _ns(getattr(t, "speaker", ""))

                if sp and sp in ct_by_speaker:
                    ct = ct_by_speaker[sp]

                    # speaker: スペース位置を復元
                    ct_sp = ct.get("speaker") or ""
                    if ct_sp and _ns(ct_sp) == sp:
                        t.speaker = ct_sp
                        ct_disp = ct.get("speaker_display") or ""
                        if ct_disp:
                            t.speaker_display = ct_disp

                    # affiliation: 内容同一ならフォーマット復元、空なら blocks から補完
                    if ct.get("affiliation"):
                        cur_aff = getattr(t, "affiliation", "") or ""
                        if cur_aff:
                            t.affiliation = _format_if_same(cur_aff, ct["affiliation"])
                        else:
                            found_aff = _find_text_in_blocks(_ns(ct["affiliation"]))
                            if found_aff:
                                t.affiliation = found_aff
                                print(f"[correct-answer-overlay] talk affiliation filled from blocks: '{found_aff}'")

                    # title: 内容同一なら改行位置を復元
                    ct_title_lines = ct.get("title_lines") or []
                    ct_title = "\n".join(ct_title_lines) or ct.get("title", "")
                    cur_title = getattr(t, "title", "") or ""
                    if ct_title_lines and _ns(ct_title.replace("\n", "")) == _ns(cur_title.replace("\n", "")):
                        t.title_lines = fix_title_lines_jp(ct_title_lines)
                        t.title = "\n".join(t.title_lines)

                elif not sp:
                    # speaker が空 → DB のいずれかのスピーカーが blocks に存在すれば補完
                    matched_speakers = {
                        _ns(getattr(tt, "speaker", ""))
                        for tt in payload.talks
                        if (not _is_program_chair_item(tt)) and getattr(tt, "speaker", "")
                    }
                    for hint_sp, ct in ct_by_speaker.items():
                        if hint_sp in matched_speakers:
                            continue
                        found_sp = _find_text_in_blocks(hint_sp)
                        if found_sp:
                            m = re.search(r"([一-龥々ぁ-んァ-ヶ]{1,5})\s*([一-龥々ぁ-んァ-ヶ]{1,5})", found_sp)
                            if m:
                                cand = norm_name(m.group(1) + m.group(2))
                                if _ns(cand) == hint_sp:
                                    t.speaker = cand
                                    t.speaker_display = build_speaker_display(cand) or cand
                                    print(f"[correct-answer-overlay] talk speaker filled from blocks: '{cand}'")
                                    break

        inline_chair_refs = {
            _ns(_json_person_name_value(ct)): ct
            for ct in correct_program_items
            if _json_is_chair_item(ct) and _json_person_name_value(ct)
        }
        if inline_chair_refs and payload.talks:
            for t in payload.talks:
                if not _is_program_chair_item(t):
                    continue
                name_key = _ns(getattr(t, "name_display", "") or getattr(t, "speaker_display", "") or "")
                ct = inline_chair_refs.get(name_key)
                if not ct:
                    continue
                if ct.get("role"):
                    t.role = ct["role"]
                ref_name = _json_person_display_value(ct)
                if ref_name and _ns(ref_name) == name_key:
                    t.name_display = ref_name
                if ct.get("affiliation"):
                    cur_aff = getattr(t, "affiliation", "") or ""
                    if cur_aff:
                        t.affiliation = _format_if_same(cur_aff, ct["affiliation"])
                    else:
                        found_aff = _find_text_in_blocks(_ns(ct["affiliation"]))
                        if found_aff:
                            t.affiliation = found_aff

        # ---- 信頼度スコアリング（正解DB + VM） ----
        _conf_scores = []

        def _field_conf(cur: str, db: str) -> float:
            """フィールドレベルの信頼度: 1.0=完全一致, 0.8=包含, 0.3=不一致, 0.0=欠損"""
            cur_n = _ns(cur or "")
            db_n = _ns(db or "")
            if not db_n:
                return 1.0  # DB に期待値なし → 常に OK
            if not cur_n:
                return 0.0  # 欠損
            if cur_n == db_n:
                return 1.0  # 完全一致
            if cur_n in db_n or db_n in cur_n:
                return 0.8  # 包含関係
            return 0.3      # 不一致

        # --- 正解DBとの照合 ---
        # chair
        if cc.get("name"):
            _conf_scores.append(("db:chair.name", _field_conf(
                getattr(payload.chair, "name", ""), cc["name"])))
        # organizer
        if correct.get("organizer"):
            _conf_scores.append(("db:organizer", _field_conf(
                payload.organizer, correct["organizer"])))
        # event_title
        if correct.get("event_title"):
            _conf_scores.append(("db:event_title", _field_conf(
                payload.event_title, correct["event_title"])))
        # talks
        payload_talks_only = [t for t in (payload.talks or []) if _is_program_talk_item(t)]
        for i, ct in enumerate(ct_list):
            if ct.get("speaker"):
                sp_val = ""
                if i < len(payload_talks_only):
                    sp_val = getattr(payload_talks_only[i], "speaker", "")
                _conf_scores.append((f"db:talk[{i}].speaker", _field_conf(sp_val, ct["speaker"])))
            if ct.get("title"):
                t_val = ""
                if i < len(payload_talks_only):
                    t_val = getattr(payload_talks_only[i], "title", "")
                _conf_scores.append((f"db:talk[{i}].title", _field_conf(t_val, ct["title"])))

        # --- VMとの照合 ---
        if vm_rows:
            _vm_speakers = []
            _vm_titles = []
            _vm_event_title = ""
            _vm_affs = []
            for _r in vm_rows:
                _d = _r if isinstance(_r, dict) and "data" not in _r else (_r.get("data") or {})
                _sp = _norm_person_name(_d.get("案内状掲載 医師名") or "")
                _vt = normalize_space(_d.get("演題") or "")
                _role = (_d.get("役職") or "").strip()
                _fac = normalize_space(_d.get("案内状掲載 施設名") or "")
                if not _vm_event_title:
                    _vm_event_title = normalize_space(_d.get("講演会名") or "")
                if _role == "演者" and _sp:
                    _vm_speakers.append(_sp)
                    _vm_titles.append(_vt)
                    _vm_affs.append(_fac)

            # VM event_title
            if _vm_event_title:
                _conf_scores.append(("vm:event_title", _field_conf(
                    payload.event_title, _vm_event_title)))

            # VM talks: speaker / title / affiliation
            _payload_sp_map = {}
            for t in (payload.talks or []):
                if _is_program_chair_item(t):
                    continue
                sp_key = _ns(_norm_person_name(getattr(t, "speaker", "") or ""))
                if sp_key:
                    _payload_sp_map[sp_key] = t

            for vi, vm_sp in enumerate(_vm_speakers):
                vm_sp_ns = _ns(vm_sp)
                if vm_sp_ns and vm_sp_ns in _payload_sp_map:
                    _conf_scores.append((f"vm:talk.speaker[{vi}]", 1.0))
                    # title
                    if vi < len(_vm_titles) and _vm_titles[vi]:
                        t_obj = _payload_sp_map[vm_sp_ns]
                        _conf_scores.append((f"vm:talk.title[{vi}]", _field_conf(
                            getattr(t_obj, "title", ""), _vm_titles[vi])))
                    # affiliation (施設名)
                    if vi < len(_vm_affs) and _vm_affs[vi]:
                        t_obj = _payload_sp_map[vm_sp_ns]
                        _conf_scores.append((f"vm:talk.affil[{vi}]", _field_conf(
                            getattr(t_obj, "affiliation", ""), _vm_affs[vi])))
                elif vm_sp_ns:
                    _conf_scores.append((f"vm:talk.speaker[{vi}]", 0.0))
                    print(f"[confidence] vm speaker not found: '{vm_sp}'")

            # VM talk_count
            if _vm_speakers:
                actual_count = sum(1 for t in (payload.talks or []) if _is_program_talk_item(t))
                expected_count = len(_vm_speakers)
                if actual_count == expected_count:
                    _conf_scores.append(("vm:talk_count", 1.0))
                elif abs(actual_count - expected_count) == 1:
                    _conf_scores.append(("vm:talk_count", 0.6))
                else:
                    _conf_scores.append(("vm:talk_count", 0.2))

        if _conf_scores:
            vals = [v for _, v in _conf_scores]
            avg_conf = sum(vals) / len(vals)
            payload.confidence = round(avg_conf, 2)

            # 低信頼フィールドをログ出力
            low_fields = [(name, v) for name, v in _conf_scores if v < 0.8]
            if low_fields:
                for name, v in low_fields:
                    print(f"[confidence] low: {name}={v:.1f}")

            print(f"[confidence] overall={payload.confidence:.2f} ({len(_conf_scores)} fields)")

            # 低信頼警告
            if avg_conf < 0.7:
                _w = list(payload.warnings or [])
                if "low_confidence" not in _w:
                    _w.append("low_confidence")
                payload.warnings = _w

    except Exception as e:
        print(f"[correct-answer-overlay] error: {e}")

    return payload


# ---------------- AI ----------------
def build_ai_prompt(
    blocks: List[TextBlock],
    draft: DesignJSON,
    speaker_map: Dict[str, str],
    time_candidates: List[str],
) -> str:
    blocks_json: List[Dict[str, Any]] = [
        {
            "text": b.text,
            "left": b.left,
            "top": b.top,
            "width": b.width,
            "height": b.height,
            "max_font_pt": round(b.max_font_pt, 2),
        }
        for b in blocks
    ]

    draft_obj = json.loads(draft.model_dump_json() if hasattr(draft, "model_dump_json") else draft.json(ensure_ascii=False))

    # 医療用語辞書を追加（精度向上のため）
    medical_terms = {
        "よくある肩書き": ["教授", "准教授", "講師", "助教", "医師", "部長", "課長", "主任", "センター長"],
        "組織名パターン": ["大学", "病院", "クリニック", "医院", "研究所", "センター", "機構", "学会"],
        "講演タイプ": ["特別講演", "招請講演", "教育講演", "シンポジウム", "ランチョンセミナー", "イブニングセミナー"],
        "時間表記パターン": ["9:00~10:00", "13:30~14:30", "18:00~19:00"],
        "講演として認識しない内容": [
            "Opening Remarks", "Closing Remarks", "開会の辞", "閉会の辞", 
            "開会挨拶", "閉会挨拶", "Welcome", "開催挨拶", "挨拶", 
            "Break", "Coffee Break", "休憩", "Reception", "懇親会",
            "Registration", "受付", "Photo Session", "集合写真",
            "Discussion", "ディスカッション", "Q&A", "質疑応答",
            "パネルディスカッション", "総合討論", "総合討議",
            "Panel Discussion", "General Discussion"
        ]
    }

    # AIには「差分だけ」返させる
    # talks は draft と同じ長さ・同じ順序の配列で返す
    return f"""あなたは、日本語の医療セミナー案内スライドから抽出済みの下書きJSONを、
根拠が明確な箇所だけ最小限修正する専門アシスタントです。

# 最重要ルール
- 出力は JSONのみ
- 推測禁止 - blocks に明示されていない内容は生成禁止
- 下書きJSON(draft)をベースに、必要最小限だけ修正すること
- draft にある情報を勝手に消さないこと
- 不確実な場合は draft の値をそのまま残すこと
- 医療用語の知識を活用して精度を向上させること

# 重要: 講演ではない内容の除外
下記のような内容は講演ではないので talks から除外してください：
- "Opening Remarks", "Closing Remarks", "開会の辞", "閉会の辞"
- "開会挨拶", "閉会挨拶", "Welcome", "開催挨拶", "挨拶"
- "Break", "Coffee Break", "休憩", "Reception", "懇親会"
- "Registration", "受付", "Photo Session", "集合写真"
これらのタイトルを持つ talks 要素は配列から削除すること

# 医療セミナー特有の知識
{json.dumps(medical_terms, ensure_ascii=False, indent=2)}

# 配列・構造の調整ルール
- talks から不適切な内容（挨拶、休憩等）は除外すること
- talks は基本的に実際の講演のみを含めること
- ただしプログラム途中に「座長/司会」行が明示されている場合は、通常講演とは別に item_type="chair" の行として talks 内の該当位置に残すこと
- item_type="chair" の行は title_lines / speaker を空にし、role / name_display / affiliation / program_index を使うこと
- item_type="talk" の行は title_lines / speaker / affiliation を使うこと
- talks の順序は時系列に従って調整可
- event_title_lines の要素数は必要に応じて調整可

# 修正可能な項目（根拠必須）
- event_title_lines: blocks に明示根拠がある場合のみ微修正可
- event_title: event_title_lines を "\\n" で結合したものにする
- chair.role / chair.name / chair.affiliation: blocks に根拠がある場合のみ修正可
- talks[i].title_lines: blocks に根拠がある場合のみ修正可
- talks[i].speaker: speaker_map のキーに完全一致する名前に限り修正可
- talks[i].affiliation: speaker_map[talks[i].speaker] の値のみ使用可
- confidence: 0.0〜1.0 の範囲で設定可（医療情報処理の確実性を反映）
- warnings: 必要なら維持・追加可

# 精度向上のための追加ガイドライン
1. 講演者名は必ず「姓 名」形式で統一
2. 所属は正式名称を優先（略称ではなく）
3. 講演タイトルは改行を適切に保持
4. 時間表記は「HH:MM~HH:MM」形式に統一
5. 敬称（先生、教授など）は適切に処理

# 修正禁止の項目
- talks[i].time は変更禁止
- datetime は変更禁止
- datetime_note は変更禁止
- organizer は変更禁止
- title_overrides は変更禁止
- datetime_parts は変更禁止
- datetime_time_newline は変更禁止
- note は変更禁止
- locked は変更禁止
- manual_override は変更禁止
- region / unit / event_id は変更禁止

# 旧字体・異体字の保持（重要）
- blocks テキストに含まれる旧字体・異体字（髙→高、﨑→崎、邉→辺、齊→斉 等）は
  絶対に新字体に置き換えないこと。blocks にある文字をそのままコピーすること。
- 例: blocks に「髙田」とあれば「髙田」のまま出力。「高田」に変換は禁止。

# speaker_map（このキーにある speaker 名だけ使用可）
{json.dumps(speaker_map, ensure_ascii=False, indent=2)}

# time_candidates（参考のみ。talk.time は変更禁止）
{json.dumps(time_candidates, ensure_ascii=False, indent=2)}

# 抽出ブロック（位置情報も参考にして文脈を理解）
{json.dumps(blocks_json, ensure_ascii=False, indent=2)}

# 下書きJSON(draft)
{json.dumps(draft_obj, ensure_ascii=False, indent=2)}

# 出力形式
- 実際の講演のみを talks に含めること（挨拶、休憩等は除外）
- 各 talks[i] は少なくとも以下を含めること:
  - item_type ("talk" または "chair")
  - time
  - title_lines
  - speaker
  - affiliation
- event_title_lines は配列で返す
- event_title は文字列で返す
- JSON以外の文は禁止
- 修正した項目については confidence を高めに設定
"""

def normalize_chair_role(role: str) -> str:
    r = normalize_space(role or "")
    r = r.replace("\n", "").replace(" ", "")

    # 総合司会だけ特別扱い
    if r == "総合司会":
        return "総合司会"

    # それ以外の司会/座長系は全部「座長」
    if any(x in r for x in ["座長", "司会"]):
        return "座長"

    return r

TIME_NORMALIZE_RE = re.compile(
    r"(\d{1,2})[:：](\d{2})\s*[〜～\-－ー−~]\s*(\d{1,2})[:：](\d{2})"
)

def normalize_time_range_talks(s: str) -> str:
    if not s:
        return ""

    s = normalize_space(s)

    m = TIME_NORMALIZE_RE.search(s)
    if not m:
        return s

    h1, m1, h2, m2 = m.groups()

    return f"{int(h1):02d}:{m1}~{int(h2):02d}:{m2}"

def clean_ai_title_lines(lines: list[str]) -> list[str]:
    out = []

    def _norm(s: str) -> str:
        return normalize_space(s or "")

    def _is_name_line(s: str) -> bool:
        s = _norm(s)
        if not s:
            return False
        if "先生" in s:
            return True
        # 鈴木勇三 / 平野勉 みたいな短い氏名だけの行
        s2 = s.replace(" ", "").replace("　", "")
        return 2 <= len(s2) <= 8 and not any(x in s for x in [
            "大学", "病院", "研究科", "医学部", "センター",
            "講師", "教授", "部長", "医長", "院長", "副会長", "幹事"
        ])

    def _is_aff_line(s: str) -> bool:
        s = _norm(s)
        return any(x in s for x in [
            "大学", "病院", "研究科", "医学部", "センター",
            "クリニック", "講師", "教授", "部長", "医長", "院長",
            "副会長", "幹事", "名誉院長"
        ])

    for ln in lines or []:
        s = _norm(ln)
        if not s:
            continue

        # 外側の引用符だけ除去対象として残す
        if s in ["「", "」", "『", "』"]:
            continue

        # 氏名行・所属行が来たらタイトル終了
        if _is_name_line(s) or _is_aff_line(s):
            break

        out.append(s)

    # 先頭末尾のカッコを剥がす
    if out:
        joined = "\n".join(out).strip()
        pairs = [("「", "」"), ("『", "』"), ("（", "）"), ("(", ")")]
        for l, r in pairs:
            if joined.startswith(l) and joined.endswith(r):
                joined = joined[len(l):-len(r)].strip()
                break
        out = [normalize_space(x) for x in joined.split("\n") if normalize_space(x)]

    return out


def clean_ai_talk_titles(payload: DesignJSON) -> DesignJSON:
    """AI処理後のタイトルクリーニング + 医療用語標準化"""
    def _clean_title_text(s: str) -> str:
        s = normalize_space(s or "")

        # 外側の引用符ペアのみ除去（片方だけ除去しない）
        for open_q, close_q in [('「', '」'), ('『', '』')]:
            if s.startswith(open_q) and s.endswith(close_q) and len(s) >= 2:
                inner = s[1:-1]
                if inner.count(open_q) == inner.count(close_q):
                    s = inner

        # 括弧のバランス補完（半角・全角混在に対応）
        open_full = s.count("（")
        close_full = s.count("）")
        open_half = s.count("(") - open_full  # 念のため重複除外(不要だが安全)
        close_half = s.count(")") - close_full
        # 半角(も全角（も開き括弧として合算して判定
        # ただし count("(") は全角を含まないので単純合算
        total_open = s.count("（") + s.count("(")
        total_close = s.count("）") + s.count(")")
        if total_open > total_close:
            # 全角開きが多い → 全角閉じで補完、半角開きが多い → 半角閉じで補完
            if s.count("（") > s.count("）"):
                s += "）"
            elif s.count("(") > s.count(")"):
                s += ")"

        # 医療用語の標準化を適用
        s = normalize_medical_terms(s)

        return s.strip()

    for t in payload.talks or []:
        lines = t.title_lines or []
        cleaned = [_clean_title_text(x) for x in lines if _clean_title_text(x)]

        # 複数行にまたがる外側カッコの除去
        # 例: ['『(仮)循環器医こそ…', 'つかいどころ』'] → カッコ除去
        if len(cleaned) >= 2:
            first = cleaned[0]
            last = cleaned[-1]
            for open_q, close_q in [('「', '」'), ('『', '』')]:
                if first.startswith(open_q) and last.endswith(close_q):
                    # 内部の開閉カッコ数をチェック
                    joined = "\n".join(cleaned)
                    inner = joined[1:-1]
                    if inner.count(open_q) == inner.count(close_q):
                        cleaned[0] = first[1:]  # 先頭行の開きカッコ除去
                        cleaned[-1] = last[:-1]  # 最終行の閉じカッコ除去
                        # 空行になった場合は除去
                        cleaned = [x for x in cleaned if x.strip()]
                        break

        if cleaned:
            t.title_lines = cleaned
            t.title = "\n".join(cleaned)
            
        # 講演者名・所属の肩書き除去と標準化
        if t.speaker:
            # 肩書き除去 → 医療用語標準化
            cleaned_speaker = normalize_person_name(t.speaker)
            t.speaker = normalize_medical_terms(cleaned_speaker).strip()
        if t.affiliation:
            # 所属に「座長」情報が混入している場合は除去
            aff_text = normalize_space(t.affiliation)
            if "座長" in aff_text:
                # 「座長 ○○ 先生 大阪刀根山…」のようなブロック全体が混入 → クリア
                t.affiliation = ""
            else:
                # 所属から肩書き除去 → 医療用語標準化
                cleaned_affiliation = normalize_affiliation(t.affiliation)
                t.affiliation = normalize_medical_terms(cleaned_affiliation).strip()
    
    # 座長情報の肩書き除去と標準化
    if payload.chair:
        if payload.chair.name:
            cleaned_chair_name = normalize_person_name(payload.chair.name)
            payload.chair.name = normalize_medical_terms(cleaned_chair_name).strip()
        if payload.chair.affiliation:
            cleaned_chair_affiliation = normalize_affiliation(payload.chair.affiliation)
            payload.chair.affiliation = normalize_medical_terms(cleaned_chair_affiliation).strip()
        if payload.chair.role:
            payload.chair.role = normalize_medical_terms(normalize_space(payload.chair.role)).strip()
            
    # 講演ではない内容を除外（Opening/Closing Remarks等）
    if payload.talks:
        non_lecture_keywords = [
            "opening remarks", "closing remarks", "開会の辞", "閉会の辞",
            "開会挨拶", "閉会挨拶", "welcome", "開催挨拶", "挨拶",
            "break", "coffee break", "休憩", "reception", "懇親会",
            "registration", "受付", "photo session", "集合写真",
            "discussion", "ディスカッション", "q&a", "質疑応答",
            "パネルディスカッション", "総合討論", "総合討議",
            "panel discussion", "general discussion"
        ]
        
        # フィルタリング: 講演として適切でない内容を除去
        filtered_talks = []
        for talk in payload.talks:
            if _is_program_chair_item(talk):
                filtered_talks.append(talk)
                continue
            talk_title = (talk.title or "").lower().strip()
            
            # 講演タイトルが除外キーワードに完全一致または実質全体がキーワードかチェック
            is_non_lecture = False
            for keyword in non_lecture_keywords:
                # タイトルがキーワードそのものか、キーワード+少量の装飾のみ
                if talk_title == keyword or talk_title.strip("　 .-–—_") == keyword:
                    is_non_lecture = True
                    break
                # タイトルの80%以上がキーワードで占められている場合
                if keyword in talk_title and len(keyword) >= len(talk_title) * 0.7:
                    is_non_lecture = True
                    break
            
            # タイトルが空の場合のみ除外（短いタイトルは医療用語で有り得る）
            if not is_non_lecture and len(talk_title.replace(" ", "")) > 0:
                filtered_talks.append(talk)
        
        payload.talks = filtered_talks
            
    return payload

async def ai_refine_json(
    blocks: List[TextBlock],
    draft: DesignJSON,
    speaker_map: Dict[str, str],
    time_candidates: List[str],
    vm_rows: list[dict],
) -> DesignJSON:
    if not OPENAI_API_KEY:
        return draft

    prompt = build_ai_prompt(blocks, draft, speaker_map, time_candidates)

    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json",
    }

    # Few-shot例: AIがよく間違えるパターンの正解を示す
    few_shot_examples = [
        # 例1: 1回目/2回目同一内容 → talksは1つだけ
        {
            "user": (
                "blocks に「1回目、2回目ともに同一の内容です」と記載。"
                "講演: 「○○治療の最前線」演者: 山田太郎 先生 / ○○大学 教授。"
                "draft の talks に同一内容が2件入っている。"
            ),
            "assistant": json.dumps({
                "talks": [
                    {
                        "time": "",
                        "title_lines": ["○○治療の最前線"],
                        "speaker": "山田太郎",
                        "affiliation": "○○大学 教授"
                    }
                ],
                "_reason": "1回目2回目同一内容なので talks は1件のみ"
            }, ensure_ascii=False)
        },
        # 例2: 開会挨拶・閉会挨拶は除外
        {
            "user": (
                "draft の talks に3件: 「開会挨拶」「糖尿病治療の新展開」「閉会の辞」。"
                "blocks を確認すると実際の講演は「糖尿病治療の新展開」のみ。"
            ),
            "assistant": json.dumps({
                "talks": [
                    {
                        "time": "",
                        "title_lines": ["糖尿病治療の新展開"],
                        "speaker": "佐藤花子",
                        "affiliation": "○○病院 内科 部長"
                    }
                ],
                "_reason": "開会挨拶・閉会の辞は講演ではないので除外"
            }, ensure_ascii=False)
        },
        # 例3: 旧字体保持 + 演者名の正確な抽出
        {
            "user": (
                "blocks に「髙橋 一郎 先生」「慶應義塾大学医学部 教授」とある。"
                "draft の speaker が「高橋一郎」（新字体に変換済み）。"
            ),
            "assistant": json.dumps({
                "talks": [
                    {
                        "time": "",
                        "title_lines": ["腎臓病の最新治療"],
                        "speaker": "髙橋一郎",
                        "affiliation": "慶應義塾大学医学部 教授"
                    }
                ],
                "_reason": "blocks の旧字体「髙」をそのまま使用"
            }, ensure_ascii=False)
        },
    ]

    few_shot_messages = []
    for ex in few_shot_examples:
        few_shot_messages.append({"role": "user", "content": ex["user"]})
        few_shot_messages.append({"role": "assistant", "content": ex["assistant"]})

    # 正解DBから類似例を動的に追加
    all_blocks_text = " ".join(b.text for b in blocks)
    event_title = draft.event_title or ""
    similar = find_similar_correct_answers(all_blocks_text, event_title, top_k=2)
    dynamic_shots = _build_dynamic_few_shot(similar)
    if dynamic_shots:
        few_shot_messages.extend(dynamic_shots)

    body = {
        "model": AI_MODEL,
        "messages": [
            {
                "role": "system",
                "content": (
                    "You are a medical seminar document processing specialist. "
                    "Return ONLY valid JSON object based on the precise instructions. "
                    "Focus on accuracy and consistency in medical terminology. "
                    "Do not add explanations. "
                    "Do not wrap in markdown."
                ),
            },
            *few_shot_messages,
            {"role": "user", "content": prompt},
        ],
        "temperature": 0.05,  # より決定論的な出力で一貫性向上
        "top_p": 0.9,        # 高品質なトークンを優先
        "presence_penalty": 0.0,   # 中立（負値はトークン繰り返しを促進し逆効果）
        "frequency_penalty": 0.1,  # 繰り返しを減らして品質向上
        "response_format": {"type": "json_object"},
        "max_tokens": 4000,  # 十分な応答長を確保
    }

    try:
        async with httpx.AsyncClient(timeout=AI_TIMEOUT) as client:
            r = await post_with_retry(
                client,
                f"{OPENAI_BASE_URL}/chat/completions",
                headers=headers,
                json_body=body,
                retries=3,
            )
            if r.status_code >= 400:
                try:
                    print("error json=", r.json())
                except Exception:
                    print("error text=", r.text)
            r.raise_for_status()
            data = r.json()
    except Exception as e:
        print(f"[ai_refine_json] request failed: {type(e).__name__}: {e}")
        draft.warnings = sorted(set((draft.warnings or []) + ["ai_request_failed"]))
        return draft

    content = (data["choices"][0]["message"]["content"] or "").strip()

    def _is_plausible_speaker_name(s: str) -> bool:
        """文字列が日本人名として妥当かチェック"""
        s = normalize_space(s or "").replace("先生", "").strip()
        s_compact = s.replace(" ", "").replace("\u3000", "")
        if not s_compact:
            return False
        # 長さチェック: 2〜8文字
        if len(s_compact) < 2 or len(s_compact) > 8:
            return False
        # 漢字が中心か（ひらがな・カタカナ名も一応許容）
        kanji_kana = sum(1 for c in s_compact if '\u4e00' <= c <= '\u9fff' or '\u3040' <= c <= '\u30ff' or '\u3400' <= c <= '\u4dbf' or '\uf900' <= c <= '\ufaff')
        if kanji_kana < len(s_compact) * 0.7:
            return False
        # 明らかに人名でないキーワード
        bad_words = [
            "大学", "病院", "センター", "クリニック", "医院", "研究", "科",
            "教授", "講師", "部長", "医長", "セミナー", "講演", "株式会社",
            "治療", "経験", "検討", "課題", "予防", "医療", "診療",
            "主催", "共催", "座長", "演者", "手術", "管理", "使用",
            "シンポジウム", "プログラム", "休憩", "質疑応答", "挨拶",
            "開会", "閉会", "司会", "後援",
        ]
        if any(w in s_compact for w in bad_words):
            return False
        return True

    def _looks_bad_talk_seed_strong(t) -> bool:
        """強化された講演データ品質チェック"""
        if _is_program_chair_item(t):
            return False
        title_lines = getattr(t, "title_lines", None) or []
        title = normalize_space(getattr(t, "title", "") or "")
        speaker = normalize_space(getattr(t, "speaker", "") or "")
        affiliation = normalize_space(getattr(t, "affiliation", "") or "")

        full_title = "\n".join([normalize_space(x) for x in title_lines if normalize_space(x)]) or title

        # title が無い
        if not full_title:
            return True

        # speaker が無い
        if not speaker:
            return True

        # speaker が人名として不正
        if not _is_plausible_speaker_name(speaker):
            return True

        # 医療セミナー特有の無効パターンをチェック
        bad_speakers = {
            "課題", "腎移植", "逐次薬物治療", "治療", "講演", "演者",
            "胆道", "肝内胆管癌治療", "肝外胆管癌治療", "休憩", "質疑応答",
            "開会", "閉会", "挨拶", "司会", "座長", "コーヒーブレイク",
            "使用経験", "検討", "株式会社", "セミナー", "シンポジウム",
            "プログラム", "主催", "共催", "後援",
        }
        if speaker in bad_speakers:
            return True

        # 講演者名の妥当性チェック（日本人名のパターン）
        if len(speaker.replace(" ", "").replace("　", "")) <= 1:
            return True
        
        # 明らかに講演者でない文字列
        if any(x in speaker for x in ["時間", "分", "場所", "会場", "Ｘ", "○", "●"]):
            return True

        # affiliation が無い or 生の演者行
        if not affiliation:
            return True
        if "演者" in affiliation and "先生" in affiliation:
            return True
            
        # 所属の妥当性チェック
        if len(affiliation.replace(" ", "").replace("　", "")) <= 2:
            return True
            
        # タイトルの妥当性チェック
        if any(x in full_title for x in ["休憩", "コーヒー", "質疑", "開会", "閉会"]):
            return True

        return False

    def _build_talks_from_parsed(parsed_talks: list[dict]) -> list[Talk]:
        out = []

        for pt in parsed_talks[:6]:
            if not isinstance(pt, dict):
                continue

            item_type = _json_item_type(pt)
            try:
                program_index = int(pt.get("program_index", len(out)) or len(out))
            except Exception:
                program_index = len(out)

            if item_type == "chair":
                name_display = normalize_person_display(
                    pt.get("name_display")
                    or pt.get("name")
                    or pt.get("speaker_display")
                    or pt.get("speaker")
                    or ""
                )
                affiliation = normalize_space(pt.get("affiliation", "") or "")
                if not (name_display or affiliation):
                    continue
                out.append(
                    Talk(
                        item_type="chair",
                        role=normalize_chair_role(pt.get("role", "") or "座長"),
                        program_index=program_index,
                        name_display=name_display,
                        speaker="",
                        speaker_display="",
                        affiliation=affiliation,
                        honorific_title=normalize_space(pt.get("honorific_title", "") or "先生"),
                    )
                )
                continue

            title_lines = normalize_lines_keep_order(_safe_list_str(pt.get("title_lines")))
            title = "\n".join(title_lines).strip() if title_lines else normalize_space(pt.get("title", "") or "")
            speaker = _norm_speaker_candidate(pt.get("speaker", "") or "")
            affiliation = normalize_space(pt.get("affiliation", "") or "")
            time = normalize_time_range_talks(pt.get("time", "") or "")

            out.append(
                Talk(
                    item_type="talk",
                    role=normalize_space(pt.get("role", "") or "演者") or "演者",
                    program_index=program_index,
                    time=time,
                    title=title,
                    title_lines=title_lines if title_lines else ([title] if title else []),
                    speaker=speaker,
                    speaker_display=build_speaker_display(speaker) if speaker else "",
                    affiliation=affiliation,
                    title_overrides=[],
                    honorific_title="先生",
                )
            )

        return out

    def _extract_json_text(s: str) -> str:
        s = (s or "").strip()
        s = re.sub(r"^\s*```json\s*", "", s, flags=re.IGNORECASE)
        s = re.sub(r"^\s*```\s*", "", s)
        s = re.sub(r"\s*```+\s*$", "", s)
        m = re.search(r"\{.*\}", s, flags=re.DOTALL)
        if m:
            s = m.group(0)
        return s.strip()

    def _escape_control_chars_in_json_strings(s: str) -> str:
        out = []
        in_string = False
        escape = False

        for ch in s:
            if in_string:
                if escape:
                    out.append(ch)
                    escape = False
                    continue

                if ch == "\\":
                    out.append(ch)
                    escape = True
                    continue

                if ch == '"':
                    out.append(ch)
                    in_string = False
                    continue

                if ch == "\n":
                    out.append("\\n")
                    continue
                if ch == "\r":
                    out.append("\\r")
                    continue
                if ch == "\t":
                    out.append("\\t")
                    continue
                if ord(ch) < 0x20:
                    out.append(f"\\u{ord(ch):04x}")
                    continue

                out.append(ch)
            else:
                out.append(ch)
                if ch == '"':
                    in_string = True
                    escape = False

        return "".join(out)

    def _fix_unterminated_string_lines(s: str) -> str:
        fixed = []

        for line in s.splitlines():
            q = 0
            escape = False

            for ch in line:
                if escape:
                    escape = False
                    continue
                if ch == "\\":
                    escape = True
                    continue
                if ch == '"':
                    q += 1

            if q % 2 == 1:
                stripped = line.rstrip()
                if stripped.endswith(","):
                    stripped = stripped[:-1] + '",'
                else:
                    stripped = stripped + '"'
                line = stripped

            fixed.append(line)

        return "\n".join(fixed)

    def parse_llm_json(content: str) -> dict:
        s = _extract_json_text(content)
        s = s.replace("“", '"').replace("”", '"').replace("’", "'")
        s = re.sub(r"\bNone\b", "null", s)
        s = re.sub(r"\bTrue\b", "true", s)
        s = re.sub(r"\bFalse\b", "false", s)
        s = _escape_control_chars_in_json_strings(s)
        s = _fix_unterminated_string_lines(s)
        s = re.sub(r",\s*([}\]])", r"\1", s)
        return json.loads(s)

    try:
        parsed = json.loads(content)
    except Exception:
        try:
            parsed = parse_llm_json(content)
        except Exception:
            print("RAW CONTENT >>>")
            print(content)
            print("<<< RAW CONTENT")
            draft.warnings = sorted(set((draft.warnings or []) + ["ai_json_parse_failed"]))
            return draft

    # -------------------------------
    # ここから「丸ごと置換」ではなく「draftに部分マージ」
    # -------------------------------
    refined = draft.model_copy(deep=True) if hasattr(draft, "model_copy") else DesignJSON(**json.loads(draft.json(ensure_ascii=False)))

    def _safe_list_str(v) -> list[str]:
        if not isinstance(v, list):
            return []
        out = []
        for x in v:
            s = normalize_space(str(x or ""))
            if s:
                out.append(s)
        return out

    def _norm_speaker_candidate(s: str) -> str:
        s = norm_name(s or "")
        return normalize_space(s)

    def _same_len_talks(parsed_talks: Any, draft_talks: List[Talk]) -> bool:
        return isinstance(parsed_talks, list) and len(parsed_talks) == len(draft_talks)

    def _looks_bad_talk_seed(t) -> bool:
        if _is_program_chair_item(t):
            return False
        title_lines = getattr(t, "title_lines", None) or []
        title = normalize_space(getattr(t, "title", "") or "")
        speaker = normalize_space(getattr(t, "speaker", "") or "")
        affiliation = normalize_space(getattr(t, "affiliation", "") or "")

        full_title = "\n".join([normalize_space(x) for x in title_lines if normalize_space(x)]) or title

        # title/speaker が空 or 明らかに壊れてる
        if not full_title:
            return True
        if not speaker:
            return True
        if any(x == speaker for x in ["胆道", "肝内胆管癌治療", "肝外胆管癌治療", "講演"]):
            return True
        if len(speaker.replace(" ", "").replace("　", "")) <= 1:
            return True
        # タイトルから誤抽出された明らかに無効な演者名
        if speaker in ["男子", "女子", "キーワード", "検査項目", "薬剤名", "症例"]:
            return True
        # 演者名がタイトルの引用符内テキストと一致（「男子」等の誤抽出）
        if full_title and f"「{speaker}」" in full_title:
            return True
        if not affiliation:
            return True

        return False

    def _is_valid_ai_patch(parsed: dict, draft: DesignJSON) -> bool:
        if not isinstance(parsed, dict):
            return False

        parsed_talks = parsed.get("talks")

        if "talks" in parsed:
            if not isinstance(parsed_talks, list):
                return False

            parsed_has_chairs = any(_json_is_chair_item(pt) for pt in parsed_talks if isinstance(pt, dict))
            draft_talk_items = [t for t in (draft.talks or []) if _is_program_talk_item(t)]
            expected_len = len(draft.talks or []) if parsed_has_chairs else len(draft_talk_items)
            parsed_len = len(parsed_talks)

            if draft.talks and parsed_len != expected_len:
                # draft側が壊れている場合はAI側の講演数を信頼する
                if any(_looks_bad_talk_seed(t) for t in draft_talk_items):
                    pass  # 壊れたseedがあるので数の不一致を許容
                # AIが講演数を減らした場合（挨拶・休憩等の除外）は許容
                elif parsed_len < expected_len:
                    pass  # non-lecture除外による減少は正当
                else:
                    return False

            draft_compare_items = list(draft.talks or []) if parsed_has_chairs else draft_talk_items
            for i, pt in enumerate(parsed_talks):
                if not isinstance(pt, dict):
                    return False
                if _json_is_chair_item(pt):
                    continue

                # draft が壊れてる talk は厳格チェックしない
                if draft_compare_items and i < len(draft_compare_items):
                    seed = draft_compare_items[i]
                    if _looks_bad_talk_seed(seed):
                        continue

                    sp = _norm_speaker_candidate(pt.get("speaker", "") or "")
                    if sp and speaker_map and sp not in speaker_map:
                        return False

                    pt_time = normalize_time_range_talks(pt.get("time", "") or "")
                    dr_time = normalize_time_range_talks(getattr(seed, "time", "") or "")
                    if pt_time and dr_time and pt_time != dr_time:
                        return False

        return True

    if not _is_valid_ai_patch(parsed, draft):
        # draft が壊れている案件は、AI全採用にフォールバック
        bad_seed = (
            not draft.talks or
            any(_looks_bad_talk_seed(t) for t in (draft.talks or []))
        )

        if bad_seed:
            try:
                refined = DesignJSON(**parsed)
                refined.warnings = sorted(set((refined.warnings or []) + ["ai_refined_fallback"]))
                return refined
            except Exception:
                pass

        draft.warnings = sorted(set((draft.warnings or []) + ["ai_patch_rejected"]))
        return draft

    def fix_spaced_english(s: str) -> str:
        s = normalize_space(s or "")

        # 完全に1文字ずつ空いてる英字
        if re.fullmatch(r"[A-Za-z](?:\s+[A-Za-z]){2,}", s):
            return s.replace(" ", "")

        # 1語ずつ判定して "U p d a t e" → "Update"
        words = s.split(" ")
        merged = []
        buf = []

        def flush_buf():
            nonlocal buf, merged
            if not buf:
                return
            token = " ".join(buf)
            if re.fullmatch(r"[A-Za-z](?:\s+[A-Za-z]){1,}", token):
                merged.append(token.replace(" ", ""))
            else:
                merged.append(token)
            buf = []

        for w in words:
            if re.fullmatch(r"[A-Za-z]", w):
                buf.append(w)
            else:
                flush_buf()
                merged.append(w)

        flush_buf()
        return " ".join(merged)
    

    # event_title_lines / event_title
    if "event_title_lines" in parsed:
        lines = normalize_lines_keep_order(_safe_list_str(parsed.get("event_title_lines")))
        draft_lines = draft.event_title_lines or []
        if lines:
            # AI がタイトル行を調整（増減両方許可）: 不要行の除去も精度向上に必要
            refined.event_title_lines = lines
            refined.event_title = "\n".join(lines).strip()

    elif getattr(refined, "event_title_lines", None):
        refined.event_title_lines = normalize_lines_keep_order(refined.event_title_lines or [])
        refined.event_title = "\n".join(refined.event_title_lines).strip()

    refined.event_title_lines = [
        fix_spaced_english(x.replace("\n", " "))
        for x in (refined.event_title_lines or [])
    ]

    refined.event_title = fix_spaced_english(
        (refined.event_title or "").replace("\n", " ")
    )

    if refined.event_title_lines:
        refined.event_title = "\n".join(refined.event_title_lines).strip()
    elif refined.event_title:
        refined.event_title_lines = [refined.event_title]

    # chair
    chair_patch = parsed.get("chair")
    _ai_swapped_chair_speaker = False  # AIが座長と演者を逆にした場合のフラグ
    if isinstance(chair_patch, dict):
        role = normalize_space(chair_patch.get("role", "") or "")
        name = normalize_person_name(chair_patch.get("name", "") or "")
        aff = normalize_space(chair_patch.get("affiliation", "") or "")

        # draftに座長名が既に抽出されているか
        current_chair_has_name = (
            refined.chair and 
            refined.chair.name and 
            len(refined.chair.name.replace(" ", "").replace("\u3000", "")) >= 2
        )

        if current_chair_has_name and name:
            draft_name_ns = refined.chair.name.replace(" ", "").replace("\u3000", "")
            ai_name_ns = name.replace(" ", "").replace("\u3000", "")
            if draft_name_ns != ai_name_ns:
                # ドラフトのspeakerが全て不正なら、座長抽出自体が間違っている可能性が高い
                # → AI変更を受け入れる
                _speaker_items = [t for t in (refined.talks or []) if _is_program_talk_item(t)]
                draft_speakers_all_bad = bool(_speaker_items) and all(
                    not _is_plausible_speaker_name(getattr(t, "speaker", "") or "")
                    for t in _speaker_items
                )
                if draft_speakers_all_bad:
                    # ドラフトの座長抽出が誤り → AI結果を採用
                    print(f"[AI REFINE DEBUG] ドラフト speaker 全不正 → AI座長名を採用: {refined.chair.name} → {name}")
                    if role:
                        refined.chair.role = normalize_chair_role(role)
                    refined.chair.name = name
                    if aff:
                        refined.chair.affiliation = aff
                else:
                    # AIが座長名を別人に変更しようとしている → 名前変更を拒否
                    # affiliationはブロックから後工程で補完される
                    print(f"[AI REFINE DEBUG] AI座長名変更を拒否: {refined.chair.name} → {name}")
                    # AIが座長と演者を逆転させた可能性が高い → talk側のspeaker/affiliationも保護
                    _ai_swapped_chair_speaker = True
            else:
                # 同一人物 → affiliation等の更新を許可
                if role:
                    refined.chair.role = normalize_chair_role(role)
                if aff:
                    refined.chair.affiliation = aff
        elif not current_chair_has_name:
            # draftに座長名がない場合のみAI結果を全面採用
            if role:
                refined.chair.role = normalize_chair_role(role)
            if name:
                refined.chair.name = name
            if aff:
                refined.chair.affiliation = aff

        if getattr(refined.chair, "name_display", ""):
            refined.chair.name_display = normalize_person_display(refined.chair.name_display)

    # talks: 件数・順序固定で index ごとにパッチ
    parsed_talks = parsed.get("talks")

    draft_regular_talks = [t for t in (refined.talks or []) if _is_program_talk_item(t)]
    parsed_regular_talks = [
        pt for pt in (parsed_talks or [])
        if isinstance(pt, dict) and _json_is_talk_item(pt)
    ] if isinstance(parsed_talks, list) else []
    parsed_has_chairs = any(
        _json_is_chair_item(pt) for pt in (parsed_talks or []) if isinstance(pt, dict)
    ) if isinstance(parsed_talks, list) else False
    draft_has_chairs = any(_is_program_chair_item(t) for t in (refined.talks or []))

    draft_is_bad = bool(draft_regular_talks) and any(_looks_bad_talk_seed_strong(t) for t in draft_regular_talks)

    # 1) draft が空、または draft が壊れてるなら AI talks を丸ごと採用
    if isinstance(parsed_talks, list) and (not refined.talks or draft_is_bad):
        # AIが座長演者を逆転させた場合、draftに有効なspeaker/affiliationがあればそれを保持
        if _ai_swapped_chair_speaker and refined.talks:
            draft_people = [(t.speaker, t.affiliation) for t in draft_regular_talks if t.speaker]
            refined.talks = _build_talks_from_parsed(parsed_talks)
            # draft側のspeaker/affiliationで上書き復元（人名として妥当な場合のみ）
            j = 0
            for i, t in enumerate(refined.talks):
                if _is_program_chair_item(t):
                    continue
                if j < len(draft_people) and draft_people[j][0] and _is_plausible_speaker_name(draft_people[j][0]):
                    t.speaker = draft_people[j][0]
                    t.speaker_display = build_speaker_display(draft_people[j][0]) or draft_people[j][0]
                    if draft_people[j][1]:
                        t.affiliation = draft_people[j][1]
                    print(f"[AI REFINE DEBUG] AI座長演者逆転: talk[{i}] speaker/affiliation復元 ({t.speaker})")
                elif j < len(draft_people) and draft_people[j][0]:
                    print(f"[AI REFINE DEBUG] AI座長演者逆転: talk[{i}] draft speaker不正のためAI値を維持 ({draft_people[j][0]})")
                j += 1
        else:
            refined.talks = _build_talks_from_parsed(parsed_talks)

    # 1.5) draft に途中座長があり、AI が通常講演だけを返した場合は座長行を保持して通常講演だけパッチ
    elif (
        isinstance(parsed_talks, list)
        and draft_has_chairs
        and not parsed_has_chairs
        and len(parsed_regular_talks) == len(draft_regular_talks)
    ):
        for i, (t, pt) in enumerate(zip(draft_regular_talks, parsed_regular_talks)):
            pt_time = normalize_time_range_talks(pt.get("time", "") or "")
            if pt_time:
                t.time = pt_time

            title_lines = normalize_lines_keep_order(_safe_list_str(pt.get("title_lines")))
            if title_lines:
                t.title_lines = title_lines
                t.title = "\n".join(title_lines).strip()

            if _ai_swapped_chair_speaker and t.speaker and t.affiliation and _is_plausible_speaker_name(t.speaker):
                print(f"[AI REFINE DEBUG] AI座長演者逆転検出: talk[{i}] speaker/affiliation保持 ({t.speaker})")
            else:
                sp = _norm_speaker_candidate(pt.get("speaker", "") or "")
                if sp:
                    t.speaker = sp
                    t.speaker_display = build_speaker_display(sp) or sp

                aff = normalize_space(pt.get("affiliation", "") or "")
                if aff:
                    t.affiliation = aff

    # 2) AI が講演数を減らした場合（挨拶・休憩等の除外）→ AI talks を採用
    elif (
        isinstance(parsed_talks, list)
        and len(parsed_talks) < len(refined.talks)
        and not (draft_has_chairs and not parsed_has_chairs)
    ):
        # AIが除外した分が妥当かチェック: AI側の全talkが有効であれば採用
        ai_talks = _build_talks_from_parsed(parsed_talks)
        if ai_talks and not any(_looks_bad_talk_seed_strong(t) for t in ai_talks):
            refined.talks = ai_talks

    # 3) draft が健全な案件だけ index patch
    elif isinstance(parsed_talks, list) and len(parsed_talks) == len(refined.talks):
        for i, pt in enumerate(parsed_talks):
            if not isinstance(pt, dict):
                continue

            t = refined.talks[i]

            if _is_program_chair_item(t) or _json_is_chair_item(pt):
                if _json_is_chair_item(pt):
                    t.item_type = "chair"
                    t.role = normalize_chair_role(pt.get("role", "") or getattr(t, "role", "") or "座長")
                    name_display = normalize_person_display(
                        pt.get("name_display")
                        or pt.get("name")
                        or getattr(t, "name_display", "")
                        or ""
                    )
                    if name_display:
                        t.name_display = name_display
                    aff = normalize_space(pt.get("affiliation", "") or "")
                    if aff:
                        t.affiliation = aff
                    t.speaker = ""
                    t.speaker_display = ""
                continue

            # time は AI を採用
            pt_time = normalize_time_range_talks(pt.get("time", "") or "")
            if pt_time:
                t.time = pt_time

            title_lines = normalize_lines_keep_order(_safe_list_str(pt.get("title_lines")))
            if title_lines:
                t.title_lines = title_lines
                t.title = "\n".join(title_lines).strip()

            # AIが座長と演者を逆転させた場合、draftのspeaker/affiliationを保持（人名として妥当な場合のみ）
            if _ai_swapped_chair_speaker and t.speaker and t.affiliation and _is_plausible_speaker_name(t.speaker):
                print(f"[AI REFINE DEBUG] AI座長演者逆転検出: talk[{i}] speaker/affiliation保持 ({t.speaker})")
            elif _ai_swapped_chair_speaker and t.speaker and not _is_plausible_speaker_name(t.speaker):
                # draft speaker が不正 → AI値を採用
                sp = _norm_speaker_candidate(pt.get("speaker", "") or "")
                if sp:
                    t.speaker = sp
                    t.speaker_display = build_speaker_display(sp) or sp
                aff = normalize_space(pt.get("affiliation", "") or "")
                if aff:
                    t.affiliation = aff
                print(f"[AI REFINE DEBUG] AI座長演者逆転検出: talk[{i}] draft speaker不正のためAI値を採用 ({t.speaker})")
            
            else:
                sp = _norm_speaker_candidate(pt.get("speaker", "") or "")
                if sp:
                    t.speaker = sp
                    t.speaker_display = build_speaker_display(sp) or sp

                aff = normalize_space(pt.get("affiliation", "") or "")
                if aff:
                    t.affiliation = aff

    # confidence
    if "confidence" in parsed:
        try:
            conf = float(parsed.get("confidence"))
            refined.confidence = max(0.0, min(1.0, conf))
        except Exception:
            pass

    # warnings
    w = set(refined.warnings or [])
    for x in (parsed.get("warnings") or []):
        sx = normalize_space(str(x or ""))
        if sx:
            w.add(sx)

    # -------------------------------
    # 以降は既存の rule-based 後処理を維持
    # -------------------------------
    DATETIME_NOTE_PAT = re.compile(r"^[※\*]\s*.+")
    DATE_PAT = re.compile(r"(20\d{2}年\s*\d{1,2}月\s*\d{1,2}日(?:\s*（[^）]+）)?)")
    TIME_RANGE_PAT2 = re.compile(r"(\d{1,2}:\d{2}\s*[～〜\-ー~]\s*\d{1,2}:\d{2})")

    def extract_datetime_note_from_blocks(blocks: List[TextBlock]) -> str:
        if not blocks:
            return ""

        def norm(s: str) -> str:
            return normalize_space(s or "")

        datetime_blocks: List[TextBlock] = []
        for b in blocks:
            txt = norm(b.text)
            txt2 = normalize_time_colon(txt)

            has_date = bool(DATE_PAT.search(txt))
            has_time = bool(TIME_RANGE_PAT2.search(txt2))
            has_datetime_label = "日時" in txt

            if has_date or has_time or has_datetime_label:
                datetime_blocks.append(b)

        if not datetime_blocks:
            return ""

        note_candidates: List[TextBlock] = []
        for b in blocks:
            txt = norm(b.text)

            if not DATETIME_NOTE_PAT.match(txt):
                continue
            if len(txt) > 60:
                continue
            if "\n" in txt and len(txt.splitlines()) >= 3:
                continue

            ng_keywords = [
                "ご視聴", "事前参加", "旅費", "ご了承ください", "担当者へご連絡",
                "医療従事者", "学生", "ご参加はご遠慮"
            ]
            if any(k in txt for k in ng_keywords):
                continue

            note_candidates.append(b)

        if not note_candidates:
            return ""

        best_text = ""
        best_score = None

        for note in note_candidates:
            note_txt = norm(note.text)

            for dt in datetime_blocks:
                dx = abs(note.left - dt.left)
                dy = abs(note.top - dt.top)
                below_bonus = 0 if note.top >= dt.top else 200000

                if dy > 1200000:
                    continue
                if dx > 5000000:
                    continue

                score = dy + dx * 0.15 + below_bonus

                if best_score is None or score < best_score:
                    best_score = score
                    best_text = note_txt

        return best_text

    def is_honsha_vm(vm_rows):
        return any(
            "VM(本社)" in (r.get("_presence_sheets") or [])
            for r in vm_rows if isinstance(r, dict)
        )

    def extract_datetime_from_blocks_v2(blocks: List[TextBlock]) -> str:
        if not blocks:
            return ""

        ordered = sorted(blocks, key=lambda b: (b.top, b.left))

        DATE_PAT = re.compile(r"(20\d{2}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日(?:\s*[（(][^）)]+[）)])?)")
        TIME_PAT = re.compile(r"(\d{1,2}[:：]\d{2}\s*[～〜\-ー−－—–~]\s*\d{1,2}[:：]\d{2})")

        # 1) 同一block内に date+time があるものを最優先
        for b in ordered:
            txt = normalize_datetime_text(b.text or "")
            md = DATE_PAT.search(txt)
            mt = TIME_PAT.search(txt)
            if md and mt:
                return normalize_space(f"{md.group(1)} {normalize_time_range(mt.group(1))}")

        # 2) date block と近傍 time block の結合
        date_block = None
        date_str = ""
        for b in ordered:
            txt = normalize_datetime_text(b.text or "")
            md = DATE_PAT.search(txt)
            if md:
                date_block = b
                date_str = md.group(1)
                break

        if date_block and date_str:
            best_time = ""
            best_score = None
            for b in ordered:
                txt = normalize_datetime_text(b.text or "")
                mt = TIME_PAT.search(txt)
                if not mt:
                    continue

                tm = normalize_time_range(mt.group(1))
                score = abs(b.top - date_block.top) + abs(b.left - date_block.left) * 0.15
                if best_score is None or score < best_score:
                    best_score = score
                    best_time = tm

            if best_time:
                return normalize_space(f"{date_str} {best_time}")
            return normalize_space(date_str)

        return ""

    # rule優先
    rule_dt = extract_datetime_from_blocks_v2(blocks)
    if rule_dt:
        # multi-session time（N回目）が既に入っている場合は単一時間で上書きしない
        current_has_multi = "回目" in (refined.datetime or "")
        rule_has_multi = "回目" in rule_dt
        if current_has_multi and not rule_has_multi:
            pass  # 既存のmulti-session datetimeを保持
        else:
            refined.datetime = normalize_space(rule_dt)

    rule_dt_note = extract_datetime_note_from_blocks(blocks)
    if is_honsha_vm(vm_rows):
        refined.datetime_note = normalize_space(rule_dt_note) if rule_dt_note else ""
    else:
        refined.datetime_note = ""

    if not refined.organizer:
        refined.organizer = extract_organizer_from_blocks(blocks)

    rule_org = extract_organizer_from_blocks(blocks)
    if rule_org:
        refined.organizer = normalize_organizer(rule_org)
    elif refined.organizer:
        refined.organizer = normalize_organizer(refined.organizer)
    refined.datetime = normalize_space(refined.datetime)
    refined.datetime_note = normalize_space(refined.datetime_note)

    refined.chair.role = normalize_chair_role(normalize_space(getattr(refined.chair, "role", "") or ""))
    refined.chair.name = normalize_person_name(getattr(refined.chair, "name", "") or "")
    if getattr(refined.chair, "name_display", ""):
        refined.chair.name_display = normalize_person_display(refined.chair.name_display)
    refined.chair.affiliation = normalize_space(getattr(refined.chair, "affiliation", "") or "")

    # --- 旧字体復元: AI が髙→高 等に置き換えた文字をブロック原文から復元 ---
    _OLD_NEW_KANJI = [
        ("髙", "高"), ("﨑", "崎"), ("邉", "辺"), ("邊", "辺"),
        ("齊", "斉"), ("齋", "斉"), ("渡邉", "渡辺"), ("渡邊", "渡辺"),
        ("廣", "広"), ("櫻", "桜"), ("國", "国"), ("壽", "寿"),
        ("眞", "真"), ("實", "実"), ("惠", "恵"), ("發", "発"),
    ]
    _all_block_text = " ".join(b.text or "" for b in blocks)
    _all_block_text_flat = _all_block_text.replace("\n", "").replace(" ", "").replace("\u3000", "")

    def _restore_old_kanji(name: str) -> str:
        """blocks 原文に旧字体があれば、AI が新字体に置換した部分を復元"""
        if not name:
            return name
        for old_ch, new_ch in _OLD_NEW_KANJI:
            if old_ch in _all_block_text and new_ch in name:
                # 名前のスペースなし版を作って照合
                name_nospace = name.replace(" ", "").replace("\u3000", "")
                restored = name_nospace.replace(new_ch, old_ch)
                if restored in _all_block_text_flat:
                    name = name.replace(new_ch, old_ch)
        return name

    refined.chair.name = _restore_old_kanji(refined.chair.name)
    refined.chair.name_display = _restore_old_kanji(refined.chair.name_display)


    def clean_speaker_text_strict(s: str) -> str:
        s = normalize_space(s or "")
        s = normalize_time_colon(s)
        s = re.sub(r"\d{1,2}[:：]\d{2}\s*[-~～〜－—–]\s*\d{1,2}[:：]\d{2}", "", s)
        s = re.sub(r"^(演者|座長|講師)\s*[:：]?\s*", "", s)
        s = s.replace("\n", " ").strip()
        s = re.sub(r"\s*先生$", "", s)
        return norm_name(s)

    def compact_talk_title_lines(lines: list[str]) -> list[str]:
        lines = [normalize_space(x) for x in (lines or []) if normalize_space(x)]
        if not lines:
            return []

        out = []
        i = 0
        while i < len(lines):
            cur = lines[i]

            # 次行と結合したいパターン
            if i + 1 < len(lines):
                nxt = lines[i + 1]

                # 英文 + 日本語の続き
                if re.search(r"[A-Za-z]$", cur) and re.search(r"^[ぁ-んァ-ヶ一-龥A-Za-z]", nxt):
                    cur = cur + nxt
                    i += 1

                # 極端に短い末尾語を前行に結合
                elif len(nxt) <= 4 and re.fullmatch(r"[A-Za-zぁ-んァ-ヶ一-龥]+", nxt):
                    cur = cur + nxt
                    i += 1

            out.append(cur)
            i += 1

        # さらに 2行目以降が短すぎるならまとめる
        if len(out) >= 3:
            merged = [out[0]]
            tail = "".join(out[1:])
            if len(tail) <= 28:
                merged.append(tail)
                return merged

        return out

    for t in refined.talks:
        t.time = normalize_time_range_talks(t.time)
        t.title_lines = compact_talk_title_lines(t.title_lines or [])
        t.title_lines = normalize_lines_keep_order(t.title_lines or [])
        if t.title_lines:
            t.title = "\n".join(t.title_lines).strip()
        t.speaker = clean_speaker_text_strict(getattr(t, "speaker", "") or "")
        t.speaker = norm_name(t.speaker)
        t.speaker = _restore_old_kanji(t.speaker)
        if speaker_map and t.speaker in speaker_map:
            t.affiliation = normalize_space(speaker_map[t.speaker] or "")
        else:
            t.affiliation = normalize_space(t.affiliation)

    refined = postprocess_refined(refined, speaker_map, time_candidates)
    refined.talks = sort_talks_by_layout(blocks, refined.talks)
    refined.talks = assign_talk_times_by_nearest_upper_time(blocks, refined.talks)

    w.add("ai_refined")
    if not refined.datetime:
        w.add("missing_datetime")
    if not refined.organizer:
        w.add("missing_organizer")
    if not refined.chair.name:
        w.add("missing_chair")
    if len(refined.talks) == 0:
        w.add("no_talks")
    if not refined.event_title_lines and not refined.event_title:
        w.add("missing_event_title")

    refined.warnings = sorted(w)
    return refined

def find_talk_anchor_top(blocks: list[TextBlock], talk: Talk) -> int:
    candidates = []

    title_lines = [normalize_space(x) for x in (talk.title_lines or []) if normalize_space(x)]
    speaker_keys = [
        normalize_space(getattr(talk, "speaker", "") or ""),
        normalize_space(getattr(talk, "speaker_display", "") or ""),
    ]
    speaker_keys = [x for x in speaker_keys if x]

    for b in blocks:
        bt = normalize_space(getattr(b, "text", "") or "")
        bt_key = bt.replace(" ", "").replace("　", "")

        score = 0

        for tl in title_lines:
            key = tl.replace(" ", "").replace("　", "")
            if key and key in bt_key:
                score += 3
            # 逆方向: ブロックテキストがタイトル行の一部（バラバラブロック対応）
            elif bt_key and len(bt_key) >= 3 and bt_key in key:
                score += 3

        for sp in speaker_keys:
            key = sp.replace(" ", "").replace("　", "").replace("先生", "")
            if key and key in bt_key.replace("先生", ""):
                score += 5

        if score > 0:
            candidates.append((score, b.top))

    if not candidates:
        return 10**18

    # アンカーは講演の「最も上にあるブロック」を使う
    # （時間ラベルは講演の上に配置されるため、上端が正確な基準になる）
    candidates.sort(key=lambda x: x[1])
    return candidates[0][1]


def sort_talks_by_layout(blocks: list[TextBlock], talks: list[Talk]) -> list[Talk]:
    items = []
    for t in talks:
        anchor_top = find_talk_anchor_top(blocks, t)
        items.append((anchor_top, t))
    items.sort(key=lambda x: x[0])
    return [t for _, t in items]

def assign_talk_times_by_nearest_upper_time(blocks: list[TextBlock], talks: list[Talk]) -> list[Talk]:
    """講演の上側にある最も近い時間を割り当て。
    イベント日時ブロック（「日時」ラベル・年月日入り）は除外。"""
    _event_dt_re = re.compile(r"20\d{2}\s*年|\d{1,2}\s*月\s*\d{1,2}\s*日|日時")

    time_blocks = []
    seen_tops = set()
    for b in blocks:
        raw = b.text or ""
        # 日時ラベル・年月日を含むブロックはイベント時間なので除外
        if _event_dt_re.search(raw):
            continue
        txt = _norm_time(raw)
        m = TIME_RE.search(txt)
        if m:
            start_time = m.group(1).replace(" ", "")
            end_time = m.group(2).replace(" ", "")
            start_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", start_time)
            end_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", end_time)
            time_blocks.append((b.top, b.left, f"{start_norm}~{end_norm}"))
            seen_tops.add(b.top)

    # バラバラブロック結合からも時間抽出（"19:00" "～" "19:30" 分割対応）
    for top, left, merged in _merge_blocks_to_rows(blocks):
        if top in seen_tops:
            continue
        if _event_dt_re.search(merged):
            continue
        txt = _norm_time(merged)
        m = TIME_RE.search(txt)
        if m:
            start_time = m.group(1).replace(" ", "")
            end_time = m.group(2).replace(" ", "")
            start_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", start_time)
            end_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", end_time)
            time_blocks.append((top, left, f"{start_norm}~{end_norm}"))
        txt = _norm_time(merged)
        m = TIME_RE.search(txt)
        if m:
            start_time = m.group(1).replace(" ", "")
            end_time = m.group(2).replace(" ", "")
            start_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", start_time)
            end_norm = re.sub(r"(\d{1,2}):(\d{2})", r"\1:\2", end_time)
            time_blocks.append((top, left, f"{start_norm}~{end_norm}"))

    for t in talks:
        if _is_program_chair_item(t):
            continue
        # 既に時間が設定されている場合はスキップ
        if normalize_space(t.time):
            continue
            
        anchor_top = find_talk_anchor_top(blocks, t)

        best = None
        best_score = None

        for top, left, tm in time_blocks:
            # 基本は talk より上の time を優先
            if top > anchor_top:
                continue

            dist = anchor_top - top
            # 距離制限を追加（あまり遠い時間は関連付けない）
            if dist > 1500000:  # 15cm相当
                continue
                
            if best_score is None or dist < best_score:
                best_score = dist
                best = tm

        if best:
            t.time = best

    return talks

# ---------------- Postprocess ----------------
def postprocess_refined(refined: DesignJSON, speaker_map: Dict[str, str], time_candidates: List[str]) -> DesignJSON:
    # event_title_lines を優先
    refined.event_title_lines = normalize_lines_keep_order(refined.event_title_lines or [])
    if refined.event_title_lines:
        refined.event_title = "\n".join(refined.event_title_lines).strip()
    else:
        # event_title しか来ない場合の救済
        if refined.event_title:
            # 1行内の ~...~ を別行化したい
            lines = []
            for raw in refined.event_title.split("\n"):
                lines.extend(split_tilde_subtitle_lines(raw))
            refined.event_title_lines = normalize_lines_keep_order(lines)
            refined.event_title = "\n".join(refined.event_title_lines).strip()

    TALK_KIND_LABELS = {
        "特別講演",
        "一般講演",
        "教育講演",
        "基調講演",
        "講演",
    }

    def clean_talk_title_lines(lines: List[str]) -> List[str]:
        out: List[str] = []

        for line in lines or []:
            s = normalize_space(line)
            if not s:
                continue

            # 単独ラベルは落とす
            if s in TALK_KIND_LABELS:
                continue

            # 先頭に付いているラベルは剥がす
            for lab in TALK_KIND_LABELS:
                if s.startswith(lab):
                    s = normalize_space(s[len(lab):])
                    break

            if s:
                out.append(s)

        return out
    
    cleaned: List[Talk] = []
    for t in refined.talks:
        if _is_program_chair_item(t):
            t.role = normalize_space(getattr(t, "role", "") or "") or "座長"
            if not normalize_space(getattr(t, "name_display", "") or ""):
                base_name = normalize_space(getattr(t, "speaker_display", "") or getattr(t, "speaker", "") or "")
                if base_name:
                    t.name_display = build_speaker_display(base_name) or base_name
            t.speaker = ""
            t.speaker_display = ""
            t.affiliation = normalize_space(getattr(t, "affiliation", "") or "")
            if t.name_display or t.affiliation:
                cleaned.append(t)
            continue

        if not (t.title_lines or t.speaker or t.affiliation or t.time):
            continue

        # title_lines 正規化（~...~ 別行化 + 重複排除）
        t.title_lines = normalize_lines_keep_order(t.title_lines or [])
        t.title_lines = clean_talk_title_lines(t.title_lines)

        sp = norm_name(t.speaker)
        t.speaker = sp
        t.speaker_display = build_speaker_display(t.speaker)

        # if sp in speaker_map:
        #     t.affiliation = speaker_map[sp] or ""
        # else:
        #     t.affiliation = ""

        if not (t.affiliation or "").strip():
            cand = (speaker_map.get(sp) or "").strip()
            # 長すぎる所属（PDFの注意文混入）を弾く
            if cand and len(cand) <= 80 and ("ご視聴" not in cand) and ("お願い" not in cand):
                t.affiliation = cand

        tm = normalize_space(t.time)
        if time_candidates and tm and tm not in time_candidates:
            tm = ""
        t.time = tm

        if t.title_lines or t.speaker or t.time:
            cleaned.append(t)
        
        if t.time and refined.datetime and normalize_space(t.time) in normalize_space(refined.datetime):
            t.time = ""
        
        

    max_items = 6 if any(_is_program_chair_item(t) for t in cleaned) else 4
    refined.talks = cleaned[:max_items]

    if refined.chair.name and not refined.chair.name_display:
        refined.chair.name_display = build_speaker_display(refined.chair.name)

    
    return refined


ROLE_WORDS2 = [
    "主任教授", "教授", "准教授", "講師", "助教",
    "副部長", "部長", "医長", "院長", "室長", "科長"
]

def looks_like_title_text(s: str) -> bool:
    s = normalize_space(s or "")
    if not s:
        return False
    if any(q in s for q in ["『", "』", "「", "」"]):
        return True
    if any(k in s for k in ["治療", "療法", "講演", "データ", "手術", "癌", "腫瘍"]):
        return True
    return len(s) >= 18

def split_affiliation_and_name_tail(s: str) -> tuple[str, str]:
    s = normalize_space(s or "")
    s = s.replace("\n", " ").strip()

    # 演者/講師ラベル除去
    s = re.sub(r"^(演者|講師)\s*[/／:：]?\s*", "", s)
    s = re.sub(r"先生$", "", s).strip()

    compact = s.replace(" ", "").replace("　", "")

    # すでに人名だけならそのまま返す
    if compact and not any(x in s for x in ["大学", "病院", "科", "センター", "教授", "部長", "医長", "院長", "医学部"]):
        if 3 <= len(compact) <= 8:
            return "", compact

    for role in ROLE_WORDS2:
        i = s.rfind(role)
        if i >= 0:
            aff = s[: i + len(role)].strip()
            name = s[i + len(role):].strip()
            compact = name.replace(" ", "").replace("　", "")

            if len(compact) < 3 or len(compact) > 8:
                continue
            if any(x in name for x in ["大学", "病院", "科", "センター", "教授", "部長", "医長", "院長", "医学部"]):
                continue

            return aff, compact

    return "", ""

def fix_title_lines_jp(lines: list[str]) -> list[str]:
    arr = [normalize_space(x) for x in (lines or []) if normalize_space(x)]
    if not arr:
        return []

    # ---- パス1: 複合助詞の不自然な分割を修正 ----
    # 「ワクチンへ\nのシフト」→「ワクチン\nへのシフト」
    COMPOUND_PARTICLES = {
        "への", "での", "との",
        "には", "では", "とは",
        "にも", "でも", "とも",
        "から", "まで",
        "ても", "ので", "のに",
        "ほど", "より",
    }
    for j in range(len(arr) - 1):
        cur_j = arr[j]
        nxt_j = arr[j + 1]
        if len(cur_j) >= 3 and len(nxt_j) >= 1:
            pair = cur_j[-1] + nxt_j[0]
            if pair in COMPOUND_PARTICLES:
                arr[j] = cur_j[:-1]
                arr[j + 1] = cur_j[-1] + nxt_j

    # ---- パス2: 短すぎる次行をマージ ----
    # 学習済みp90が取得できればそれを上限に、なければ28文字をデフォルト
    _llen_fix = _get_title_line_len_cache()
    MAX_MERGE_LEN = int(max(20, min(_llen_fix.get("talk_title_p90") or 28, 35)))
    out = []
    i = 0
    while i < len(arr):
        cur = arr[i]

        if i + 1 < len(arr):
            nxt = arr[i + 1]
            merged_len = len(cur) + len(nxt)

            # 1文字だけの次行
            if len(nxt) == 1:
                out.append(cur + nxt)
                i += 2
                continue

            # ひらがな始まりの短い次行（る / が / に / を 等）
            if re.match(r"^[ぁ-ん]", nxt) and len(nxt) <= 4 and merged_len <= MAX_MERGE_LEN:
                out.append(cur + nxt)
                i += 2
                continue

            # 漢字1字だけ残った不自然分割（性 などの短いサフィックス）のみマージ
            # nxt が長い独立した行（例:「実臨床データについて（仮）」）はマージしない
            if re.match(r"^[一-龠々]", nxt) and len(cur) >= 6 and len(nxt) <= 3 and merged_len <= MAX_MERGE_LEN:
                out.append(cur + nxt)
                i += 2
                continue

        out.append(cur)
        i += 1

    return out

def repair_chair_from_multiline_block(payload: DesignJSON, blocks: list[TextBlock]) -> DesignJSON:
    """
    座長情報の修復を行う関数
    ただし、既に正しい座長情報が抽出されている場合は修復をスキップする
    """
    if not getattr(payload, "chair", None):
        return payload

    # 座長情報が既に存在する場合は修復をスキップ
    # extract_chair_by_blocksで正しく抽出されている場合を優先
    if payload.chair.name and payload.chair.affiliation:
        print(f"[CHAIR REPAIR DEBUG] 座長情報が既に存在するため修復をスキップ: {payload.chair.name}")
        return payload

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    def _norm(s: str) -> str:
        return normalize_space(s or "")

    def _key(s: str) -> str:
        return _norm(s).replace(" ", "").replace("　", "")

    def _person_key(s: str) -> str:
        s = _norm(s)
        s = s.replace("先生", "")
        return s.replace(" ", "").replace("　", "")

    role_words = ["大学", "病院", "研究科", "医学部", "センター", "科", "講師", "教授", "部長", "医長", "副部長"]

    def _is_name_line(s: str) -> bool:
        s = _norm(s)
        if not s or "先生" not in s:
            return False
        if any(x in s for x in role_words):
            return False

        name = s.replace("先生", "").replace(" ", "").replace("　", "")
        return 2 <= len(name) <= 8

    def _looks_like_valid_person_name(s: str) -> bool:
        s = _norm(s)
        if not s:
            return False
        if any(x in s for x in role_words):
            return False
        s2 = s.replace("先生", "").replace(" ", "").replace("　", "")
        return 2 <= len(s2) <= 8

    def _looks_like_bad_affiliation(s: str) -> bool:
        s = _norm(s)
        if not s:
            return True
        # 役職だけで施設名がない
        has_role = any(x in s for x in ["講師", "教授", "部長", "医長", "副部長"])
        has_org = any(x in s for x in ["大学", "病院", "研究科", "医学部", "センター"])
        return has_role and not has_org

    current_name = _norm(payload.chair.name or "")
    current_name_key = _person_key(current_name)
    current_name_is_valid = _looks_like_valid_person_name(current_name)

    current_aff = _norm(payload.chair.affiliation or "")
    current_aff_is_bad = _looks_like_bad_affiliation(current_aff)

    first_talk_top = None
    for b in ordered:
        if "講演" in _key(b.text):
            first_talk_top = b.top
            break

    for i, b in enumerate(ordered):
        raw = _norm(b.text)
        if "座長" not in _key(raw):
            continue

        role = "座長"
        name = ""
        aff_lines = []

        lines = [_norm(x) for x in str(b.text).split("\n") if _norm(x)]
        for ln in lines:
            k = _key(ln)
            if "座長" in k:
                continue

            # 既存chair名が妥当なときだけ優先
            if current_name_is_valid and current_name_key and current_name_key in _person_key(ln):
                name = current_name
                continue

            if _is_name_line(ln) and not name:
                name = norm_name(ln.replace("先生", ""))
            else:
                if any(x in ln for x in role_words):
                    aff_lines.append(ln)

        allow_nearby_name_search = not current_name_is_valid and not name

        for j in range(max(0, i - 3), min(len(ordered), i + 4)):
            bj = ordered[j]
            if first_talk_top is not None and bj.top >= first_talk_top:
                break

            ln = _norm(bj.text)
            if not ln:
                continue

            key = _key(ln)
            if "講演" in key or "演者" in key:
                break

            if allow_nearby_name_search and _is_name_line(ln):
                name = norm_name(ln.replace("先生", ""))
                allow_nearby_name_search = False
                continue

            if any(x in ln for x in role_words):
                aff_lines.append(ln)

        org_words = ["大学", "病院", "研究科", "医学部", "センター"]
        role_words = ["講師", "教授", "部長", "医長", "副部長"]

        org_lines = []
        role_lines = []

        for ln in lines:
            if any(x in ln for x in org_words):
                org_lines.append(ln)
            elif any(x in ln for x in role_words):
                role_lines.append(ln)

        aff_out = []

        # まず施設
        if org_lines:
            aff_out.extend(org_lines)

        # 次に役職
        if role_lines:
            aff_out.extend(role_lines)


        # 既存nameは「妥当なときだけ」優先
        if current_name_is_valid:
            name = current_name

        # aff_out = []
        # seen = set()
        # for x in aff_lines:
        #     x2 = _norm(x)
        #     if not x2:
        #         continue
        #     if "座長" in _key(x2):
        #         continue
        #     if "先生" in x2:
        #         continue
        #     if x2 not in seen:
        #         aff_out.append(x2)
        #         seen.add(x2)

        if name:
            payload.chair.role = role
            payload.chair.name = norm_name(name)
            payload.chair.name_display = build_speaker_display(payload.chair.name) or payload.chair.name

            # affiliation は空欄 or 明らかに壊れている時は上書き
            if aff_out and (not current_aff or current_aff_is_bad):
                payload.chair.affiliation = "\n".join(aff_out).strip()

            payload.chair.honorific_title = "先生"
            return payload

    return payload

# ---------------- Parse (Rule + AI) ----------------
def parse_blocks_to_design_json(blocks: List[TextBlock], vm_rows: Optional[List[dict]] = None) -> DesignJSON:
    warnings: List[str] = []
    confidence = 0.78

    event_title_lines = extract_event_title_lines_from_blocks(blocks)
    event_title = "\n".join(event_title_lines).strip()

    # VM に「講演会名」があれば、抽出タイトルが案内文ヘッダー等の場合に優先使用
    if vm_rows:
        _vm_title = ""
        for _r in vm_rows:
            _d = _r if isinstance(_r, dict) and "data" not in _r else (_r.get("data") or {})
            _vt = normalize_space(_d.get("講演会名") or "")
            if _vt:
                _vm_title = _vt
                break
        if _vm_title:
            _et_norm = event_title.replace(" ", "").replace("\u3000", "")
            _vm_norm = _vm_title.replace(" ", "").replace("\u3000", "")
            # VM 講演会名と抽出タイトルの Jaccard 類似度で判断
            # Jaccard < 0.5（かなり違う）なら VM を優先（画像タイトル誤抽出・案内文ヘッダー誤採用等に対応）
            # → 正解DBがない初回からも有効
            _char_jaccard = len(set(_et_norm) & set(_vm_norm)) / max(len(set(_et_norm) | set(_vm_norm)), 1)
            if not _et_norm or _char_jaccard < 0.5:
                print(f"[vm-event-title] use VM 講演会名: '{_vm_title}' (was: '{event_title[:40]}', jaccard={_char_jaccard:.2f})")
                event_title_lines = [_vm_title]
                event_title = _vm_title

    print("event_title_lines:", event_title_lines)

    dt = extract_datetime_from_blocks(blocks)
    print("datetime:", dt)

    org = extract_organizer_from_blocks(blocks)  # ←主催: を含めたいなら別途調整（必要なら次で直す）

    speaker_map = extract_speaker_affil_map_by_blocks(blocks)
    # VMデータがあれば、VMの医師名をblocks内で検索してspeaker_mapを強化
    if vm_rows:
        speaker_map = enrich_speaker_map_with_vm(speaker_map, blocks, vm_rows)
    chair = extract_chair_by_blocks(blocks, speaker_map)

    talks = extract_talks_by_blocks(blocks, speaker_map, chair)

    if not event_title:
        warnings.append("missing_event_title"); confidence -= 0.2
    if not dt:
        warnings.append("missing_datetime"); confidence -= 0.15
    if not org:
        warnings.append("missing_organizer"); confidence -= 0.1
    if not chair.name:
        warnings.append("missing_chair"); confidence -= 0.1
    if len(talks) == 0:
        warnings.append("no_talks"); confidence -= 0.35


    confidence = float(min(max(confidence, 0.0), 1.0))

    payload = DesignJSON(
        event_title_lines=event_title_lines,
        event_title=event_title,
        datetime=normalize_space(dt),
        organizer=normalize_organizer(org),
        chair=Chair(role=chair.role, name=chair.name, name_display=chair.name_display, affiliation=chair.affiliation),
        talks=talks[:6],
        warnings=sorted(set(warnings)),
        confidence=confidence,
    )
    return apply_inline_program_extraction(payload, blocks)





# ---------------- Render (HTML→PNG) ----------------

# ---------------- VMヒント（演題演者）: PPTX探索精度UP + 欠損のみ補完 ----------------

def _norm_person_name(v: str) -> str:
    s = str(v or "").replace("\u3000", " ")
    s = " ".join(s.split())
    return s.replace(" ", "")

def _vm_aff_str(vm: dict) -> str:
    facility = (vm.get("案内状掲載 施設名") or "").strip()
    dept = (vm.get("案内状掲載 所属科") or "").strip()
    role = (vm.get("案内状掲載 役職") or "").strip()
    parts = [p for p in [facility, dept, role] if p]
    return " ".join(parts).strip()

def _norm_title_key(s: str) -> str:
    s = normalize_space(str(s or ""))
    # 記号ゆらぎを減らす
    s = s.replace("～", "〜").replace("−", "-").replace("—", "-").replace("–", "-").replace("－", "-")
    # かっこ/引用符などを除去（マッチ安定）
    for ch in ['"', "“", "”", "「", "」", "’", "‘", "（", "）", "(", ")", "【", "】", "[", "]", "『", "』"]:
        s = s.replace(ch, "")
    # スペース除去
    return s.replace(" ", "").replace("\u3000", "")

_UNWANTED_TITLE_WORDS = [
    "開会", "閉会", "開会の辞", "閉会の辞",
    "挨拶", "ご挨拶", "総合司会", "司会",
    "休憩", "休", "intermission",
    "動画上映", "ビデオ", "Video", "上映",
    "事務連絡", "諸連絡", "注意事項", "ご案内",
    "総合討論", "討論", "質疑", "Q&A",
    "オープニング", "エンディング",
]

def _is_unwanted_talk(title: str) -> bool:
    t = normalize_space(title)
    if not t:
        return True
    k = _norm_title_key(t)
    return any(_norm_title_key(w) in k for w in _UNWANTED_TITLE_WORDS)

def looks_like_real_talk(t: Talk) -> bool:
    title = normalize_space("\n".join(t.title_lines or []).strip() or (t.title or ""))
    if not title:
        return False
    if "演題" in title:
        return True
    if len(_norm_title_key(title)) >= 12:
        return True
    if (t.time or "").strip() and (t.speaker or t.speaker_display or "").strip():
        return True
    return False

def _vm_speaker_titles(vm_rows: list[dict]) -> list[str]:
    titles = []
    for r in (vm_rows or []):
        if (r.get("役職") or "") != "演者":
            continue
        v = (r.get("演題") or "").strip()
        if v:
            titles.append(v)
    return titles

def _time_start_minutes(t: str) -> int:
    m = re.search(r"(\d{1,2}):(\d{2})", str(t or ""))
    if not m:
        return 10**9
    hh, mm = map(int, m.groups())
    return hh * 60 + mm

def _strip_outer_quotes(s: str) -> str:
    s2 = str(s or "").strip()
    # 外側を囲む引用符ペアのみ除去（中身の引用符は残す）
    if s2.startswith("「") and s2.endswith("」"):
        # 内部に追加の「」ペアがない、または最外ペアが対応している場合のみ
        inner = s2[1:-1]
        if inner.count("「") == inner.count("」"):
            s2 = inner
    if s2.startswith("『") and s2.endswith("』"):
        inner = s2[1:-1]
        if inner.count("『") == inner.count("』"):
            s2 = inner
    return s2.strip()

def strip_outer_quotes_loose(s: str) -> str:
    s = normalize_space(s or "")

    # まず完全ペアを剥がす
    pairs = [
        ("「", "」"),
        ("『", "』"),
        ("（", "）"),
        ("(", ")"),
    ]
    changed = True
    while changed and s:
        changed = False
        for l, r in pairs:
            if s.startswith(l) and s.endswith(r):
                s = s[len(l):-len(r)].strip()
                changed = True

    # 片側だけ残った外カッコも落とす（カウントが不均衡な場合のみ）
    for l, r in pairs:
        while s.startswith(l) and s.count(l) > s.count(r):
            s = s[len(l):].strip()
        while s.endswith(r) and s.count(r) > s.count(l):
            s = s[:-len(r)].strip()

    return s

def _clean_title_lines(t):
    if getattr(t, "title_lines", None):
        t.title_lines = [_strip_outer_quotes(x) for x in t.title_lines]
    if getattr(t, "title", None):
        t.title = _strip_outer_quotes(t.title)
    return t

TIME_RANGE_PAT = re.compile(
    r"\d{1,2}[:：]\d{2}\s*[~\-–—−－〜～]\s*\d{1,2}[:：]\d{2}"
)

ROLE_PAT = re.compile(r"(演者|総合司会|座長)")

def clean_speaker_text(s: str) -> str:
    s = str(s or "")
    s = TIME_RANGE_PAT.sub("", s)      # 時間帯を除去
    s = ROLE_PAT.sub("", s)            # 演者/司会/座長系を除去
    s = re.sub(r"\s+", " ", s).strip()

    # 漢字間の変なスペースは潰す（前 田潤 → 前田潤）
    s = re.sub(r"(?<=[一-龥])\s+(?=[一-龥])", "", s)

    # 最後に姓名スペースを付け直す（あなたの build_speaker_display を使う）
    s = build_speaker_display(s)
    return s


def normalize_talk_speakers(payload: DesignJSON) -> DesignJSON:
    for t in (payload.talks or []):
        if _is_program_chair_item(t):
            continue
        base = (t.speaker_display or "").strip() or (t.speaker or "").strip()
        cleaned = clean_speaker_text(base)

        # display は姓名スペースあり
        t.speaker_display = cleaned
        # speaker はスペースなしに統一（検索/キー用）
        t.speaker = cleaned.replace(" ", "")
    return payload

def _talk_title_text(t) -> str:
    return normalize_space("\n".join(t.title_lines or []).strip() or (t.title or ""))

def _same_person(a: str, b: str) -> bool:
    return normalize_key(a or "").replace("先生", "") == normalize_key(b or "").replace("先生", "")

def _has_any_talk_signal(t) -> bool:
    return any([
        normalize_space(getattr(t, "title", "") or ""),
        any(normalize_space(x) for x in (getattr(t, "title_lines", None) or [])),
        normalize_space(getattr(t, "speaker", "") or ""),
        normalize_space(getattr(t, "affiliation", "") or ""),
        normalize_space(getattr(t, "time", "") or ""),
    ])

def _is_obviously_bad_talk(t, chair_name: str = "") -> bool:
    if _is_program_chair_item(t):
        return False

    title = _talk_title_text(t)
    speaker = normalize_space(getattr(t, "speaker", "") or "")
    affiliation = normalize_space(getattr(t, "affiliation", "") or "")

    if not _has_any_talk_signal(t):
        return True

    if _is_unwanted_talk(title):
        return True

    # event title や注意書きだけ
    if title.startswith("⚫") or "事前参加登録" in title or "担当者へご連絡" in title:
        return True

    # 座長行そのもの
    if title.startswith("座長") and not speaker:
        return True

    # chair本人で、titleも弱いなら落とす
    if chair_name and speaker and _same_person(speaker, chair_name):
        # 本当に chair が講演している可能性もあるので title が薄い時だけ
        if not title or "講演" in title or title == normalize_space(getattr(t, "title", "") or ""):
            return True

    # speakerもaffiliationもなく、titleだけ短い
    if not speaker and not affiliation and len(title) < 12:
        return True

    return False


def prune_talks_using_vm_titles(payload: DesignJSON, vm_rows: list[dict]) -> DesignJSON:
    talks = list(payload.talks or [])
    if not talks:
        return payload

    # 1) 明らかな不要物だけ落とす
    filtered = []
    for t in talks:
        if _is_program_chair_item(t):
            filtered.append(t)
            continue
        tt = normalize_space("\n".join(t.title_lines or []).strip() or (t.title or ""))
        if _is_unwanted_talk(tt):
            continue
        filtered.append(t)

    # 2) title重複だけ除去
    seen = set()
    dedup = []
    for t in filtered:
        if _is_program_chair_item(t):
            dedup.append(t)
            continue
        key = _norm_title_key("\n".join(t.title_lines or []).strip() or (t.title or ""))
        if not key:
            dedup.append(t)
            continue
        if key in seen:
            continue
        seen.add(key)
        dedup.append(t)

    # 3) VMでは削らない。最大4件だけ残す
    payload.talks = dedup[:4]
    payload.warnings = sorted(set((payload.warnings or []) + ["talks_pruned_by_vm_hint"]))
    return payload





def fill_chair_affiliation_from_vm_hint(payload: DesignJSON, blocks: list[TextBlock], vm_rows: list[dict]) -> DesignJSON:
    if not vm_rows or not payload.chair or (payload.chair.affiliation or "").strip() == "":
        pass
    else:
        return payload  # 既に入ってる

    # name -> vm（役職で絞らない。座長行が来る前提）
    vm_by_name = {}
    for vm in vm_rows:
        n = _norm_person_name(vm.get("案内状掲載 医師名", ""))
        if n:
            vm_by_name[n] = vm

    chair_key = _norm_person_name(payload.chair.name or payload.chair.name_display or "")
    vm = vm_by_name.get(chair_key)
    if not vm:
        return payload

    # 1) 施設名ヒントで blocks から拾い直す（talkと同じ戦略）
    facility = (vm.get("案内状掲載 施設名") or "").strip()
    if facility:
        pptx_aff = _join_affiliation_near_facility(blocks, facility)
        if pptx_aff:
            payload.chair.affiliation = pptx_aff
            return payload

    # 2) それでもダメなら VM 文字列を入れる（最後の保険）
    aff = _vm_aff_str(vm)
    if aff:
        payload.chair.affiliation = aff

    return payload

def fill_chair_affiliation_from_blocks(payload: DesignJSON, blocks: list[TextBlock]) -> DesignJSON:
    if not getattr(payload, "chair", None):
        return payload
    # chair.nameがアルファベット系や'PROGRAM'なら絶対に補完しない
    name = (payload.chair.name or "").strip()
    # 見出しワード一覧
    heading_words = {"PROGRAM", "P R O G R A M", "AGENDA", "SCHEDULE", "TIME TABLE", "タイムテーブル", "プログラム"}
    if not name or name.upper() in heading_words:
        print(f"[DEBUG] Chair name '{name}' is heading word, skip affiliation fill.")
        return payload
    # ★ 既に所属情報がある場合はスキップ（上流のAI結果を尊重）
    current_affiliation = (payload.chair.affiliation or "").strip()
    if current_affiliation and len(current_affiliation) >= 4:
        return payload

    print(f"[DEBUG] Chair affiliation check: current='{current_affiliation}' (length={len(current_affiliation)})")

    def _nospace(s: str) -> str:
        return normalize_space(s).replace(" ", "").replace("\u3000", "")

    key_ns = _nospace(payload.chair.name)
    print(f"[DEBUG] Looking for chair name: '{payload.chair.name}' (no space: '{key_ns}')")

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    target = None
    print(f"[DEBUG] Checking {len(ordered)} blocks for chair...")
    
    for i, b in enumerate(ordered):
        t = normalize_space(b.text)
        t_ns = _nospace(t)
        print(f"[DEBUG] Block {i}: '{t[:50]}...' (contains 座長: {'座長' in normalize_key(t)}, contains chair name: {key_ns in t_ns if key_ns else False})")
        
        if "座長" in normalize_key(t) and key_ns and key_ns in t_ns:
            print(f"[DEBUG] ★ FOUND chair block {i}: '{t}'")
            target = b
            break

    if not target:
        # 座長ラベルと座長名が別ブロックにある場合: 隣接ブロックから検索
        if key_ns:
            chair_label_block = None
            chair_name_block = None
            for b in ordered:
                t = normalize_space(b.text)
                t_ns = _nospace(t)
                if "座長" in normalize_key(t) and not chair_label_block:
                    chair_label_block = b
                if key_ns in t_ns and not chair_name_block:
                    chair_name_block = b
            if chair_label_block and chair_name_block and chair_label_block is not chair_name_block:
                vertical_dist = abs(chair_label_block.top - chair_name_block.top)
                if vertical_dist < 500000:
                    print(f"[DEBUG] Chair label and name in adjacent blocks (dist={vertical_dist})")
                    # ブロック全体から名前を除去して所属を抽出
                    full_text = normalize_space(chair_name_block.text.replace("\n", " "))
                    # 名前+先生を除去
                    aff_text = full_text
                    for pat in [payload.chair.name + "先生", payload.chair.name]:
                        aff_text = aff_text.replace(pat, "")
                    name_with_space = (payload.chair.name_display or "")
                    if name_with_space:
                        for pat in [name_with_space + "先生", name_with_space]:
                            aff_text = aff_text.replace(pat, "")
                    # 先頭末尾のカッコ・空白を整理
                    aff_text = re.sub(r'^[\s（(]+', '', aff_text)
                    aff_text = re.sub(r'[\s）)]+$', '', aff_text)
                    aff_text = normalize_space(aff_text).strip()
                    if aff_text and len(aff_text) >= 3:
                        print(f"[DEBUG] Extracted affiliation from adjacent block: '{aff_text}'")
                        payload.chair.affiliation = aff_text
                        return payload
                    target = chair_name_block

    if not target:
        print("[DEBUG] No chair block found")
        return payload

    def looks_like_affil(s: str) -> bool:
        s = normalize_space(s).replace("\n", " ")
        if not s:
            return False
        # アルファベットのみやPROGRAMは除外（ただし英語所属ワードは許可）
        if s.upper() in {"PROGRAM", "AGENDA", "SCHEDULE", "TIME TABLE"}:
            return False
        if "先生" in s or "座長" in normalize_key(s):
            return False
        if any(w in s for w in ["日時", "会場", "共催", "主催", "提供", "視聴", "登録"]):
            return False
        # 日本語・英語の所属キーワード
        kw = [
            "大学", "病院", "クリニック", "センター", "科", "部", "教授", "講師", "医師", "部長", "院長",
            "university", "hospital", "clinic", "center", "department", "faculty", "school", "institute", "division", "professor", "doctor", "dr.", "md", "phd"
        ]
        # 英語所属っぽいパターン（カンマ区切りやofを含む）
        if any(w in s.lower() for w in kw):
            return True
        if "," in s and "of" in s.lower():
            return True
        return any(w in s for w in kw)

    # ★ まず同一ブロック内から所属を抽出
    target_text = normalize_space(target.text)
    lines = [line.strip() for line in target_text.split('\n') if line.strip()]
    
    print(f"[DEBUG] Chair block found. Target text: {repr(target_text)}")
    print(f"[DEBUG] Split lines: {lines}")
    print(f"[DEBUG] Chair name key (no space): {repr(key_ns)}")
    
    # 「座長」と名前以外の行で所属っぽいものを探す
    for i, line in enumerate(lines):
        print(f"[DEBUG] Processing line {i}: {repr(line)}")
        
        # 座長が含まれる行でも、1行内に所属情報があるかチェック
        if "座長" in normalize_key(line):
            print(f"[DEBUG] Line {i} contains 座長, checking for inline affiliation")
            # 座長名の後に所属があるかチェック (パターン: "座長 名前先生 所属")
            if key_ns and key_ns in _nospace(line):
                # 「座長」を除去
                remaining = re.sub(r'座長\s*', '', line).strip()
                
                # 人名部分（漢字+ひらがな+先生）を除去する汎用的な処理
                # 「[漢字・ひらがな・カタカナ・空白]+先生」パターンを除去
                remaining = re.sub(r'^[一-龠ぁ-ゔァ-ヴー々〆〤\s]*先生\s*', '', remaining)
                
                remaining = remaining.strip()
                print(f"[DEBUG] After removing chair info, remaining: '{remaining}'")
                
                if remaining and looks_like_affil(remaining):
                    print(f"[DEBUG] Found inline affiliation in 座長 line {i}: '{remaining}'")
                    payload.chair.affiliation = remaining  
                    return payload
                else:
                    print(f"[DEBUG] Remaining part doesn't look like affiliation")
            
            print(f"[DEBUG] Skipping line {i}: contains 座長 but no inline affiliation found")
            continue
        # 名前（先生付き）を含む行をスキップ
        if key_ns and key_ns in _nospace(line):
            print(f"[DEBUG] Skipping line {i}: contains chair name")
            continue
        if "先生" in line:
            print(f"[DEBUG] Skipping line {i}: contains 先生")
            continue
        # 残った行で所属っぽければ採用
        if looks_like_affil(line):
            print(f"[DEBUG] Found affiliation in same block: {repr(line)}")
            payload.chair.affiliation = line
            return payload
        else:
            print(f"[DEBUG] Line {i} doesn't look like affiliation")

    print("[DEBUG] No affiliation found in same block, trying horizontal search...")

    # ★ 横並び優先（高さ帯一致）- フォールバック
    cand = []
    for b in ordered:
        if b is target:
            continue

        # 右側
        if b.left <= target.left:
            continue

        # 高さ帯が重なる（重要）
        if not (abs(b.top - target.top) < 400000):
            continue

        s = normalize_space(b.text.replace("\n", " "))
        if not looks_like_affil(s):
            continue

        dx = b.left - target.left
        dy = abs(b.top - target.top)
        score = dx + dy * 0.3
        cand.append((score, s))

    if cand:
        cand.sort(key=lambda x: x[0])
        payload.chair.affiliation = cand[0][1]
        return payload

    # fallback: 下方向
    for b in ordered:
        if b.top <= target.top:
            continue
        s = normalize_space(b.text.replace("\n", " "))
        if looks_like_affil(s):
            payload.chair.affiliation = s
            break

    return payload

def _find_best_block_idx_by_hint(blocks: List[TextBlock], hint: str, *, min_sim: float = 0.60) -> int:
    if not hint:
        return -1
    hint = str(hint).strip()
    if not hint:
        return -1

    # 部分一致優先
    for i, b in enumerate(blocks):
        t = (b.text or "").strip()
        if t and hint in t:
            return i

    # 類似一致
    best_i = -1
    best_sc = 0.0
    for i, b in enumerate(blocks):
        t = (b.text or "").strip()
        if not t:
            continue
        sc = sim(hint, t)
        if sc > best_sc:
            best_sc = sc
            best_i = i
    return best_i if best_sc >= min_sim else -1

def _join_affiliation_near_facility(blocks: List[TextBlock], facility_hint: str, *, max_follow: int = 4, y_limit: int | None = None) -> str:
    """PPTX上の所属表記を優先して作る: 施設名ブロック + 近傍（科/役職）を結合。
    Layer-1: 正解DBから学習した位置パターン（affil_rel_to_speaker_y）を使って y_limit を自動調整。"""
    if not facility_hint:
        return ""
    idx = _find_best_block_idx_by_hint(blocks, facility_hint, min_sim=0.60)
    if idx < 0:
        return ""

    b0 = blocks[idx]
    base_top = b0.top
    base_left = b0.left

    # Layer-1: 学習した affil_rel_to_speaker_y からブロック間隔の目安を算出
    if y_limit is None:
        lpc = _get_layout_pattern_cache()
        aff_stats = lpc.get("affil_rel_to_speaker_y") or {}
        if aff_stats.get("count", 0) >= 5:
            # 学習済みパターンの四分位範囲の絶対最大値 x2 を許容範囲とする
            learned_limit = max(abs(aff_stats.get("q25", 0)), abs(aff_stats.get("q75", 0))) * 2
            y_limit = max(int(learned_limit), 300000)  # 最低 300000
            print(f"[layout-pattern] _join_affil y_limit={y_limit} (learned from {aff_stats['count']} samples)")
        else:
            y_limit = 500000  # デフォルト（学習データ不足時）

    parts = [(b0.text or "").strip()]
    taken = 0

    # 施設名の直後に並ぶ「科/部/役職」っぽい行を拾う
    for j in range(idx + 1, min(len(blocks), idx + 1 + 30)):
        bj = blocks[j]
        tj = (bj.text or "").strip()
        if not tj:
            continue

        # 近傍制約（縦位置/横位置）
        if abs(bj.top - base_top) > y_limit:
            continue
        if abs(bj.left - base_left) > y_limit // 2:
            # 横が大きくズレるものは別カラムの可能性
            continue

        if not any(k in tj for k in ["科", "部", "センター", "内科", "外科", "教授", "准教授", "講師", "医長", "部長"]):
            continue

        parts.append(tj)
        taken += 1
        if taken >= max_follow:
            break

    return "\n".join([p for p in parts if p]).strip()

def _pick_vm_row_by_talk(vm_by_name: dict, t) -> dict | None:
    # speaker を優先してVMを引く（displayが壊れてても耐える）
    cand = []
    sp1 = _norm_person_name(getattr(t, "speaker", "") or "")
    sp2 = _norm_person_name(getattr(t, "speaker_display", "") or "")
    if sp1: cand.append(sp1)
    if sp2 and sp2 != sp1: cand.append(sp2)

    for key in cand:
        vm = vm_by_name.get(key)
        if vm:
            return vm
    return None

def apply_vm_hints_from_blocks(blocks: List[TextBlock], payload: DesignJSON, vm_rows: List[dict]) -> DesignJSON:
    """VMはヒントとしてのみ使用。最終値はPPTX(blocks)から取得して埋める（上書きは欠損/矛盾時のみ）。"""
    if not vm_rows or not getattr(payload, "talks", None):
        return payload

    # name -> vm
    vm_by_name: Dict[str, dict] = {}
    for vm in vm_rows:
        n = _norm_person_name(vm.get("案内状掲載 医師名", ""))
        if n:
            vm_by_name[n] = vm

    # facilities list (to detect mismatch)
    facilities = [ (vm.get("案内状掲載 施設名") or "").strip() for vm in vm_rows if (vm.get("案内状掲載 施設名") or "").strip() ]

    for t in payload.talks:
        if _is_program_chair_item(t):
            continue
        # sp = _norm_person_name(getattr(t, "speaker_display", "") or getattr(t, "speaker", ""))
        # vm = vm_by_name.get(sp)
        vm = _pick_vm_row_by_talk(vm_by_name, t)
        if not vm:
            continue

        facility = (vm.get("案内状掲載 施設名") or "").strip()
        if not facility:
            continue

        pptx_aff = _join_affiliation_near_facility(blocks, facility)

        if not pptx_aff:
            continue

        # 座長情報が混入した結果を採用しない
        if "座長" in pptx_aff or "演者" in pptx_aff:
            continue

        cur = (t.affiliation or "").strip()

        # 欠損なら入れる。入ってるが別施設っぽければPPTX値で修正（PPTX由来なのでOK）
        if not cur:
            t.affiliation = pptx_aff
        else:
            if facility not in cur:
                # もしcurが他の施設名を含んでいたら矛盾とみなす
                if facility not in cur and any(f and f in cur for f in facilities):
                    t.affiliation = pptx_aff

    return payload

def fill_missing_from_vm(payload: DesignJSON, vm_rows: List[dict]) -> DesignJSON:
    if not vm_rows or not payload.talks:
        return payload

    vm_by_name = {
        _norm_person_name(r.get("案内状掲載 医師名", "")): r
        for r in vm_rows
        if _norm_person_name(r.get("案内状掲載 医師名", ""))
    }

    for t in payload.talks:
        if _is_program_chair_item(t):
            continue
        sp = _norm_person_name(t.speaker or t.speaker_display or "")
        vm = vm_by_name.get(sp)
        if not vm:
            continue

        cur_title = normalize_space(getattr(t, "title", "") or "")
        cur_title_lines = [
            normalize_space(x)
            for x in (getattr(t, "title_lines", None) or [])
            if normalize_space(x)
        ]

        # タイトル補完（title も title_lines も空のときのみ）
        if not cur_title and not cur_title_lines:
            v = normalize_space(vm.get("演題") or "")
            if v:
                t.title = v
                t.title_lines = [ln for ln in v.split("\n") if normalize_space(ln)]

        # 所属補完（完全空のときのみ）
        if not normalize_space(getattr(t, "affiliation", "") or ""):
            aff = _vm_aff_str(vm)
            if aff:
                t.affiliation = aff

    return payload

def build_vm_title_map(vm_rows: list[dict]) -> dict[str, dict]:
    """
    return: { normalized_title: {"speaker": "...", "affiliation": "...", "title": "..."} }
    """
    def norm_title(s: str) -> str:
        s = normalize_space(s or "")
        # 記号ゆれ吸収（必要なら増やす）
        s = s.replace("〜", "～")
        s = re.sub(r"[‐-–—−]", "-", s)
        s = s.replace(" ", "").replace("\u3000", "")
        return s

    title_map: dict[str, dict] = {}
    for r in (vm_rows or []):
        if (r.get("役職") or "") != "演者":
            continue
        title = (r.get("演題") or "").strip()
        sp = norm_name(r.get("案内状掲載 医師名") or "")
        fac = normalize_space(r.get("案内状掲載 施設名") or "")
        dept = normalize_space(r.get("案内状掲載 所属科") or "")
        pos = normalize_space(r.get("案内状掲載 役職") or "")

        aff = " ".join([x for x in [fac, dept, pos] if x]).strip()
        if not title or not sp:
            continue

        key = norm_title(title)
        title_map[key] = {"speaker": sp, "affiliation": aff, "title": title}

    return title_map

def normalize_speaker_display(payload: DesignJSON) -> DesignJSON:
    if getattr(payload, "chair", None):
        name = normalize_space(payload.chair.name or "")
        payload.chair.name_display = build_speaker_display(name)

    for t in getattr(payload, "talks", []) or []:
        if _is_program_chair_item(t):
            if getattr(t, "name_display", ""):
                t.name_display = build_speaker_display(t.name_display)
            continue
        name = normalize_space(t.speaker or "")
        t.speaker_display = build_speaker_display(name)

    return payload

# def normalize_speaker_display(payload: DesignJSON) -> DesignJSON:
#     if not getattr(payload, "talks", None):
#         return payload

#     for t in payload.talks:
#         if (t.speaker or "").strip():
#             # speaker_display が空 or 不正なら再生成
#             if not (t.speaker_display or "").strip():
#                 t.speaker_display = build_speaker_display(t.speaker)

#     # chair も同様
#     if getattr(payload, "chair", None):
#         ch = payload.chair
#         if (ch.name or "").strip():
#             if not (ch.role or "").strip():
#                 ch.role = detect_chair_role((ch.name_display or "") + " " + (ch.name or ""))
#             if not (ch.name_display or "").strip():
#                 ch.name_display = build_speaker_display(ch.name)

#     return payload


def prune_talks_heuristic_only(payload: DesignJSON) -> DesignJSON:
    talks = list(payload.talks or [])
    if not talks:
        return payload

    chair_name = getattr(payload.chair, "name", "") or ""

    filtered = []
    for t in talks:
        if _is_program_chair_item(t):
            filtered.append(t)
            continue
        if _is_obviously_bad_talk(t, chair_name=chair_name):
            continue
        filtered.append(t)

    grouped = {}
    for t in filtered:
        if _is_program_chair_item(t):
            grouped.setdefault(f"__chair__{len(grouped)}", []).append(t)
            continue
        key = _norm_title_key(_talk_title_text(t))
        if not key:
            key = f"__idx__{len(grouped)}"
        grouped.setdefault(key, []).append(t)

    dedup = []
    for _, group in grouped.items():
        if len(group) == 1:
            dedup.append(group[0])
            continue

        def score_talk(x):
            return (
                1 if normalize_space(getattr(x, "speaker", "") or "") else 0,
                1 if normalize_space(getattr(x, "affiliation", "") or "") else 0,
                1 if normalize_space(getattr(x, "time", "") or "") else 0,
                len(_talk_title_text(x)),
            )

        group = sorted(group, key=score_talk, reverse=True)
        dedup.append(group[0])

    payload.talks = dedup[:4]
    payload.warnings = sorted(set((payload.warnings or []) + ["talks_pruned_heuristic_only"]))
    return payload

ROLE_WORDS = [
    "教授", "准教授", "講師", "助教", "医長", "部長",
    "院長", "副院長", "主任", "センター長"
]

def append_vm_role_to_talk_affiliation(payload, vm_rows: list[dict]) -> None:
    if not vm_rows or not getattr(payload, "talks", None):
        return

    def norm_key(s: str) -> str:
        return normalize_space(s or "").replace(" ", "").replace("　", "")

    def clean_affiliation_text(s: str) -> str:
        x = normalize_space(s or "")
        if not x:
            return ""

        # よくある重複だけ軽く掃除
        x = x.replace("病院病院", "病院")

        for role in ROLE_WORDS:
            x = re.sub(rf"({re.escape(role)})\s*{re.escape(role)}", role, x)

        return x

    def append_role_to_affiliation(affiliation: str, role: str) -> str:
        aff = clean_affiliation_text(affiliation)
        role = normalize_space(role or "")

        if not role:
            return affiliation  # \n を保持するため元の文字列を返す
        if not aff:
            return role

        # 同じ role が既に入っていれば追加しない
        if role in aff:
            return affiliation  # \n を保持するため元の文字列を返す

        # 既に別の役職語が入っているなら、むやみに追加しない
        if any(rw in aff for rw in ROLE_WORDS):
            return affiliation  # \n を保持するため元の文字列を返す

        return f"{aff} {role}".strip()

    # 医師名ベースで VM を引けるようにする
    vm_by_name = {}
    for r in vm_rows:
        d = r if isinstance(r, dict) and "data" not in r else (r.get("data") or {})
        doctor = norm_key(d.get("案内状掲載 医師名", ""))
        if doctor:
            vm_by_name[doctor] = d

    for t in payload.talks:
        if _is_program_chair_item(t):
            continue
        sp = norm_key(getattr(t, "speaker", ""))
        if not sp:
            continue

        vm = vm_by_name.get(sp)
        if not vm:
            continue

        role = normalize_space(vm.get("案内状掲載 役職", "") or "")
        if not role:
            continue

        aff = getattr(t, "affiliation", "") or ""
        t.affiliation = append_role_to_affiliation(aff, role)

def looks_like_talk_title_text(s: str) -> bool:
    s = normalize_space(s)
    if not s:
        return False
    if len(s) < 8:
        return False
    if any(k in s for k in ["ワクチン", "治療", "講演", "セミナー", "Up to date", "Update"]):
        return True
    if "～" in s or "〜" in s:
        return True
    if s.startswith(("(仮）", "（仮）", "(仮)", "（仮)")):
        return True
    return False


def clean_chair_fields(payload: DesignJSON) -> DesignJSON:
    if not getattr(payload, "chair", None):
        return payload

    name = normalize_space(payload.chair.name or "")
    aff = normalize_space(payload.chair.affiliation or "")

    # name から明らかなノイズを削る
    name = re.sub(r"^(座長|司会)\s*", "", name)
    name = re.sub(r"(座長|司会)", "", name)
    name = name.replace("先生", "").strip()

    # 所属語が混ざっていたら切り落とす
    m = re.search(r"(大学|病院|研究科|医学部|センター|科|講師|教授|部長|医長|副部長)", name)
    if m:
        cut = m.start()
        # 手前に人名が残っていればそこだけ採用
        maybe_name = name[:cut].strip()
        if maybe_name:
            name = maybe_name

    name = norm_name(name)

    if name:
        payload.chair.name = name
        payload.chair.name_display = build_speaker_display(name) or name

    payload.chair.affiliation = aff
    payload.chair.role = normalize_chair_role(payload.chair.role)
    return payload

def split_name_and_affiliation_strict(s: str) -> tuple[str, str]:
    s = normalize_space(s)

    # 伊東直哉先生（名古屋市立大学大学院医学研究科
    m = re.match(r"^(?P<name>[^（(]+?)(?:先生)?\s*[（(](?P<aff>.+)$", s)
    if m:
        return normalize_space(m.group("name")), normalize_space(m.group("aff"))

    # 中西重清先生（中西内科院長）
    m = re.match(r"^(?P<name>[^（(]+?)(?:先生)?\s*[（(](?P<aff>.+)[）)]\s*$", s)
    if m:
        return normalize_space(m.group("name")), normalize_space(m.group("aff"))

    # 伊東直哉先生
    m = re.match(r"^(?P<name>.+?)(?:先生)\s*$", s)
    if m:
        return normalize_space(m.group("name")), ""

    return "", ""

def repair_talks_from_blocks(payload: DesignJSON, blocks: list[TextBlock]) -> DesignJSON:
    def is_time_line(s: str) -> str:
        s2 = normalize_time_colon(normalize_space(s))
        m = TIME_RANGE_RE.search(s2)
        return normalize_space(m.group(1)) if m else ""
    
    if not getattr(payload, "talks", None):
        return payload

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    def _clean_speaker_line(s: str) -> str:
        s = normalize_space(s)
        s = normalize_time_colon(s)  # ← 追加

        # 時間削除
        s = re.sub(r"\d{1,2}:\d{2}\s*[～〜~\-－—–]\s*\d{1,2}:\d{2}", "", s)

        # ラベル削除
        s = re.sub(r"(演者|座長)\s*[:：]?", "", s)

        return s.strip()


    def strip_label(prefixes, s: str) -> str:
        s2 = normalize_space(s or "")
        s2_key = normalize_key(s2)

        for p in prefixes:
            p_key = normalize_key(p)
            if s2_key.startswith(p_key):
                chars = list(p_key)
                pat = r"^" + r"\s*".join(map(re.escape, chars)) + r"\s*(?:[:：/／]\s*)?"
                s2 = re.sub(pat, "", s2).strip()
                return s2

        return s2

    def _norm(s: str) -> str:
        return normalize_space(s or "").replace("　", " ")

    def _find_talk_anchor(no: int) -> Optional[TextBlock]:
        labels_by_no = {
            1: ["講演1", "講演１", "講演①", "講演Ⅰ", "教育講演", "演題1", "演題１", "演題①"],
            2: ["講演2", "講演２", "講演②", "講演Ⅱ", "演題2", "演題２", "演題②"],
            3: ["講演3", "講演３", "講演③", "講演Ⅲ", "特別講演", "演題3", "演題３", "演題③"],
            4: ["講演4", "講演４", "講演④", "講演Ⅳ", "演題4", "演題４", "演題④"],
        }

        no_zen = str(no).translate(str.maketrans('1234567890', '１２３４５６７８９０'))
        labels = labels_by_no.get(no, [f"講演{no}", f"講演{no_zen}", f"演題{no}", f"演題{no_zen}"])

        for b in ordered:
            # 改行で分割してから各行を検索（時間が連結されて誤マッチを防ぐ）
            lines = (b.text or "").split("\n")
            for line in lines:
                key = normalize_key(line)
                for lb in labels:
                    lb_key = normalize_key(lb)
                    # 数字末尾ラベルは後続数字があると誤マッチする
                    # 例: "講演1" が "特別講演18:30" にヒットしないよう
                    if lb_key and lb_key[-1].isdigit():
                        if re.search(re.escape(lb_key) + r'(?!\d)', key):
                            return b
                    else:
                        if lb_key in key:
                            return b
        return None

    def _looks_like_title(s: str) -> bool:
        s = _norm(s)
        if not s:
            return False
        if "講演" in s or "演者" in normalize_key(s) or "座長" in normalize_key(s):
            return False
        if re.match(r'^\s*(?:主催|共催)\s*[：:]', s):
            return False
        if looks_like_datetime_text(s):
            return False
        if is_time_line(s):
            return False
        if "先生" in s:
            return False
        return len(s) >= 6

    def _extract_person_near_enja(
    seg: list[TextBlock],
    chair_name: str = "",
    chair_aff: str = "",
) -> tuple[str, str]:
        speaker = ""
        affiliation = ""

        chair_name_key = normalize_key(chair_name or "").replace("先生", "")
        chair_aff_key = normalize_key(chair_aff or "")

        bad_words = ["大学", "病院", "研究科", "医学部", "センター", "講師", "教授", "部長", "医長", "院長", "副部長"]

        def _is_name_only(s: str) -> bool:
            s = _norm(s).replace("先生", "").strip()
            s2 = s.replace(" ", "").replace("　", "")
            if not s2:
                return False
            if any(x in s for x in bad_words):
                return False
            # 講演ラベル（講演1, 講演2, 演題1 等）は名前ではない
            if re.match(r'^(講演|演題)[0-9０-９①-⑩ⅠⅡⅢⅣⅤ]*$', s2):
                return False
            return 2 <= len(s2) <= 8

        def _clean_name_candidate(s: str) -> str:
            s = clean_speaker_text(s)
            s = s.replace("先生", "").strip()
            return norm_name(s)

        def _same_as_chair_name(s: str) -> bool:
            if not chair_name_key:
                return False
            return normalize_key(s or "").replace("先生", "") == chair_name_key

        def _same_as_chair_aff(s: str) -> bool:
            if not chair_aff_key:
                return False
            return normalize_key(s or "") == chair_aff_key

        def _score_name(block_idx: int, line_idx: int, anchor_idx: int, line: str) -> tuple:
            dist = abs(block_idx - anchor_idx)
            key = normalize_key(line or "")

            # 演者と同一行なら最優先
            same_line_bonus = 0 if "演者" in key else 10

            # 「先生」付きは少し優先
            sensei_bonus = 0 if "先生" in line else 1

            return (dist, same_line_bonus, sensei_bonus, line_idx)

        def _score_aff(block_idx: int, line_idx: int, anchor_idx: int, line: str) -> tuple:
            dist = abs(block_idx - anchor_idx)
            return (dist, line_idx)

        for i, b in enumerate(seg):
            key = normalize_key(b.text or "")
            if "演者" not in key:
                continue

            # 演者ブロックの前後を見る
            start = max(0, i - 2)
            end = min(len(seg), i + 4)
            cand_blocks = seg[start:end]

            name_cands = []
            aff_cands = []

            for rel_idx, cb in enumerate(cand_blocks):
                abs_idx = start + rel_idx
                lines = [_norm(x) for x in str(cb.text).split("\n") if _norm(x)]

                for li, ln in enumerate(lines):
                    ln_key = normalize_key(ln)

                    # 座長行は除外
                    if "座長" in ln_key:
                        continue

                    # 名前候補
                    nm = _clean_name_candidate(ln)
                    if _is_name_only(nm) and not _same_as_chair_name(nm):
                        name_cands.append((_score_name(abs_idx, li, i, ln), nm))

                    # 所属候補
                    if looks_like_affil_line(ln) and not _same_as_chair_aff(ln):
                        aff_cands.append((_score_aff(abs_idx, li, i, ln), ln))

                # 複数行ブロック: 1行目が所属なら後続の科名・役職行を結合
                if len(lines) > 1 and looks_like_affil_line(lines[0]) and not _same_as_chair_aff(lines[0]):
                    role_dept_kw = ["科", "内科", "外科", "教授", "准教授", "講師", "部長",
                                    "医長", "院長", "部", "室", "課", "主任"]
                    joined = [lines[0]]
                    for sub_ln in lines[1:]:
                        sub_ln_s = sub_ln.strip()
                        if sub_ln_s and any(k in sub_ln_s for k in role_dept_kw) and len(sub_ln_s) <= 20:
                            joined.append(sub_ln_s)
                    if len(joined) > 1:
                        joined_aff = " ".join(joined)
                        aff_cands.append((_score_aff(abs_idx, 0, i, joined_aff), joined_aff))

                # ブロック全体が「演者 石井 康隆先生 ...」型のときの補助
                block_text = _norm(cb.text)
                block_key = normalize_key(block_text)

                if "演者" in block_key:
                    # 演者の後ろの名前（日本語）
                    m = re.search(r"演者\s*([一-龥々]{1,4}\s*[一-龥々]{1,4})\s*先生?", block_text)
                    if m:
                        nm = norm_name(m.group(1))
                        if nm and not _same_as_chair_name(nm):
                            name_cands.append(((abs(abs_idx - i), 0, 0, 0), nm))

                    # 演者の後ろにラテン文字名（Prof. Harm Jan Bogaard 等）
                    if not m:
                        m_lat = re.search(r"演者\s*(?:Prof\.?\s*)?([A-Za-z][A-Za-z .\-]+[A-Za-z])", block_text)
                        if m_lat:
                            lat_name = normalize_space(m_lat.group(1))
                            if len(lat_name) >= 3 and not _same_as_chair_name(lat_name):
                                name_cands.append(((abs(abs_idx - i), 0, 0, 0), lat_name))

                    # 演者の後ろに所属も続くケース
                    m2 = re.search(
                        r"演者\s*[一-龥々]{1,4}\s*[一-龥々]{1,4}\s*先生?\s+(.+)$",
                        block_text
                    )
                    if m2:
                        cand_aff = _norm(m2.group(1))
                        if looks_like_affil_line(cand_aff) and not _same_as_chair_aff(cand_aff):
                            aff_cands.append(((abs(abs_idx - i), 0), cand_aff))

            if name_cands:
                name_cands.sort(key=lambda x: x[0])
                speaker = name_cands[0][1]

            if aff_cands:
                aff_cands.sort(key=lambda x: x[0])
                affiliation = aff_cands[0][1]

            if speaker:
                return speaker, affiliation

        return "", ""

    def clean_title_lines2(lines: list[str]) -> list[str]:
        out = []

        for ln in lines:
            s = normalize_space(ln)

            # 単独カッコ除去
            if s in ["「", "」", "『", "』"]:
                continue

            # 明確な名前行は除外
            if re.fullmatch(r"[一-龥々]{1,4}\s*[一-龥々]{1,4}\s*先生?", s):
                continue

            # 所属行は除外
            if any(x in s for x in [
                "大学", "病院", "研究科", "医学部", "センター",
                "クリニック", "講師", "教授", "部長", "医長", "院長", "名誉院長"
            ]):
                continue

            s = strip_outer_quotes_loose(s).strip()
            if s:
                out.append(s)

        return out


    def _extract_title_near_anchor(anchor: TextBlock, seg: list[TextBlock]) -> list[str]:
        title_lines: list[str] = []

        # anchorの少し上も含める
        around = [b for b in seg if b.top >= anchor.top - 450000]
        around = sorted(around, key=lambda b: (b.top, b.left))

        def _is_name_line(s: str) -> bool:
            s = _norm(s)
            if not s or "先生" not in s:
                return False
            bad = ["大学", "病院", "研究科", "医学部", "センター", "科", "講師", "教授", "部長", "医長", "院長", "副部長"]
            return not any(x in s for x in bad)

        def strip_outer_quotes(s: str) -> str:
            s = normalize_space(s or "")

            PAIRS = [
                ("「", "」"),
                ("『", "』"),
                ("（", "）"),
                ("(", ")"),
            ]

            for l, r in PAIRS:
                if s.startswith(l) and s.endswith(r):
                    return s[len(l):-len(r)].strip()

            return s

        def _clean_title_piece(s: str) -> str:
            s = _norm(s)
            s = re.sub(r"^\d{1,2}:\d{2}\s*[～〜~\-－]\s*\d{1,2}:\d{2}\s*", "", s)
            s = re.sub(r"^講\s*演\s*[0-9０-９①②③④⑤⑥⑦⑧⑨⑩ⅠⅡⅢⅣⅤIVX]+\s*", "", s)
            return s.strip()

        

        def _looks_like_title_piece(s: str) -> bool:
            s = _clean_title_piece(s)

            if not s:
                return False

            if is_non_talk_heading(s):
                return False

            if any(x in s for x in ["講演", "演者", "座長"]):
                return False

            # 主催・共催行は演題ではない
            if re.match(r'^\s*(?:主催|共催)\s*[：:]', s):
                return False

            if looks_like_affil_line(s):
                return False

            if "先生" in s:
                return False

            # 時間文字列は除外
            if is_time_line(s):
                return False

            # ★追加：短すぎる or 記号だけは除外
            if len(s) < 8:
                return False

            return True

        # まず time と同じブロックにタイトル前半があるケースを優先
        for i, b in enumerate(around):
            raw = _norm(b.text)
            if not is_time_line(raw):
                continue

            parts = [_clean_title_piece(x) for x in raw.split("\n")]
            parts = [x for x in parts if _looks_like_title_piece(x)]
            if parts:
                title_lines.extend(parts)

                # ← ここは for i の内側
                for j in range(i + 1, min(i + 3, len(around))):
                    raw2 = _norm(around[j].text)
                    nxt = _clean_title_piece(raw2)

                    if _is_name_line(raw2) or looks_like_affil_line(raw2):
                        break

                    if _looks_like_title_piece(nxt):
                        title_lines.append(nxt)
                        break
                break

        # timeブロックから取れなければ従来fallback
        if not title_lines:
            for i, b in enumerate(around):
                s = _clean_title_piece(b.text)
                if not _looks_like_title_piece(s):
                    continue

                title_lines.append(s)

                # 後続ブロックを探索（講演ラベル・時間ブロックはスキップ）
                for j in range(i + 1, min(i + 4, len(around))):
                    raw_j = _norm(around[j].text)
                    nxt = _clean_title_piece(around[j].text)

                    # 名前・所属・演者ラベルに到達したら終了
                    if _is_name_line(raw_j) or looks_like_affil_line(raw_j):
                        break
                    if "演者" in normalize_key(raw_j):
                        break

                    # 講演ラベル・時間ブロックはスキップして先を見る
                    if any(x in raw_j for x in ["講演", "講 演"]) or is_time_line(raw_j):
                        continue

                    if _looks_like_title_piece(nxt):
                        title_lines.append(nxt)
                    break

                break

        # 重複除去
        out = []
        seen = set()
        for x in title_lines:
            x = _norm(x)
            if x and x not in seen:
                out.append(x)
                seen.add(x)

        # 各行ごとに外側引用符を除去（改行を潰さない）
        out = [strip_outer_quotes_loose(x) for x in out if x]

        return out[:4]


    
    def is_bad_speaker(s: str) -> bool:
        s = normalize_space(s or "")
        if not s:
            return True

        # 所属ワードが入ってたらアウト
        bad = ["大学", "病院", "研究科", "医学部", "センター", "講師", "教授", "部長", "医長", "院長"]
        if any(x in s for x in bad):
            return True

        # 明らかに人名ではない一般語
        s_ns = s.replace(" ", "").replace("　", "")
        non_person_words = {
            "遠慮", "参加", "視聴", "登録", "配信", "質問", "回答",
            "講演", "演題", "座長", "司会", "開催", "案内", "申請",
            "治療", "診断", "手術", "検査", "予防", "感染",
        }
        if s_ns in non_person_words:
            return True

        # 講演ラベル（講演1, 演題2 等）は人名ではない
        if re.match(r'^(講演|演題)[0-9０-９①-⑩ⅠⅡⅢⅣⅤ]*$', s_ns):
            return True

        # 長すぎる（ラテン文字名は20文字まで許容）
        is_latin = bool(re.search(r'[A-Za-z]', s))
        max_len = 25 if is_latin else 10
        if len(s_ns) > max_len:
            return True

        return False

    def clean_speaker_text(s: str) -> str:
        s = normalize_space(s or "")

        # 時間削除
        s = re.sub(r"\d{1,2}[:：]\d{2}\s*[-~～〜]\s*\d{1,2}[:：]\d{2}", "", s)

        # ラベル削除
        s = re.sub(r"(演者|座長)\s*", "", s)

        # 「先生」を先に除去してから人名パターンを取る
        s = re.sub(r"\s*先生\s*$", "", s).strip()

        # 最後の人名だけ取る
        m = re.search(r"([一-龥々]{1,4}\s*[一-龥々]{1,4})$", s)
        if m:
            return norm_name(m.group(1))

        return norm_name(s)
    
    def clean_affiliation_text(s: str) -> str:
        s = normalize_space(s or "")

        # 座長情報が混入していたら空にする
        if "座長" in s:
            return ""

        # 演者とか混ざってたら除去
        s = re.sub(r"(演者|座長)\s*", "", s)

        # 「ご所属：」プレフィックスを除去
        s = re.sub(r'^ご?所属\s*[:：]\s*', '', s)

        return s.strip()


    # talks[0], talks[1]... を各「講演N」アンカーから拾い直す
    talks = list(payload.talks or [])

    # ラベル語が所属に入っていたらクリア（"演者", "座長" 等）
    _label_aff_pat = re.compile(r'^(演者|座長|講演\d*|演題\d*)$')
    for t in talks:
        aff_ns = re.sub(r'[\s\u3000]+', '', getattr(t, 'affiliation', '') or '')
        if aff_ns and _label_aff_pat.match(aff_ns):
            print(f"[repair-talks] clearing label-like affiliation: '{t.affiliation}'")
            t.affiliation = ''

    for idx, t in enumerate(talks, start=1):
        anchor = _find_talk_anchor(idx)

        # フォールバック: 番号なし「演題」「演者」ブロックをアンカーにする（1講演のみ）
        if not anchor and len(talks) == 1:
            for _b in ordered:
                for _line in (_b.text or "").split("\n"):
                    _k = normalize_key(_line).rstrip("：:")
                    if _k in ("演題", "演者"):
                        anchor = _b
                        break
                if anchor:
                    break

        if not anchor:
            continue

        # 次の講演アンカーまでをこの講演の範囲にする
        next_anchor = _find_talk_anchor(idx + 1)

        x0 = anchor.left - 3000000
        x1 = anchor.left + 4500000
        y0 = anchor.top - 100000
        y1 = (next_anchor.top - 150000) if next_anchor else (anchor.top + 3500000)

        seg = [b for b in ordered if in_region(b, x0, y0, x1, y1)]
        seg = sorted(seg, key=lambda b: (b.top, b.left))

        # title
        def has_meaningful_title(t) -> bool:
            lines = [normalize_space(x) for x in (getattr(t, "title_lines", None) or []) if normalize_space(x)]
            title = normalize_space(getattr(t, "title", "") or "")
            full = "\n".join(lines) if lines else title

            if not full:
                return False

            if re.search(r"[^\s　]+\s+[^\s　]+\s*先生?$", full):
                return False

            if any(x in full for x in ["大学", "病院", "研究科", "医学部", "センター", "講師", "教授", "部長", "医長", "院長"]):
                return False

            return len(full) >= 6

        title_lines = _extract_title_near_anchor(anchor, seg)
        title_lines = clean_title_lines2(title_lines)
        # title_lines = normalize_title_lines(title_lines)

        if title_lines:
            if not has_meaningful_title(t):
                t.title_lines = title_lines
                t.title = "\n".join(title_lines)
            else:
                # ブロック抽出タイトルが既存タイトルを包含する場合は拡張適用
                existing_key = (getattr(t, "title", "") or "").replace("\n", "").replace(" ", "").replace("\u3000", "")
                extracted_key = "\n".join(title_lines).replace("\n", "").replace(" ", "").replace("\u3000", "")
                if existing_key and existing_key in extracted_key and len(extracted_key) > len(existing_key):
                    t.title_lines = title_lines
                    t.title = "\n".join(title_lines)

        # speaker / affiliation
        speaker, affiliation = _extract_person_near_enja(seg,chair_name=getattr(payload.chair, "name", "") or "",chair_aff=getattr(payload.chair, "affiliation", "") or "",)

        # フォールバック: "演者"ラベルが無い場合、"先生"ブロックから演者名を抽出
        if not speaker:
            chair_name_key = normalize_key(getattr(payload.chair, "name", "") or "").replace("先生", "")
            chair_aff_key = normalize_key(getattr(payload.chair, "affiliation", "") or "")
            speaker_block = None
            for b in sorted(seg, key=lambda x: (x.top, x.left)):
                bt = normalize_space(b.text or "")
                if "先生" not in bt:
                    continue
                # 所属っぽいブロックは除外（大学/病院等が先頭にある）
                bt_key = normalize_key(bt)
                if any(bt_key.startswith(normalize_key(w)) for w in ["大学", "病院", "センター", "クリニック"]):
                    continue
                # "姓\n名先生" or "姓 名 先生" パターンから名前を抽出
                bt_joined = bt.replace("\n", " ")
                m_name = re.search(r"([一-龥々ぁ-んァ-ヶ]{1,5})\s*([一-龥々ぁ-んァ-ヶ]{1,5})\s*先生", bt_joined)
                if m_name:
                    cand = norm_name(m_name.group(1) + m_name.group(2))
                    cand_key = normalize_key(cand).replace("先生", "")
                    if cand_key and cand_key != chair_name_key and is_valid_person_name(cand):
                        speaker = cand
                        speaker_block = b
                        break

            # "先生"ブロックから演者を見つけた場合、近傍ブロックから所属を探す
            if speaker_block and not affiliation:
                for b in sorted(seg, key=lambda x: (x.top, x.left)):
                    # 名前ブロックの近く（上下 300,000 EMU 以内）の所属候補
                    if abs(b.top - speaker_block.top) > 300000:
                        continue
                    if b is speaker_block:
                        continue
                    bt = normalize_space(b.text or "").replace("\n", " ")
                    if looks_like_affil_line(bt) and normalize_key(bt) != chair_aff_key:
                        affiliation = bt
                        break

        # 演者が座長と同一なら再抽出対象にする
        _chair_nm_key = normalize_key(getattr(payload.chair, "name", "") or "").replace("先生", "")
        _t_sp_key = normalize_key(t.speaker or "").replace("先生", "")
        _speaker_is_chair = bool(_chair_nm_key and _t_sp_key and _t_sp_key == _chair_nm_key)

        if speaker and (not t.speaker or is_bad_speaker(t.speaker) or _speaker_is_chair):
            sp = clean_speaker_text(speaker)

            chair_name_key = normalize_key(getattr(payload.chair, "name", "") or "").replace("先生", "")
            sp_key = normalize_key(sp or "").replace("先生", "")

            if sp_key and sp_key != chair_name_key:
                t.speaker = sp
                t.speaker_display = build_speaker_display(sp) or sp

        def _affiliation_score(s: str) -> tuple[int, int]:
            s = normalize_space(s or "")
            key = normalize_key(s)

            dept_words = [
                "科", "内科", "外科", "産婦人科", "婦人科", "泌尿器科", "呼吸器内科",
                "腫瘍内科", "循環器内科", "消化器内科", "教室", "講座", "部"
            ]

            has_dept = 1 if any(w in s for w in dept_words) else 0
            return (has_dept, len(key))

        if affiliation:
            aff = clean_affiliation_text(affiliation)
            cur_aff = normalize_space(getattr(t, "affiliation", "") or "")

            if aff:
                if not cur_aff:
                    t.affiliation = aff
                else:
                    # 既存に候補が含まれるなら既存優先
                    if normalize_key(aff) in normalize_key(cur_aff):
                        t.affiliation = cur_aff
                    # 候補の方が明らかに情報量が多い時だけ更新
                    elif _affiliation_score(aff) > _affiliation_score(cur_aff):
                        t.affiliation = aff

        # 最後の保険: speaker だけ取れて affiliation がない場合は speaker_map 相当を使いたいならここで補完
        # if t.speaker and not t.affiliation:
        #     t.affiliation = aff_from_speaker_map(t.speaker)

    # ── セカンドパス: アンカー順序と talk 順序のズレで拾い漏れた所属を補完 ──
    chair_aff_key = normalize_key(getattr(payload.chair, "affiliation", "") or "")
    def _is_bad_affiliation(s: str) -> bool:
        """ラベル語だけ・明らかに不正な所属"""
        s_ns = re.sub(r'[\s\u3000]+', '', s or '')
        if not s_ns:
            return True
        if _label_aff_pat.match(s_ns):
            return True
        return False

    for t in talks:
        if not t.speaker or (normalize_space(t.affiliation or "") and not _is_bad_affiliation(t.affiliation)):
            continue
        speaker_key = normalize_key(t.speaker).replace("先生", "")
        # blocks 全体から speaker 名を含むブロックを探す
        speaker_block = None
        for b in ordered:
            bt_flat = (b.text or "").replace("\n", "").replace(" ", "").replace("\u3000", "").replace("先生", "")
            if speaker_key and speaker_key in normalize_key(bt_flat):
                speaker_block = b
                break
        if not speaker_block:
            continue
        # speaker ブロックの近傍から所属候補を探す
        for b in ordered:
            if abs(b.top - speaker_block.top) > 400000:
                continue
            if b is speaker_block:
                continue
            bt = normalize_space(b.text or "").replace("\n", " ")
            if looks_like_affil_line(bt) and normalize_key(bt) != chair_aff_key:
                t.affiliation = bt
                print(f"[repair-talks] second-pass affiliation for {t.speaker}: {bt}")
                break

    payload.talks = talks
    return payload

def strip_outer_parens_suffix(s: str) -> str:
    s = normalize_space(s)
    s = re.sub(r"[）)\]]+\s*$", "", s).strip()
    return s

def split_person_and_affiliation(s: str) -> tuple[str, str]:
    s = normalize_space(s)

    # 例: 中西重清先生（中西内科院長）
    m = re.match(r"^(?P<name>[^（(]+?)(?:先生)?\s*[（(](?P<aff>.+?)[）)]\s*$", s)
    if m:
        return normalize_space(m.group("name")), normalize_space(m.group("aff"))

    # 例: 伊東直哉先生（名古屋市立大学大学院医学研究科
    m = re.match(r"^(?P<name>[^（(]+?)(?:先生)?\s*[（(](?P<aff>.+?)\s*$", s)
    if m:
        return normalize_space(m.group("name")), normalize_space(m.group("aff"))

    # 例: 伊東直哉先生
    m = re.match(r"^(?P<name>.+?)(?:先生)\s*$", s)
    if m:
        return normalize_space(m.group("name")), ""

    return "", ""

def finalize_people_fields(payload: DesignJSON) -> DesignJSON:
    # chair
    if getattr(payload, "chair", None):
        raw_name = normalize_space(payload.chair.name or "")
        raw_aff = normalize_space(payload.chair.affiliation or "")

        # affiliation側に「名前（所属）」が残っていたら分解し直す
        # ただし所属テキスト自体に括弧が含まれる場合（○○科（リウマチ）部長 等）は
        # 誤分解しないよう、name部分が所属機関を含まないことを確認する
        n, a = split_person_and_affiliation(raw_aff)
        _ORG_KEYWORDS = ("大学", "病院", "医院", "クリニック", "センター", "科", "厚生", "診療所")
        if a and not any(kw in (n or "") for kw in _ORG_KEYWORDS):
            payload.chair.affiliation = strip_outer_parens_suffix(a)
            if not raw_name and n:
                payload.chair.name = n
                payload.chair.name_display = n
        
        # 座長所属から人名を除去（デバッグ強化版）
        if raw_name and payload.chair.affiliation:
            print(f"[DEBUG] Chair name: '{raw_name}'")
            print(f"[DEBUG] Chair affiliation before: '{payload.chair.affiliation}'")
            
            # 改行をスペースに変換して正規化
            clean_affiliation = re.sub(r'[\n\r]+', ' ', payload.chair.affiliation)
            clean_affiliation = re.sub(r'\s+', ' ', clean_affiliation).strip()
            
            # シンプルな名前除去（「髙田 慶应」「髙田慶応」両方対応）
            name_without_space = raw_name.replace(' ', '')
            name_with_space = raw_name if ' ' in raw_name else f"{raw_name[:2]} {raw_name[2:]}" if len(raw_name) >= 4 else raw_name
            
            patterns_to_remove = [
                raw_name,           # 元の名前
                name_without_space, # スペースなし版
                name_with_space,    # スペースあり版
            ]
            
            print(f"[DEBUG] Patterns to remove: {patterns_to_remove}")
            
            # 各パターンを除去
            for pattern in patterns_to_remove:
                if pattern and pattern in clean_affiliation:
                    print(f"[DEBUG] Removing: '{pattern}'")
                    clean_affiliation = clean_affiliation.replace(pattern, "")
            
            # 「先生」も除去
            clean_affiliation = clean_affiliation.replace("先生", "")
            clean_affiliation = re.sub(r'\s+', ' ', clean_affiliation).strip()
            
            payload.chair.affiliation = clean_affiliation
            print(f"[DEBUG] Chair affiliation after: '{clean_affiliation}'")

        payload.chair.affiliation = strip_outer_parens_suffix(payload.chair.affiliation or "")
        # 「ご所属：」プレフィックスを除去
        payload.chair.affiliation = re.sub(r'^ご?所属\s*[:：]\s*', '', payload.chair.affiliation)

    # talks - 不正な演者名を修正（講演自体は残す）
    for t in getattr(payload, "talks", []) or []:
        speaker = normalize_space(t.speaker or "")
        title_lines = getattr(t, "title_lines", []) or []
        title_text = " ".join(title_lines)
        
        # 「男子」かつキーワード文脈 → タイトルから誤抽出された演者名をクリア
        if speaker == "男子" and 'キーワード' in title_text:
            print(f"[DEBUG] Clearing invalid speaker: '{speaker}' (keyword context)")
            t.speaker = ""
            t.speaker_display = ""
        
        # 「ご所属：」プレフィックスを除去
        aff = normalize_space(getattr(t, "affiliation", "") or "")
        aff = re.sub(r'^ご?所属\s*[:：]\s*', '', aff)
        if aff != normalize_space(getattr(t, "affiliation", "") or ""):
            t.affiliation = aff

        # 所属に座長情報が混入している場合の最終クリーニング
        aff = normalize_space(getattr(t, "affiliation", "") or "")
        if aff and "座長" in aff:
            t.affiliation = ""
        
        # 所属に「演題名」が混入している場合の除去
        # 例: 「大阪大学 ... 助教 「 癌と腸内細菌叢について （仮）」」
        aff = getattr(t, "affiliation", "") or ""
        if aff and re.search(r'[「「]', aff):
            aff_clean = re.sub(r'\s*[「「][^」」]*[」」]?\s*$', '', aff).strip()
            if aff_clean:
                t.affiliation = aff_clean
        
        t.affiliation = strip_outer_parens_suffix(t.affiliation or "")
        
        # 演題 title / title_lines から主催・共催行を除去
        _tl = getattr(t, "title_lines", []) or []
        _tl_clean = [ln for ln in _tl if not re.match(r'^\s*(?:主催|共催)\s*[：:]', ln)]
        if _tl_clean and len(_tl_clean) < len(_tl):
            t.title_lines = _tl_clean
            t.title = "\n".join(_tl_clean)

    return payload

ROLE_KEYWORDS = [
    "主任教授", "教授", "准教授", "講師", "助教",
    "病院長", "院長", "部長", "医長", "センター長",
    "科長", "所長", "責任者"
]

def extract_role_only(s: str) -> str:
    s = normalize_space(s or "")

    # 完全一致優先（主任教授 → 教授より優先）
    for role in ROLE_KEYWORDS:
        if role in s:
            return role

    return ""

def fill_chair_role_from_blocks(payload: DesignJSON, blocks: list[TextBlock]) -> DesignJSON:
    if not getattr(payload, "chair", None):
        return payload

    name = normalize_space(payload.chair.name or "")
    if not name:
        return payload

    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    # 名前含むブロックを探す
    for i, b in enumerate(ordered):
        t = normalize_space(b.text)
        if name in t:
            # 周辺を見に行く（上下2〜3ブロック）
            for j in range(max(0, i - 2), min(len(ordered), i + 3)):
                cand = normalize_space(ordered[j].text)

                # 役職っぽいものを探す
                role = extract_role_only(cand)
                if role:
                    payload.chair.affiliation = role
                    return payload

    return payload

def repair_talk_speaker_tail_split(payload: DesignJSON) -> DesignJSON:
    for t in getattr(payload, "talks", []) or []:
        sp = normalize_space(getattr(t, "speaker", "") or "")
        aff = normalize_space(getattr(t, "affiliation", "") or "")

        if not sp:
            continue

        # 演者/講師ラベル除去
        sp_clean = re.sub(r"^(演者|講師)\s*[/／:：]?\s*", "", sp).strip()
        sp_compact = sp_clean.replace(" ", "").replace("　", "")

        # すでに人名だけならそのまま採用して終了
        if sp_compact and not any(x in sp_clean for x in ["大学", "病院", "科", "センター", "教授", "部長", "医長", "院長", "医学部"]):
            if 3 <= len(sp_compact) <= 8:
                t.speaker = sp_compact
                continue

        aff2, name2 = split_affiliation_and_name_tail(sp_clean)
        if not name2:
            continue

        compact = name2.replace(" ", "").replace("　", "")
        if len(compact) < 3:
            continue

        t.speaker = compact
        if aff2 and not aff:
            t.affiliation = aff2

    return payload


def rebuild_talks_from_anchors(blocks: list[TextBlock], payload: DesignJSON) -> DesignJSON:
    ordered = sorted(blocks, key=lambda b: (b.top, b.left))

    def _find_talk_anchor(no: int):
        labels_by_no = {
            1: ["講演1", "講演１", "講演①", "講演Ⅰ"],
            2: ["講演2", "講演２", "講演②", "講演Ⅱ"],
            3: ["講演3", "講演３", "講演③", "講演Ⅲ"],
            4: ["講演4", "講演４", "講演④", "講演Ⅳ"],
        }
        labels = labels_by_no.get(no, [f"講演{no}"])
        for b in ordered:
            key = normalize_key(b.text or "")
            if any(normalize_key(lb) in key for lb in labels):
                return b
        return None

    def _extract_time_from_seg(seg):
        for b in seg:
            txt = normalize_time_colon(normalize_space(b.text or ""))
            m = re.search(r"(\d{1,2}:\d{2})\s*[～〜~\-－—–]\s*(\d{1,2}:\d{2})", txt)
            if m:
                return f"{m.group(1)}~{m.group(2)}"
        return ""

    def _extract_title_from_seg(seg):
        lines = []
        for b in seg:
            txt = normalize_space(b.text or "")
            key = normalize_key(txt)

            if not txt:
                continue
            if "講演" in key:
                continue
            if "演者" in key:
                continue
            if "座長" in key:
                continue
            if looks_like_affil_line(txt):
                continue
            if re.search(r"\d{1,2}:\d{2}\s*[～〜~\-－—–]\s*\d{1,2}:\d{2}", txt):
                continue
            if len(txt) < 6:
                continue

            lines.append(txt)

        lines = normalize_lines_keep_order(lines)
        return lines[:3]

    def _extract_speaker_aff_from_seg(seg):
        speaker = ""
        affiliation = ""

        for i, b in enumerate(seg):
            txt = normalize_space(b.text or "")
            key = normalize_key(txt)

            if "演者" not in key:
                continue

            # 同一行: 演者 市田 晃佑 先生
            m = re.search(r"演\s*者\s*([一-龥々]{1,4}\s*[一-龥々]{1,4})\s*先生?", txt)
            if m:
                speaker = norm_name(m.group(1))

                # 同一行の所属
                m2 = re.search(
                    r"演\s*者\s*[一-龥々]{1,4}\s*[一-龥々]{1,4}\s*先生?\s*(.+)$",
                    txt
                )
                if m2:
                    cand_aff = normalize_space(m2.group(1))
                    if looks_like_affil_line(cand_aff):
                        affiliation = cand_aff
                if speaker:
                    return speaker, affiliation

            # 改行型
            lines = [normalize_space(x) for x in txt.split("\n") if normalize_space(x)]
            for li, ln in enumerate(lines):
                if "演者" not in normalize_key(ln):
                    continue
                if li + 1 < len(lines):
                    cand = re.sub(r"\s*先生$", "", lines[li + 1]).strip()
                    speaker = norm_name(cand)
                if li + 2 < len(lines):
                    cand_aff = normalize_space(lines[li + 2])
                    if looks_like_affil_line(cand_aff):
                        affiliation = cand_aff
                if speaker:
                    return speaker, affiliation

            # 次ブロック fallback
            if i + 1 < len(seg):
                nxt = normalize_space(seg[i + 1].text or "")
                cand = re.sub(r"\s*先生$", "", nxt).strip()
                cand = norm_name(cand)
                if re.fullmatch(r"[一-龥々]{2,8}", cand.replace(" ", "")):
                    speaker = cand

            if i + 2 < len(seg):
                cand_aff = normalize_space(seg[i + 2].text or "")
                if looks_like_affil_line(cand_aff):
                    affiliation = cand_aff

            if speaker:
                return speaker, affiliation

        return "", ""

    new_talks = []

    for no in [1, 2, 3, 4]:
        anchor = _find_talk_anchor(no)
        if not anchor:
            continue

        next_anchor = _find_talk_anchor(no + 1)

        x0 = anchor.left - 3000000
        x1 = anchor.left + 7000000
        y0 = anchor.top - 200000
        y1 = (next_anchor.top - 100000) if next_anchor else (anchor.top + 2500000)

        seg = [b for b in ordered if x0 <= b.left <= x1 and y0 <= b.top <= y1]
        seg = sorted(seg, key=lambda b: (b.top, b.left))

        title_lines = _extract_title_from_seg(seg)
        speaker, affiliation = _extract_speaker_aff_from_seg(seg)
        time = _extract_time_from_seg(seg)

        if not (title_lines or speaker or affiliation or time):
            continue

        title = "\n".join(title_lines).strip()

        new_talks.append(
            Talk(
                time=time,
                title=title,
                title_lines=title_lines,
                speaker=speaker,
                speaker_display=build_speaker_display(speaker) if speaker else "",
                affiliation=affiliation,
                title_overrides=[],
                honorific_title="先生",
            )
        )

    if new_talks:
        payload.talks = new_talks[:4]

    return payload

def dump_titles(tag, payload):
    print(f"--- {tag} ---")
    print(payload)
    for i, t in enumerate(payload.talks or []):
        print("idx=", i)
        print("title=", repr(getattr(t, "title", "")))
        print("title_lines=", getattr(t, "title_lines", []))
        print("speaker=", repr(getattr(t, "speaker", "")))
        print("affiliation=", repr(getattr(t, "affiliation", "")))
        
def is_non_talk_heading(s: str) -> bool:
    s = normalize_space(s or "")
    key = normalize_key(s)

    if not s:
        return True

    exact_ng = {
        "program", "p r o g r a m",
        "一般講演", "特別講演", "講演", "座長", "演者",
        "主催", "共催", "日時", "形式", "開催",
    }
    if s.lower() in exact_ng or key in {normalize_key(x) for x in exact_ng}:
        return True

    # PROGRAM のような英字ばらし
    if re.fullmatch(r"(?:[A-Za-z]\s*){4,}", s):
        return True

    # 開催形式系
    if "live配信" in s.lower() or "web" in s.lower() and "開催" in s:
        return True

    # 注意書き系
    if "ご視聴" in s or "事前参加予約" in s or "旅費の負担" in s:
        return True

    return False

async def pptx_to_json_vm_hint(pptx_path: Path, vm_rows: List[dict], debug_blocks_path: Optional[Path] = None) -> DesignJSON:
    """PPTX優先。VMは精度を上げるヒントとして blocks からの拾い直しにのみ使用し、欠損時のみVMで補完する。"""
    blocks = extract_blocks_any(pptx_path, first_only=True)
    blocks = merge_event_title_blocks_strict(blocks)

    if debug_blocks_path:
        dbg = [
            {
                "text": b.text,
                "left": b.left,
                "top": b.top,
                "width": b.width,
                "height": b.height,
                "max_font_pt": round(b.max_font_pt, 2),
            }
            for b in blocks
        ]
        debug_blocks_path.write_text(json.dumps(dbg, ensure_ascii=False, indent=2), encoding="utf-8")

    speaker_map = extract_speaker_affil_map_by_blocks(blocks)
    # VM医師名でspeaker_mapを強化（ブロック位置ベースで所属も取得）
    speaker_map = enrich_speaker_map_with_vm(speaker_map, blocks, vm_rows)
    time_candidates = extract_time_candidates_from_blocks(blocks)

    draft = parse_blocks_to_design_json(blocks, vm_rows=vm_rows)
    print("draft", draft)

    refined = await ai_refine_json(blocks, draft, speaker_map, time_candidates, vm_rows)
    dump_titles("after ai_refine_json", refined)

    # AI が VM 講演会名を上書きした場合に再適用
    if vm_rows:
        _vm_title_post = ""
        for _r in vm_rows:
            _d = _r if isinstance(_r, dict) and "data" not in _r else (_r.get("data") or {})
            _vt = normalize_space(_d.get("講演会名") or "")
            if _vt:
                _vm_title_post = _vt
                break
        if _vm_title_post:
            _ai_et_norm = (refined.event_title or "").replace(" ", "").replace("\u3000", "")
            _vm_post_norm = _vm_title_post.replace(" ", "").replace("\u3000", "")
            _post_jaccard = len(set(_ai_et_norm) & set(_vm_post_norm)) / max(len(set(_ai_et_norm) | set(_vm_post_norm)), 1)
            if _post_jaccard < 0.5:
                # blocks に存在しないテキストは導入しない（ログのみ）
                print(f"[vm-event-title] mismatch: vm='{_vm_title_post[:40]}' ai='{refined.event_title[:40]}' jaccard={_post_jaccard:.2f} (skip overwrite)")

    # AI が VM 演題を上書きした場合 → blocks に存在しないテキストは導入しない
    # VM 演題は参考情報としてのみ使用（上書き廃止）
    if vm_rows and refined.talks:
        _vm_by_speaker: dict[str, str] = {}
        for _r in vm_rows:
            _d = _r if isinstance(_r, dict) and "data" not in _r else (_r.get("data") or {})
            _sp = _norm_person_name(_d.get("案内状掲載 医師名") or "")
            _vt = normalize_space(_d.get("演題") or "")
            if _sp and _vt:
                _vm_by_speaker[_sp] = _vt
        for _t in refined.talks:
            _sp_key = _norm_person_name(getattr(_t, "speaker", "") or getattr(_t, "speaker_display", "") or "")
            _vm_enden = _vm_by_speaker.get(_sp_key, "")
            if not _vm_enden:
                continue
            _ai_title_norm = (getattr(_t, "title", "") or "").replace(" ", "").replace("\u3000", "")
            _vm_enden_norm = _vm_enden.replace(" ", "").replace("\u3000", "")
            _t_jaccard = len(set(_ai_title_norm) & set(_vm_enden_norm)) / max(len(set(_ai_title_norm) | set(_vm_enden_norm)), 1)
            if _t_jaccard < 0.4:
                print(f"[vm-talk-title] mismatch for '{_sp_key}': vm='{_vm_enden[:40]}' ai='{_ai_title_norm[:40]}' jaccard={_t_jaccard:.2f} (skip overwrite)")

    def has_chair_shift_pattern(payload: DesignJSON) -> bool:
        talks = list(payload.talks or [])
        if len(talks) < 2:
            return False

        chair_name = normalize_key(getattr(payload.chair, "name", "") or "").replace("先生", "")
        if not chair_name:
            return False

        first_sp = normalize_key(getattr(talks[0], "speaker", "") or "").replace("先生", "")
        return bool(first_sp and first_sp == chair_name)

    if has_chair_shift_pattern(refined):
        refined = rebuild_talks_from_anchors(blocks, refined)


    refined = clean_ai_talk_titles(refined)
    dump_titles("after clean_ai_talk_titles", refined)

    refined = repair_chair_from_multiline_block(refined, blocks)
    dump_titles("after repair_chair_from_multiline_block", refined)

    refined = clean_chair_fields(refined)
    dump_titles("after clean_chair_fields", refined)

    refined = repair_talks_from_blocks(refined, blocks)
    dump_titles("after repair_talks_from_blocks", refined)

    # 正解DBヒントによる空フィールド補完（blocks 内に実在するテキストのみ）
    _ca_hints = compute_correct_answer_hints(blocks, refined.event_title or "")
    refined = fill_empty_fields_from_blocks_with_hints(refined, blocks, _ca_hints)
    dump_titles("after fill_empty_fields_from_blocks_with_hints", refined)

    refined = assign_talk_times_by_anchor(blocks, refined)
    dump_titles("after assign_talk_times_by_anchor", refined)

    # 近接性による時間割り当て（既存の時間をリセットしない）
    refined = assign_talk_times_by_proximity(blocks, refined)
    dump_titles("after assign_talk_times_by_proximity", refined)
    
    # 最終手段：上位の時間による割り当て
    refined.talks = assign_talk_times_by_nearest_upper_time(blocks, refined.talks)
    dump_titles("after assign_talk_times_by_nearest_upper_time", refined)

    # 時間が割り当てられた talks を時間順にソート（演題番号順に整列）
    def _sort_talks_by_time(talks):
        if not talks or len(talks) <= 1:
            return talks
        # 全 talk に時間が設定されている場合のみソート
        if all(normalize_space(getattr(t, "time", "") or "") for t in talks):
            def _time_sort_key(t):
                tm = normalize_space(getattr(t, "time", "") or "")
                m = re.match(r"(\d{1,2}):(\d{2})", tm)
                if m:
                    return int(m.group(1)) * 60 + int(m.group(2))
                return 9999
            return sorted(talks, key=_time_sort_key)
        return talks

    refined.talks = _sort_talks_by_time(list(refined.talks or []))
    dump_titles("after sort_talks_by_time", refined)

    def _same_person(a: str, b: str) -> bool:
        return normalize_key(a or "").replace("先生", "") == normalize_key(b or "").replace("先生", "")

    def _title_key(t) -> str:
        lines = [normalize_space(x) for x in (getattr(t, "title_lines", None) or []) if normalize_space(x)]
        title = "\n".join(lines) if lines else normalize_space(getattr(t, "title", "") or "")
        return normalize_key(title)

    def drop_chair_duplicate_talks(payload: DesignJSON) -> DesignJSON:
        talks = list(payload.talks or [])
        if len(talks) <= 1:
            return payload

        chair_name = getattr(payload.chair, "name", "") or ""
        if not chair_name:
            return payload

        grouped = {}
        for t in talks:
            if _is_program_chair_item(t):
                grouped.setdefault(f"__chair__{len(grouped)}", []).append(t)
                continue
            grouped.setdefault(_title_key(t), []).append(t)

        kept = []
        for _, group in grouped.items():
            if len(group) == 1:
                kept.append(group[0])
                continue

            # 同じタイトルが複数あるなら、chair本人を除外
            non_chair = [
                t for t in group
                if not _same_person(getattr(t, "speaker", "") or "", chair_name)
            ]

            if non_chair:
                kept.extend(non_chair)
            else:
                kept.extend(group[:1])

        payload.talks = kept
        return payload

    refined = drop_chair_duplicate_talks(refined)

    vm_titles = _vm_speaker_titles(vm_rows)
    if vm_titles:
        refined = prune_talks_using_vm_titles(refined, vm_rows)
    else:
        refined = prune_talks_heuristic_only(refined)
    dump_titles("after prune_talks", refined)

    refined = apply_vm_hints_from_blocks(blocks, refined, vm_rows)
    dump_titles("after apply_vm_hints_from_blocks", refined)

    refined = fill_missing_from_vm(refined, vm_rows)
    dump_titles("after fill_missing_from_vm", refined)

    append_vm_role_to_talk_affiliation(refined, vm_rows)
    dump_titles("after append_vm_role_to_talk_affiliation", refined)

    refined = fill_chair_affiliation_from_vm_hint(refined, blocks, vm_rows)
    dump_titles("after fill_chair_affiliation_from_vm_hint", refined)

    refined = fill_chair_affiliation_from_blocks(refined, blocks)
    dump_titles("after fill_chair_affiliation_from_blocks", refined)

    refined = normalize_speaker_display(refined)
    dump_titles("after normalize_speaker_display", refined)

    refined = finalize_people_fields(refined)
    dump_titles("after finalize_people_fields", refined)

    refined = apply_learned_text_roles(refined)
    dump_titles("after apply_learned_text_roles", refined)

    refined = apply_learned_affiliation_format(refined)
    dump_titles("after apply_learned_affiliation_format", refined)

    refined = fill_datetime_parts(refined, blocks)

    refined = apply_inline_program_extraction(refined, blocks)
    dump_titles("after apply_inline_program_extraction", refined)

    # NOTE: apply_correct_answer_overlay はバッチフロー側で
    # apply_precise_typeset_initial の「後」に呼ぶ（typeset が改行位置を上書きするため）

    # --- warnings を実態に合わせて再計算 ---
    _w = list(refined.warnings or [])
    if refined.chair and refined.chair.name:
        _w = [w for w in _w if w != "missing_chair"]
    if refined.organizer:
        _w = [w for w in _w if w != "missing_organizer"]
    if refined.talks:
        _w = [w for w in _w if w != "no_talks"]
    if refined.event_title or refined.event_title_lines:
        _w = [w for w in _w if w != "missing_event_title"]
    if refined.datetime:
        _w = [w for w in _w if w != "missing_datetime"]
    refined.warnings = sorted(set(_w))

    return refined



def trim_last_pixel(path: str):
    img = Image.open(path)
    w, h = img.size
    if h > 1:
        img = img.crop((0, 0, w, h - 1))
        img.save(path, quality=100)


async def render_png_bytes(payload: DesignJSON) -> tuple[bytes, str]:
    global _cached_template, _browser

    if _cached_template is None:
        _cached_template = TEMPLATE_PATH.read_text(encoding="utf-8")

    if _browser is None:
        raise RuntimeError("Playwright browser is not initialized")

    context = await _browser.new_context(
        viewport=BASE_VIEWPORT,
        device_scale_factor=1,
    )
    page = await context.new_page()

    try:
        page.on("pageerror", lambda e: print("[pageerror]", e))
        page.on("console", lambda m: print("[console]", m.type, m.text))

        await page.goto(TEMPLATE_PATH.resolve().as_uri(), wait_until="domcontentloaded")
        await page.evaluate("() => document.fonts && document.fonts.ready")

        data_json = (
            payload.model_dump_json()
            if hasattr(payload, "model_dump_json")
            else payload.json(ensure_ascii=False)
        )
        data_obj = json.loads(data_json)

        await page.evaluate(
            """(data) => {
                window.__DATA__ = data;
                if (typeof window.__render === "function") window.__render();
            }""",
            data_obj,
        )

        await page.wait_for_selector('html[data-ready="1"]', timeout=30000)
        await page.wait_for_selector(".wrap", timeout=30000)

        await page.evaluate("""
        () => {
        const footer = document.querySelector(".footer");
        if (!footer) return;
        const r = footer.getBoundingClientRect();
        const dy = Math.round(r.top) - r.top;
        footer.style.transform = `translateY(${dy}px)`;
        }
        """)
        await page.wait_for_timeout(50)

        await page.evaluate("""
        () => {
            document.documentElement.style.margin = "0";
            document.body.style.margin = "0";
            document.body.style.padding = "0";
            const wrap = document.querySelector(".wrap");
            if (wrap) {
                wrap.style.margin = "0";
                wrap.style.display = "block";
            }
        }
        """)

        wrap = page.locator(".wrap")
        footer = page.locator(".footer")

        for _ in range(60):
            wrap_box = await wrap.bounding_box()
            footer_box = await footer.bounding_box()
            if wrap_box and footer_box and wrap_box["height"] > 10:
                break
            await page.wait_for_timeout(100)
        else:
            raise RuntimeError("layout not ready")

        wrap_box = await wrap.bounding_box()
        footer_box = await footer.bounding_box()

        if not wrap_box or not footer_box:
            raise RuntimeError("bounding box is None")

        clip_x = math.floor(wrap_box["x"])
        clip_y = math.floor(wrap_box["y"])
        clip_w = math.ceil(wrap_box["width"])

        footer_bottom = footer_box["y"] + footer_box["height"]
        clip_h = min(math.ceil(footer_bottom - wrap_box["y"]), MAX_HEIGHT)

        await page.set_viewport_size({
            "width": max(clip_x + clip_w, 1),
            "height": max(clip_y + clip_h, 1),
        })

        jpg_bytes = await page.screenshot(
            type="jpeg",
            quality=100,
            clip={
                "x": clip_x,
                "y": clip_y,
                "width": clip_w,
                "height": clip_h,
            },
        )

        html = await page.content()
        return jpg_bytes, html

    finally:
        await context.close()

async def render_png(payload: DesignJSON, out_path: Path, debug_html_path: Path):
    global _cached_template
    if _cached_template is None:
        _cached_template = TEMPLATE_PATH.read_text(encoding="utf-8")

    async with async_playwright() as p:
        browser = await p.chromium.launch(
            args=["--no-sandbox", "--disable-dev-shm-usage"],
        )

        context = await browser.new_context(
            viewport=BASE_VIEWPORT,
            device_scale_factor=1,
        )
        page = await context.new_page()

        page.on("pageerror", lambda e: print("[pageerror]", e))
        page.on("console", lambda m: print("[console]", m.type, m.text))

        await page.goto(TEMPLATE_PATH.resolve().as_uri(), wait_until="domcontentloaded")
        await page.evaluate("() => document.fonts && document.fonts.ready")

        data_json = (
            payload.model_dump_json()
            if hasattr(payload, "model_dump_json")
            else payload.json(ensure_ascii=False)
        )
        data_obj = json.loads(data_json)

        await page.evaluate(
            """(data) => {
                window.__DATA__ = data;
                if (typeof window.__render === "function") window.__render();
            }""",
            data_obj,
        )

        await page.wait_for_selector('html[data-ready="1"]', timeout=30000)
        await page.wait_for_selector(".wrap", timeout=30000)

        # 念のため余白系を潰す
        await page.evaluate("""
        () => {
            document.documentElement.style.margin = "0";
            document.body.style.margin = "0";
            document.body.style.padding = "0";
            const wrap = document.querySelector(".wrap");
            if (wrap) {
                wrap.style.margin = "0";
                wrap.style.display = "block";
            }
        }
        """)

        wrap = page.locator(".wrap")

        for _ in range(60):
            box = await wrap.bounding_box()
            if box and box["height"] and box["height"] > 10:
                break
            await page.wait_for_timeout(100)
        else:
            html = await page.content()
            debug_html_path.write_text(html, encoding="utf-8")
            raise RuntimeError(f"wrap bounding box not ready; wrote {debug_html_path}")

        box = await wrap.bounding_box()
        if not box:
            raise RuntimeError("wrap bounding box is None")

        clip_x = math.floor(box["x"])
        clip_y = math.floor(box["y"])
        clip_w = math.ceil(box["width"])
        clip_h = min(math.ceil(box["height"]), MAX_HEIGHT)

        print(f"wrap bounding box: {box}")
        print(f"clip: x={clip_x}, y={clip_y}, w={clip_w}, h={clip_h}")

        await page.set_viewport_size({
            "width": max(clip_x + clip_w, 1),
            "height": max(clip_y + clip_h, 1),
        })

        await page.screenshot(
            path=str(out_path),
            type="jpeg",
            quality=100,
            clip={
                "x": clip_x,
                "y": clip_y,
                "width": clip_w,
                "height": clip_h,
            },
        )

        # Chromium の 1px 余り対策
        trim_last_pixel(str(out_path))

        await context.close()
        await browser.close()



def guess_content_type(path: Path) -> str:
    ctype, _ = mimetypes.guess_type(str(path))
    return ctype or "application/octet-stream"


def upload_to_storage(local_path: Path, remote_path: str, upsert: bool = True):
    url = f"{SUPABASE_URL}/storage/v1/object/{SUPABASE_BUCKET}/{remote_path}"

    headers = {
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Content-Type": guess_content_type(local_path),
        "x-upsert": "true" if upsert else "false",
    }

    with local_path.open("rb") as f:
        res = requests.post(url, headers=headers, data=f)

    if not res.ok:
        raise RuntimeError(f"storage upload failed: {res.status_code} {res.text}")

    return res.json()


def upload_bytes_to_storage(
    data: bytes,
    remote_path: str,
    content_type: str,
    upsert: bool = True,
):
    url = f"{SUPABASE_URL}/storage/v1/object/{SUPABASE_BUCKET}/{remote_path}"
    headers = {
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Content-Type": content_type,
        "x-upsert": "true" if upsert else "false",
        "Cache-Control": "no-cache",
    }

    res = requests.post(url, headers=headers, data=data, timeout=60)
    if not res.ok:
        raise RuntimeError(f"storage upload failed: {res.status_code} {res.text}")
    return res.json()


def upload_json_to_storage(data: dict, remote_path: str, upsert: bool = True):
    raw = json.dumps(data, ensure_ascii=False).encode("utf-8")
    return upload_bytes_to_storage(
        raw,
        remote_path=remote_path,
        content_type="application/json; charset=utf-8",
        upsert=upsert,
    )


def upload_text_to_storage(text: str, remote_path: str, content_type: str, upsert: bool = True):
    return upload_bytes_to_storage(
        text.encode("utf-8"),
        remote_path=remote_path,
        content_type=content_type,
        upsert=upsert,
    )

def _storage_auth_headers() -> dict:
    return {
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
    }

def _authenticated_storage_url(remote_path: str) -> str:
    return f"{SUPABASE_URL}/storage/v1/object/{SUPABASE_BUCKET}/{remote_path}"

def download_storage_file(remote_path: str) -> bytes:
    url = _authenticated_storage_url(remote_path)
    r = requests.get(url, headers=_storage_auth_headers(), timeout=30)
    if r.status_code >= 400:
        raise HTTPException(status_code=404, detail=f"storage file not found: {remote_path}")
    return r.content

def download_storage_json(remote_path: str) -> dict:
    url = _authenticated_storage_url(remote_path)
    r = requests.get(url, headers=_storage_auth_headers(), timeout=10)
    if r.status_code >= 400:
        raise HTTPException(status_code=404, detail="job json not found")
    return r.json()

def _normalize_signed_url(signed: str) -> str:
    if signed.startswith("http"):
        return signed
    if signed.startswith("/") and not signed.startswith("/storage/v1"):
        signed = f"/storage/v1{signed}"
    return f"{SUPABASE_URL}{signed}"

# --- signed URL cache (in-memory, TTL-based) ---
_signed_url_cache: Dict[str, tuple] = {}  # path -> (url, expires_at)
_SIGNED_URL_TTL = 2400  # 40 min (URLs expire in 60 min)

def _get_cached_signed_urls(paths: list[str]) -> tuple[dict[str, str], list[str]]:
    """Return (cached_map, missing_paths)."""
    now = time.time()
    cached = {}
    missing = []
    for p in paths:
        entry = _signed_url_cache.get(p)
        if entry and entry[1] > now:
            cached[p] = entry[0]
        else:
            missing.append(p)
    return cached, missing

def _put_signed_url_cache(mapping: dict[str, str]):
    exp = time.time() + _SIGNED_URL_TTL
    for p, url in mapping.items():
        _signed_url_cache[p] = (url, exp)
    # evict expired entries when cache grows
    if len(_signed_url_cache) > 5000:
        now = time.time()
        expired = [k for k, v in _signed_url_cache.items() if v[1] <= now]
        for k in expired:
            del _signed_url_cache[k]

def create_signed_url(remote_path: str, expires_in: int = 3600) -> str:
    cached, _ = _get_cached_signed_urls([remote_path])
    if remote_path in cached:
        return cached[remote_path]
    url = f"{SUPABASE_URL}/storage/v1/object/sign/{SUPABASE_BUCKET}/{remote_path}"
    headers = {**_storage_auth_headers(), "Content-Type": "application/json"}
    r = requests.post(url, headers=headers, json={"expiresIn": expires_in}, timeout=10)
    if r.status_code >= 400:
        raise RuntimeError(f"signed url failed: {r.status_code} {r.text}")
    body = r.json()
    signed = body.get("signedURL") or body.get("signedUrl") or ""
    if not signed:
        raise RuntimeError(f"signed url empty: {body}")
    result = _normalize_signed_url(signed)
    _put_signed_url_cache({remote_path: result})
    return result

def create_signed_urls_batch(remote_paths: list[str], expires_in: int = 3600) -> dict[str, str]:
    if not remote_paths:
        return {}
    cached, missing = _get_cached_signed_urls(remote_paths)
    if not missing:
        return cached
    url = f"{SUPABASE_URL}/storage/v1/object/sign/{SUPABASE_BUCKET}"
    headers = {**_storage_auth_headers(), "Content-Type": "application/json"}
    r = requests.post(url, headers=headers, json={"expiresIn": expires_in, "paths": missing}, timeout=15)
    if r.status_code >= 400:
        logger.warning(f"batch signed url failed: {r.status_code} {r.text}")
        return cached
    fresh = {}
    for item in r.json():
        path = item.get("path", "")
        signed = item.get("signedURL") or item.get("signedUrl") or ""
        if path and signed:
            fresh[path] = _normalize_signed_url(signed)
    _put_signed_url_cache(fresh)
    return {**cached, **fresh}

def delete_storage_files(paths: list[str]):
    if not paths:
        return

    url = f"{SUPABASE_URL}/storage/v1/object/{SUPABASE_BUCKET}"
    headers = {
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Content-Type": "application/json",
    }

    # Supabase Storage delete は「paths の配列」で消す
    res = requests.delete(
        url,
        headers=headers,
        json={"prefixes": paths},  # ← ここは環境によって paths ではなく prefixes でなく files list endpoint差異あり得る
        timeout=30,
    )

    if not res.ok:
        raise RuntimeError(f"storage delete failed: {res.status_code} {res.text}")

def storage_paths(job_id: str) -> dict:
    return {
        "preview": f"{job_id}/preview.jpg",
        "json": f"{job_id}/latest.json",
        "debug_blocks": f"{job_id}/debug_blocks.json",
        "debug_html": f"{job_id}/debug.html",
    }


def upload_required_assets(
    job_id: str,
    payload_dict: dict,
    jpg_bytes: bytes,
):
    sp = storage_paths(job_id)

    # 必須
    upload_bytes_to_storage(
        jpg_bytes,
        sp["preview"],
        content_type="image/jpeg",
        upsert=True,
    )

    upload_json_to_storage(
        payload_dict,
        sp["json"],
        upsert=True,
    )


def upload_optional_assets(
    job_id: str,
    debug_html: Optional[str] = None,
    debug_blocks_path: Optional[Path] = None,
):
    sp = storage_paths(job_id)

    # 任意
    if debug_html:
        upload_text_to_storage(
            debug_html,
            sp["debug_html"],
            content_type="text/html; charset=utf-8",
            upsert=True,
        )

    if debug_blocks_path and debug_blocks_path.exists():
        upload_to_storage(
            debug_blocks_path,
            sp["debug_blocks"],
            upsert=True,
        )


# --- 非同期並列版 Storage アップロード ---
_upload_executor = ThreadPoolExecutor(max_workers=4)

async def upload_all_assets_async(
    job_id: str,
    payload_dict: dict,
    jpg_bytes: bytes,
    debug_html: Optional[str] = None,
    debug_blocks_path: Optional[Path] = None,
):
    """必須＋任意アセットを並列で Storage にアップロード"""
    sp = storage_paths(job_id)
    loop = asyncio.get_event_loop()

    def _upload_jpg():
        upload_bytes_to_storage(jpg_bytes, sp["preview"], content_type="image/jpeg", upsert=True)

    def _upload_json():
        upload_json_to_storage(payload_dict, sp["json"], upsert=True)

    tasks = [
        loop.run_in_executor(_upload_executor, _upload_jpg),
        loop.run_in_executor(_upload_executor, _upload_json),
    ]

    if debug_html:
        _dh = debug_html

        def _upload_html():
            upload_text_to_storage(_dh, sp["debug_html"], content_type="text/html; charset=utf-8", upsert=True)

        tasks.append(loop.run_in_executor(_upload_executor, _upload_html))

    if debug_blocks_path and debug_blocks_path.exists():
        _dbp = debug_blocks_path

        def _upload_blocks():
            upload_to_storage(_dbp, sp["debug_blocks"], upsert=True)

        tasks.append(loop.run_in_executor(_upload_executor, _upload_blocks))

    await asyncio.gather(*tasks)

# ---------------- App ----------------
app = FastAPI(title="PPTX → JSON → HTML → jpg (Keep Newlines + Split ~...~)")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://fragment-design-tool.vercel.app",
        "http://localhost:5173",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.mount("/fonts", StaticFiles(directory=str(APP_DIR / "fonts")), name="fonts")

# データベース接続エラーのグローバルハンドラー
@app.exception_handler(OperationalError)
async def database_exception_handler(request: Request, exc: OperationalError):
    error_msg = str(exc)
    logger.error(f"Database operational error: {error_msg}")
    
    if "timed out" in error_msg.lower():
        return JSONResponse(
            status_code=503,
            content={
                "error": "database_timeout",
                "message": "データベースへの接続がタイムアウトしました。しばらく待ってから再試行してください。",
                "details": "Database connection timeout. Please try again later."
            }
        )
    elif "connection" in error_msg.lower():
        return JSONResponse(
            status_code=503, 
            content={
                "error": "database_connection_failed",
                "message": "データベースに接続できません。システム管理者にお問い合わせください。",
                "details": "Database connection failed. Please contact system administrator."
            }
        )
    else:
        return JSONResponse(
            status_code=500,
            content={
                "error": "database_error",
                "message": "データベースエラーが発生しました。",
                "details": f"Database error: {error_msg}"
            }
        )

@app.on_event("startup")
async def startup():
    global _cached_template, _pw, _browser

    # データベース初期化（エラーハンドリング付き）
    try:
        logger.info("Starting database initialization...")
        init_db()
        logger.info("Database initialization completed successfully")
    except OperationalError as e:
        error_msg = f"Database connection failed during startup: {e}"
        logger.error(error_msg)
        if "timed out" in str(e).lower():
            logger.error("This appears to be a connection timeout issue. Please check:")
            logger.error("1. Network connectivity to the database server")
            logger.error("2. Database server status")
            logger.error("3. Firewall settings")
            logger.error("4. DATABASE_URL configuration")
        raise RuntimeError(error_msg) from e
    except Exception as e:
        error_msg = f"Unexpected error during database initialization: {e}"
        logger.error(error_msg)
        raise RuntimeError(error_msg) from e

    if not TEMPLATE_PATH.exists():
        raise RuntimeError(f"template.html not found: {TEMPLATE_PATH}")
    _cached_template = TEMPLATE_PATH.read_text(encoding="utf-8")

    # Playwright browser path (Render disk)
    browsers_path = os.getenv("PLAYWRIGHT_BROWSERS_PATH")
    if browsers_path:
        Path(browsers_path).mkdir(parents=True, exist_ok=True)

    # Playwright を1回だけ起動して使い回す
    _pw = await async_playwright().start()
    try:
        os.environ.setdefault("LECTURE_TOOL_CHROME_EXECUTABLE_PATH", _pw.chromium.executable_path)
    except Exception as e:
        print("[startup chromium executable path warning]", e)
    _browser = await _pw.chromium.launch(
        args=["--no-sandbox", "--disable-dev-shm-usage"],
    )

    # Ensure Chromium exists
    # try:
    #     subprocess.check_call(["python", "-m", "playwright", "install", "chromium"])
    # except Exception as e:
    #     raise RuntimeError(f"Playwright install failed: {e}")

@app.on_event("shutdown")
async def shutdown():
    global _browser, _pw

    try:
        if _browser is not None:
            await _browser.close()
    except Exception as e:
        print("[shutdown browser close error]", e)
    finally:
        _browser = None

    try:
        if _pw is not None:
            await _pw.stop()
    except Exception as e:
        print("[shutdown playwright stop error]", e)
    finally:
        _pw = None


def _sse(event: str, data: dict) -> str:
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"


def _lecture_normalize_drive_folder_id(value: str) -> str:
    text = (value or "").strip()
    if not text:
        return ""
    folder_match = re.search(r"/folders/([^/?#]+)", text)
    if folder_match:
        return folder_match.group(1)
    query_match = re.search(r"[?&]id=([^&#]+)", text)
    if query_match:
        return query_match.group(1)
    return text


LECTURE_TOOL_SPREADSHEET_KEY = os.getenv(
    "LECTURE_TOOL_SPREADSHEET_KEY",
    "1BA4e8UpnC9MSA7vIaX6nDKwbLnJjSlR4WGfnSilycQQ",
)
LECTURE_TOOL_SHEET_GID = int(os.getenv("LECTURE_TOOL_SHEET_GID", "156086772"))
LECTURE_TOOL_OUTPUT_DIR = DATA_DIR / "lecture_search_guide"
LECTURE_TOOL_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
LECTURE_TOOL_IMAGE_WIDTH = 2048
LECTURE_TOOL_IMAGE_QUALITY = 100
LECTURE_TOOL_DRIVE_IMAGE_QUALITY = int(os.getenv("LECTURE_TOOL_DRIVE_IMAGE_QUALITY", "95"))
LECTURE_TOOL_DRIVE_CONFIG_SIZES = [(453, 640), (480, 679)]
LECTURE_TOOL_TARGET_IMAGE_BYTES = int(os.getenv("LECTURE_TOOL_TARGET_IMAGE_BYTES", str(100 * 1024)))
LECTURE_TOOL_TARGET_IMAGE_TOLERANCE = int(os.getenv("LECTURE_TOOL_TARGET_IMAGE_TOLERANCE", str(15 * 1024)))
LECTURE_TOOL_DRIVE_FOLDER_ID = _lecture_normalize_drive_folder_id(os.getenv("LECTURE_TOOL_DRIVE_FOLDER_ID") or "")
LECTURE_TOOL_DEFAULT_VAULT_COMMAND = f"node {shlex.quote(str(APP_DIR / 'tools' / 'lecture_vault_register.js'))}"
LECTURE_TOOL_VAULT_COMMAND = (os.getenv("LECTURE_TOOL_VAULT_COMMAND") or LECTURE_TOOL_DEFAULT_VAULT_COMMAND).strip()
LECTURE_TOOL_VAULT_ACCOUNTS = [
    account.strip()
    for account in (
        os.getenv(
            "LECTURE_TOOL_VAULT_ACCOUNTS",
            "mika.hirawatari@msd.com,maika.mori@msd.com,yura.fukuhara@msd.com,hidenori.sonohata@msd.com,Hayato.Seto@vv-agency.com",
        )
    ).split(",")
    if account.strip()
]


LECTURE_TOOL_COLUMN_ALIASES = {
    "lecture_id": ["講演会ID", "システムID", "event_id", "Event ID", "講演会 id"],
    "presentation_id": ["プレゼンテーションID", "Presentation ID", "presentation_id", "crmPresentationId_b"],
    "product": ["Product", "プロダクト", "製品", "製品名"],
    "reception_date": ["受付日", "受付日付", "開催受付日", "reception_date"],
    "category": ["区分", "分類", "category"],
    "event_date": ["開催日", "開催日時", "event_date"],
    "event_time": ["時間", "開催時間", "event_time"],
    "event_name": ["講演会名", "event_name"],
    "image_name": ["画像名", "画像ファイル名", "image", "image_name"],
    "media_file_name": ["メディアファイル名", "メディア名", "media", "media_file_name"],
    "presentation_name": [
        "プレゼンテーション/キーメッセージ名",
        "プレゼンテーション名",
        "キーメッセージ名",
        "presentation",
    ],
}


class LectureToolCancelled(Exception):
    pass


LECTURE_TOOL_CANCEL_LOCK = threading.Lock()
LECTURE_TOOL_CANCEL_CONTROLLERS: dict[str, dict[str, Any]] = {}
LECTURE_TOOL_VAULT_MAX_PARALLEL = max(
    1,
    int(os.getenv("LECTURE_TOOL_VAULT_MAX_PARALLEL", str(max(1, len(LECTURE_TOOL_VAULT_ACCOUNTS))))),
)
LECTURE_TOOL_VAULT_SEMAPHORE = threading.BoundedSemaphore(LECTURE_TOOL_VAULT_MAX_PARALLEL)
LECTURE_TOOL_VAULT_LOCKS_LOCK = threading.Lock()
LECTURE_TOOL_VAULT_ACCOUNT_LOCKS: dict[str, threading.Lock] = {}


def _lecture_register_cancel_controller(session_id: str) -> dict[str, Any]:
    controller = {
        "event": threading.Event(),
        "process": None,
        "createdAt": datetime.now(timezone.utc).isoformat(),
    }
    with LECTURE_TOOL_CANCEL_LOCK:
        LECTURE_TOOL_CANCEL_CONTROLLERS[session_id] = controller
    return controller


def _lecture_unregister_cancel_controller(session_id: str) -> None:
    with LECTURE_TOOL_CANCEL_LOCK:
        LECTURE_TOOL_CANCEL_CONTROLLERS.pop(session_id, None)


def _lecture_set_cancel_process(session_id: str, proc: subprocess.Popen | None) -> None:
    should_stop = False
    with LECTURE_TOOL_CANCEL_LOCK:
        controller = LECTURE_TOOL_CANCEL_CONTROLLERS.get(session_id)
        if controller is not None:
            controller["process"] = proc
            should_stop = bool(proc and controller["event"].is_set())
    if should_stop and proc:
        _lecture_stop_process(proc)


def _lecture_stop_process(proc: subprocess.Popen) -> None:
    if proc.poll() is not None:
        return
    try:
        os.killpg(proc.pid, signal.SIGTERM)
    except Exception:
        try:
            proc.terminate()
        except Exception:
            pass


def _lecture_cancel_session(session_id: str) -> bool:
    proc = None
    with LECTURE_TOOL_CANCEL_LOCK:
        controller = LECTURE_TOOL_CANCEL_CONTROLLERS.get(session_id)
        if not controller:
            return False
        controller["event"].set()
        proc = controller.get("process")
    if proc:
        _lecture_stop_process(proc)
    return True


def _lecture_raise_if_cancelled(cancel_event: threading.Event | None, message: str = "処理を中断しました。") -> None:
    if cancel_event and cancel_event.is_set():
        raise LectureToolCancelled(message)


def _lecture_vault_account_key(vault_account: str) -> str:
    return normalize_space(vault_account or "").lower()


def _lecture_get_vault_account_lock(vault_account: str) -> threading.Lock:
    key = _lecture_vault_account_key(vault_account) or "__blank__"
    with LECTURE_TOOL_VAULT_LOCKS_LOCK:
        lock = LECTURE_TOOL_VAULT_ACCOUNT_LOCKS.get(key)
        if lock is None:
            lock = threading.Lock()
            LECTURE_TOOL_VAULT_ACCOUNT_LOCKS[key] = lock
        return lock


def _lecture_header_key(value: str) -> str:
    return re.sub(r"[\s\u3000]+", "", normalize_space(value or "")).lower()


def _lecture_sheet_text(value: Any) -> str:
    return str(value or "").strip(" \t\r\n\f\v\u3000")


def _lecture_pick_column(headers: list[str], field: str) -> str:
    normalized = {_lecture_header_key(h): h for h in headers}
    for alias in LECTURE_TOOL_COLUMN_ALIASES[field]:
        hit = normalized.get(_lecture_header_key(alias))
        if hit:
            return hit
    return ""


def _lecture_find_header(values: list[list[str]]) -> tuple[int, list[str], dict[str, str]]:
    for idx, row in enumerate(values[:20], start=1):
        headers = make_unique(row)
        columns = {
            field: _lecture_pick_column(headers, field)
            for field in LECTURE_TOOL_COLUMN_ALIASES.keys()
        }
        if columns["image_name"] and columns["media_file_name"] and columns["presentation_name"]:
            return idx, headers, columns
    raise RuntimeError("必要な列（画像名、メディアファイル名、プレゼンテーション/キーメッセージ名）が見つかりません。")


def _lecture_safe_name(value: str, fallback: str) -> str:
    text = normalize_space(value or "")
    text = re.sub(r'[\\/:*?"<>|]+', "_", text)
    text = re.sub(r"\s+", "", text)
    text = text.strip("._ ")
    return text or fallback


def _lecture_normalize_filename(value: str) -> str:
    name = Path(normalize_space(value or "")).name
    stem = Path(name).stem if Path(name).suffix else name
    return re.sub(r"[\s\u3000]+", "", stem).lower()


def _lecture_sheet_image_filename(row: dict[str, Any], fallback: str) -> str:
    raw = _lecture_sheet_text(row.get("imageName") or fallback or "image.jpg")
    name = Path(raw).name
    stem = _lecture_safe_name(Path(name).stem if Path(name).suffix else name, "image")
    return f"{stem}.jpg"


def _lecture_spreadsheet_url() -> str:
    return f"https://docs.google.com/spreadsheets/d/{LECTURE_TOOL_SPREADSHEET_KEY}/edit?gid={LECTURE_TOOL_SHEET_GID}"


def _lecture_open_spreadsheet_and_worksheet():
    scope = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive",
    ]
    credentials = get_gsa_credentials(scope)
    gc = gspread.authorize(credentials)
    workbook = _retry_gspread(lambda: gc.open_by_key(LECTURE_TOOL_SPREADSHEET_KEY))
    worksheets = _retry_gspread(lambda: workbook.worksheets())
    for ws in worksheets:
        if int(getattr(ws, "id", -1)) == LECTURE_TOOL_SHEET_GID:
            return workbook, ws
    raise RuntimeError(f"gid={LECTURE_TOOL_SHEET_GID} のシートが見つかりません。")


def _lecture_open_worksheet():
    return _lecture_open_spreadsheet_and_worksheet()[1]


def _lecture_drive_credentials():
    scope = ["https://www.googleapis.com/auth/drive"]
    credentials = get_gsa_credentials(scope)
    credentials.refresh(GoogleAuthRequest())
    return credentials


def _lecture_drive_escape_query_value(value: str) -> str:
    return value.replace("\\", "\\\\").replace("'", "\\'")


def _lecture_drive_error(resp: requests.Response, action: str) -> RuntimeError:
    message = resp.text
    try:
        body = resp.json()
        message = body.get("error", {}).get("message") or message
    except Exception:
        pass
    return RuntimeError(f"Google Drive {action} failed ({resp.status_code}): {message}")


def _lecture_drive_find_files(folder_id: str, filename: str) -> list[dict[str, Any]]:
    credentials = _lecture_drive_credentials()
    q = (
        f"name = '{_lecture_drive_escape_query_value(filename)}' "
        f"and '{_lecture_drive_escape_query_value(folder_id)}' in parents "
        "and trashed = false"
    )
    resp = requests.get(
        "https://www.googleapis.com/drive/v3/files",
        headers={"Authorization": f"Bearer {credentials.token}"},
        params={
            "q": q,
            "fields": "files(id,name,modifiedTime)",
            "pageSize": 20,
            "orderBy": "modifiedTime desc",
            "supportsAllDrives": "true",
            "includeItemsFromAllDrives": "true",
            "corpora": "allDrives",
        },
        timeout=30,
    )
    if not resp.ok:
        raise _lecture_drive_error(resp, "file lookup")
    return resp.json().get("files") or []


def _lecture_drive_upload_image(filename: str, data: bytes) -> dict[str, Any]:
    if not LECTURE_TOOL_DRIVE_FOLDER_ID:
        return {}

    credentials = _lecture_drive_credentials()
    headers = {"Authorization": f"Bearer {credentials.token}"}
    existing_files = _lecture_drive_find_files(LECTURE_TOOL_DRIVE_FOLDER_ID, filename)
    existing_id = existing_files[0]["id"] if existing_files else ""
    metadata = {"name": filename, "mimeType": "image/jpeg"}
    if not existing_id:
        metadata["parents"] = [LECTURE_TOOL_DRIVE_FOLDER_ID]

    files = {
        "metadata": ("metadata", json.dumps(metadata, ensure_ascii=False), "application/json; charset=UTF-8"),
        "file": (filename, data, "image/jpeg"),
    }

    if existing_id:
        url = f"https://www.googleapis.com/upload/drive/v3/files/{existing_id}"
        resp = requests.patch(
            url,
            headers=headers,
            params={"uploadType": "multipart", "supportsAllDrives": "true", "fields": "id,name,webViewLink"},
            files=files,
            timeout=60,
        )
    else:
        resp = requests.post(
            "https://www.googleapis.com/upload/drive/v3/files",
            headers=headers,
            params={"uploadType": "multipart", "supportsAllDrives": "true", "fields": "id,name,webViewLink"},
            files=files,
            timeout=60,
        )
    if not resp.ok:
        raise _lecture_drive_error(resp, "upload")
    body = resp.json()
    return {
        "id": body.get("id", existing_id),
        "name": body.get("name", filename),
        "webViewLink": body.get("webViewLink", ""),
        "size": len(data),
        "updated": bool(existing_id),
        "duplicateExistingCount": max(0, len(existing_files) - 1),
    }


def _lecture_bool(value: Any, default: bool = True) -> bool:
    if value is None:
        return default
    return str(value).strip().lower() in {"1", "true", "yes", "on", "checked"}


def _lecture_progress(progress, step: str, message: str, **data) -> None:
    if progress:
        progress({"step": step, "message": message, **data})


def _lecture_vault_register_packages(
    packages: list[dict[str, Any]],
    vault_account: str,
    progress=None,
    *,
    session_id: str = "",
    cancel_event: threading.Event | None = None,
) -> dict[str, Any]:
    if not LECTURE_TOOL_VAULT_COMMAND:
        raise RuntimeError("Vault登録コマンドが未設定です。LECTURE_TOOL_VAULT_COMMAND を設定してください。")

    _lecture_raise_if_cancelled(cancel_event, "Vault登録を中断しました。")
    payload = json.dumps(
        {
            "vaultAccount": vault_account,
            "packages": packages,
        },
        ensure_ascii=False,
    )
    _lecture_progress(
        progress,
        "vault",
        f"Vault登録をまとめて開始しています: {len(packages)} 件",
        count=len(packages),
        vaultAccount=vault_account,
    )
    proc = subprocess.Popen(
        shlex.split(LECTURE_TOOL_VAULT_COMMAND),
        stdin=subprocess.PIPE,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        start_new_session=True,
    )
    if session_id:
        _lecture_set_cancel_process(session_id, proc)
    timeout_seconds = int(os.getenv("LECTURE_TOOL_VAULT_TIMEOUT", "1800"))
    timed_out = threading.Event()

    def kill_on_timeout():
        if proc.poll() is None:
            timed_out.set()
            _lecture_stop_process(proc)

    watchdog = threading.Timer(timeout_seconds, kill_on_timeout)
    watchdog.daemon = True
    watchdog.start()
    assert proc.stdin is not None
    assert proc.stdout is not None
    assert proc.stderr is not None
    try:
        proc.stdin.write(payload)
        proc.stdin.close()

        stderr_lines: list[str] = []
        marker_results: list[dict[str, Any]] = []
        for line in proc.stderr:
            _lecture_raise_if_cancelled(cancel_event, "Vault登録を中断しました。")
            text = line.strip()
            if not text:
                continue
            if text.startswith("__VAULT_RESULT__"):
                try:
                    result_item = json.loads(text.removeprefix("__VAULT_RESULT__"))
                    if isinstance(result_item, dict):
                        marker_results.append(result_item)
                        if progress:
                            progress({"event": "vault-result", "result": result_item})
                except Exception:
                    pass
                continue
            stderr_lines.append(text)
            _lecture_progress(progress, "vault", text, vaultAccount=vault_account)

        stdout = proc.stdout.read()
        try:
            return_code = proc.wait(timeout=10)
        except subprocess.TimeoutExpired:
            _lecture_stop_process(proc)
            return_code = proc.wait(timeout=10)
        _lecture_raise_if_cancelled(cancel_event, "Vault登録を中断しました。")
        if timed_out.is_set():
            raise RuntimeError(f"Vault登録に失敗しました。timeout={timeout_seconds}s: プロセスが終了しなかったため停止しました。")
        if return_code != 0:
            message = "\n".join(stderr_lines) or stdout.strip()
            raise RuntimeError(f"Vault登録に失敗しました。exit={return_code}: {message}")

        stdout = (stdout or "").strip()
        try:
            body = json.loads(stdout) if stdout else {}
        except Exception:
            body = {"output": stdout}
        if marker_results and (not isinstance(body, dict) or not isinstance(body.get("results"), list)):
            return {"results": marker_results, "vaultAccount": vault_account}
        return body if isinstance(body, dict) else {"output": body}
    finally:
        watchdog.cancel()
        if session_id:
            _lecture_set_cancel_process(session_id, None)


def _lecture_fetch_sheet_rows() -> dict[str, Any]:
    workbook, ws = _lecture_open_spreadsheet_and_worksheet()
    values = _retry_gspread(lambda: ws.get_all_values())
    if not values:
        return {
            "spreadsheetTitle": getattr(workbook, "title", ""),
            "spreadsheetUrl": _lecture_spreadsheet_url(),
            "sheetTitle": ws.title,
            "sheetGid": getattr(ws, "id", LECTURE_TOOL_SHEET_GID),
            "headerRow": 0,
            "columns": {},
            "rows": [],
        }

    header_row, headers, columns = _lecture_find_header(values)
    rows: list[dict[str, Any]] = []

    for row_number, raw in enumerate(values[header_row:], start=header_row + 1):
        data = {headers[i]: raw[i] if i < len(raw) else "" for i in range(len(headers))}
        lecture_id = normalize_space(data.get(columns.get("lecture_id", ""), ""))
        presentation_id = normalize_space(data.get(columns.get("presentation_id", ""), ""))
        product = normalize_space(data.get(columns.get("product", ""), ""))
        reception_date = normalize_space(data.get(columns.get("reception_date", ""), ""))
        category = normalize_space(data.get(columns.get("category", ""), ""))
        event_name = _lecture_sheet_text(data.get(columns.get("event_name", ""), ""))
        event_date = normalize_space(data.get(columns.get("event_date", ""), ""))
        event_time = normalize_space(data.get(columns.get("event_time", ""), ""))
        image_name = _lecture_sheet_text(data.get(columns["image_name"], ""))
        media_file_name = _lecture_sheet_text(data.get(columns["media_file_name"], ""))
        presentation_name = _lecture_sheet_text(data.get(columns["presentation_name"], ""))

        if not (image_name or media_file_name or presentation_name):
            continue

        rows.append(
            {
                "id": str(row_number),
                "rowNumber": row_number,
                "lectureId": lecture_id,
                "presentationId": presentation_id,
                "product": product,
                "receptionDate": reception_date,
                "category": category,
                "imageName": image_name,
                "imageKey": _lecture_normalize_filename(image_name),
                "mediaFileName": media_file_name,
                "presentationName": presentation_name,
                "eventName": event_name,
                "eventDate": event_date,
                "eventTime": event_time,
            }
        )

    return {
        "spreadsheetTitle": getattr(workbook, "title", ""),
        "spreadsheetUrl": _lecture_spreadsheet_url(),
        "sheetTitle": ws.title,
        "sheetGid": getattr(ws, "id", LECTURE_TOOL_SHEET_GID),
        "headerRow": header_row,
        "columns": columns,
        "rows": rows,
    }


def _lecture_session_root(session_id: str) -> Path:
    if not re.fullmatch(r"[a-f0-9]{32}", session_id or ""):
        raise HTTPException(status_code=400, detail="invalid session id")
    return LECTURE_TOOL_OUTPUT_DIR / session_id


def _lecture_result_files(session_id: str, limit: int = 300) -> list[dict[str, Any]]:
    result_dir = _lecture_session_root(session_id) / "result"
    if not result_dir.exists():
        return []

    files: list[dict[str, Any]] = []
    for path in result_dir.rglob("*"):
        if not path.is_file():
            continue
        try:
            stat = path.stat()
        except OSError:
            continue
        files.append(
            {
                "path": path.relative_to(result_dir).as_posix(),
                "name": path.name,
                "size": stat.st_size,
                "modified": datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(),
            }
        )

    files.sort(key=lambda item: item["modified"], reverse=True)
    return files[:limit]


def _lecture_jpeg_bytes(
    img: Image.Image,
    *,
    quality: int,
    subsampling: int = 2,
    optimize: bool = True,
) -> bytes:
    out = io.BytesIO()
    img.save(out, format="JPEG", quality=quality, subsampling=subsampling, optimize=optimize)
    return out.getvalue()


def _lecture_jpeg_near_target(img: Image.Image) -> tuple[bytes, int]:
    target = LECTURE_TOOL_TARGET_IMAGE_BYTES
    tolerance = LECTURE_TOOL_TARGET_IMAGE_TOLERANCE
    lower = max(0, target - tolerance)
    upper = target + tolerance

    candidates: list[tuple[int, int, bytes]] = []
    for subsampling in (2, 1, 0):
        lo, hi = 1, 100
        best_under: tuple[int, bytes] | None = None
        while lo <= hi:
            q = (lo + hi) // 2
            data = _lecture_jpeg_bytes(img, quality=q, subsampling=subsampling)
            size = len(data)
            candidates.append((abs(size - target), q, data))
            if size <= upper:
                best_under = (q, data)
                lo = q + 1
            else:
                hi = q - 1

        if best_under:
            q, data = best_under
            if len(data) >= lower:
                return data, q

    candidates.sort(key=lambda item: (item[0], -item[1]))
    _, quality, data = candidates[0]
    return data, quality


def _lecture_resize_width_jpeg(data: bytes) -> tuple[bytes, dict[str, int]]:
    with Image.open(io.BytesIO(data)) as img:
        img = ImageOps.exif_transpose(img)
        original_width, original_height = img.size
        next_height = max(1, round(original_height * (LECTURE_TOOL_IMAGE_WIDTH / original_width)))
        img = img.convert("RGB")
        if img.size != (LECTURE_TOOL_IMAGE_WIDTH, next_height):
            img = img.resize((LECTURE_TOOL_IMAGE_WIDTH, next_height), Image.Resampling.LANCZOS)

        output = _lecture_jpeg_bytes(
            img,
            quality=LECTURE_TOOL_IMAGE_QUALITY,
            subsampling=0,
            optimize=False,
        )
        return output, {
            "originalWidth": original_width,
            "originalHeight": original_height,
            "width": LECTURE_TOOL_IMAGE_WIDTH,
            "height": next_height,
            "quality": LECTURE_TOOL_IMAGE_QUALITY,
            "bytes": len(output),
        }


def _lecture_drive_jpeg_near_target(data: bytes) -> tuple[bytes, dict[str, int]]:
    with Image.open(io.BytesIO(data)) as img:
        img = ImageOps.exif_transpose(img).convert("RGB")
        original_width, original_height = img.size

        target = LECTURE_TOOL_TARGET_IMAGE_BYTES

        candidates: list[tuple[int, bytes, int, int]] = []
        for width, height in LECTURE_TOOL_DRIVE_CONFIG_SIZES:
            work = img if img.size == (width, height) else img.resize((width, height), Image.Resampling.LANCZOS)
            output = _lecture_jpeg_bytes(
                work,
                quality=LECTURE_TOOL_DRIVE_IMAGE_QUALITY,
                subsampling=2,
                optimize=True,
            )
            candidates.append((abs(len(output) - target), output, width, height))

        candidates.sort(key=lambda item: (item[0], item[1]))
        _, output, width, height = candidates[0]
        return output, {
            "originalWidth": original_width,
            "originalHeight": original_height,
            "width": width,
            "height": height,
            "quality": LECTURE_TOOL_DRIVE_IMAGE_QUALITY,
            "bytes": len(output),
            "targetBytes": LECTURE_TOOL_TARGET_IMAGE_BYTES,
        }


def _lecture_contain_jpeg(data: bytes, size: tuple[int, int]) -> bytes:
    with Image.open(io.BytesIO(data)) as img:
        img = ImageOps.exif_transpose(img).convert("RGB")
        img.thumbnail(size, Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", size, (0, 0, 0))
        x = (size[0] - img.width) // 2
        y = (size[1] - img.height) // 2
        canvas.paste(img, (x, y))

        out = io.BytesIO()
        canvas.save(out, format="JPEG", quality=LECTURE_TOOL_IMAGE_QUALITY, subsampling=0)
        return out.getvalue()


def _lecture_write_html_assets(result_dir: Path, media_name: str, count: int) -> None:
    (result_dir / "css").mkdir(parents=True, exist_ok=True)
    (result_dir / "js").mkdir(parents=True, exist_ok=True)

    reset_css = "@charset \"UTF-8\";html,body,div,span,object,iframe,h1,h2,h3,h4,h5,h6,p,blockquote,pre,abbr,address,cite,code,del,dfn,em,img,ins,kbd,q,samp,small,strong,sub,sup,var,b,i,dl,dt,dd,ol,ul,li,fieldset,form,label,legend,table,caption,tbody,tfoot,thead,tr,th,td,article,aside,canvas,details,figcaption,figure,footer,header,hgroup,menu,nav,section,summary,time,mark,audio,video{margin:0;padding:0;border:0;outline:0;font-size:100%;vertical-align:baseline;background:transparent}body{line-height:1}article,aside,details,figcaption,figure,footer,header,hgroup,menu,nav,section{display:block}blockquote,q{quotes:none}blockquote:before,blockquote:after,q:before,q:after{content:none}a{margin:0;padding:0;font-size:100%;vertical-align:baseline;background:transparent}ins{background-color:#ff9;color:#000;text-decoration:none}mark{background-color:#ff9;color:#000;font-style:italic;font-weight:bold}del{text-decoration:line-through}abbr[title],dfn[title]{border-bottom:1px dotted;cursor:help}table{border-collapse:collapse;border-spacing:0}hr{display:block;height:1px;border:0;border-top:1px solid #ccc;margin:1em 0;padding:0}input,select{vertical-align:middle}ul{list-style:none}ol{list-style:none}img{vertical-align:top;font-size:0;line-height:0}body,button,input,select,textarea{font-family:sans-serif}em{font-style:normal}\n"
    index_css = """@charset "UTF-8";
.wrapper {
\twidth: 100vw;
\t\theight: 100vh;
\t\tposition: relative;
\t\toverflow: hidden;
}

.wrapper img {
\twidth: 100%;
}





"""
    script_index_js = "(() => {\r\n\r\n\t\r\n\r\n})();\r\n\r\n\r\n"

    (result_dir / "css" / "reset.css").write_text(
        reset_css,
        encoding="utf-8",
    )
    (result_dir / "css" / "index.css").write_text(
        index_css,
        encoding="utf-8",
    )
    with (result_dir / "js" / "script_index.js").open("w", encoding="utf-8", newline="") as f:
        f.write(script_index_js)

    image_tags = "\n".join(
        f'\t\t\t<img id="link{i}" src="images/{i}.jpg" width="100%">'
        for i in range(1, count + 1)
    )
    html = f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0,minimum-scale=1.0,maximum-scale=1.0,user-scalable=no">
<meta name="format-detection" content="telephone=no">
<title></title>
<link rel="stylesheet" type="text/css" href="css/reset.css" />
<link rel="stylesheet" type="text/css" href="css/index.css">
<link rel="stylesheet" type="text/css" href="shared/css/shared.css">
</head>
<body>
\t<section class="wrapper">
\t\t<div class="inner">
{image_tags}
\t\t</div>
\t</section>
</body>
<script src="shared/vendor/jquery.min.js" ></script>
<script src="shared/vendor/iscroll.js"></script>
<script src="shared/vendor/veeva-library.js"></script>
<script src="js/script_index.js"></script>
<script src="shared/js/shared.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/jquery.qrcode/1.0/jquery.qrcode.min.js"></script>
</html>
"""
    (result_dir / f"{media_name}.html").write_text(html, encoding="utf-8")


def _lecture_zip_directory(source_dir: Path, zip_path: Path, root_name: str) -> None:
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for path in source_dir.rglob("*"):
            if path.is_file():
                arcname = Path(root_name) / path.relative_to(source_dir)
                zf.write(path, arcname.as_posix())


def _lecture_result_url(session_id: str, rel_path: str) -> str:
    encoded = "/".join(quote(part) for part in rel_path.split("/"))
    return f"/lecture-tool/results/{session_id}/{encoded}"


def _lecture_drive_folder_url() -> str:
    if not LECTURE_TOOL_DRIVE_FOLDER_ID:
        return ""
    return f"https://drive.google.com/drive/folders/{quote(LECTURE_TOOL_DRIVE_FOLDER_ID)}"


def _lecture_generation_result(
    session_id: str,
    packages: list[dict[str, Any]],
    small_images: list[dict[str, Any]],
    drive_uploads: list[dict[str, Any]],
    drive_errors: list[dict[str, str]],
    vault_registrations: list[dict[str, Any]],
    vault_errors: list[dict[str, str]],
) -> dict[str, Any]:
    return {
        "sessionId": session_id,
        "packages": packages,
        "smallImages": small_images,
        "driveFolderConfigured": bool(LECTURE_TOOL_DRIVE_FOLDER_ID),
        "driveFolderId": LECTURE_TOOL_DRIVE_FOLDER_ID,
        "driveFolderUrl": _lecture_drive_folder_url(),
        "googleServiceAccountEmail": get_gsa_client_email(),
        "driveUploads": drive_uploads,
        "driveErrors": drive_errors,
        "vaultConfigured": bool(LECTURE_TOOL_VAULT_COMMAND),
        "vaultRegistrations": vault_registrations,
        "vaultErrors": vault_errors,
        "resultFiles": _lecture_result_files(session_id),
    }


def _lecture_generate_packages(
    records: list[dict[str, Any]],
    vault_account: str = "",
    progress=None,
    *,
    session_id: str | None = None,
    cancel_event: threading.Event | None = None,
) -> dict[str, Any]:
    session_id = session_id or new_session_id()
    result_root = _lecture_session_root(session_id) / "result"
    result_root.mkdir(parents=True, exist_ok=True)
    _lecture_progress(progress, "start", "生成処理を開始しました。", sessionId=session_id)
    _lecture_raise_if_cancelled(cancel_event)

    groups: dict[tuple[str, str], list[dict[str, Any]]] = {}
    for record in records:
        row = record["row"]
        presentation_name = row["presentationName"] or "presentation"
        media_file_name = row["mediaFileName"] or Path(record["filename"]).stem
        groups.setdefault((presentation_name, media_file_name), []).append(record)

    packages: list[dict[str, Any]] = []
    small_images = [
        {
            "filename": record["filename"],
            "width": record["imageInfo"]["originalWidth"],
            "height": record["imageInfo"]["originalHeight"],
            "row": record["row"],
        }
        for record in records
        if record["imageInfo"]["originalWidth"] < LECTURE_TOOL_IMAGE_WIDTH
    ]

    drive_uploads: list[dict[str, Any]] = []
    drive_errors: list[dict[str, str]] = []
    vault_registrations: list[dict[str, Any]] = []
    vault_errors: list[dict[str, str]] = []

    for (presentation_name, media_file_name), items in groups.items():
        _lecture_raise_if_cancelled(cancel_event)
        _lecture_progress(
            progress,
            "package",
            f"生成物を作成しています: {media_file_name}",
            presentationName=presentation_name,
            mediaFileName=media_file_name,
        )
        presentation_safe = _lecture_safe_name(presentation_name, "presentation")
        media_safe = _lecture_safe_name(media_file_name, "media")
        result_dir = result_root / presentation_safe / media_safe
        images_dir = result_dir / "images"
        drive_images_dir = result_dir / "drive_images"
        images_dir.mkdir(parents=True, exist_ok=True)

        for idx, item in enumerate(items, start=1):
            _lecture_raise_if_cancelled(cancel_event)
            (images_dir / f"{idx}.jpg").write_bytes(item["resizedBytes"])

        _lecture_raise_if_cancelled(cancel_event)
        _lecture_progress(progress, "zip", f"HTMLとサムネイルを作成しています: {media_file_name}", mediaFileName=media_file_name)
        first_data = items[0]["originalBytes"]
        (result_dir / f"{media_safe}-full.jpg").write_bytes(_lecture_contain_jpeg(first_data, (1024, 768)))
        (result_dir / f"{media_safe}-thumb.jpg").write_bytes(_lecture_contain_jpeg(first_data, (200, 150)))
        _lecture_write_html_assets(result_dir, media_safe, len(items))

        zip_path = result_root / presentation_safe / f"{media_safe}.zip"
        _lecture_raise_if_cancelled(cancel_event)
        _lecture_progress(progress, "zip", f"ZIPを作成しています: {zip_path.name}", mediaFileName=media_file_name)
        _lecture_zip_directory(result_dir, zip_path, media_safe)
        _lecture_progress(progress, "zip", f"ZIP作成完了: {zip_path.name}", mediaFileName=media_file_name)

        row_payloads = [item["row"] for item in items]
        categories = [item["row"].get("category") for item in items if item["row"].get("category")]
        category = next((value for value in categories if "修正" in value), categories[0] if categories else "")
        package = {
            "presentationName": presentation_name,
            "presentationId": next((item["row"].get("presentationId") for item in items if item["row"].get("presentationId")), ""),
            "product": next((item["row"].get("product") for item in items if item["row"].get("product")), ""),
            "category": category,
            "mediaFileName": media_file_name,
            "path": str(result_dir.relative_to(result_root)),
            "absolutePath": str(result_dir),
            "zipPath": str(zip_path.relative_to(result_root)),
            "absoluteZipPath": str(zip_path),
            "zipUrl": _lecture_result_url(session_id, str(zip_path.relative_to(result_root)).replace(os.sep, "/")),
            "count": len(items),
            "rows": row_payloads,
            "driveEnabled": any(item.get("driveEnabled") for item in items),
            "vaultEnabled": any(item.get("vaultEnabled") for item in items),
            "vaultAccount": vault_account,
        }

        packages.append(package)
        if progress:
            progress(
                {
                    "event": "partial",
                    "result": {
                        "ok": True,
                        "imageWidth": LECTURE_TOOL_IMAGE_WIDTH,
                        "imageQuality": LECTURE_TOOL_IMAGE_QUALITY,
                        "targetImageBytes": LECTURE_TOOL_TARGET_IMAGE_BYTES,
                        **_lecture_generation_result(
                            session_id,
                            packages,
                            small_images,
                            drive_uploads,
                            drive_errors,
                            vault_registrations,
                            vault_errors,
                        ),
                    },
                }
            )

        for item in items:
            _lecture_raise_if_cancelled(cancel_event)
            drive_filename = _lecture_sheet_image_filename(item["row"], item["filename"])
            download_url = ""
            try:
                _lecture_raise_if_cancelled(cancel_event)
                _lecture_progress(progress, "drive", f"Drive用画像を作成しています: {drive_filename}", filename=drive_filename)
                drive_bytes, drive_info = _lecture_drive_jpeg_near_target(item["originalBytes"])
                drive_images_dir.mkdir(parents=True, exist_ok=True)
                local_drive_path = drive_images_dir / drive_filename
                local_drive_path.write_bytes(drive_bytes)
                download_url = _lecture_result_url(session_id, str(local_drive_path.relative_to(result_root)).replace(os.sep, "/"))
                local_drive_item = {
                    "filename": drive_filename,
                    "rowNumber": item["row"].get("rowNumber"),
                    "mediaFileName": media_file_name,
                    "imageInfo": drive_info,
                    "downloadUrl": download_url,
                    "uploadRequested": bool(item.get("driveEnabled")),
                    "uploaded": False,
                }

                if item.get("driveEnabled") and LECTURE_TOOL_DRIVE_FOLDER_ID:
                    _lecture_raise_if_cancelled(cancel_event)
                    _lecture_progress(progress, "drive", f"Google Driveへアップロードしています: {drive_filename}", filename=drive_filename)
                    drive_uploads.append(
                        {
                            **local_drive_item,
                            "uploaded": True,
                            **_lecture_drive_upload_image(drive_filename, drive_bytes),
                        }
                    )
                    _lecture_progress(progress, "drive", f"Google Driveアップロード完了: {drive_filename}", filename=drive_filename)
                else:
                    local_drive_item["uploadSkipped"] = True
                    if item.get("driveEnabled") and not LECTURE_TOOL_DRIVE_FOLDER_ID:
                        local_drive_item["uploadSkipReason"] = "Google Driveフォルダが未設定です。"
                    elif not item.get("driveEnabled"):
                        local_drive_item["uploadSkipReason"] = "Drive格納のチェックが外れています。Drive格納は未実行です。"
                    drive_uploads.append(local_drive_item)
                    _lecture_progress(progress, "drive", f"Drive用画像作成完了: {drive_filename}", filename=drive_filename)
            except Exception as exc:
                drive_errors.append(
                    {
                        "filename": drive_filename,
                        "mediaFileName": media_file_name,
                        "downloadUrl": download_url,
                        "error": str(exc),
                    }
                )
                _lecture_progress(progress, "drive", f"Drive用画像処理失敗: {drive_filename}", filename=drive_filename, error=str(exc))

    if progress:
        progress(
            {
                "event": "partial",
                "result": {
                    "ok": True,
                    "imageWidth": LECTURE_TOOL_IMAGE_WIDTH,
                    "imageQuality": LECTURE_TOOL_IMAGE_QUALITY,
                    "targetImageBytes": LECTURE_TOOL_TARGET_IMAGE_BYTES,
                    **_lecture_generation_result(
                        session_id,
                        packages,
                        small_images,
                        drive_uploads,
                        drive_errors,
                        vault_registrations,
                        vault_errors,
                    ),
                },
            }
        )
        _lecture_progress(progress, "drive", "Drive処理が全件完了しました。")

    vault_packages = [package for package in packages if package["vaultEnabled"]]
    if vault_packages:
        _lecture_raise_if_cancelled(cancel_event, "Vault登録の開始前に中断しました。")
        vault_slot_acquired = False
        vault_account_lock_acquired = False
        vault_account_lock = _lecture_get_vault_account_lock(vault_account)
        try:
            _lecture_progress(
                progress,
                "vault",
                f"Vault登録の並列実行枠を確保しています。（最大 {LECTURE_TOOL_VAULT_MAX_PARALLEL} 件）",
            )
            while not vault_slot_acquired:
                _lecture_raise_if_cancelled(cancel_event, "Vault登録の待機中に中断しました。")
                vault_slot_acquired = LECTURE_TOOL_VAULT_SEMAPHORE.acquire(timeout=1)
            _lecture_progress(progress, "vault", "Vault登録の並列実行枠を確保しました。")

            _lecture_progress(progress, "vault", f"Vaultアカウントの利用枠を確保しています: {vault_account}")
            while not vault_account_lock_acquired:
                _lecture_raise_if_cancelled(cancel_event, "Vault登録の待機中に中断しました。")
                vault_account_lock_acquired = vault_account_lock.acquire(timeout=1)
            _lecture_progress(progress, "vault", f"Vaultアカウントの利用枠を確保しました: {vault_account}")

            vault_body = _lecture_vault_register_packages(
                vault_packages,
                vault_account,
                progress=progress,
                session_id=session_id,
                cancel_event=cancel_event,
            )
            vault_results = vault_body.get("results") if isinstance(vault_body, dict) else None
            if not isinstance(vault_results, list):
                vault_results = [vault_body]
            for idx, vault_result in enumerate(vault_results):
                if not isinstance(vault_result, dict):
                    vault_result = {"output": vault_result}
                package = vault_packages[idx] if idx < len(vault_packages) else {}
                common = {
                    "presentationName": package.get("presentationName", ""),
                    "presentationId": package.get("presentationId", ""),
                    "category": package.get("category", ""),
                    "mediaFileName": package.get("mediaFileName", ""),
                    "zipPath": package.get("zipPath", ""),
                    "vaultAccount": vault_account,
                }
                if vault_result.get("error"):
                    vault_errors.append({**common, "error": str(vault_result.get("error"))})
                    _lecture_progress(progress, "vault", f"Vault登録失敗: {common['mediaFileName']}", mediaFileName=common["mediaFileName"], error=str(vault_result.get("error")))
                else:
                    vault_registrations.append({**common, **vault_result})
                    _lecture_progress(progress, "vault", f"Vault登録完了: {common['mediaFileName']}", mediaFileName=common["mediaFileName"])
        except LectureToolCancelled:
            raise
        except Exception as exc:
            for package in vault_packages:
                vault_errors.append(
                    {
                        "presentationName": package.get("presentationName", ""),
                        "presentationId": package.get("presentationId", ""),
                        "category": package.get("category", ""),
                        "mediaFileName": package.get("mediaFileName", ""),
                        "zipPath": package.get("zipPath", ""),
                        "vaultAccount": vault_account,
                        "error": str(exc),
                    }
                )
            _lecture_progress(progress, "vault", "Vault登録処理が停止しました。", error=str(exc))
        finally:
            if vault_account_lock_acquired:
                vault_account_lock.release()
            if vault_slot_acquired:
                LECTURE_TOOL_VAULT_SEMAPHORE.release()

    return _lecture_generation_result(
        session_id,
        packages,
        small_images,
        drive_uploads,
        drive_errors,
        vault_registrations,
        vault_errors,
    )


@app.get("/lecture-tool/status")
async def lecture_tool_status():
    try:
        sheet = _lecture_fetch_sheet_rows()
        return JSONResponse(
            {
                "ok": True,
                "mode": "spreadsheet_upload",
                "spreadsheetKey": LECTURE_TOOL_SPREADSHEET_KEY,
                "spreadsheetTitle": sheet["spreadsheetTitle"],
                "spreadsheetUrl": sheet["spreadsheetUrl"],
                "sheetTitle": sheet["sheetTitle"],
                "sheetGid": sheet["sheetGid"],
                "rowCount": len(sheet["rows"]),
                "imageWidth": LECTURE_TOOL_IMAGE_WIDTH,
                "imageQuality": LECTURE_TOOL_IMAGE_QUALITY,
                "driveImageQuality": LECTURE_TOOL_DRIVE_IMAGE_QUALITY,
                "driveConfigSizes": [{"width": width, "height": height} for width, height in LECTURE_TOOL_DRIVE_CONFIG_SIZES],
                "targetImageBytes": LECTURE_TOOL_TARGET_IMAGE_BYTES,
                "driveFolderConfigured": bool(LECTURE_TOOL_DRIVE_FOLDER_ID),
                "driveFolderId": LECTURE_TOOL_DRIVE_FOLDER_ID,
                "googleServiceAccountEmail": get_gsa_client_email(),
                "vaultConfigured": bool(LECTURE_TOOL_VAULT_COMMAND),
                "vaultAccounts": LECTURE_TOOL_VAULT_ACCOUNTS,
                "vaultMaxParallel": LECTURE_TOOL_VAULT_MAX_PARALLEL,
                "checkedAt": datetime.now(timezone.utc).isoformat(),
            }
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"sheet status failed: {exc}") from exc


@app.get("/lecture-tool/sheet-rows")
async def lecture_tool_sheet_rows():
    try:
        sheet = _lecture_fetch_sheet_rows()
        return JSONResponse({"ok": True, "vaultAccounts": LECTURE_TOOL_VAULT_ACCOUNTS, **sheet})
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"spreadsheet fetch failed: {exc}") from exc


async def _lecture_records_from_uploads(
    files: List[UploadFile] = File(...),
    rowIds: List[str] = Form(...),
    driveEnabled: Optional[List[str]] = Form(None),
    vaultEnabled: Optional[List[str]] = Form(None),
) -> list[dict[str, Any]]:
    if not files:
        raise HTTPException(status_code=400, detail="画像をアップロードしてください。")
    if len(files) != len(rowIds):
        raise HTTPException(status_code=400, detail="画像数と選択行数が一致しません。")
    sheet = _lecture_fetch_sheet_rows()
    rows_by_id = {row["id"]: row for row in sheet["rows"]}

    records: list[dict[str, Any]] = []
    for idx, upload in enumerate(files):
        row_id = normalize_space(rowIds[idx])
        row = rows_by_id.get(row_id)
        if not row:
            raise HTTPException(
                status_code=400,
                detail=f"{upload.filename}: スプレッドシート行が選択されていません。行番号を指定してください。",
            )

        data = await upload.read()
        try:
            resized, image_info = _lecture_resize_width_jpeg(data)
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"{upload.filename}: 画像として読み込めません。") from exc

        records.append(
            {
                "filename": upload.filename or f"image_{idx + 1}",
                "row": row,
                "originalBytes": data,
                "resizedBytes": resized,
                "imageInfo": image_info,
                "driveEnabled": _lecture_bool(driveEnabled[idx] if driveEnabled and idx < len(driveEnabled) else None, True),
                "vaultEnabled": _lecture_bool(vaultEnabled[idx] if vaultEnabled and idx < len(vaultEnabled) else None, True),
            }
        )
    return records


@app.post("/lecture-tool/generate")
async def lecture_tool_generate(
    files: List[UploadFile] = File(...),
    rowIds: List[str] = Form(...),
    driveEnabled: Optional[List[str]] = Form(None),
    vaultEnabled: Optional[List[str]] = Form(None),
    vaultAccount: str = Form(""),
):
    records = await _lecture_records_from_uploads(files, rowIds, driveEnabled, vaultEnabled)
    small_images = [
        {
            "filename": record["filename"],
            "width": record["imageInfo"]["originalWidth"],
            "height": record["imageInfo"]["originalHeight"],
            "row": record["row"],
        }
        for record in records
        if record["imageInfo"]["originalWidth"] < LECTURE_TOOL_IMAGE_WIDTH
    ]
    vault_account = normalize_space(vaultAccount)
    if any(record.get("vaultEnabled") for record in records) and not vault_account:
        raise HTTPException(status_code=400, detail="Vault登録対象がある場合はVaultアカウントを選択してください。")
    result = _lecture_generate_packages(records, vault_account)
    return JSONResponse(
        {
            "ok": True,
            "imageWidth": LECTURE_TOOL_IMAGE_WIDTH,
            "imageQuality": LECTURE_TOOL_IMAGE_QUALITY,
            "targetImageBytes": LECTURE_TOOL_TARGET_IMAGE_BYTES,
            **result,
        }
    )


@app.post("/lecture-tool/generate-stream")
async def lecture_tool_generate_stream(
    request: Request,
    files: List[UploadFile] = File(...),
    rowIds: List[str] = Form(...),
    driveEnabled: Optional[List[str]] = Form(None),
    vaultEnabled: Optional[List[str]] = Form(None),
    vaultAccount: str = Form(""),
):
    records = await _lecture_records_from_uploads(files, rowIds, driveEnabled, vaultEnabled)
    vault_account = normalize_space(vaultAccount)
    if any(record.get("vaultEnabled") for record in records) and not vault_account:
        raise HTTPException(status_code=400, detail="Vault登録対象がある場合はVaultアカウントを選択してください。")

    session_id = new_session_id()
    controller = _lecture_register_cancel_controller(session_id)

    def worker(progress_queue: queue.Queue):
        def progress(payload: dict[str, Any]):
            progress_queue.put({"event": "progress", **payload})

        try:
            result = _lecture_generate_packages(
                records,
                vault_account,
                progress=progress,
                session_id=session_id,
                cancel_event=controller["event"],
            )
            progress_queue.put(
                {
                    "event": "done",
                    "ok": True,
                    "imageWidth": LECTURE_TOOL_IMAGE_WIDTH,
                    "imageQuality": LECTURE_TOOL_IMAGE_QUALITY,
                    "targetImageBytes": LECTURE_TOOL_TARGET_IMAGE_BYTES,
                    **result,
                }
            )
        except LectureToolCancelled as exc:
            progress_queue.put({"event": "cancelled", "sessionId": session_id, "message": str(exc)})
        except Exception as exc:
            progress_queue.put({"event": "error", "message": str(exc)})
        finally:
            _lecture_unregister_cancel_controller(session_id)
            progress_queue.put(None)

    async def event_stream():
        progress_queue: queue.Queue = queue.Queue()
        thread = threading.Thread(target=worker, args=(progress_queue,), daemon=True)
        thread.start()
        completed = False
        try:
            while True:
                try:
                    item = await asyncio.to_thread(progress_queue.get, True, 0.5)
                except queue.Empty:
                    if await request.is_disconnected():
                        _lecture_cancel_session(session_id)
                        break
                    continue
                if item is None:
                    completed = True
                    break
                event_name = item.pop("event", "progress")
                yield _sse(event_name, item)
        finally:
            if not completed:
                _lecture_cancel_session(session_id)

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@app.post("/lecture-tool/cancel/{session_id}")
async def lecture_tool_cancel(session_id: str):
    try:
        _lecture_session_root(session_id)
    except HTTPException:
        raise
    if not _lecture_cancel_session(session_id):
        raise HTTPException(status_code=404, detail="対象の処理が見つからないか、すでに完了しています。")
    return JSONResponse({"ok": True, "sessionId": session_id, "message": "中断リクエストを送信しました。"})


@app.get("/lecture-tool/results/{session_id}/{rel_path:path}")
async def lecture_tool_result_file(session_id: str, rel_path: str, request: Request):
    result_dir = (_lecture_session_root(session_id) / "result").resolve()
    target = (result_dir / rel_path).resolve()

    try:
        target.relative_to(result_dir)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="invalid result path") from exc

    if target == result_dir:
        raise HTTPException(status_code=400, detail="invalid result path")
    if not target.is_file():
        raise HTTPException(status_code=404, detail="result file not found")

    media_type = mimetypes.guess_type(target.name)[0] or "application/octet-stream"
    force_download = str(request.query_params.get("download") or "").lower() in {"1", "true", "yes"}
    if force_download:
        return FileResponse(
            target,
            media_type=media_type,
            filename=target.name,
            content_disposition_type="attachment",
        )
    return FileResponse(target, media_type=media_type)

@app.post("/upload/simple/stream")
async def upload_simple_stream(
    files: List[UploadFile] = File(...),
    regions: List[str] = Form(...),
    units: List[str] = Form(None),
):
    """スプレッドシートを読み込まない簡易アップロード。region(必須) + unit(任意) のみ。"""
    session_id = new_session_id()

    if not files:
        raise HTTPException(400, "files is empty")
    if len(regions) != len(files):
        raise HTTPException(400, f"regions length mismatch: {len(regions)} != {len(files)}")
    if units is None:
        units = [""] * len(files)
    if len(units) != len(files):
        raise HTTPException(400, f"units length mismatch: {len(units)} != {len(files)}")

    VALID_REGIONS = {"VP", "PH", "ONC"}
    total = len(files)

    session_dir = Path("jobs") / f"session_{session_id}"
    session_dir.mkdir(parents=True, exist_ok=True)

    buffered: List[Dict[str, Any]] = []
    for i, f in enumerate(files):
        filename = f.filename or f"file_{i}"
        suffix = Path(filename).suffix.lower()
        region = (regions[i] or "").strip().upper()
        unit = (units[i] or "").strip()

        item = {"index": i, "filename": filename, "suffix": suffix, "region": region, "unit": unit}

        if suffix not in [".pptx", ".pdf"]:
            item["precheck"] = {"ok": False, "error": "not_supported_file"}
            buffered.append(item)
            continue
        if region not in VALID_REGIONS:
            item["precheck"] = {"ok": False, "error": f"invalid_region: {region}（VP/PH/ONC のいずれかを選択してください）"}
            buffered.append(item)
            continue

        try:
            data = await f.read()
            in_path = session_dir / f"{i}_{uuid.uuid4().hex}{suffix}"
            in_path.write_bytes(data)
            item["precheck"] = {"ok": True}
            item["in_path"] = str(in_path)
        except Exception as e:
            item["precheck"] = {"ok": False, "error": f"upload_read_failed: {e}"}
        finally:
            try:
                await f.close()
            except Exception:
                pass

        buffered.append(item)

    async def gen():
        yield _sse("start", {"sessionId": session_id, "total": total})

        try:
            yield _sse("phase", {"phase": "processing", "message": "生成を開始します…（スプレッドシート不使用）"})
            out: List[Dict[str, Any]] = []

            for it in buffered:
                i = it["index"]
                filename = it["filename"]
                region = it["region"]
                unit = it["unit"]

                yield _sse("item_start", {"index": i, "filename": filename})

                if not it["precheck"]["ok"]:
                    err = it["precheck"]["error"]
                    out.append({"filename": filename, "ok": False, "error": err})
                    yield _sse("item_done", {"index": i, "filename": filename, "ok": False, "error": err})
                    continue

                in_path = Path(it["in_path"])
                job_id = uuid.uuid4().hex
                p = job_paths(job_id)

                try:
                    import time as _time_mod
                    _t0 = _time_mod.monotonic()

                    payload = await pptx_to_json_vm_hint(
                        in_path,
                        [],  # VM rows なし
                        debug_blocks_path=p.get("debug_blocks"),
                    )
                    _t1 = _time_mod.monotonic()
                    print(f"[TIMING][simple] pptx_to_json_vm_hint: {_t1 - _t0:.2f}s")

                    payload = normalize_for_render(payload)
                    payload = post_format_design_initial(payload)
                    _t2 = _time_mod.monotonic()
                    print(f"[TIMING][simple] normalize+post_format: {_t2 - _t1:.2f}s")

                    payload = await apply_precise_typeset_initial(payload)
                    _t3 = _time_mod.monotonic()
                    print(f"[TIMING][simple] apply_precise_typeset: {_t3 - _t2:.2f}s")

                    payload = ensure_display_fields(payload)

                    _blocks_for_overlay = []
                    try:
                        _dbp = p.get("debug_blocks")
                        if _dbp and _dbp.exists():
                            _raw = json.loads(_dbp.read_text(encoding="utf-8"))
                            _blocks_for_overlay = _raw
                    except Exception:
                        pass
                    payload = apply_correct_answer_overlay(payload, _blocks_for_overlay)
                    payload = ensure_display_fields(payload)
                    dump_titles("after apply_correct_answer_overlay", payload)
                    _t4 = _time_mod.monotonic()
                    print(f"[TIMING][simple] overlay+display: {_t4 - _t3:.2f}s")

                    payload.region = region
                    payload.unit = unit
                    payload.event_id = ""

                    payload.talks = sorted(
                        payload.talks or [],
                        key=lambda x: (
                            getattr(x, "program_index", 10**9),
                            getattr(x, "_talk_index", 10**9),
                            _time_start_minutes(getattr(x, "time", "")),
                        )
                    )

                    payload_dict = (
                        payload.model_dump(exclude_none=True)
                        if hasattr(payload, "model_dump")
                        else json.loads(payload.json(ensure_ascii=False))
                    )

                    jpg_bytes, debug_html = await render_png_bytes(payload)
                    _t5 = _time_mod.monotonic()
                    print(f"[TIMING][simple] render_png_bytes: {_t5 - _t4:.2f}s")

                    # ローカルにJSON・JPGを保存（一覧/編集画面の高速化）
                    try:
                        p["json"].write_text(json.dumps(payload_dict, ensure_ascii=False), encoding="utf-8")
                        p["jpg"].write_bytes(jpg_bytes)
                    except Exception:
                        pass

                    upsert_job_ok(job_id, filename, payload, session_id, "")
                    _t6 = _time_mod.monotonic()
                    print(f"[TIMING][simple] local_save+upsert: {_t6 - _t5:.2f}s")
                    print(f"[TIMING][simple] TOTAL: {_t6 - _t0:.2f}s")

                    # Storage アップロードはバックグラウンドで実行（ローカル保存済みなので遅延OK）
                    asyncio.create_task(upload_all_assets_async(
                        job_id,
                        payload_dict=payload_dict,
                        jpg_bytes=jpg_bytes,
                        debug_html=debug_html,
                        debug_blocks_path=p.get("debug_blocks"),
                    ))

                    out.append({"filename": filename, "jobId": job_id, "ok": True})
                    yield _sse("item_done", {"index": i, "filename": filename, "ok": True, "jobId": job_id})

                except Exception as e:
                    tb = traceback.format_exc()
                    print("[upload/simple error]", filename, job_id)
                    print(tb)

                    out.append({"filename": filename, "jobId": job_id, "ok": False, "error": str(e)})
                    yield _sse("item_done", {"index": i, "filename": filename, "ok": False, "jobId": job_id, "error": str(e)})

            ok_count = sum(1 for r in out if r.get("ok"))
            yield _sse("done", {"sessionId": session_id, "count": ok_count, "results": out})

        except Exception as e:
            tb = traceback.format_exc()
            print(tb)
            yield _sse("fatal", {"message": str(e)})
        finally:
            try:
                shutil.rmtree(session_dir, ignore_errors=True)
            except Exception:
                pass

    return StreamingResponse(
        gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/upload/batch/stream")
async def upload_batch_stream(
    files: List[UploadFile] = File(...),
    eventIds: List[str] = Form(...),
):
    session_id = new_session_id()

    if not files:
        raise HTTPException(400, "files is empty")
    if len(eventIds) != len(files):
        raise HTTPException(400, f"eventIds length mismatch: {len(eventIds)} != {len(files)}")

    total = len(files)

    # ✅ ここで UploadFile を全部 “生きてるうちに” 退避する
    buffered: List[Dict[str, Any]] = []
    # sessionごとの temp dir（好きな場所でOK）
    session_dir = Path("jobs") / f"session_{session_id}"
    session_dir.mkdir(parents=True, exist_ok=True)

    for i, f in enumerate(files):
        filename = f.filename or f"file_{i}"
        suffix = Path(filename).suffix.lower()

        item = {"index": i, "filename": filename, "suffix": suffix, "eventId": (eventIds[i] or "").strip()}

        if suffix not in [".pptx", ".pdf"]:
            item["precheck"] = {"ok": False, "error": "not_supported_file"}
            buffered.append(item)
            continue
        if not item["eventId"]:
            item["precheck"] = {"ok": False, "error": "event_id_required"}
            buffered.append(item)
            continue

        try:
            data = await f.read()              # ✅ return前に読む
            in_path = session_dir / f"{i}_{uuid.uuid4().hex}{suffix}"
            in_path.write_bytes(data)
            item["precheck"] = {"ok": True}
            item["in_path"] = str(in_path)
        except Exception as e:
            item["precheck"] = {"ok": False, "error": f"upload_read_failed: {e}"}
        finally:
            # 任意：明示的に閉じておく（なくてもOK）
            try:
                await f.close()
            except Exception:
                pass

        buffered.append(item)

    async def gen():
        yield _sse("start", {"sessionId": session_id, "total": total})

        try:
            yield _sse("phase", {"phase": "sheet_open", "message": "スプレッドシート接続中…"})
            # ---- Spreadsheet open (once) ----
            scope = [
                "https://www.googleapis.com/auth/spreadsheets",
                "https://www.googleapis.com/auth/drive",
            ]
            credentials = get_gsa_credentials(scope)
            gc = gspread.authorize(credentials)

            SPREADSHEET_KEY = "1hiV0Ve2cnYyrPkBuZcZIcLWeAnJ-ucNiB0P4owZpXug"
            workbook = gc.open_by_key(SPREADSHEET_KEY)

            PRESENCE_SHEETS = ["VM(GWET)", "VM(例外)", "VM(本社)"]
            VM_SHEET = "演題演者（VM）"
            PRESENCE_HEADER_ROW = 2
            VM_HEADER_ROW = 1
            PRESENCE_ID_COL = "システムID"

            yield _sse("phase", {"phase": "precheck", "message": "事前チェック中…"})

            valid_event_ids = []
            for it in buffered:
                if it["precheck"]["ok"]:
                    valid_event_ids.append(it["eventId"])

            ws_map = _retry_gspread(lambda: {ws.title: ws for ws in workbook.worksheets()})

            yield _sse("phase", {"phase": "batch_fetch", "message": "VM/Presence一括取得中…"})
            presence_rows_by_event, vm_rows_by_event, _ = batch_fetch_system_and_vm_rows(
                workbook,
                ws_map=ws_map,
                event_ids=valid_event_ids,
                presence_sheets=PRESENCE_SHEETS,
                presence_header_row=PRESENCE_HEADER_ROW,
                presence_id_col=PRESENCE_ID_COL,
                vm_sheet=VM_SHEET,
                vm_header_row=VM_HEADER_ROW,
                vm_id_col_candidates=["講演会ID"],
                col_end="Z",
            )

            # （あなたの _parse_ymd / presence_rows_by_file_index / vm_rows_by_file_index のロジックは
            #    buffered を元に組み直すのが一番安全。ここでは “生成部分” の直しだけ見せます）

            yield _sse("phase", {"phase": "processing", "message": "生成を開始します…"})
            out: List[Dict[str, Any]] = []

            for it in buffered:
                i = it["index"]
                filename = it["filename"]
                event_id = it["eventId"]

                yield _sse("item_start", {"index": i, "filename": filename, "eventId": event_id})

                if not it["precheck"]["ok"]:
                    err = it["precheck"]["error"]
                    out.append({"filename": filename, "ok": False, "error": err})
                    yield _sse("item_done", {"index": i, "filename": filename, "ok": False, "error": err})
                    continue

                # ✅ もう UploadFile は触らない。退避したパスだけ使う
                in_path = Path(it["in_path"])

                # presence/vm は event_id から引く（ここはあなたの既存ロジックに合わせてOK）
                presence_rows = presence_rows_by_event.get(event_id, []) or []
                if not presence_rows:
                    out.append({"filename": filename, "ok": False, "error": "event_id_not_found"})
                    yield _sse("item_done", {"index": i, "filename": filename, "ok": False, "error": "event_id_not_found"})
                    continue
                vm_rows = vm_rows_by_event.get(event_id, []) or []

                job_id = uuid.uuid4().hex
                p = job_paths(job_id)

                try:
                    payload = await pptx_to_json_vm_hint(
                        in_path,
                        vm_rows,
                        debug_blocks_path=p.get("debug_blocks"),
                    )
                    payload = normalize_for_render(payload)
                    payload = post_format_design_initial(payload)
                    payload = await apply_precise_typeset_initial(payload)
                    payload = ensure_display_fields(payload)

                    # 正解DBで後処理の上書きを復元（typeset後に適用して改行位置を保持）
                    _blocks_for_overlay = []
                    try:
                        _dbp = p.get("debug_blocks")
                        if _dbp and _dbp.exists():
                            _raw = json.loads(_dbp.read_text(encoding="utf-8"))
                            _blocks_for_overlay = _raw  # list[dict] – overlay は dict 対応済み
                    except Exception:
                        pass
                    payload = apply_correct_answer_overlay(payload, _blocks_for_overlay, vm_rows=vm_rows)
                    payload = ensure_display_fields(payload)
                    dump_titles("after apply_correct_answer_overlay", payload)

                    payload.region = presence_rows[0].get("VP/PH/ONC", "")
                    payload.unit = presence_rows[0].get("取得単位：フラグメントデザインへの内容記載", "")
                    payload.event_id = event_id

                    payload.talks = sorted(
                        payload.talks or [],
                        key=lambda x: (
                            getattr(x, "program_index", 10**9),
                            getattr(x, "_talk_index", 10**9),
                            _time_start_minutes(getattr(x, "time", "")),
                        )
                    )

                    payload_dict = (
                        payload.model_dump(exclude_none=True)
                        if hasattr(payload, "model_dump")
                        else json.loads(payload.json(ensure_ascii=False))
                    )

                    jpg_bytes, debug_html = await render_png_bytes(payload)

                    # ローカルにJSON・JPGを保存（一覧/編集画面の高速化）
                    try:
                        p["json"].write_text(json.dumps(payload_dict, ensure_ascii=False), encoding="utf-8")
                        p["jpg"].write_bytes(jpg_bytes)
                    except Exception:
                        pass

                    upsert_job_ok(job_id, filename, payload, session_id, event_id)

                    # Storage アップロードはバックグラウンドで実行（ローカル保存済みなので遅延OK）
                    asyncio.create_task(upload_all_assets_async(
                        job_id,
                        payload_dict=payload_dict,
                        jpg_bytes=jpg_bytes,
                        debug_html=debug_html,
                        debug_blocks_path=p.get("debug_blocks"),
                    ))

                    out.append({
                        "filename": filename,
                        "jobId": job_id,
                        "ok": True,
                    })
                    yield _sse("item_done", {
                        "index": i,
                        "filename": filename,
                        "ok": True,
                        "jobId": job_id,
                    })

                except Exception as e:
                    tb = traceback.format_exc()
                    print("[upload/batch error]", filename, job_id)
                    print(tb)

                    out.append({
                        "filename": filename,
                        "jobId": job_id,
                        "ok": False,
                        "error": str(e),
                    })
                    yield _sse("item_done", {
                        "index": i,
                        "filename": filename,
                        "ok": False,
                        "jobId": job_id,
                        "error": str(e),
                    })

            ok_count = sum(1 for r in out if r.get("ok"))
            yield _sse("done", {"sessionId": session_id, "count": ok_count, "results": out})

        except Exception as e:
            tb = traceback.format_exc()
            print(tb)
            yield _sse("fatal", {"message": str(e)})
        finally:
            try:
                shutil.rmtree(session_dir, ignore_errors=True)
            except Exception:
                pass

    return StreamingResponse(
        gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/render")
async def render(req: RenderReq, background_tasks: BackgroundTasks):
    with db_connect() as con:
        row = con.execute(
            "SELECT locked, filename, session_id, event_id FROM jobs WHERE job_id=%s",
            (req.jobId,),
        ).fetchone()

    if not row:
        raise HTTPException(404, "job not found")

    if bool(row.get("locked")):
        raise HTTPException(400, "This job is locked.")

    filename = row.get("filename") or ""
    session_id = row.get("session_id") or ""
    event_id = row.get("event_id") or ""

    payload = req.design
    # payloadにevent_idが含まれていればそちらを優先（編集画面から変更可能）
    if (getattr(payload, "event_id", "") or "").strip():
        event_id = payload.event_id.strip()
    # /render はエディタからの手動保存でのみ呼ばれるので常に manual_override=True
    payload.manual_override = True

    payload_dict = (
        payload.model_dump(exclude_none=True)
        if hasattr(payload, "model_dump")
        else json.loads(payload.json(ensure_ascii=False))
    )

    jpg_bytes, debug_html = await render_png_bytes(payload)

    # ローカルにもJSON・JPGを保存（レスポンス前に完了 — 高速）
    try:
        rp = job_paths(req.jobId)
        rp["json"].write_text(json.dumps(payload_dict, ensure_ascii=False), encoding="utf-8")
        rp["jpg"].write_bytes(jpg_bytes)
    except Exception:
        pass

    # invalidate signed URL cache for this preview (image changed)
    preview_cache_key = f"{req.jobId}/preview.jpg"
    _signed_url_cache.pop(preview_cache_key, None)

    # Storage アップロード + DB更新はバックグラウンドで実行
    # （レスポンスには previewDataUrl が含まれるので待つ必要なし）
    _job_id_bg = req.jobId
    _payload_dict_bg = payload_dict
    _jpg_bytes_bg = jpg_bytes
    _debug_html_bg = debug_html
    _filename_bg = filename
    _payload_bg = payload
    _session_id_bg = session_id
    _event_id_bg = event_id

    async def _upload_and_upsert_bg():
        try:
            await upload_all_assets_async(
                _job_id_bg,
                payload_dict=_payload_dict_bg,
                jpg_bytes=_jpg_bytes_bg,
                debug_html=_debug_html_bg,
            )
        except Exception as e:
            print("[storage upload error][/render bg]", _job_id_bg, e)
        try:
            upsert_job_ok(_job_id_bg, _filename_bg, _payload_bg, _session_id_bg, _event_id_bg)
        except Exception as e:
            print("[upsert error][/render bg]", _job_id_bg, e)

    def _upload_and_upsert_bg_sync():
        """BackgroundTasks はワーカースレッドで実行されるため、
        asyncio の event loop が存在しない。新しいループを作って実行する。"""
        loop = asyncio.new_event_loop()
        try:
            loop.run_until_complete(_upload_and_upsert_bg())
        finally:
            loop.close()

    background_tasks.add_task(_upload_and_upsert_bg_sync)

    # manual_override=True → 正解DBに自動登録（バックグラウンドで実行）
    _job_id_for_bg = req.jobId
    _payload_for_bg = payload
    _payload_dict_for_bg = payload_dict

    def _save_correct_answer_bg():
        try:
            # blocksをローカル → Storage の順で取得
            blocks_text = ""
            _blocks_json_for_save = []

            # ローカルファイルを先に試す
            _local_blocks = job_paths(_job_id_for_bg).get("debug_blocks")
            if _local_blocks and _local_blocks.exists():
                try:
                    blocks_data = json.loads(_local_blocks.read_text(encoding="utf-8"))
                    _blocks_json_for_save = blocks_data if isinstance(blocks_data, list) else []
                    blocks_text = " ".join(
                        (b.get("text", "") if isinstance(b, dict) else "")
                        for b in _blocks_json_for_save
                    )
                except Exception:
                    pass

            if not blocks_text.strip():
                # ローカルになければ Storage fallback
                sp = storage_paths(_job_id_for_bg)
                blocks_url = _authenticated_storage_url(sp['debug_blocks'])
                headers_s = _storage_auth_headers()
                resp = requests.get(blocks_url, headers=headers_s, timeout=30)
                if resp.status_code == 200:
                    blocks_data = resp.json()
                    _blocks_json_for_save = blocks_data if isinstance(blocks_data, list) else []
                    blocks_text = " ".join(
                        (b.get("text", "") if isinstance(b, dict) else "")
                        for b in _blocks_json_for_save
                    )
                else:
                    print(f"[correct-answer] blocks fetch failed: {resp.status_code} for {_job_id_for_bg}")

            if not blocks_text.strip():
                parts = []
                etl = getattr(_payload_for_bg, "event_title_lines", None) or []
                if etl:
                    parts.extend(etl)
                elif getattr(_payload_for_bg, "event_title", ""):
                    parts.append(_payload_for_bg.event_title)
                if getattr(_payload_for_bg, "organizer", ""):
                    parts.append(_payload_for_bg.organizer)
                if getattr(_payload_for_bg, "datetime", ""):
                    parts.append(_payload_for_bg.datetime)
                for t in (_payload_for_bg.talks or []):
                    if getattr(t, "title", ""):
                        parts.append(t.title)
                    if getattr(t, "speaker", ""):
                        parts.append(t.speaker)
                    if getattr(t, "affiliation", ""):
                        parts.append(t.affiliation)
                if getattr(_payload_for_bg, "chair", None):
                    if getattr(_payload_for_bg.chair, "name", ""):
                        parts.append(_payload_for_bg.chair.name)
                blocks_text = " ".join(parts)
                print(f"[correct-answer] using payload fallback for blocks_text ({len(blocks_text)} chars)")

            etl = getattr(_payload_for_bg, "event_title_lines", None) or []
            full_event_title = " ".join(etl) if etl else (getattr(_payload_for_bg, "event_title", "") or "")

            save_correct_answer(
                blocks_text=blocks_text,
                correct_json=_payload_dict_for_bg,
                event_title=full_event_title,
                job_id=_job_id_for_bg,
                blocks_json=_blocks_json_for_save or None,
            )
            print(f"[correct-answer] saved for {_job_id_for_bg} (blocks_text={len(blocks_text)} chars)")
        except Exception as e:
            print(f"[correct-answer][auto-register] {_job_id_for_bg}: {e}")

    background_tasks.add_task(_save_correct_answer_bg)

    return JSONResponse({
        "jobId": req.jobId,
        "json": payload_dict,
        "warnings": getattr(payload, "warnings", None),
        "previewUrl": f"/preview/{req.jobId}.jpg",
        "downloadUrl": f"/download/{req.jobId}.jpg",
        "previewDataUrl": f"data:image/jpeg;base64,{base64.b64encode(jpg_bytes).decode('ascii')}",
    })

def _parse_date_start(s: str) -> Optional[datetime]:
    s = (s or "").strip()
    if not s:
        return None
    # "YYYY-MM-DD" をUTC 00:00:00として扱う（必要ならJSTに変更）
    return datetime.strptime(s, "%Y-%m-%d").replace(tzinfo=timezone.utc)

def _parse_date_end(s: str) -> Optional[datetime]:
    s = (s or "").strip()
    if not s:
        return None
    # inclusive end にしたいなら 23:59:59.999999
    dt = datetime.strptime(s, "%Y-%m-%d").replace(tzinfo=timezone.utc)
    return dt.replace(hour=23, minute=59, second=59, microsecond=999999)

@app.post("/jobs/restore/batch")
async def restore_from_json_batch(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(400, "files is empty")

    session_id = new_session_id()
    results = []

    for f in files:
        job_id = uuid.uuid4().hex
        try:
            raw = await f.read()
            text = raw.decode("utf-8")
            data = json.loads(text)

            payload = DesignJSON.model_validate(data)

            payload_dict = (
                payload.model_dump(exclude_none=True)
                if hasattr(payload, "model_dump")
                else json.loads(payload.json(ensure_ascii=False))
            )

            jpg_bytes, debug_html = await render_png_bytes(payload)

            try:
                await upload_all_assets_async(
                    job_id,
                    payload_dict=payload_dict,
                    jpg_bytes=jpg_bytes,
                    debug_html=debug_html,
                )
            except Exception as e:
                print("[storage upload error][/jobs/restore/batch]", job_id, e)
                raise HTTPException(500, f"storage upload failed: {e}")

            # ローカルにもJSON・JPGを保存（一覧/編集画面の高速化）
            try:
                rp = job_paths(job_id)
                rp["json"].write_text(json.dumps(payload_dict, ensure_ascii=False), encoding="utf-8")
                rp["jpg"].write_bytes(jpg_bytes)
            except Exception:
                pass

            event_id = getattr(payload, "event_id", "") or ""
            upsert_job_ok(job_id, f.filename or "restore.json", payload, session_id, event_id)

            results.append({
                "ok": True,
                "filename": f.filename,
                "jobId": job_id,
                "eventId": getattr(payload, "event_id", None),
                "previewUrl": f"/preview/{job_id}.jpg",
            })

        except Exception as e:
            print("[restore/batch error]", f.filename, job_id)
            print(traceback.format_exc())
            results.append({
                "ok": False,
                "filename": f.filename,
                "jobId": job_id,
                "error": str(e),
            })

    ok_count = sum(1 for r in results if r.get("ok"))
    return JSONResponse({
        "ok": True,
        "count": ok_count,
        "results": results,
        "sessionId": session_id,
    })


@app.get("/jobs")
async def list_jobs(
    q: str = "",
    status: Optional[Literal["ok", "error"]] = None,
    warning: str = "",
    manual: Optional[bool] = None,
    locked: Optional[bool] = None,
    min_conf: Optional[float] = None,
    max_conf: Optional[float] = None,
    created_from: str = "",
    created_to: str = "",
    page: int = 1,
    page_size: int = 30,
    order: Literal["updated_desc", "created_desc"] = "updated_desc",
):
    page = max(page, 1)
    page_size = min(max(page_size, 1), 200)
    offset = (page - 1) * page_size

    where: List[str] = []
    params: List[Any] = []

    cf = _parse_date_start(created_from)
    ct = _parse_date_end(created_to)
    if cf:
        where.append("created_at >= %s")
        params.append(cf)
    if ct:
        where.append("created_at <= %s")
        params.append(ct)

    if status:
        where.append("status = %s")
        params.append(status)

    if q and q.strip():
        where.append("(filename ILIKE %s OR title ILIKE %s OR organizer ILIKE %s OR event_id ILIKE %s)")
        like = f"%{q.strip()}%"
        params.extend([like, like, like, like])

    if manual is not None:
        where.append("manual_override = %s")
        params.append(bool(manual))

    if locked is not None:
        where.append("locked = %s")
        params.append(bool(locked))

    if min_conf is not None:
        where.append("confidence >= %s")
        params.append(float(min_conf))

    if max_conf is not None:
        where.append("confidence <= %s")
        params.append(float(max_conf))

    # jsonb array contains: warnings_json @> '["missing_x"]'
    if warning and warning.strip():
        where.append("warnings_json @> %s::jsonb")
        params.append(json.dumps([warning.strip()], ensure_ascii=False))

    where_sql = ("WHERE " + " AND ".join(where)) if where else ""

    order_sql = "ORDER BY updated_at DESC" if order == "updated_desc" else "ORDER BY created_at DESC"

    with db_connect() as con:
        total = con.execute(
            f"SELECT COUNT(*) AS c FROM jobs {where_sql}",
            params,
        ).fetchone()["c"]

        rows = con.execute(
            f"""
            SELECT job_id, filename, session_id, event_id, status,
                   created_at, updated_at, title, organizer, datetime,
                   confidence, warnings_json, manual_override, note,
                   locked, error_message
            FROM jobs
            {where_sql}
            {order_sql}
            LIMIT %s OFFSET %s
            """,
            params + [page_size, offset],
        ).fetchall()

    items = [row_to_job_item(r) for r in rows]

    # プレビュー画像は /preview/{jobId}.jpg 経由で取得するので
    # ここでは署名付き URL を生成しない（1.6s → ~50ms に高速化）

    return {
        "page": page,
        "pageSize": page_size,
        "total": total,
        "items": items,
    }




# ------------------------------------------------------------
# ジョブメタ更新（manual_override / note / locked）
# ------------------------------------------------------------

class JobPatch(BaseModel):
    manual_override: Optional[bool] = None
    note: Optional[str] = None
    locked: Optional[bool] = None

@app.get("/job/{job_id}")
async def get_job(job_id: str):
    p = job_paths(job_id)

    # 1) ローカル json 優先
    if p["json"].exists():
        data = json.loads(p["json"].read_text(encoding="utf-8"))
        data = ensure_display_fields_in_dict(data)
        return JSONResponse({"jobId": job_id, "json": data})

    # 2) DBに存在確認
    with db_connect() as con:
        row = con.execute(
            "SELECT job_id FROM jobs WHERE job_id=%s",
            (job_id,)
        ).fetchone()

    if not row:
        raise HTTPException(status_code=404, detail="job not found")

    # 3) Storage fallback → ダウンロード後ローカルにキャッシュ
    try:
        data = download_storage_json(f"{job_id}/latest.json")
        data = ensure_display_fields_in_dict(data)
        try:
            p["json"].write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
        except Exception:
            pass  # キャッシュ書き込み失敗は無視
        return JSONResponse({"jobId": job_id, "json": data})
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(status_code=404, detail="job json not found")



@app.get("/preview/{job_id}.jpg")
async def preview(job_id: str, background_tasks: BackgroundTasks):
    # ローカルファイルがあれば高速に返す（リダイレクト不要）
    # no-cache: ブラウザは毎回確認するが ETag/Last-Modified で 304 を返すので高速
    local = DATA_DIR / job_id / "preview.jpg"
    if local.exists() and local.stat().st_size > 0:
        return FileResponse(local, media_type="image/jpeg",
                            headers={"Cache-Control": "no-cache"})

    sp = storage_paths(job_id)
    try:
        signed = create_signed_url(sp["preview"], expires_in=600)
        # バックグラウンドでローカルキャッシュ（次回から高速配信）
        def _cache_preview():
            try:
                data = download_storage_file(sp["preview"])
                local.parent.mkdir(parents=True, exist_ok=True)
                local.write_bytes(data)
            except Exception:
                pass
        background_tasks.add_task(_cache_preview)
        return RedirectResponse(url=signed, status_code=307,
                                headers={"Cache-Control": "no-cache"})
    except Exception:
        data = download_storage_file(sp["preview"])
        try:
            local.parent.mkdir(parents=True, exist_ok=True)
            local.write_bytes(data)
        except Exception:
            pass
        return Response(content=data, media_type="image/jpeg",
                        headers={"Cache-Control": "no-cache"})


@app.get("/export/{job_id}.zip")
async def export_zip_single(job_id: str, background_tasks: BackgroundTasks):
    event_id = resolve_event_id(job_id)
    p = job_paths(job_id)
    sp = storage_paths(job_id)

    # ローカル優先 → Storage fallback
    if p["jpg"].exists() and p["jpg"].stat().st_size > 0:
        jpg_bytes = p["jpg"].read_bytes()
    else:
        jpg_bytes = download_storage_file(sp["preview"])
    if p["json"].exists() and p["json"].stat().st_size > 0:
        json_bytes = p["json"].read_bytes()
    else:
        json_bytes = download_storage_file(sp["json"])

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".zip")
    tmp_path = Path(tmp.name)
    tmp.close()

    with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as z:
        z.writestr(f"{event_id}_招聘.jpg", jpg_bytes)
        z.writestr(f"{event_id}_backup.json", json_bytes)

    background_tasks.add_task(lambda: os.remove(tmp_path) if tmp_path.exists() else None)
    return FileResponse(
        tmp_path,
        media_type="application/zip",
        filename=f"{event_id}_export.zip",
    )

@app.get("/download/{job_id}.jpg")
async def download(job_id: str, background_tasks: BackgroundTasks):
    event_id = resolve_event_id(job_id)
    filename = f"{event_id}_招聘.jpg"

    # ローカル優先 → Storage fallback
    local = DATA_DIR / job_id / "preview.jpg"
    if local.exists() and local.stat().st_size > 0:
        return FileResponse(
            local,
            media_type="image/jpeg",
            filename=filename,
        )

    sp = storage_paths(job_id)
    jpg_bytes = download_storage_file(sp["preview"])

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    tmp_path = Path(tmp.name)
    tmp.write(jpg_bytes)
    tmp.close()

    background_tasks.add_task(lambda: os.remove(tmp_path) if tmp_path.exists() else None)

    return FileResponse(
        tmp_path,
        media_type="image/jpeg",
        filename=filename,
    )

@app.get("/debug/{job_id}/latest.json")
async def debug_latest(job_id: str):
    # ローカル優先 → Storage fallback
    local = DATA_DIR / job_id / "latest.json"
    if local.exists() and local.stat().st_size > 0:
        try:
            data = json.loads(local.read_text(encoding="utf-8"))
            data = ensure_display_fields_in_dict(data)
            return JSONResponse(content=data, headers={"Cache-Control": "no-cache"})
        except Exception:
            pass
    sp = storage_paths(job_id)
    data = download_storage_json(sp["json"])
    data = ensure_display_fields_in_dict(data)
    return JSONResponse(content=data, headers={"Cache-Control": "no-cache"})



@app.get("/debug/{job_id}/blocks.json")
async def debug_blocks(job_id: str):
    # ローカル優先 → Storage fallback
    local = DATA_DIR / job_id / "blocks.json"
    if local.exists() and local.stat().st_size > 0:
        try:
            data = json.loads(local.read_text(encoding="utf-8"))
            return JSONResponse(content=data, headers={"Cache-Control": "no-cache"})
        except Exception:
            pass
    sp = storage_paths(job_id)
    data = download_storage_json(sp["debug_blocks"])
    return JSONResponse(content=data, headers={"Cache-Control": "no-cache"})





# ------------------------------------------------------------
# 選択ジョブのPNGをまとめてZIP（納品用）
# ------------------------------------------------------------
def resolve_event_id(job_id: str) -> str:
    try:
        with db_connect() as con:
            row = con.execute(
                "SELECT event_id FROM jobs WHERE job_id=%s",
                (job_id,),
            ).fetchone()
        if row and (row.get("event_id") or "").strip():
            return sanitize_basename(row["event_id"].strip())
    except Exception:
        pass

    return sanitize_basename(job_id)

def sanitize_basename(s: str) -> str:
    s = (s or "").strip()
    s = re.sub(r"[\\/]", "_", s)
    s = re.sub(r"\.pptx$", "", s, flags=re.I)
    s = re.sub(r"\s+", " ", s).strip()
    return s or "file"

def unique_name(base: str, used: dict[str, int]) -> str:
    # baseは拡張子なし
    n = used.get(base, 0) + 1
    used[base] = n
    return base if n == 1 else f"{base} ({n})"


class ExportReq(BaseModel):
    jobIds: List[str] = Field(default_factory=list)
    nameMode: Literal["jobId", "filename"] = "filename"  # zip内ファイル名
    includeJson: bool = False

@app.post("/jobs/export.zip")
async def export_zip(req: ExportReq, background_tasks: BackgroundTasks):
    if not req.jobIds:
        raise HTTPException(400, "jobIds is empty")

    # job_id -> filename（Postgres）
    with db_connect() as con:
        rows = con.execute(
            "SELECT job_id, filename FROM jobs WHERE job_id = ANY(%s)",
            (req.jobIds,),
        ).fetchall()
    mp = {r["job_id"]: (r.get("filename") or "") for r in rows}

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)

    export_id = f"export_{int(time.time())}_{uuid.uuid4().hex}"
    zip_path = EXPORT_DIR / f"{export_id}.zip"

    def pick_base(job_id: str, used_names: dict[str, int]) -> str:
        # eventId を最優先
        base0 = resolve_event_id(job_id)

        # fallback
        if not base0:
            base0 = mp.get(job_id) or job_id

        base0 = sanitize_basename(base0)

        # 重複対応
        used_names[base0] = used_names.get(base0, 0) + 1
        return base0 if used_names[base0] == 1 else f"{base0} ({used_names[base0]})"

    used_names: dict[str, int] = {}
    added = 0

    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as z:
        for job_id in req.jobIds:
            try:
                base = pick_base(job_id, used_names)

                # preview.jpg は Storage から取得
                jpg_bytes = download_storage_file(f"{job_id}/preview.jpg")
                z.writestr(f"{base}_招聘.jpg", jpg_bytes)
                added += 1

                # # latest.json も必要なら Storage から取得
                # if req.includeJson:
                #     try:
                #         json_bytes = download_storage_file(f"{job_id}/latest.json")
                #         z.writestr(f"{base}_backup.json", json_bytes)
                #     except HTTPException:
                #         # json が無い個体は jpg だけ入れて続行
                #         pass

            except HTTPException:
                # preview.jpg が無い job はスキップ
                continue
            except Exception as e:
                print("[jobs/export.zip error]", job_id, e)
                continue

    if added == 0:
        if zip_path.exists():
            os.remove(zip_path)
        raise HTTPException(404, "no exportable jobs found")

    # 送信後にzip削除
    background_tasks.add_task(lambda: os.remove(zip_path) if zip_path.exists() else None)

    return FileResponse(
        str(zip_path),
        media_type="application/zip",
        filename="export.zip",
    )

class ExportPdfReq(BaseModel):
    jobIds: List[str] = Field(default_factory=list)
    pageSize: Literal["fit", "a4"] = "fit"   # fit: 画像サイズそのまま / a4: A4に載せる
    orientation: Literal["portrait", "landscape", "auto"] = "auto"


def _fit_to_a4(img: Image.Image, orientation: str = "auto") -> Image.Image:
    # 150dpiくらいのA4
    PORTRAIT = (1240, 1754)
    LANDSCAPE = (1754, 1240)

    if orientation == "portrait":
        canvas_size = PORTRAIT
    elif orientation == "landscape":
        canvas_size = LANDSCAPE
    else:
        canvas_size = LANDSCAPE if img.width > img.height else PORTRAIT

    canvas = Image.new("RGB", canvas_size, "white")

    ratio = min(canvas_size[0] / img.width, canvas_size[1] / img.height)
    new_w = max(1, int(img.width * ratio))
    new_h = max(1, int(img.height * ratio))

    resized = img.resize((new_w, new_h), Image.LANCZOS)
    x = (canvas_size[0] - new_w) // 2
    y = (canvas_size[1] - new_h) // 2
    canvas.paste(resized, (x, y))
    return canvas


@app.post("/jobs/export.pdf")
async def export_pdf(req: ExportPdfReq, background_tasks: BackgroundTasks):
    if not req.jobIds:
        raise HTTPException(400, "jobIds is empty")

    images: list[Image.Image] = []

    for job_id in req.jobIds:
        sp = storage_paths(job_id)
        jpg_bytes = download_storage_file(sp["preview"])

        try:
            img = Image.open(io.BytesIO(jpg_bytes)).convert("RGB")
        except Exception as e:
            raise HTTPException(500, f"failed to open preview jpg: {job_id}: {e}")

        if req.pageSize == "a4":
            img = _fit_to_a4(img, req.orientation)

        images.append(img)

    if not images:
        raise HTTPException(400, "no images")

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    tmp_path = Path(tmp.name)
    tmp.close()

    first = images[0]
    rest = images[1:]

    first.save(
        tmp_path,
        "PDF",
        resolution=150.0,
        save_all=True,
        append_images=rest,
    )

    background_tasks.add_task(lambda: os.remove(tmp_path) if tmp_path.exists() else None)

    return FileResponse(
        tmp_path,
        media_type="application/pdf",
        filename=f"selected_{len(req.jobIds)}items.pdf",
    )


def cleanup_old_vm_diff_previews(ttl_sec: int = 60 * 100) -> None:
    """
    念のための掃除。
    10分以上前の preview を削除。
    """
    now = time.time()
    for f in VM_DIFF_PREVIEW_DIR.glob("*"):
        try:
            if f.is_file() and (now - f.stat().st_mtime > ttl_sec):
                f.unlink()
        except Exception:
            pass


def delete_file_quietly(path: Path | str) -> None:
    try:
        p = Path(path)
        if p.exists():
            p.unlink()
    except Exception:
        pass





# =========================
# VM取得
# =========================

def get_endai_enja_vm_rows_by_event_id(event_id: str) -> tuple[list[str], list[dict]]:
    event_id = normalize_space(event_id or "")
    if not event_id:
        return [], []

    scope = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive",
    ]

    credentials = get_gsa_credentials(scope)
    gc = gspread.authorize(credentials)

    spreadsheet_key = "1hiV0Ve2cnYyrPkBuZcZIcLWeAnJ-ucNiB0P4owZpXug"
    workbook = gc.open_by_key(spreadsheet_key)
    ws = workbook.worksheet("演題演者（VM）")

    values = ws.get_all_values()
    if not values or len(values) < 2:
        return [], []

    # header row = 1
    headers = make_unique(values[0])

    id_col_candidates = ["講演会ID", "システムID", "event_id", "Event ID"]
    id_col = next((c for c in id_col_candidates if c in headers), None)
    if not id_col:
        return headers, []

    rows: list[dict] = []
    for raw in values[1:]:
        row = {
            headers[i]: raw[i] if i < len(raw) else ""
            for i in range(len(headers))
        }
        rid = normalize_space(row.get(id_col, ""))
        if rid == event_id:
            rows.append(row)

    return headers, rows


def shape_vm_rows_for_diff(rows: list[dict], headers_in_sheet_order: list[str]) -> list[dict]:
    out: list[dict] = []

    for r in rows or []:
        r = r or {}
        shaped: dict[str, str] = {}

        for header in headers_in_sheet_order:
            shaped[header] = normalize_space(r.get(header, ""))

        out.append(shaped)

    return out



# =========================
# blocks抽出 wrapper
# ここは既存関数名に合わせて差し替え
# =========================

def extract_text_blocks_for_vm_diff(file_path: str) -> list[Any]:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        if not pdf_has_extractable_text(file_path):
            raise ValueError(
                "このPDFはテキスト抽出できません。"
                "スキャンPDFや画像PDFの可能性があります。"
                "テキスト抽出可能なPDFをご使用ください。"
            )

        return extract_blocks_from_pdf2(file_path)

    raise ValueError(
        "このファイル形式は対応していません。"
        "テキスト抽出可能な PDF を使用してください。"
    )

PDF_PREVIEW_ZOOM = 2.0

def render_first_pdf_page_to_image(file_path: str) -> tuple[str, int, int]:
    doc = fitz.open(file_path)
    try:
        if len(doc) == 0:
            raise ValueError("pdf has no pages")

        page = doc.load_page(0)
        mat = fitz.Matrix(PDF_PREVIEW_ZOOM, PDF_PREVIEW_ZOOM)
        pix = page.get_pixmap(matrix=mat, alpha=False)

        out_path = DATA_DIR / f"pdf_preview_{uuid.uuid4().hex}.jpg"
        pix.save(str(out_path))

        return str(out_path), pix.width, pix.height
    finally:
        doc.close()


def render_first_slide_to_image(file_path: str) -> tuple[str, int, int]:
    """
    PPTX -> PDF (LibreOffice) -> first page JPG
    """
    src = Path(file_path)
    work_dir = DATA_DIR / f"pptx_preview_{uuid.uuid4().hex}"
    work_dir.mkdir(parents=True, exist_ok=True)

    pdf_path = work_dir / (src.stem + ".pdf")

    try:
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(work_dir),
            str(src),
        ]
        proc = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=120,
        )
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice convert failed: {proc.stderr or proc.stdout}")

        if not pdf_path.exists():
            raise RuntimeError("converted pdf not found")

        preview_img_path, width, height = render_first_pdf_page_to_image(str(pdf_path))
        return preview_img_path, width, height

    finally:
        # pdf は後で不要なので掃除
        try:
            if pdf_path.exists():
                pdf_path.unlink()
        except Exception:
            pass
        try:
            if work_dir.exists():
                work_dir.rmdir()
        except Exception:
            pass

def normalize_blocks_to_preview_pixels(
    blocks: list[Any],
    preview_width: int,
    preview_height: int,
    coord_unit: str | None = None,
    slide_width: float | None = None,
    slide_height: float | None = None,
) -> list[dict]:
    out = []

    for b in blocks or []:
        if isinstance(b, dict):
            getv = b.get
        else:
            getv = lambda k, d=None: getattr(b, k, d)

        unit = coord_unit or getv("_coord_unit", "px")

        left_raw = float(getv("left", 0) or 0)
        top_raw = float(getv("top", 0) or 0)
        width_raw = float(getv("width", 0) or 0)
        height_raw = float(getv("height", 0) or 0)

        if unit == "emu":
            sw = float(slide_width or getv("_slide_width", 1) or 1)
            sh = float(slide_height or getv("_slide_height", 1) or 1)

            scale_x = preview_width / sw
            scale_y = preview_height / sh

            left = left_raw * scale_x
            top = top_raw * scale_y
            width = width_raw * scale_x
            height = height_raw * scale_y

        elif unit == "pdf_page":
            pw = float(getv("_page_width", 1) or 1)
            ph = float(getv("_page_height", 1) or 1)

            scale_x = preview_width / pw
            scale_y = preview_height / ph

            left = left_raw * scale_x
            top = top_raw * scale_y
            width = width_raw * scale_x
            height = height_raw * scale_y

        else:
            left = left_raw
            top = top_raw
            width = width_raw
            height = height_raw

        out.append({
            "text": getv("text", "") or "",
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "max_font_pt": getv("max_font_pt", 0) or 0,
        })

    return out
# =========================
# preview画像生成 wrapper
# ここも既存関数名に合わせて差し替え
# 戻り値: (preview_image_path, width, height)
# =========================

def pdf_has_extractable_text(file_path: str) -> bool:
    try:
        doc = fitz.open(file_path)
        try:
            for page in doc:
                text = page.get_text("text")
                if text and text.strip():
                    return True
            return False
        finally:
            doc.close()
    except Exception:
        return False



def extract_text_blocks_for_vm_diff(file_path: str) -> list[Any]:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        if not pdf_has_extractable_text(file_path):
            raise ValueError(
                "このPDFはテキスト抽出できません。"
                "スキャンPDFや画像PDFの可能性があります。"
                "テキスト抽出可能なPDFをご使用ください。"
            )

        return extract_blocks_from_pdf2(file_path)

    raise ValueError(
        "このファイル形式は対応していません。"
        "テキスト抽出可能な PDF を使用してください。"
    )




# =========================
# endpoint: event_id -> VM rows
# =========================
@app.post("/vm-diff/by-event-id")
async def vm_diff_by_event_id(event_id: str = Form(...)):
    event_id = normalize_space(event_id or "")
    if not event_id:
        raise HTTPException(status_code=400, detail="講演会IDを入力してください。")

    try:
        headers, vm_rows = get_endai_enja_vm_rows_by_event_id(event_id)

        return {
            "ok": True,
            "event_id": event_id,
            "headers": headers,
            "vm_rows": shape_vm_rows_for_diff(vm_rows, headers),
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"VM取得に失敗しました: {e}")


@app.post("/vm-diff/extract-text-blocks")
async def extract_text_blocks_endpoint(
    file: UploadFile = File(...),
    eventId: str = Form(""),
):
    suffix = os.path.splitext(file.filename or "")[1].lower()
    if suffix != ".pdf":
        raise HTTPException(
            status_code=400,
            detail="このファイル形式は対応していません。テキスト抽出可能な PDF を使用してください。"
        )

    tmp_path = ""

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix, dir=str(DATA_DIR)) as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name

        event_id = normalize_space(eventId or "")

        headers: list[str] = []
        vm_rows: list[dict] = []
        if event_id:
            try:
                headers, vm_rows = get_endai_enja_vm_rows_by_event_id(event_id)
            except Exception:
                headers, vm_rows = [], []

        blocks = extract_text_blocks_for_vm_diff(tmp_path)

        return {
            "ok": True,
            "event_id": event_id,
            "headers": headers,
            "vm_rows": shape_vm_rows_for_diff(vm_rows, headers),
            "blocks": blocks_to_dicts(blocks),
        }

    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("extract_text_blocks_endpoint failed")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"failed to extract text blocks: {type(e).__name__}: {e}")
    finally:
        if tmp_path:
            delete_file_quietly(tmp_path)

# =========================
# endpoint: preview 取得後に削除
# =========================

@app.get("/vm-diff/preview/{filename}")
def get_vm_diff_preview(filename: str, background_tasks: BackgroundTasks):
    path = VM_DIFF_PREVIEW_DIR / filename
    if not path.exists():
        raise HTTPException(status_code=404, detail="preview not found")

    # 返却後に削除
    background_tasks.add_task(delete_file_quietly, path)

    return FileResponse(path)


class JobDeleteReq(BaseModel):
    delete_files: bool = True
    force: bool = False  # lockedでも消したい場合だけTrue

def delete_storage_file(remote_path: str):
    url = f"{SUPABASE_URL}/storage/v1/object/{SUPABASE_BUCKET}/{remote_path}"
    headers = {
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
    }

    res = requests.delete(url, headers=headers, timeout=30)
    if res.status_code not in (200, 204, 404):
        raise RuntimeError(f"storage delete failed: {res.status_code} {res.text}")

@app.delete("/job/{job_id}")
async def delete_job(job_id: str, req: JobDeleteReq = JobDeleteReq()):
    with db_connect() as con:
        row = con.execute(
            "SELECT locked FROM jobs WHERE job_id=%s",
            (job_id,),
        ).fetchone()

        if not row:
            raise HTTPException(404, "job not found")

        if bool(row.get("locked")) and not req.force:
            raise HTTPException(400, "This job is locked.")

        con.execute("DELETE FROM jobs WHERE job_id=%s", (job_id,))

    deleted_storage = False
    storage_errors = []

    if req.delete_files:
        sp = storage_paths(job_id)
        for remote_path in sp.values():
            try:
                delete_storage_file(remote_path)
            except Exception as e:
                storage_errors.append(f"{remote_path}: {e}")

        deleted_storage = len(storage_errors) == 0

    return {
        "ok": True,
        "jobId": job_id,
        "deletedStorage": deleted_storage,
        "storageErrors": storage_errors,
    }


# ---------------- 正解DB API ----------------

class CorrectAnswerReq(BaseModel):
    job_id: str
    design: dict  # 確定済みDesignJSON


@app.post("/correct-answer/register")
async def register_correct_answer(req: CorrectAnswerReq):
    """ユーザーが確定した結果を正解DBに登録する"""
    job_id = req.job_id
    design = req.design

    # blocksをストレージから取得
    blocks_text = ""
    _blocks_json_for_save = []
    try:
        sp = storage_paths(job_id)
        blocks_url = f"{SUPABASE_URL}/storage/v1/object/authenticated/{sp['debug_blocks']}"
        headers = {
            "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
            "apikey": SUPABASE_SERVICE_ROLE_KEY,
        }
        resp = requests.get(blocks_url, headers=headers, timeout=30)
        if resp.status_code == 200:
            blocks_data = resp.json()
            _blocks_json_for_save = blocks_data if isinstance(blocks_data, list) else []
            blocks_text = " ".join(
                (b.get("text", "") if isinstance(b, dict) else "")
                for b in _blocks_json_for_save
            )
    except Exception as e:
        print(f"[correct-answer] blocks fetch failed: {e}")

    event_title = design.get("event_title", "")

    save_correct_answer(
        blocks_text=blocks_text,
        correct_json=design,
        event_title=event_title,
        job_id=job_id,
        blocks_json=_blocks_json_for_save or None,
    )

    answers = _load_correct_answers()
    return {"ok": True, "total_answers": len(answers), "job_id": job_id}


@app.get("/correct-answer/list")
async def list_correct_answers():
    """登録済みの正解データ一覧"""
    answers = _load_correct_answers()
    return {
        "total": len(answers),
        "answers": [
            {
                "job_id": a.get("job_id", ""),
                "event_title": a.get("event_title", ""),
                "created_at": a.get("created_at", ""),
            }
            for a in answers
        ],
    }


@app.delete("/correct-answer/{job_id}")
async def delete_correct_answer(job_id: str):
    """正解データを削除"""
    try:
        with db_connect() as con:
            result = con.execute(
                "DELETE FROM correct_answers WHERE job_id = %s", (job_id,)
            )
            deleted = result.rowcount
            con.commit()
        return {"ok": True, "deleted": deleted}
    except Exception as e:
        raise HTTPException(500, f"delete failed: {e}")
